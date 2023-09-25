from pathlib import Path
from enum import Enum
import pandas as pd
import numpy as np
import streamlit as st
from plotly.subplots import make_subplots
import plotly.graph_objects as go
import plotly.figure_factory as ff
import plotly.express as px
import plotly.io as pio
import pip
from PIL import Image
import openpyxl
import lasio
import time
import io
from io import BytesIO
from io import StringIO
import pymongo
from st_aggrid import AgGrid
from fpdf import FPDF
import os
from pylatex import Document, Section, Subsection, Command, MiniPage, LargeText, Figure
import pdflatex
from streamlit_plotly_events import plotly_events
import asyncio
import altair as alt
import yagmail
import requests
import re
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import chardet 
import csv
from streamlit_toggle import st_toggle_switch
import traceback
import csv as cf
import gc
from tempfile import NamedTemporaryFile

# https://github.com/predict-idlab/plotly-resampler
# plotly-resampler==0.8.3.2
# Automatic resampler to aggregate data in graph view
# from plotly_resampler import register_plotly_resampler
# Manually aggregate data for more speedup, configurability
# from plotly_resampler import FigureResampler, FigureWidgetResampler

# https://bithub.ulterra.com/get_bit_details/50118
# returns {"serial_number": "50118", "type": "RPS613", "size": "6.750", "item_no": "U03776", "body_material": "MATRIX", "data_found": true}

import base64
from xhtml2pdf import pisa
from weasyprint import HTML

#  Azure blob setup
# from azure.storage.blob import *
from azure.storage.blob import BlobServiceClient
# import dotenv

userdata = requests.get('https://ipapi.co/json').json()

# st.secrets['blobname']
# st.secrets['blobkey']
# st.secrets['container']
connection_string = 'DefaultEndpointsProtocol=https;AccountName=' + st.secrets['blobname'] + ';AccountKey=' + st.secrets['blobkey'] + ';EndpointSuffix=core.windows.net'

blob_service = BlobServiceClient.from_connection_string(connection_string)

def blob_upload(df, df_name):
    blob_client = blob_service.get_blob_client(container=st.secrets['container'], blob=df_name)
    try:
        output = df.to_csv(index=False, encoding="utf-8")
    except Exception as e:
        print(e)
        st.error(f"error blob df-to-csv: {e}")
    
    try:
        blob_exists = blob_client.exists()
        if not blob_exists:
            blob_client.upload_blob(output, overwrite=False)
            status = 'Loaded'
        else:            
            status = 'Exists'
        # , blob_type="BlockBlob")
    except Exception as e:
        print(e)
        st.error(f"error blob upload: {e}")

    return status 

# yag = yagmail.SMTP()
# yag = yagmail.SMTP('mygmailusername', 'mygmailpassword')
yag = yagmail.SMTP(st.secrets['yagmailusername'], st.secrets['yagmailpassword'])
# contents = [
#     "This is the body, and here is just text http://somedomain/image.png",
#     "You can find an audio file attached.", '/local/path/to/song.mp3'
# ]
# yag.send('to@someone.com', 'subject', contents)

# # Alternatively, with a simple one-liner:
# yagmail.SMTP('mygmailusername').send('to@someone.com', 'subject', contents)

# setup for dataframe cleaning
class Hdr(str, Enum):
    date = 'YYYY/MM/DD'
    time = 'HH:MM:SS'
    bit = 'Bit Diameter'
    diff = 'Differential Pressure (psi)'
    hole_depth = 'Hole Depth (feet)'
    bit_depth = 'Bit Depth'
    bit_size = 'Bit Size'
    rop = 'Rate Of Penetration (ft_per_hr)'
    arop = 'Adira ROP (ft_per_hr)'
    crop = 'Rate Of Penetration (ft_per_hr)'
    rpm = 'Rotary RPM (RPM)'
    wob = 'Weight on Bit (klbs)'
    trq = 'Rotary Torque'
    pump = 'Pump Pressure'
    flowrate = "Flow Rate"
    dt = 'DateTime'
    dtn = 'DateTimeNormalized'
    dte = 'DateTimeElapsed'
    gr = 'Gamma Ray'
    hook = 'Hook Load (klbs)'
    magtf = 'Magnetic Toolface (degrees)'
    gravtf = 'Gravity Toolface (degrees)'
    dtc = 'DTC'
    dts = 'DTS'
    sprstate = 'Superstate'
    substate = 'Substate'
    az = 'Azimuth'
    inc = 'Inclinitation'
    memo = 'Memos'
    nan = 'Empty'


class DropOptions:
    yes = 'Yes'
    no = 'No'

# ROP, RotRPM, WOB, Diff, TotalRPM, FlowRate, PumpPress, TD Torque, Gamma, DoC, MSE
# for Chee, Canada 	True Vertical Depth (meters)	Standpipe Pressure (kPa)	On Bottom Hours (hrs)	DNMX_SHKA (G) (no_unit)	DNMX_SHKL (G) (no_unit)	DNMX_SSR ((0-3) Units) (no_unit)	PD STEER (Units) (no_unit)	PDAZM (Degrees) (no_unit)	PDINC (Degrees) (no_unit)	AutoDriller Throttle (percent)	Tool Face (degrees)	Rig Sub State (unitless)	Time Of Penetration (min_per_m)
# another chee, DateTime	Hole Depth(m)	Bit Position(m)	Bit Hours(hr)			Mud Weight In(kg/mÂ³)	Toolface Grav(deg)	Toolface Mag(deg)		Block Height(m)	Gamma Ray(API)	Bit TVD(m)	Axial Shock Level ()	Axial Shock RMS()	Ensign Bit Weight(kdaN)	Ensign Pump Pressure(kPa)	Ensign TD Actual Torque(kftÂ·lbf)	Ensign TD Actual RPM(RPM)	Lateal Shock Level ()	Lateral Shock RMS ()	Stick Slip Level (deg)

final_column_name_map = {
    'YYYY/MM/DD': Hdr.date,
    'YYYYMMDD': Hdr.date,
    'Date': Hdr.date,
    'DATE': Hdr.date,
    'HH:MM:SS': Hdr.time,
    'HHMMSS': Hdr.time,
    'Date Time': Hdr.dt,
    'Time': Hdr.dt,
    'DateTime': Hdr.dt,
    'DATE TIME': Hdr.dt,
    'TIME': Hdr.dt,
    'DATETIME': Hdr.dt,
    'YYYY/MM/DD HH:MM:SS': Hdr.dt,
    'AQ. TIME': Hdr.dt,
    'Depth': Hdr.hole_depth,
    'Depth (Feet)': Hdr.hole_depth,
    'Depth(Feet)': Hdr.hole_depth,
    'Depth (feet)': Hdr.hole_depth,
    'Depth(feet)': Hdr.hole_depth,
    'DEPTH': Hdr.hole_depth,
    'DEPTH (FEET)': Hdr.hole_depth,
    'DEPTH(FEET)': Hdr.hole_depth,
    'DEPTH (FT)': Hdr.hole_depth,
    'DEPTH(FT)': Hdr.hole_depth,
    'DEPT': Hdr.hole_depth,
    'Hole Depth (feet)': Hdr.hole_depth,
    'Hole Depth (meters)': Hdr.hole_depth,
    'Hole Depth(m)': Hdr.hole_depth,
    'Hole Depth': Hdr.hole_depth,
    'HoleDepth': Hdr.hole_depth,
    'hole depth': Hdr.hole_depth,
    'holedepth': Hdr.hole_depth,
    'HOLE DEPTH': Hdr.hole_depth,
    'HOLE DEPTH (FEET)': Hdr.hole_depth,
    'HOLE DEPTH (FT)': Hdr.hole_depth,
    'HOLE DEPTH(FT)': Hdr.hole_depth,
    'HOLE DEPTH (METERS)': Hdr.hole_depth,
    'HOLE DEPTH (M)': Hdr.hole_depth,
    'HOLE DEPTH(M)': Hdr.hole_depth,
    'HOLEDEPTH': Hdr.hole_depth,
    'Bit Size': Hdr.bit_size,
    'Bit Size (inches)': Hdr.bit_size,
    'BIT SIZE': Hdr.bit_size,
    'BIT SIZE (INCHES)': Hdr.bit_size,
    'BIT SIZE (IN)': Hdr.bit_size,
    'BITSIZE': Hdr.bit_size,
    'BITSIZE (INCHES)': Hdr.bit_size,
    'BITSIZE (IN)': Hdr.bit_size,    
    'Bit Size': Hdr.bit_size,
    'Bit Diameter': Hdr.bit_size,
    'Bit Dia': Hdr.bit_size,
    'BIT SIZE': Hdr.bit_size,
    'BIT DIAMETER': Hdr.bit_size,
    'BIT DIA': Hdr.bit_size,
    'Bit Depth': Hdr.bit_depth,
    'Bit Depth (feet)': Hdr.bit_depth,
    'Bit Depth (meters)': Hdr.bit_depth,
    'Bit Depth (Feet)': Hdr.bit_depth,
    'Bit Depth (Meters)': Hdr.bit_depth,
    'BIT DEPTH': Hdr.bit_depth,
    'BIT DEPTH (FEET)': Hdr.bit_depth,
    'BIT DEPTH (FT)': Hdr.bit_depth,
    'BIT DEPTH (METERS)': Hdr.bit_depth,
    'BIT DEPTH (M)': Hdr.bit_depth,
    'Bit Position(m)': Hdr.bit_depth,
    'Bit Position': Hdr.bit_depth,
    'BIT POSITION (M)': Hdr.bit_depth,
    'BIT POSITION (FT)': Hdr.bit_depth,
    'BIT POSITION(M)': Hdr.bit_depth,
    'BIT POSITION(FT)': Hdr.bit_depth,
    'BIT POSITION': Hdr.bit_depth,
    'WOB': Hdr.wob,   
    'WOB (SURFACE)': Hdr.wob,   
    'WOB (K#)': Hdr.wob,  
    'WOB (k#)': Hdr.wob,  
    'Weight on Bit (klbs)': Hdr.wob,
    'Weight on Bit (kDaN)': Hdr.wob,
    'Weight on Bit(kdaN)': Hdr.wob,
    'Weight on Bit': Hdr.wob,   
    'Bit Weight': Hdr.wob,    
    'BIT WEIGHT': Hdr.wob, 
    'WEIGHT ON BIT (KLBS)': Hdr.wob,
    'WEIGHT ON BIT (KDAN)': Hdr.wob,
    'WEIGHT ON BIT(KDAN)': Hdr.wob,
    'WEIGHT ON BIT': Hdr.wob,
    'Rate Of Penetration': Hdr.rop,
    'Rate Of Penetration (ft_per_hr)': Hdr.rop,
    'Rate Of Penetration (m_per_hr)': Hdr.rop,
    'ROP - Average': Hdr.rop,
    'ROP - Average(m/hr)': Hdr.rop,
    'ROP': Hdr.rop,
    'ROP (fph)': Hdr.rop,
    'ROP/5 AVG (ft/hr)': Hdr.rop,
    'RATE OF PENETRATION': Hdr.rop,
    'RATE OF PENETRATION (FT_PER_HR)': Hdr.rop,
    'RATE OF PENETRATION (M_PER_HR)': Hdr.rop,
    'ROP - AVERAGE': Hdr.rop,
    'ROP - AVERAGE(M/HR)': Hdr.rop,
    'ROP': Hdr.rop,
    'ROP (FPH)': Hdr.rop,
    'ROP/5 AVG (FT/HR)': Hdr.rop,
    'Rotary RPM': Hdr.rpm,
    'Rotary RPM (RPM)': Hdr.rpm,
    'Top Drive RPM': Hdr.rpm,
    'Top Drive RPM(RPM)': Hdr.rpm,
    'Top Drive RPM (RPM)': Hdr.rpm,
    'RPM': Hdr.rpm,
    'RPM (rev/mins)': Hdr.rpm,
    'RPM (REV/MINS)': Hdr.rpm,
    'RPM(rev/mins)': Hdr.rpm,
    'RPM(REV/MINS)': Hdr.rpm,
    'ROTARY RPM': Hdr.rpm,
    'ROTARY RPM (RPM)': Hdr.rpm,
    'TOP DRIVE RPM': Hdr.rpm,
    'TOP DRIVE RPM(RPM)': Hdr.rpm,
    'TOP DRIVE RPM (RPM)': Hdr.rpm,
    'Differential Pressure': Hdr.diff,
    'Differential Pressure (psi)': Hdr.diff,
    'Differential Pressure (kPa)': Hdr.diff,
    'Diff Press(kPa)': Hdr.diff,
    'Diff Press (kPa)': Hdr.diff,
    'Diff Press(PSI)': Hdr.diff,
    'Diff Press (PSI)': Hdr.diff,
    'DIFF': Hdr.diff,
    'Diff Press': Hdr.diff,
    'DIFFERENTIAL PRESSURE': Hdr.diff,
    'DIFFERENTIAL PRESSURE (PSI)': Hdr.diff,
    'DIFFERENTIAL PRESSURE (KPA)': Hdr.diff,
    'DIFF PRESS(KPA)': Hdr.diff,
    'DIFF PRESS (KPA)': Hdr.diff,
    'DIFF PRESS(PSI)': Hdr.diff,
    'DIFF PRESS (PSI)': Hdr.diff,
    'DIFF PRESS': Hdr.diff,
    'Rotary Torque': Hdr.trq,
    'Rotary Torque (unitless)': Hdr.trq,
    'Convertible Torque (kft_lb)': Hdr.trq,
    'Top Drive Torque': Hdr.trq,
    'Top Drive Torque(ft·lbf)': Hdr.trq,
    'Torque': Hdr.trq,
    'TORQUE': Hdr.trq,
    'Torque (FtLB)': Hdr.trq,
    'ROTARY TORQUE': Hdr.trq,
    'ROTARY TORQUE (UNITLESS)': Hdr.trq,
    'CONVERTIBLE TORQUE (KFT_LB)': Hdr.trq,
    'TOP DRIVE TORQUE': Hdr.trq,
    'TOP DRIVE TORQUE(FT·LBF)': Hdr.trq,
    'TORQUE (KFT.LBF)': Hdr.trq,
    'TORQUE (FTLB)': Hdr.trq,
    'Flow Rate': Hdr.flowrate,
    'Flowrate': Hdr.flowrate,
    'Total GPM (GPM)': Hdr.flowrate,
    'Total Pump Output': Hdr.flowrate,
    'Total Pump Output (gal_per_min)': Hdr.flowrate,
    'Total Pump Output (m3_per_min)': Hdr.flowrate,
    'TOTAL GPM (GPM)': Hdr.flowrate,
    'TOTAL PUMP OUTPUT': Hdr.flowrate,
    'TOTAL PUMP OUTPUT (GAL_PER_MIN)': Hdr.flowrate,
    'TOTAL PUMP OUTPUT (M3_PER_MIN)': Hdr.flowrate,
    'FLOW IN (GPM)': Hdr.flowrate,
    'Pump Pressure': Hdr.pump,
    'Pump Pressure(kPa)': Hdr.pump,
    'Pump Output TTL(m³/min)': Hdr.pump,
    'Pump Press 1 (PSI)': Hdr.pump,
    'PUMP PRESSURE': Hdr.pump,
    'PUMP PRESSURE(KPA)': Hdr.pump,
    'PUMP OUTPUT TTL(M³/MIN)': Hdr.pump,
    'PUMP PRESS 1 (PSI)': Hdr.pump,
    'SPP (PSI)': Hdr.pump,
    'Rig Super State (unitless)': Hdr.sprstate,
    'Rig Super State': Hdr.sprstate,
    'Rig Sub State (unitless)': Hdr.substate,
    'Rig Sub State': Hdr.substate,
    'RIG SUPER STATE (UNITLESS)': Hdr.sprstate,
    'RIG SUPER STATE': Hdr.sprstate,
    'RIG SUB STATE (UNITLESS)': Hdr.substate,
    'RIG SUB STATE': Hdr.substate,
    'Azimuth': Hdr.az,
    'Azimuth (degrees)': Hdr.az,
    'Svy Azimuth': Hdr.az,
    'Svy Azimuth(deg)': Hdr.az,
    'Svy Azimuth (deg)': Hdr.az,
    'Svy Azimuth(Deg)': Hdr.az,
    'Svy Azimuth (Deg)': Hdr.az,
    'AZIMUTH': Hdr.az,
    'AZIMUTH (DEGREES)': Hdr.az,
    'SVY AZIMUTH': Hdr.az,
    'SVY AZIMUTH(DEG)': Hdr.az,
    'SVY AZIMUTH (DEG)': Hdr.az,
    'SVY AZIMUTH(DEG)': Hdr.az,
    'SVY AZIMUTH (DEG)': Hdr.az,
    'Inclination (degrees)': Hdr.inc,
    'Svy Inclination': Hdr.inc,
    'Svy Inclination(deg)': Hdr.inc,
    'Svy Inclination (deg)': Hdr.inc,
    'Svy Inclination (DEG)': Hdr.inc,
    'INCLINATION (DEGREES)': Hdr.inc,
    'SVY INCLINATION': Hdr.inc,
    'SVY INCLINATION(DEG)': Hdr.inc,
    'SVY INCLINATION (DEG)': Hdr.inc,
    'SVY INCLINATION (DEG)': Hdr.inc,
    'Gamma': Hdr.gr,
    'Gamma at Bit (api)': Hdr.gr,
    'Gamma (api)': Hdr.gr,
    'Gamma Ray': Hdr.gr, 
    'GammaRay': Hdr.gr, 
    'GR': Hdr.gr,
    'Gamma Ray 1 (Borehole Corr) (API)': Hdr.gr,
    'GAMMA': Hdr.gr,
    'GAMMA AT BIT (API)': Hdr.gr,
    'GAMMA (API)': Hdr.gr,
    'GAMMA RAY': Hdr.gr, 
    'GAMMARAY': Hdr.gr, 
    'GR': Hdr.gr,
    'GAMMA RAY 1 (BOREHOLE CORR) (API)': Hdr.gr,
    'Hookload': Hdr.hook,
    'HOOKLOAD': Hdr.hook,
    'Hook Load (klbs)': Hdr.hook,
    'Hook Load': Hdr.hook,
    'Hook Load(kdaN)': Hdr.hook,
    'HOOK LOAD (KLBS)': Hdr.hook,
    'HOOK LOAD': Hdr.hook,
    'HOOK LOAD(KDAN)': Hdr.hook,
    'HOOKLOAD (KLBF)': Hdr.hook,
    'Magnetic Toolface (degrees)': Hdr.magtf,
    'MAGNETIC TOOLFACE (DEGREES)': Hdr.magtf,
    'Gravity Toolface (degrees)': Hdr.gravtf,
    'GRAVITY TOOLFACE (DEGREES)': Hdr.gravtf,
    'DTC': Hdr.dtc,
    'DTS': Hdr.dts,
    'XDTS': Hdr.dts,
    'XDT': Hdr.dtc,
    'Memos': Hdr.memo,  
    'MEMOS': Hdr.memo, 
    '': Hdr.nan, 
}

# df.columns.str.match('Unnamed'): Hdr.rop

regex1_column_name_map = {    
    'Rate Of Pe': Hdr.arop,
}

regex2_column_name_map = {
    re.compile("^['rate of pen']"): Hdr.arop,
    
}


class adiraPDF(FPDF):
    def header(self):
        # # Select Arial bold 15
        # self.set_font('Arial', 'B', 15)       
        self.set_font('Arial', '', 12)  
        self.image('images/adira-logo-name-lrg.png', 10, 10, 100)
        # self.cell(30, 30, f'header-REPORT',0,1)
        
    def footer(self):
        # Go to 1.5 cm from bottom
        self.set_y(-15)
        # Select Arial italic 8
        self.set_font('Arial', 'I', 8)
        # Print centered page number
        self.cell(0, 10, 'Page %s' % self.page_no(), 0, 0, 'C')
        
        self.image('Ulterra_teal_250px.png', self.set_x(-20), self.set_y(-35), 70)
        # self.cell(40, 20, f'Created by ADIRA v2.34')
        self.text(40, 605, f'Created by ADIRA v2.34')
        

def figure_to_base64(figures):
    images_html = ""
    for figure in figures:
        image = str(base64.b64encode(figure.to_image(format="png", scale=2)))[2:-1]
        images_html += (f'<img src="data:image/png;base64,{image}"><br>')
    return images_html
 
def create_html_report(template_file, images_html):
    with open(template_file,'r') as f:
        template_html = f.read()
    report_html = template_html.replace("{{ FIGURES }}", images_html)
    return report_html
 
def convert_html_to_pdf(source_html, output_filename):
    with open(f"{output_filename}", "w+b") as f:
        pisa_status = pisa.CreatePDF(source_html, dest=f)
    return pisa_status.err

def combine_plotly_figs_to_html(plotly_figs, html_fname, include_plotlyjs='cdn', 
                                separator=None, auto_open=False):
    with open(html_fname, 'w') as f:
        f.write(plotly_figs[0].to_html(include_plotlyjs=include_plotlyjs))
        for fig in plotly_figs[1:]:
            if separator:
                f.write(separator)
            f.write(fig.to_html(full_html=False, include_plotlyjs=False))

    if auto_open:
        import pathlib, webbrowser
        uri = pathlib.Path(html_fname).absolute().as_uri()
        webbrowser.open(uri)

def convert_df(df):
   return df.to_csv(index=False).encode('utf-8')

# def spherical_to_cartesian(theta, phi):
#     x = math.cos(phi) * math.sin(theta)
#     y = math.sin(phi) * math.sin(theta)
#     z = math.cos(theta)
#     return x, y, z

# disabled cache because it broken pptx creation
# @st.cache_data(hash_funcs={dict: lambda _: None}) # hash_funcs because dict can't be hashed
# @st.cache_data
def make_traces(well_num, def_ref_df, def_df, def_dfsvy, well_color, chnl_num, interp, fig_meter, fig_ch, fig_dvd, fig_dvdi, fig_rotsld, fig_bx, fig_svymap):
    #  Single Well Comparison mode
    if well_num == 0:
        # if chnl_num == 5:
        #     fig_ch.add_trace(go.Scatter(legendgroup="dt", x=def_df['dte'], y=def_df[Hdr.hole_depth], line= {'color':'#121212', 'width':0.4}), row=1, col=1)
        #     fig_ch.add_trace(go.Scatter(legendgroup="rop", x=def_df[Hdr.rop].rolling(interp).mean(), y=def_df[Hdr.hole_depth], line= {'color':'#008C9A', 'width':0.4}), row=1, col=2)
        #     fig_ch.add_trace(go.Scatter(legendgroup="wob", x=def_df[Hdr.wob].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line= {'color':'#8ACFDD', 'width':0.4}), row=1, col=3)
        #     fig_ch.add_trace(go.Scatter(legendgroup="rpm", x=def_df[Hdr.rpm].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line= {'color':'#7A7D81', 'width':0.4}), row=1, col=4)
        #     fig_ch.add_trace(go.Scatter(legendgroup="diff", x=def_df[Hdr.diff].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line= {'color':'#98C007', 'width':0.4}), row=1, col=5)
        # else:
        
        if Hdr.hole_depth in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup="dt",name='Hours',x=def_df['dte'], y=def_df[Hdr.hole_depth], line= {'color':'#121212', 'width':0.4}, fill='tozerox'), row=1, col=1)
        if Hdr.rop in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup="rop",name='ROP', x=def_df[Hdr.rop].rolling(interp).mean(), y=def_df[Hdr.hole_depth], line= {'color':'#008C9A', 'width':0.4}), row=1, col=2)
            fig_ch.update_xaxes(title_text=f'<span style="color:#121212"> <b>AVG:</b> {round(np.mean(def_df[Hdr.rop]),1)}</span>',title_standoff=1, row=1, col=2)
        if Hdr.wob in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup="wob",name='WOB', x=def_df[Hdr.wob].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line= {'color':'#8ACFDD', 'width':0.4}), row=1, col=3)
            fig_ch.update_xaxes(title_text=f'<span style="color:#8ACFDD"> <b>AVG:</b> {round(np.mean(def_df[Hdr.wob]),1)}</span>',title_standoff=1, row=1, col=3)
        if Hdr.rpm in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup="rpm",name='RPM', x=def_df[Hdr.rpm].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line= {'color':'#7A7D81', 'width':0.4}), row=1, col=4)
            fig_ch.update_xaxes(title_text=f'<span style="color:#7A7D81"> <b>AVG:</b> {round(np.mean(def_df[Hdr.rpm]),1)}</span>',title_standoff=1, row=1, col=4)
        if Hdr.diff in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup="diff",name='DIFF', x=def_df[Hdr.diff].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line= {'color':'#98C007', 'width':0.4}), row=1, col=5)
            fig_ch.update_xaxes(title_text=f'<span style="color:#98C007"> <b>AVG:</b> {round(np.mean(def_df[Hdr.diff]),1)}</span>',title_standoff=1, row=1, col=5)
        if Hdr.trq in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup="trq",name='TRQ', x=def_df[Hdr.trq].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line= {'color':'#008C9A', 'width':0.4}), row=1, col=6)
            fig_ch.update_xaxes(title_text=f'<span style="color:#008C9A"> <b>AVG:</b> {round(np.mean(def_df[Hdr.trq]),1)}</span>',title_standoff=1, row=1, col=6)
        if Hdr.flowrate in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup="flow",name='Flowrate', x=def_df[Hdr.flowrate].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line= {'color':'#8ACFDD', 'width':0.4}), row=1, col=7)
            fig_ch.update_xaxes(title_text=f'<span style="color:#8ACFDD"> <b>AVG:</b> {round(np.mean(def_df[Hdr.flowrate]),1)}</span>',title_standoff=1, row=1, col=7)
        if Hdr.gr in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup="gr",name='GR', x=def_df[Hdr.gr].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line= {'color':'#7A7D81', 'width':0.4}), row=1, col=8)
            fig_ch.update_xaxes(title_text=f'<span style="color:#7A7D81"> <b>AVG:</b> {round(np.mean(def_df[Hdr.gr]),1)}</span>',title_standoff=1, row=1, col=8)
        if 'ads-spr' in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup="State",name='Super', x=def_df['ads-spr'], y=def_df[Hdr.hole_depth],line= {'color':'#98C007', 'width':2}), row=1, col=9)
            if 'ads-sub' in def_df.columns:
                fig_ch.add_trace(go.Scatter(legendgroup="State",name='Sub', x=def_df['ads-sub'], y=def_df[Hdr.hole_depth],line= {'color':'#008C9A', 'width':0.4}), row=1, col=9)
                fig_ch.update_xaxes(title_text=f'<span style="color:#98C007"> <b>Super</b></span>       <span style="color:#008C9A"> <b>Sub</b></span>',title_standoff=1, row=1, col=9)
        if 'MSE' in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup="MSE",name='MSE', x=def_df['MSE'], y=def_df[Hdr.hole_depth],line= {'color':'#8ACFDD', 'width':0.4}), row=1, col=10)
            fig_ch.update_xaxes(title_text=f'<span style="color:#8ACFDD"> <b>AVG:</b> {round(np.mean(def_df["MSE"]),1)}</span>',title_standoff=1, row=1, col=10)
        
        # fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.25,y=1.045,text=(f'<b>ROP:</b> {round(np.mean(df[Hdr.rop]),1)}'),textangle=0,showarrow=False,font=dict(color='#008C9A',size=12))
        # fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.45,y=1.045,text=(f'<b>WOB:</b> {round(np.mean(df[Hdr.wob]),1)}'),textangle=0,showarrow=False,font=dict(color='#8ACFDD',size=12))
        # fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.65,y=1.045,text=(f'<b>RPM:</b> {round(df[Hdr.rpm].mean(),0)}'),textangle=0,showarrow=False,font=dict(color='#7A7D81',size=12))
        # fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.85,y=1.045,text=(f'<b>DIFF:</b> {round(df[Hdr.diff].mean(),1)}'),textangle=0,showarrow=False,font=dict(color='#98C007',size=12))
        
        # if Hdr.rop in def_df.columns:
        #     fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.14,y=1.045,text=(f'<b>AVG:</b> {round(np.mean(df[Hdr.rop]),1)}'),textangle=0,showarrow=False,font=dict(color='#008C9A',size=12))
        # if Hdr.wob in def_df.columns:
        #     fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.25,y=1.045,text=(f'<b>AVG:</b> {round(np.mean(df[Hdr.wob]),1)}'),textangle=0,showarrow=False,font=dict(color='#8ACFDD',size=12))
        # if Hdr.rpm in def_df.columns:
        #     fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.37,y=1.045,text=(f'<b>AVG:</b> {round(df[Hdr.rpm].mean(),0)}'),textangle=0,showarrow=False,font=dict(color='#7A7D81',size=12))
        # if Hdr.diff in def_df.columns:
        #     fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.48,y=1.045,text=(f'<b>AVG:</b> {round(df[Hdr.diff].mean(),1)}'),textangle=0,showarrow=False,font=dict(color='#98C007',size=12))
        # if Hdr.trq in def_df.columns:
        #     fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.6,y=1.045,text=(f'<b>AVG:</b> {round(df[Hdr.trq].mean(),1)}'),textangle=0,showarrow=False,font=dict(color='#008C9A',size=12))
        # if Hdr.flowrate in def_df.columns:
        #     fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.7,y=1.045,text=(f'<b>AVG:</b> {round(df[Hdr.flowrate].mean(),1)}'),textangle=0,showarrow=False,font=dict(color='#8ACFDD',size=12))
        # if Hdr.gr in def_df.columns:
        #     fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.81,y=1.045,text=(f'<b>AVG:</b> {round(df[Hdr.gr].mean(),1)}'),textangle=0,showarrow=False,font=dict(color='#7A7D81',size=12))
        # if Hdr.sprstate in def_df.columns:
        #     fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.93,y=1.045,text=(f'Super'),textangle=0,showarrow=False,font=dict(color='#98C007',size=10))
        # if Hdr.substate in def_df.columns:
        #     fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.97,y=1.045,text=(f'Sub'),textangle=0,showarrow=False,font=dict(color='#008C9A',size=10))
        # if 'MSE' in def_df.columns:
        #     fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.97,y=1.045,text=(f'<b>AVG:</b> {round(df["MSE"].mean(),1)}'),textangle=0,showarrow=False,font=dict(color='#7A7D81',size=12))
        
        fig_ch.update_layout(showlegend = False)
        
        if not len(def_dfsvy) == 0:
            fig_svymap.add_trace(go.Scatter3d(name='Well_1',x=def_dfsvy['3dx'],y=def_dfsvy['3dy'],z=def_dfsvy['3dz'],mode="lines",line= {'color':df[Hdr.rop],'colorscale':'Viridis','width':2, 'showscale':True}))
            # fig_svymap.update_layout(showscale=True)
        # marker=dict(size=4,color=z,colorscale='plasma',),
        
        # if 'dte' in def_df.columns:  
        if set(['dte', Hdr.hole_depth]).issubset(def_df.columns):                     
            fig_dvd.add_trace(go.Scatter(x=def_df['dte'], y=def_df[Hdr.hole_depth], line= {'color':well_color}, name='Well 1', fill='tozerox'))
            fig_dvdi.add_trace(go.Scatter(x=def_df['dte'], y=def_df[Hdr.hole_depth], line= {'color':well_color}, name=f'Well_1'), row=1,col=1)  
                  
        if Hdr.inc in def_df.columns:
            fig_dvdi.add_trace(go.Scatter(x=def_df[Hdr.inc], y=def_df[Hdr.hole_depth], line= {'color':well_color}, line_shape='spline', line_smoothing=1.3, name=f'Well_1'), row=1,col=2)
        # 'line': {'shape': 'spline', 'smoothing': 1.3}
        try:
            # number={"font":{"size":20}}
            fig_meter.add_trace(go.Indicator(title='ROP',mode="number+delta",value=round(np.mean(df[Hdr.rop]),1),number={'suffix': ' ft/hr'},delta={'reference': round(np.mean(df[Hdr.rop]), 1),'relative':True,'valueformat': '.1%'},number_font_color="#008C9A",domain={'row': 0, 'column': 1}))
        except Exception as e:
            st.error(f"error ind delta size1: {e}")
        try: 
            fig_meter.add_trace(go.Indicator(title='WOB',mode="number+delta",value=round(np.mean(df[Hdr.wob]),1),number={'suffix': ' kbs'},delta={'reference': round(np.mean(df[Hdr.wob]), 1),'relative':True,'valueformat': '.1%'},number_font_color="#008C9A",domain={'row': 0, 'column': 2}))
        except Exception as e:
            st.error(f"error ind delta size2: {e}")
        
        
        if Hdr.rop in def_df.columns:   
            fig_bx.add_trace(go.Box(y=def_df[Hdr.rop],name='Well_1',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color=well_color), row=1, col=1)        
        if Hdr.wob in def_df.columns:
            fig_bx.add_trace(go.Box(y=def_df[Hdr.wob],name='Well_1',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color=well_color), row=1, col=2)
        if Hdr.rpm in def_df.columns:
            fig_bx.add_trace(go.Box(y=def_df[Hdr.rpm],name='Well_1',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color=well_color), row=2, col=1)
        if Hdr.diff in def_df.columns:
            fig_bx.add_trace(go.Box(y=def_df[Hdr.diff],name='Well_1',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color=well_color), row=2, col=2)
        
        if Hdr.rop in def_df.columns:
            fig_bx.add_annotation(x='Well_1',y=def_df[Hdr.rop].mean(),text=f'{round(df[Hdr.rop].mean(),1)}',showarrow=False,font_size=12,font_color='black',xshift=0,yshift=(round(df[Hdr.rop].mean(),1) * 0.1), row=1, col=1)
        if Hdr.wob in def_df.columns:
            fig_bx.add_annotation(x='Well_1',y=def_df[Hdr.wob].mean(),text=f'{round(df[Hdr.wob].mean(),1)}',showarrow=False,font_size=12,font_color='black',xshift=0,yshift=(round(df[Hdr.wob].mean(),1) * 0.1), row=1, col=2)
        if Hdr.rpm in def_df.columns:
            fig_bx.add_annotation(x='Well_1',y=def_df[Hdr.rpm].mean(),text=f'{round(df[Hdr.rpm].mean(),1)}',showarrow=False,font_size=12,font_color='black',xshift=0,yshift=(round(df[Hdr.rpm].mean(),1) * 0.1), row=2, col=1)
        if Hdr.diff in def_df.columns:
            fig_bx.add_annotation(x='Well_1',y=def_df[Hdr.diff].mean(),text=f'{round(df[Hdr.diff].mean(),0)}',showarrow=False,font_size=12,font_color='black',xshift=0,yshift=(round(df[Hdr.diff].mean(),0) * 0.1), row=2, col=2)
        
        def_dfrot = def_df.loc[def_df[Hdr.substate] == 1]
        def_dfslide = def_df.loc[def_df[Hdr.substate] == 2]
        # with st.expander(f'Rotate Data {well_num}'):
        #     st.write(def_dfrot)
        # with st.expander(f'Slide Data {well_num}'):
        #     st.write(def_dfslide)
        
        try:        
            # if [Hdr.hole_depth, Hdr.rop] in def_df.columns:                
            if set([Hdr.hole_depth, Hdr.rop]).issubset(def_df.columns): 
                # Col1 Text Averages
                fig_rotsld.add_annotation(xref="paper",yref="paper",x=0.01, y=1.01, showarrow=False,text=f'<b>Well 1:</b> ')            
                fig_rotsld.add_annotation(xref="paper",yref="paper",x=0.01, y=0.97, showarrow=False,text=f'Sliding ROP: {round(def_dfslide[Hdr.rop].mean(),1)}')
                fig_rotsld.add_annotation(xref="paper",yref="paper",x=0.01, y=0.94, showarrow=False,text=f'Rotating ROP: {round(def_dfrot[Hdr.rop].mean(),1)}')
                fig_rotsld.add_annotation(xref="paper",yref="paper",x=0.01, y=0.91, showarrow=False,text=f'Total ROP: {round(def_df[Hdr.rop].mean(),1)}')
                # Col 2 Pie Breakdown  
                fig_rotsld.add_annotation(xref="paper",yref="paper",x=0.25, y=1.2, showarrow=False,text="Time (%)")
                # fig_rotsld.add_annotation(xref="x domain",yref="y domain",x=0.5, y=1.2, showarrow=False,text='Times (%)', row=1, col=2)
                fig_rotsld.add_trace(go.Pie(values=[len(def_dfrot[Hdr.rop]), len(def_dfslide[Hdr.rop])], labels=['Rotating', 'Sliding'], marker=dict(colors=['black','teal']), name=''), row=1, col=2)
                # Col 3 Box/whisker
                fig_rotsld.add_annotation(xref="x domain",yref="y domain",x=0.5, y=1.2, showarrow=False,text='ROP', row=1, col=3)
                fig_rotsld.add_trace(go.Box(y=def_dfrot[Hdr.rop],name=f'Rotating, {round(def_dfrot[Hdr.rop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color='black'), row=1, col=3)
                fig_rotsld.add_trace(go.Box(y=def_dfslide[Hdr.rop],name=f'Sliding, {round(def_dfslide[Hdr.rop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color='teal'), row=1, col=3)
                fig_rotsld.add_trace(go.Box(y=def_df[Hdr.rop],name=f'Total, {round(def_df[Hdr.rop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color='teal'), row=1, col=3)
                
                # fig_rotsld.add_trace(go.Box(y=def_dfrot[Hdr.rop],name=f'Rotating, {round(def_dfrot[Hdr.rop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color='black'), row=1, col=2)
                # fig_rotsld.add_trace(go.Box(y=def_dfslide[Hdr.rop],name=f'Sliding, {round(def_dfpslide[Hdr.rop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color='teal'), row=1, col=2)
                # fig_rotsld.add_annotation(xref="x domain",yref="y domain",x=0.5, y=1.2, showarrow=False,text="ROTATING", row=1, col=2)
                
                # Col 4 spec chart 
                # fig_rotsld.add_trace(go.Scatter(name="Rotate", x=def_dfrot[Hdr.hole_depth], y=def_dfrot[Hdr.rop], mode='markers', marker= {'color':'black', 'size':0.6}), row=1, col=4)
                # fig_rotsld.add_trace(go.Scatter(name="Slide", x=def_dfslide[Hdr.hole_depth], y=def_dfslide[Hdr.rop], mode='markers', marker= {'color':"teal", 'size':0.6}), row=1, col=4)
                fig_rotsld.add_trace(go.Scatter(name="Rotate", x=def_dfrot[Hdr.hole_depth], y=def_dfrot[Hdr.rop], mode='lines', connectgaps=False, line= {'color':'black', 'width':0.6}), row=1, col=4)
                fig_rotsld.add_trace(go.Scatter(name="Slide", x=def_dfslide[Hdr.hole_depth], y=def_dfslide[Hdr.rop], mode='lines', connectgaps=False, line= {'color':"teal", 'width':0.6}), row=1, col=4)
                # fig_rotsld.add_trace(go.Scatter(name="pRotate", x=def_dfprot[Hdr.hole_depth], y=def_dfprot[Hdr.rop], mode='markers', marker= {'color':'black', 'size':0.6}), row=1, col=3)
                # fig_rotsld.add_trace(go.Scatter(name="pSlide", x=def_dfpslide[Hdr.hole_depth], y=def_dfpslide[Hdr.rop], mode='markers', marker= {'color':"teal", 'size':0.6}), row=1, col=3)
        except Exception as e:
            print(e)
            st.info(f"error rotslide1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")

    # Multiwell comparison mode
    else:             
        # fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name=f"Well_{well_num}", x=def_df['dte'], y=def_df[Hdr.hole_depth], line= {'color':well_color, 'width':0.4}), row=1, col=1)
        # fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df[Hdr.rop].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=2)
        # fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df[Hdr.wob].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=3)
        # fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df[Hdr.rpm].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=4)
        # fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df[Hdr.diff].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=5)
        
        
        if Hdr.hole_depth in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name=f"Well_{well_num}", x=def_df['dte'], y=def_df[Hdr.hole_depth], line= {'color':well_color, 'width':0.4}), row=1, col=1)        
        if Hdr.rop in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df[Hdr.rop].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=2)        
        if Hdr.wob in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df[Hdr.wob].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}, showlegend=False), row=1, col=3)        
        if Hdr.rpm in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df[Hdr.rpm].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=4)        
        if Hdr.diff in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df[Hdr.diff].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=5)        
        if Hdr.trq in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df[Hdr.trq].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=6)        
        if Hdr.flowrate in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df[Hdr.flowrate].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=7)        
        if Hdr.gr in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df[Hdr.gr].rolling(interp).mean(), y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=8)        
        if 'ads-spr' in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df['ads-spr'], y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':2}), row=1, col=9)                   
        if 'ads-sub'in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df['ads-sub'], y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=9)
        if 'MSE' in def_df.columns:
            fig_ch.add_trace(go.Scatter(legendgroup=f'Well_{well_num}', name='', x=def_df['MSE'], y=def_df[Hdr.hole_depth],line={'color':well_color, 'width':0.4}), row=1, col=10)
        
        #  Annotaiton offset for each well trace average            
        if chnl_num is not None:
            plotxpos = 0.18 + (80/chnl_num)
        else:
            plotxpos = 0.2
            
        if well_num is not None:
            xpos = 0.01 + (0.05*well_num)
        else:
            xpos = 0.01        
            
        # fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'{round(def_df[Hdr.rop].mean(),1)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=12))
        # fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'{round(def_df[Hdr.wob].mean(),1)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=12))
        # fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'{round(def_df[Hdr.rpm].mean(),0)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=12))
        # fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'{round(def_df[Hdr.diff].mean(),0)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=12))
        
        try:
            if Hdr.rop in def_df.columns:
                fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'{round(def_df[Hdr.rop].mean(),1)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=16))
            if Hdr.wob in def_df.columns:
                fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'{round(def_df[Hdr.wob].mean(),1)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=16))
            if Hdr.rpm in def_df.columns:
                fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'{round(def_df[Hdr.rpm].mean(),0)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=16))
            if Hdr.diff in def_df.columns:
                fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'{round(def_df[Hdr.diff].mean(),0)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=16))
            if Hdr.trq in def_df.columns:
                fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'{round(def_df[Hdr.trq].mean(),0)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=16))
            if Hdr.flowrate in def_df.columns:
                fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'{round(def_df[Hdr.flowrate].mean(),0)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=16))
            if Hdr.gr in def_df.columns:
                fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'{round(def_df[Hdr.gr].mean(),0)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=16))
            if 'ads-spr' in def_df.columns:
                fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos,y=1.045,text=(f'Super-{well_num}'),textangle=0,showarrow=False,font=dict(color=well_color,size=16))
            if 'ads-sub' in def_df.columns:
                fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos+(xpos/2),y=1.045,text=(f'Sub-{well_num}'),textangle=0,showarrow=False,font=dict(color=well_color,size=16))
            if 'MSE' in def_df.columns:
                fig_ch.add_annotation(xref='paper',yref='paper',xanchor='left',x=(plotxpos*well_num)+xpos+(xpos/2),y=1.045,text=(f'{round(def_df["MSE"].mean(),0)}'),textangle=0,showarrow=False,font=dict(color=well_color,size=16))
        
        except Exception as e:
            print(e)
            st.error(f"error chan annot: {e}  //// {traceback.print_exc()} //// {traceback.format_exc()}")
        
        
        # if ['dte', Hdr.hole_depth] in def_df.columns:      
        if set(['dte', Hdr.hole_depth]).issubset(def_df.columns):
            # fig_dvd.add_trace(go.Scatter(x=df['dte'], y=df[Hdr.hole_depth], line= {'color':'#008C9A'}, name='Well 1', fill='tozerox'))
            fig_dvd.add_trace(go.Scatter(x=def_df['dte'], y=def_df[Hdr.hole_depth], line= {'color':well_color}, name=f'Well_{well_num}', fill='tonextx'))        
            fig_dvdi.add_trace(go.Scatter(x=def_df['dte'], y=def_df[Hdr.hole_depth], line= {'color':well_color}, name=f'Well_{well_num}'), row=1,col=1)
        
        if Hdr.inc in def_df.columns:
            fig_dvdi.add_trace(go.Scatter(x=def_df[Hdr.inc], y=def_df[Hdr.hole_depth], line= {'color':well_color}, line_shape='spline', line_smoothing=1.3, name=f'Well_{well_num}'), row=1,col=2)
        
        if Hdr.rop in def_df.columns:
            fig_meter.add_trace(go.Indicator(title='ROP',mode="number+delta",value=round(def_df[Hdr.rop].mean(),1),number={'suffix': ' ft/hr'},delta={'reference': round(def_ref_df[Hdr.rop].mean()),'relative':True,'valueformat': '.1%'},number_font_color="#8ACFDD" ,domain={'row': 0, 'column': 1}))
        if Hdr.wob in def_df.columns:
            fig_meter.add_trace(go.Indicator(title='WOB',mode="number+delta",value=round(def_df[Hdr.wob].mean(),1),number={'suffix': ' kbs'},delta={'reference': round(def_ref_df[Hdr.wob].mean()),'relative':True,'valueformat': '.1%'},number_font_color="#8ACFDD" ,domain={'row': 0, 'column': 2}))
                            
        if Hdr.rop in def_df.columns:
            fig_bx.add_trace(go.Box(y=def_df[Hdr.rop],name=f'Well_{well_num}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color=well_color,line_color=well_color), row=1, col=1)
        if Hdr.wob in def_df.columns:
            fig_bx.add_trace(go.Box(y=def_df[Hdr.wob],name=f'Well_{well_num}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color=well_color,line_color=well_color), row=1, col=2)
        if Hdr.rpm in def_df.columns:
            fig_bx.add_trace(go.Box(y=def_df[Hdr.rpm],name=f'Well_{well_num}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color=well_color,line_color=well_color), row=2, col=1)
        if Hdr.diff in def_df.columns:
            fig_bx.add_trace(go.Box(y=def_df[Hdr.diff],name=f'Well_{well_num}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color=well_color,line_color=well_color), row=2, col=2)
        
        if Hdr.rop in def_df.columns:
            fig_bx.add_annotation(x=f'Well_{well_num}',y=def_df[Hdr.rop].mean(),text=f'{round(def_df[Hdr.rop].mean(),1)}',showarrow=False,font_size=12,font_color='black',xshift=0,yshift=30, row=1, col=1)
        if Hdr.wob in def_df.columns:
            fig_bx.add_annotation(x=f'Well_{well_num}',y=def_df[Hdr.wob].mean(),text=f'{round(def_df[Hdr.wob].mean(),1)}',showarrow=False,font_size=12,font_color='black',xshift=0,yshift=30, row=1, col=2)
        if Hdr.rpm in def_df.columns:
            fig_bx.add_annotation(x=f'Well_{well_num}',y=def_df[Hdr.rpm].mean(),text=f'{round(def_df[Hdr.rpm].mean(),1)}',showarrow=False,font_size=12,font_color='black',xshift=0,yshift=30, row=2, col=1)
        if Hdr.diff in def_df.columns:
            fig_bx.add_annotation(x=f'Well_{well_num}',y=def_df[Hdr.diff].mean(),text=f'{round(def_df[Hdr.diff].mean(),0)}',showarrow=False,font_size=12,font_color='black',xshift=0,yshift=30, row=2, col=2)
        
        if not len(def_dfsvy) == 0:
            fig_svymap.add_trace(go.Scatter3d(name=f'Well_{well_num}',x=def_dfsvy['3dx'],y=def_dfsvy['3dy'],z=def_dfsvy['3dz'],mode="lines",line= {'color':well_color, 'width':2}))
            # fig_svymap.update_layout(showscale=False)

        # def_dfrot = def_df.loc[def_df['ads-rot'] == 1]
        # def_dfslide = def_df.loc[def_df['ads-slide'] == 1]
        
        def_dfrot = def_df.loc[def_df[Hdr.substate] == 1]
        def_dfslide = def_df.loc[def_df[Hdr.substate] == 2]
        
        # with st.expander(f'Rotate Data {well_num}'):
        #     st.write(def_dfrot)
        # with st.expander(f'Slide Data {well_num}'):
        #     st.write(def_dfslide)
            
        try:
            # if [Hdr.hole_depth, Hdr.rop] in def_df.columns:
            if set([Hdr.hole_depth, Hdr.rop]).issubset(def_df.columns):
                # # fig_rotsld.add_trace(go.Indicator(title='R%',mode="number+delta",value=round(def_dfrot[Hdr.rop].mean(),1),number={'suffix': ' ft/hr'},delta={'reference': round(len(def_dfrot[Hdr.rop])/len(def_df[Hdr.rop]),1),'relative':True,'valueformat': '.1%'},number_font_color="#8ACFDD" ,domain={'row': 0, 'column': 1}), row=well_num, col=2)

                # fig_rotsld.add_trace(go.Box(y=def_dfrot[Hdr.rop],name=f'Avg ROP, {round(def_dfrot[Hdr.rop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color='black'), row=well_num, col=2)
                # fig_rotsld.add_annotation(xref="x domain",yref="y domain",x=0.5, y=1.2, showarrow=False,text="ROTATING", row=well_num, col=1)
                # fig_rotsld.add_trace(go.Box(y=def_dfslide[Hdr.rop],name=f'Avg ROP, {round(def_dfslide[Hdr.rop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color='green'), row=well_num, col=2)
                # fig_rotsld.add_annotation(xref="x domain",yref="y domain",x=0.5, y=1.2, showarrow=False,text="SLIDING", row=well_num, col=2)
                # fig_rotsld.add_trace(go.Scatter(name="Rotate", x=def_dfrot[Hdr.hole_depth], y=def_dfrot[Hdr.rop], mode='markers', marker= {'color':'black', 'size':0.4}), row=well_num, col=3)
                # fig_rotsld.add_trace(go.Scatter(name="Slide", x=def_dfslide[Hdr.hole_depth], y=def_dfslide[Hdr.rop], mode='markers', marker= {'color':"green", 'size':0.4}), row=well_num, col=3)
                # # fig_rotsld.add_trace(go.Scatter(name="Rotate", x=def_dfrot[Hdr.hole_depth], y=def_dfrot[Hdr.rop], mode='markers', marker= {'color':'black', 'width':0.4}), row=well_num, col=3)
                # # fig_rotsld.add_trace(go.Scatter(name="Slide", x=def_dfslide[Hdr.hole_depth], y=def_dfslide[Hdr.rop], mode='markers', marker= {'color':"gray", 'width':0.4}), row=well_num, col=3)    
                
                # Col 1 Averages
                fig_rotsld.add_annotation(xref="paper",yref="paper",x=0.01, y=(1-(well_num/4)+0.25), showarrow=False,text=f'<b>Well {well_num}:</b> {round(def_df[Hdr.rop].mean(),1)}')
                fig_rotsld.add_annotation(xref="paper",yref="paper",x=0.01, y=(1-(well_num/4)+0.22), showarrow=False,text=f'Sliding ROP: {round(def_dfslide[Hdr.rop].mean(),1)}')
                fig_rotsld.add_annotation(xref="paper",yref="paper",x=0.01, y=(1-(well_num/4)+0.18), showarrow=False,text=f'Rotating ROP: {round(def_dfrot[Hdr.rop].mean(),1)}')
                fig_rotsld.add_annotation(xref="paper",yref="paper",x=0.01, y=(1-(well_num/4)+0.15), showarrow=False,text=f'Total ROP: {round(def_df[Hdr.rop].mean(),1)}')
                # Col 2 Pie             
                fig_rotsld.add_trace(go.Pie(values=[len(def_dfrot[Hdr.rop]), len(def_dfslide[Hdr.rop])], labels=['Rotating', 'Sliding'], marker=dict(colors=['black','teal']), name=''), row=well_num, col=2)
                # Col 3 Box/whisker 
                fig_rotsld.add_annotation(xref="x domain",yref="y domain",x=0.5, y=1.2, showarrow=False,text="ROP", row=well_num, col=3)
                fig_rotsld.add_trace(go.Box(y=def_dfrot[Hdr.rop],name=f'Rotating, {round(def_dfrot[Hdr.rop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color='black'), row=well_num, col=3)
                fig_rotsld.add_trace(go.Box(y=def_dfslide[Hdr.rop],name=f'Sliding, {round(def_dfslide[Hdr.rop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color='teal'), row=well_num, col=3)
                fig_rotsld.add_trace(go.Box(y=def_df[Hdr.rop],name=f'Total, {round(def_df[Hdr.rop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color='teal'), row=well_num, col=3)
                # Col 4 Chart 
                # fig_rotsld.add_trace(go.Scatter(name="Rotate", x=def_dfrot[Hdr.hole_depth], y=def_dfrot[Hdr.rop], mode='markers', marker= {'color':'black', 'size':0.6}), row=well_num, col=4)
                # fig_rotsld.add_trace(go.Scatter(name="Slide", x=def_dfslide[Hdr.hole_depth], y=def_dfslide[Hdr.rop], mode='markers', marker= {'color':"teal", 'size':0.6}), row=well_num, col=4)
                fig_rotsld.add_trace(go.Scatter(name="Rotate", x=def_dfrot[Hdr.hole_depth], y=def_dfrot[Hdr.rop], mode='lines', connectgaps=False, line= {'color':'black', 'width':0.6}), row=well_num, col=4)
                fig_rotsld.add_trace(go.Scatter(name="Slide", x=def_dfslide[Hdr.hole_depth], y=def_dfslide[Hdr.rop], mode='lines', connectgaps=False, line= {'color':"teal", 'width':0.6}), row=well_num, col=4)

        except Exception as e:
            print(e)
            st.error(f"error rotslide0: {e}  //// {traceback.print_exc()} //// {traceback.format_exc()}")
        
    
# # Need to implement * FIX* 
# def miracle_cleanup(ref_df, def_qu_lim):
#     # ref_df.columns = ref_df.columns.str.lower() # make lowercase to unify dictionary application next line
#     ref_df.rename(columns=final_column_name_map, inplace=True)
#     # ref_df.columns.str.lower() = set(final_column_name_map) # another method to apply column rename dictionary
#     num = ref_df._get_numeric_data()
#     # Null negative values
#     num[num < 0] = np.nan
#     ref_df.reset_index(inplace=True)
#     # count = 1
#     # with stcol0.spinner('Finding time domain...'):
#     try:
#         # df[Hdr.dt] = df['YYYY/MM/DD']+df['HH:MM:SS']
#         ref_df[Hdr.dt] = pd.to_datetime(ref_df['YYYY/MM/DD'] + ' ' + ref_df['HH:MM:SS'])
#         # avg_rop1 = np.mean(ref_df[Hdr.rop])
#         # avg_wob1 = np.mean(ref_df[Hdr.wob])
                
#         try: 
#             dfst = ref_df[Hdr.dt].values[1]
#             # df.loc[df[Hdr.rop]!='0','A'].index[0]
#             ref_df['dtst'] = (dfst)
#             ref_df['dte'] = ((ref_df[Hdr.dt] - ref_df['dtst']) / np.timedelta64(1, "h"))
#         except Exception as e:
#             print(e)
#             st.error(f"error calculating elapsed time0: {e}")
        
#         # Check for rigstates                
#         if not set([Hdr.sprstate,Hdr.substate]).issubset(df.columns):
#             # st.warning(f"rigstate present") 
#         # else:
#             st.warning(f"rigstate not present")
#             with st.expander ('W1 pre-miracle rigstate'):
#                 st.write(ref_df)
#             df = miracle_rigstate(ref_df)
#             with st.expander ('W1 post-miracle rigstate'):
#                 st.write(ref_df)
        
#         # Check for kilo ft_lb torque, correct
#         ref_df[Hdr.trq] = np.where(np.mean(ref_df[Hdr.trq]) < 1, ref_df[Hdr.trq]*1000, ref_df[Hdr.trq])
        
#     except Exception as e:
#         print(e)
#         st.error(f"error cleanup counter: {e}")
    
#     return ref_df


def miracle_rigstate(ref_df):
    
    # Pason super States
    # 0Uknown
    # 1Drilling
    # 4Transition
    # 5Out of Hole
    # 6Pre-Spud
    # 7Tripping
    
    # Pason Sub States
    # 0Unknown 
    # 1Rotary drilling 
    # 2Sliding 
    # 3Connection 
    # 4Static 
    # 5Circulating 
    # 6Washing 
    # 7Reaming 
    # 10Pulling out of hole (POOH)
    # 11Running in hole (RIH)

    # Rig State logic chain
    #  change in hole Depth, if df[Hdr.hole_depth] - df[Hdr.hole_depth].shift(1) = 0 then deltaHD = 0
    #  change in bit depth, if df[Hdr.bit_depth] - df[Hdr.bit_depth].shift(1) = 0 then deltaBD = 0
    #  calc off-bottom value, df[Hdr.hole_depth] - df[Hdr.bit_depth] = off-bottom value
    #  on-bottom check, if off-bottom value is <=0.03 then true, else false
    
    # Superstates
    # Drilling, If on bottom is true and (rpm > 0 or flow rate > 0)
    
    # substates
    # State1 - Rotating, if On-Bottom is true and top-drive RPM > 30, then rotating true
    # State2 - Sliding, if On-Bottom is true and top-drive RPM < 30, then sliding true
    # State3 - in Slips, if on-bottom is false and HookLoad is < 50 and Flow Rate < 100 then slips is true
    # State - Trip In, if off-bottom is > 100 and delta bit depth is > 0 and hookload >= 50 then tripping in
    # State - Trip Out, if off-bottom is > 100 and delta bit depth is < 0 and hookload >= 50 then tripping in
    # State5  - Circulating, if on-bottom is false and flow rate > 100 and delta bit depth = 0 then circulating
    
    # If states are not present
    if not set([Hdr.sprstate, Hdr.substate]).issubset(ref_df.columns):
        ref_df[Hdr.sprstate] = 80
        ref_df[Hdr.substate] = 80
        
        if set([Hdr.hole_depth, Hdr.bit_depth]).issubset(ref_df.columns):
            # Running change in hole depth vs prior line
            ref_df['ads-dHD'] = (ref_df[Hdr.hole_depth] - ref_df[Hdr.hole_depth].shift(1))
            #  running change in bit depth vs prior line
            ref_df['ads-dBD'] = (ref_df[Hdr.bit_depth] - ref_df[Hdr.bit_depth].shift(1))
            # comparison of delta between hole depth and bit depth
            ref_df['ads-offbtm'] = (ref_df[Hdr.hole_depth] - ref_df[Hdr.bit_depth])
            # if depth comps are negligible, on bottom is true
            ref_df['ads-onbtm'] = np.where((ref_df['ads-offbtm'] <= 0.03), 1,0)
        
            # determine rotating drilling if on bottom and RPM over 30
            if set([Hdr.rpm]).issubset(ref_df.columns):
                ref_df['ads-rot'] = np.where(((ref_df['ads-onbtm']==1) & (ref_df[Hdr.rpm]>=30)), 1,0) 
                # determine sliding drilling if on bottom and RPM under 30
                ref_df['ads-slide'] = np.where(((ref_df['ads-onbtm']==1) & (ref_df[Hdr.rpm]<30)), 1,0) 
            else:
                ref_df['ads-rot'] = 0
                ref_df['ads-slide'] = 0
        
            if set([Hdr.wob, Hdr.flowrate]).issubset(ref_df.columns):
                #  Req hookload to calc these sub states, used WOB for temporary solution *FIX*
                ref_df['ads-slips'] = np.where(((ref_df['ads-onbtm']==0) & (ref_df[Hdr.wob]<50) & (ref_df[Hdr.flowrate]<100)), 1,0) 
                ref_df['ads-trpin'] = np.where(((ref_df['ads-offbtm']>100) & (ref_df['ads-dBD']>0) & (ref_df[Hdr.wob]>=50)), 1,0)
                ref_df['ads-trpout'] = np.where(((ref_df['ads-offbtm']>100) & (ref_df['ads-dBD']<0) & (ref_df[Hdr.wob]>=50)), 1,0)
                ref_df['ads-circ'] = np.where(((ref_df['ads-onbtm']==0) & (ref_df[Hdr.flowrate]>100) & (ref_df['ads-dBD']==0)), 1,0)
            else:
                ref_df['ads-slips'] = 0
                ref_df['ads-trpin'] = 0
                ref_df['ads-trpout'] = 0
                ref_df['ads-circ'] = 0
        
            # Roll up into Pason states
            ref_df['ads-spr'] = np.where((ref_df['ads-rot'] == 1) | (ref_df['ads-slide'] == 1), 1, 0)
            ref_df['ads-sub'] = np.where((ref_df['ads-rot'] == 1), 1, np.where((ref_df['ads-slide'] == 1), 2, np.where((ref_df['ads-circ'] == 1), 5, 0)))
        
    else:
        # do some shwifty
        if set([Hdr.hole_depth, Hdr.bit_depth]).issubset(ref_df.columns):
            ref_df['ads-dHD'] = (ref_df[Hdr.hole_depth] - ref_df[Hdr.hole_depth].shift(1))
            ref_df['ads-dBD'] = (ref_df[Hdr.bit_depth] - ref_df[Hdr.bit_depth].shift(1))
            ref_df['ads-offbtm'] = (ref_df[Hdr.hole_depth] - ref_df[Hdr.bit_depth])
            ref_df['ads-onbtm'] = np.where((ref_df['ads-offbtm'] > 0.03), 1,0)
        
            if set([Hdr.rpm]).issubset(ref_df.columns):
                ref_df['ads-rot'] = np.where(((ref_df['ads-onbtm']==1) & (ref_df[Hdr.rpm]>=30)), 1,0) 
                ref_df['ads-slide'] = np.where(((ref_df['ads-onbtm']==1) & (ref_df[Hdr.rpm]<30)), 1,0) 
            else:
                ref_df['ads-rot'] = 0
                ref_df['ads-slide'] = 0
                
            if set([Hdr.wob, Hdr.flowrate]).issubset(ref_df.columns):
                #  Req hookload
                ref_df['ads-slips'] = np.where(((ref_df['ads-onbtm']==0) & (ref_df[Hdr.wob]<50) & (ref_df[Hdr.flowrate]<100)), 1,0) 
                ref_df['ads-trpin'] = np.where(((ref_df['ads-offbtm']>100) & (ref_df['ads-dBD']>0) & (ref_df[Hdr.wob]>=50)), 1,0)
                ref_df['ads-trpout'] = np.where(((ref_df['ads-offbtm']>100) & (ref_df['ads-dBD']<0) & (ref_df[Hdr.wob]>=50)), 1,0)
                ref_df['ads-circ'] = np.where(((ref_df['ads-onbtm']==0) & (ref_df[Hdr.flowrate]>100) & (ref_df['ads-dBD']==0)), 1,0) 
            else:
                ref_df['ads-slips'] = 0
                ref_df['ads-trpin'] = 0
                ref_df['ads-trpout'] = 0
                ref_df['ads-circ'] = 0
        
        # Roll up into Pason states
        ref_df['ads-spr'] = ref_df[Hdr.sprstate]
        ref_df['ads-sub'] = ref_df[Hdr.substate]
         
    return ref_df 
    
def miracle_survey(ref_df):
    # https://github.com/andymcdgeo/Andys_YouTube_Notebooks/blob/main/13%20-%20Welly%20Location.ipynb
    # https://www.youtube.com/watch?v=mcCa5oxN1Ys
    
    if set([Hdr.hole_depth, 'ads-spr', Hdr.inc, Hdr.az, Hdr.rop, Hdr.wob]).issubset(ref_df.columns):
        try:
            dfsvy = ref_df[[Hdr.hole_depth, 'ads-spr', Hdr.inc, Hdr.az, Hdr.rop, Hdr.wob]].copy()
            dfsvy['theta'] = 0
            dfsvy['thetaRF'] = 0
            # dfsvy['tvd'] = dfsvy[Hdr.hole_depth].values[0]
            dfsvy['tvd'] = 0
            # dfsvy['tvd'].iloc[0] = dfsvy[Hdr.hole_depth].iloc[0]
            dfsvy['3dx2'] = 0
            dfsvy['3dy2'] = 0
            
            dfsvy['3dx'] = 0
            dfsvy['3dy'] = 0
            dfsvy['3dz'] = dfsvy[Hdr.hole_depth] 
            # .loc[Hdr.sprstate == 1]
            dfsvy['svy'] = 0                    
            # df = df[df[Hdr.sprstate] == 1]
            # dfsvy.iloc[df[Hdr.sprstate == 1].argmax()] = 1
            # dfsvy['svy'] = np.where(((dfsvy[Hdr.inc] != dfsvy[Hdr.inc].shift(1)) | (dfsvy[Hdr.az] != dfsvy[Hdr.az].shift(1))), 1, 0)    
            
            dfsvy['svy'] = np.where(((dfsvy[Hdr.inc]==-999.25) | (dfsvy[Hdr.az]==-999.25)), 0, np.where(((dfsvy[Hdr.inc] != dfsvy[Hdr.inc].shift(1)) | (dfsvy[Hdr.az] != dfsvy[Hdr.az].shift(1))), np.where(((dfsvy[Hdr.inc] == dfsvy[Hdr.inc].shift(-1)) & (dfsvy[Hdr.az] == dfsvy[Hdr.az].shift(-1))), 1, 0) ,0))
            # dfsvy['svy'] = np.where(((dfsvy[Hdr.inc]==-999.25) | (dfsvy[Hdr.az]==-999.25)), 0, np.where(((dfsvy[Hdr.hole_depth]-dfsvy[Hdr.hole_depth].query('svy == svy.max()')) < 85), 0, np.where(((dfsvy[Hdr.inc] != dfsvy[Hdr.inc].shift(1)) | (dfsvy[Hdr.az] != dfsvy[Hdr.az].shift(1))), np.where(((dfsvy[Hdr.inc] == dfsvy[Hdr.inc].shift(-1)) & (dfsvy[Hdr.az] == dfsvy[Hdr.az].shift(-1))), dfsvy['sysvy'].cumcount+1, 0) ,0)))
            
            dfsvy = dfsvy[dfsvy['svy'] >= 1].copy()              
            
            try:
                dfsvy['svy'] = dfsvy.groupby(['svy']).cumcount() + 1
            except Exception as e:
                print(e)
                st.error(f"error svy counter: {e}")
                                        
            dfsvy = dfsvy.reset_index(drop=True)
            
    
        # df['sysvy'].iloc[df[Hdr.sprstate == 1].argmax()] = 1
        # =IF(OR(B3=-999.25,C3=-999.25),0,IF(A3-INDEX(A$2:A2,MATCH(MAX(D$2:D2),D$2:D2,0))<85,0,IF(OR(B3<>B2,C3<>C2),IF(AND(B3=B4,C3=C4),MAX(D$2:D2)+1,0),0)))
        # df['sysvy'] = np.where((df[Hdr.inc]=-999.25 | df[Hdr.az]=-999.25), 0, np.where((df[Hdr.hole_depth]-df[Hdr.hole_depth].query('sysvy == sysvy.max()')) < 85 ,0, np.where((df[Hdr.inc] <> df[Hdr.inc].shift(1) | df[Hdr.az]<>df[Hdr.az].shift(1)), np.where((df[Hdr.inc]=df[Hdr.inc].shift(-1) & df[Hdr.az]=df[Hdr.az].shift(-1)), df['sysvy'].cumcount+1, 0) ,0)))
        # df['sysvy'] = np.where((df[Hdr.inc]==-999.25 | df[Hdr.az]==-999.25), 0, np.where((df[Hdr.hole_depth]-df[Hdr.hole_depth].query('sysvy == sysvy.max()')) < 85 ,0, np.where((df[Hdr.inc] != df[Hdr.inc].shift(1) | df[Hdr.az] != df[Hdr.az].shift(1)), np.where((df[Hdr.inc] == df[Hdr.inc].shift(-1) & df[Hdr.az] == df[Hdr.az].shift(-1)), df['sysvy'].cumcount+1, 0) ,0)))
        # df['sysvy'] = np.where((df[Hdr.inc] != df[Hdr.inc].shift(1) | df[Hdr.az] != df[Hdr.az].shift(1)), np.where((df[Hdr.inc] == df[Hdr.inc].shift(-1) & df[Hdr.az] == df[Hdr.az].shift(-1)), 1, 0) ,0)
        
    
            try:                                                
                dfsvy['theta'] = np.arccos((np.cos(np.radians(dfsvy[Hdr.inc] - dfsvy[Hdr.inc].shift(1))) - np.sin(np.radians(dfsvy[Hdr.inc])) * np.sin(np.radians(dfsvy[Hdr.inc].shift(1))) * (1 - np.cos(np.radians(dfsvy[Hdr.az] - dfsvy[Hdr.az].shift(1))))))                    
                dfsvy['thetaRF'] = np.where(dfsvy['theta'] == 0, 1, 2 / dfsvy['theta'] * np.tan(dfsvy['theta'] / 2))
                # =IF(ISBLANK($B8),"",      IF(     AND($C8=$C7,$D8=$D7)   ,  ($B8-$B7)*COS(RADIANS($C8)),                                                                                                                                                                          ($B8-$B7)/                                                  (2*(COS(RADIANS($C7)) + COS(RADIANS($C8)))*$L8))  )
                dfsvy['tvd'] = np.where((dfsvy[Hdr.inc] == dfsvy[Hdr.inc].shift(1)) & (dfsvy[Hdr.az] == dfsvy[Hdr.az].shift(1)), ((dfsvy[Hdr.hole_depth] - dfsvy[Hdr.hole_depth].shift(1)) * (np.cos(np.radians(dfsvy[Hdr.inc])))), (((dfsvy[Hdr.hole_depth] - dfsvy[Hdr.hole_depth].shift(1)) / 2) * (np.cos(np.radians(dfsvy[Hdr.inc].shift(1))) + np.cos(np.radians(dfsvy[Hdr.inc])))) * dfsvy['thetaRF'])
                dfsvy['tvd'].iloc[0] = dfsvy[Hdr.hole_depth].iloc[0]
                dfsvy['3dz'] = dfsvy['tvd'].cumsum()
                                        
                dfsvy['3dx2'] = np.where((dfsvy[Hdr.inc] == dfsvy[Hdr.inc].shift(1)) & (dfsvy[Hdr.az] == dfsvy[Hdr.az].shift(1)), ((dfsvy[Hdr.hole_depth] - dfsvy[Hdr.hole_depth].shift(1)) * (np.sin(np.radians(dfsvy[Hdr.az])) * np.sin(np.radians(dfsvy[Hdr.inc])))), ((dfsvy[Hdr.hole_depth] - dfsvy[Hdr.hole_depth].shift(1)) / 2) * (np.sin(np.radians(dfsvy[Hdr.inc].shift(1))) * np.sin(np.radians(dfsvy[Hdr.az].shift(1))) + np.sin(np.radians(dfsvy[Hdr.inc])) * np.sin(np.radians(dfsvy[Hdr.az]))) * dfsvy['thetaRF'])
                dfsvy['3dy2'] = np.where((dfsvy[Hdr.inc] == dfsvy[Hdr.inc].shift(1)) & (dfsvy[Hdr.az] == dfsvy[Hdr.az].shift(1)), ((dfsvy[Hdr.hole_depth] - dfsvy[Hdr.hole_depth].shift(1)) * (np.cos(np.radians(dfsvy[Hdr.az])) * np.sin(np.radians(dfsvy[Hdr.inc])))), ((dfsvy[Hdr.hole_depth] - dfsvy[Hdr.hole_depth].shift(1)) / 2) * (np.sin(np.radians(dfsvy[Hdr.inc].shift(1))) * np.cos(np.radians(dfsvy[Hdr.az].shift(1))) + np.sin(np.radians(dfsvy[Hdr.inc])) * np.cos(np.radians(dfsvy[Hdr.az]))) * dfsvy['thetaRF'])
                                            
                dfsvy['NSe'] = dfsvy[Hdr.hole_depth].shift(1)
                
                try:
                    dfsvy['3dx'] = dfsvy['3dx2'].cumsum() 
                    dfsvy['3dy'] = dfsvy['3dy2'].cumsum() 
                except Exception as e:
                    print(e)
                    st.error(f"error adding deltaXY: {e}")
            
            except Exception as e:
                print(e)
                st.error(f"error calculating cartesianRF from az/inc: {e}")
                
        except Exception as e:
            print(e)
            st.error(f"error svy calc: {e}")
        # st.write(dfsvy)  
    else:        
        try:
            dfsvy = ref_df[[Hdr.hole_depth]].copy()
            dfsvy[Hdr.az] = 0
            dfsvy[Hdr.inc] = 0
            
            dfsvy['3dz'] = 0
            dfsvy['3dx'] = 0 
            dfsvy['3dy'] = 0 
        except Exception as e:
            print(e)
            st.error(f"error dfsvy no az/inc: {e}")
    
    return dfsvy


@st.cache_data
def get_chart_47479(def_df, use_container_width: bool):
    import altair as alt
    # from vega_datasets import data

    source = def_df
    # data.sp500.url

    brush = alt.selection(type='interval', encodings=['x'])

    base = alt.Chart(source).mark_area().encode(x = 'holedepth:N',y = 'dte:T').properties(width=600,height=200)

    upper = base.encode(alt.X('holedepth:N', scale=alt.Scale(domain=brush)))

    lower = base.properties(height=60).add_selection(brush)

    chart = upper & lower

    # tab1, tab2 = st.tabs(["Streamlit theme (default)", "Altair native theme"])
    # with tab1:
    st.altair_chart(chart, use_container_width=True)

def edrleg():
    
    # BDIA Bit Diameter:	12.25” or maybe input from elsewhere
    # HKLDT Hookload Threshold:	80 klbs
    # MTC Motor Torque Constant:	9.0 lb-ft/psi
    # RPG Revs per Gallon Constant:	0.33
    # RPMT String RPM Threshold:	15

    global df
    global dfsvy
    
    # Standard Color Palette
    # class upal(str, Enum):
    #     teal1 = '#05929f'
    #     teal2 = '#3F7077'
    #     gray1 = '#CBC9C9'
    #     gray2 = '#9B9DA0'
    #     gray3 = '#7A7D81'
    #     green1 = '#98C21F'
    #     orig1 = '#008080'
    #     orig2 = '#40E0D0'
    #     orig3 = '#32CD32'
    #     orig4 = '#5D3FD3'
    #     orig5 = '#121212' 
    color0 = '#121212'
    color1 = '#008C9A'
    color3 = '#8ACFDD'
    color2 = '#7A7D81'
    color4 = '#98C007'
    gbg = '#FBFBFB'
    ggrid = '#EEEEF0'
    
    # Plotly Resampler automatic mode
    # register_plotly_resampler(mode='auto')
    
    # if not edr_file:
    #     st.warning('No file has been uploaded')
    
    # else:        
    count = 0
    pip.main(['install', 'openpyxl'])

    dbmsg = st.sidebar.empty()
    
    # Initialize monogodb connection.
    # Uses st.cache_resource to only run once.
    @st.cache_resource()
    def init_connection():
        return pymongo.MongoClient("mongodb+srv://" + st.secrets['username'] + ":" + st.secrets['password'] + "@cluster0.bteptlv.mongodb.net/?retryWrites=true&w=majority")
    
    client = init_connection()
    db = client.edrtest1
    # Pull data from the collection.
    # Uses st.cache_data to only rerun when the query changes or after 10 min.
    st.cache_data(ttl=600)
    def get_data():
        db = client.mydb
        items = db.mycollection.find()
        items = list(items)  # make hashable for st.cache_data
        return items

    items = get_data()
    # # Print results.
    for item in items:
        dbmsg.success(f"Status :{item['pet']}:")
    
    # pet: "white_check_mark"
    
    
    # define tab names
    tab_titles = [
        'Overview',
        'Channels',
        'Days vs. Depth',
        'Survey (3D)',
        'Sliding',
        'Toolface',
        'Analysis',
        'Data Quality',
        'Reports'
    ]
        
    # Set metric columns
    hed1, hed2, hed3, hed4, hed5, hed6 = st.columns(6)
    
    # Enable tabs
    tabs = st.tabs(tab_titles)
    
    with tabs[0]:
        ovcol0, ovcol1 = st.columns(2)        
        
    with tabs[1]:        
        # Set metric columns
        w0, well1, well2, well3, well4, w5 = st.columns(6)
        # set plot container
        # stcol0, stcol1, stcol2 = st.columns([0,8,1])
        stcol1, stcol2 = st.columns([8,1])

    hide_st_style = """
                    <style>
                    footer {visibility:hidden;}
                    </style>
                    """
    
    st.markdown(hide_st_style, unsafe_allow_html=True)


    # 7.1.13 sidebar content display
    st.sidebar.subheader('Import file here and Select data')

    # ##uploaded_lasfile = st.sidebar.file_uploader(label='Upload las file here', type= ['las', 'LAS'])
    # uploaded_file = st.sidebar.file_uploader(label='Upload CSV or Excel File here', type= ['csv', 'xlsx'])
    # uploaded_file2 = None
    # uploaded_file3 = None
    # uploaded_file4 = None
    # try:
    #     uploaded_file2 = st.sidebar.file_uploader(label='Upload second CSV or Excel File here', type= ['csv', 'xlsx'])
    #     try:
    #         uploaded_file3 = st.sidebar.file_uploader(label='Upload third CSV or Excel File here', type= ['csv', 'xlsx'])
    #         try:
    #             uploaded_file4 = st.sidebar.file_uploader(label='Upload fourth CSV or Excel File here', type= ['csv', 'xlsx'])
    #         except Exception as e:
    #             print(e)
    #     except Exception as e:
    #         print(e)
    # except Exception as e:
    #     print(e)
    
    
    uploaded_file = None
    uploaded_file2 = None
    uploaded_file3 = None
    uploaded_file4 = None
    uploaded_fs = st.sidebar.file_uploader(label='Upload CSV or Excel File here', type= ['csv', 'xlsx'], accept_multiple_files=True)   
    
    clearall = st.sidebar.button("Clear All")
    if clearall or uploaded_fs is None:
        # Delete all the items in Session state
        for key in st.session_state.keys():
            del st.session_state[key]        
        uploaded_file = None
        uploaded_file2 = None
        uploaded_file3 = None
        uploaded_file4 = None
        count = 0
        uploaded_fs = None
        try:
            if df:
                df.empty
                df = None
                del df
            if df2:
                df2.empty
                df2 = None
                del df2
            if df3:
                df3.empty
                df3 = None
                del df3
            if df4:
                df4.empty
                df4 = None  
                del df4          
            gc.collect()
        except Exception as e:
            print(e)
            
        # Clear values from *all* all in-memory and on-disk data caches:
        # i.e. clear values from both square and cube
        st.cache_data.clear() 
        st.cache_resource.clear()
    
        
    well1name = None 
    well2name = None 
    well3name = None 
    well4name = None
    try:
        if uploaded_fs[0]:
            uploaded_file = uploaded_fs[0]
            well1name = uploaded_fs[0].name
            # st.write(uploaded_fs[0].name)
    except Exception as e:
        print(e)  
        st.info(f"**Upload file** to start, in the left sidebar.")
    
    try:
        if uploaded_fs[1]:
            uploaded_file2 = uploaded_fs[1]
            well2name = uploaded_fs[1].name
            # st.write(uploaded_fs[1].name)
    except Exception as e:
        print(e)
    
    try:
        if uploaded_fs[2]:
            uploaded_file3 = uploaded_fs[2] 
            well3name = uploaded_fs[2].name  
            # st.write(uploaded_fs[2].name) 
    except Exception as e:
        print(e)  
        
    try:
        if uploaded_fs[3]:
            uploaded_file4 = uploaded_fs[3] 
            well4name = uploaded_fs[3].name
            # st.write(uploaded_fs[3].name)  
    except Exception as e:
        print(e)  

    # @st.cache_data
    def load_csv_df(file_in):
        # #  Test code to try to parse out metadata header rows
        # # https://stackoverflow.com/questions/60434664/automatically-determine-header-row-when-reading-csv-in-pandas                    
        # https://stackoverflow.com/questions/18039057/python-pandas-error-tokenizing-data                      
        # # https://www.reddit.com/r/CodingHelp/comments/u2lwnn/how_to_read_the_lines_from_a_streamlit_file/                    
        # # https://stackoverflow.com/questions/51530785/pandas-read-csv-skiprows-determine-rows-to-skip                                            
        
        try:
            df_head = pd.read_csv(file_in, encoding_errors='ignore', error_bad_lines=False, header=None)
            head_len = len(df_head)                                               
            # st.write(df_head)
            # st.write(f'head1 line count: {len(df_head)}')                                                
        except Exception as e:
            print(e)  
            st.info(f"error df head: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                     
        # skiprows = csvcount -1
        # st.write(f'Skiprow count: {skiprows}')
        if head_len >= 100:
            # majority of file comes thru, header likely at top
            try:
                file_in.seek(0)
                loader_df = pd.read_csv(file_in, encoding_errors='ignore')
                # df = pd.read_csv(uploaded_file, index=False)                            
                # https://sparkbyexamples.com/pandas/pandas-remove-duplicate-columns-from-dataframe/
                # df.T.drop_duplicates().T
            except Exception as e:
                print(e)  
                # st.info(f"error df load: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                                
                try:
                    file_in.seek(0)
                    loader_df = pd.read_csv(file_in, encoding_errors='ignore', error_bad_lines=False)
                except Exception as e:
                    print(e)  
                    st.info(f"error df errorload: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                    
        else: 
            # Majority of file does not come thru, header likely under some extra meta lines
            try:                               
                # df = pd.read_csv(f, delimiter=",")
                # df.index += 1  
                file_in.seek(0)                      
                loader_df = pd.read_csv(file_in, skiprows=head_len, encoding_errors='ignore', error_bad_lines=False)
                with st.expander('No Header Data'):
                    st.write(loader_df)
            except Exception as e:
                print(e)  
                st.info(f"error df skiprow: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")   
        
        return loader_df
    
    
    # @st.cache_data
    def load_xls_df(file_in):
        try:               
            loader_df2 = pd.read_excel(file_in, engine= 'openpyxl')                                                  
        except Exception as e:
            print(e)  
            loader_df2 = None
            st.info(f"error df head: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                     
        
        return loader_df2
        
    # create a list of dataframes from files uploaded by the user.
    # global df
    # with st.spinner('Constructing Dataframe...'):
    dfs_to_compare = []
    uploaded_files = [uploaded_file, uploaded_file2, uploaded_file3, uploaded_file4]
        
    # handling for session state df https://discuss.streamlit.io/t/how-to-use-session-state-to-save-file-uploads-and-filters/36443
    # if 'ssdf' not in st.session_state:
    #     st.session_state.ssdf = pd.DataFrame()
    # if 'ssdf2' not in st.session_state:
    #     st.session_state.ssdf2 = pd.DataFrame()
    # if 'ssdf3' not in st.session_state:
    #     st.session_state.ssdf3 = pd.DataFrame()
    # if 'ssdf4' not in st.session_state:
    #     st.session_state.ssdf4 = pd.DataFrame()
    
    try:
        for user_upload in uploaded_files:
        # attempt to read uploaded file into pandas
            
                
            # df = pd.DataFrame
            if uploaded_file is not None:
                collectionname = uploaded_file.name 
                
                # try:                    
                #     with st.expander ('write check'):
                #         st.write(uploaded_file)  
                #         filedec = st.write(chardet.detect(uploaded_file.read()))
                # except Exception as e:
                #     print(e)  
                #     st.error(f"error writetest: {e}")  
                    
                if uploaded_file.name.lower().endswith(('.csv')):                    
                    # #  Test code to try to parse out metadata header rows
                    # # https://stackoverflow.com/questions/60434664/automatically-determine-header-row-when-reading-csv-in-pandas                    
                    # https://stackoverflow.com/questions/18039057/python-pandas-error-tokenizing-data                      
                    # # https://www.reddit.com/r/CodingHelp/comments/u2lwnn/how_to_read_the_lines_from_a_streamlit_file/                    
                    # # https://stackoverflow.com/questions/51530785/pandas-read-csv-skiprows-determine-rows-to-skip                                            
                    
                    # sessions state for df https://discuss.streamlit.io/t/apps-keeps-resetting-and-loosing-imported-files-need-to-start-again/13837/10
                    df = load_csv_df(uploaded_file)
                    # st.sessions_state.ssdf = df
                    # st.write(st.sessions_state.ssdf)
                    
                    # df = load_csv_df(uploaded_file)
                    
                    # try:
                    #     df_head = pd.read_csv(uploaded_file, encoding_errors='ignore', error_bad_lines=False, header=None)
                    #     head1_len = len(df_head)                                               
                    #     # st.write(df_head)
                    #     # st.write(f'head1 line count: {len(df_head)}')                                                
                    # except Exception as e:
                    #     print(e)  
                    #     st.info(f"error df head: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                     
                    # # skiprows = csvcount -1
                    # # st.write(f'Skiprow count: {skiprows}')
                    # if head1_len >= 100:
                    #     # majority of file comes thru, header likely at top
                    #     try:
                    #         uploaded_file.seek(0)
                    #         df = pd.read_csv(uploaded_file, encoding_errors='ignore')
                    #         # df = pd.read_csv(uploaded_file, index=False)                            
                    #         # https://sparkbyexamples.com/pandas/pandas-remove-duplicate-columns-from-dataframe/
                    #         # df.T.drop_duplicates().T
                    #     except Exception as e:
                    #         print(e)  
                    #         st.info(f"error df load: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                                
                    #         try:
                    #             uploaded_file.seek(0)
                    #             df = pd.read_csv(uploaded_file, encoding_errors='ignore', error_bad_lines=False)
                    #         except Exception as e:
                    #             print(e)  
                    #             st.info(f"error df errorload: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                    
                    # else: 
                    #     # Majority of file does not come thru, header likely under some extra meta lines
                    #     try:                               
                    #         # df = pd.read_csv(f, delimiter=",")
                    #         # df.index += 1  
                    #         uploaded_file.seek(0)                      
                    #         df = pd.read_csv(uploaded_file, skiprows=head1_len, encoding_errors='ignore', error_bad_lines=False)
                    #         with st.expander('No Header Data'):
                    #             st.write(df)
                    #     except Exception as e:
                    #         print(e)  
                    #         st.info(f"error df skiprow: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                
                #  https://sparkbyexamples.com/pandas/pandas-remove-duplicate-columns-from-dataframe/                     
                if uploaded_file.name.lower().endswith(('.xls','.xlsx')):
                    df = load_xls_df(uploaded_file)                    
                    # try:
                    #     df = pd.read_excel(uploaded_file, engine= 'openpyxl')
                    # except Exception as e:
                    #     print(e)
                    #     st.error(f"error df1 excel: {e}")
            # else:
            #     df = pd.DataFrame(None)         
                                    
            if uploaded_file2 is not None:
                if uploaded_file2.name.lower().endswith(('.csv')):
                    
                    df2 = load_csv_df(uploaded_file2)
                    # try:               
                    #     df2_head = pd.read_csv(uploaded_file2, encoding_errors='ignore', error_bad_lines=False, header=None)
                    #     head2_len = len(df2_head)                                                 
                    # except Exception as e:
                    #     print(e)  
                    #     st.info(f"error df2 head: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                    # if head2_len >= 20:
                    #     # majority of file comes thru, header likely at top
                    #     try:
                    #         uploaded_file2.seek(0)
                    #         df2 = pd.read_csv(uploaded_file2, encoding_errors='ignore')
                    #     except Exception as e:
                    #         print(e)  
                    #         st.info(f"error df2 load: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")   
                                                
                    #         try:
                    #             uploaded_file2.seek(0)
                    #             df2 = pd.read_csv(uploaded_file2, encoding_errors='ignore', error_bad_lines=False)
                    #         except Exception as e:
                    #             print(e)  
                    #             st.info(f"error df2 errorload: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                    
                    # else: 
                    #     # Majority of file does not come thru, header likely under some extra meta lines
                    #     try:
                    #         uploaded_file2.seek(0)                      
                    #         df2 = pd.read_csv(uploaded_file2, skiprows=head2_len, encoding_errors='ignore', error_bad_lines=False)
                    #     except Exception as e:
                    #         print(e)  
                    #         st.info(f"error df2 skiprow: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                
                                     
                if uploaded_file2.name.lower().endswith(('.xls','.xlsx')):
                    df2 = load_xls_df(uploaded_file2)                    
                    # try:
                    #     df2 = pd.read_excel(uploaded_file2, engine= 'openpyxl')
                    # except Exception as e:
                    #     print(e)
                    #     st.error(f"error df2 excel: {e}")
                     
            # else:
            #     df2 = pd.DataFrame(None)        
                                    
            if uploaded_file3 is not None:
                if uploaded_file3.name.lower().endswith(('.csv')):
                    
                    df3 = load_csv_df(uploaded_file3)
                    # try:               
                    #     df3_head = pd.read_csv(uploaded_file3, encoding_errors='ignore', error_bad_lines=False, header=None)
                    #     head3_len = len(df3_head)                                                 
                    # except Exception as e:
                    #     print(e)  
                    #     st.info(f"error df3 head: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                    # if head3_len >= 30:
                    #     # majority of file comes thru, header likely at top
                    #     try:
                    #         uploaded_file3.seek(0)
                    #         df3 = pd.read_csv(uploaded_file3, encoding_errors='ignore')
                    #     except Exception as e:
                    #         print(e)  
                    #         st.info(f"error df3 load: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")   
                                                
                    #         try:
                    #             uploaded_file3.seek(0)
                    #             df3 = pd.read_csv(uploaded_file3, encoding_errors='ignore', error_bad_lines=False)
                    #         except Exception as e:
                    #             print(e)  
                    #             st.info(f"error df3 errorload: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                    
                    # else: 
                    #     # Majority of file does not come thru, header likely under some extra meta lines
                    #     try:
                    #         uploaded_file3.seek(0)                      
                    #         df = pd.read_csv(uploaded_file3, skiprows=head3_len, encoding_errors='ignore', error_bad_lines=False)
                    #     except Exception as e:
                    #         print(e)  
                    #         st.info(f"error df3 skiprow: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                
                if uploaded_file3.name.lower().endswith(('.xls','.xlsx')):
                    df3 = load_xls_df(uploaded_file3)                    
                    # try:
                    #     df3 = pd.read_excel(uploaded_file3, engine= 'openpyxl')
                    # except Exception as e:
                    #     print(e)
                    #     st.error(f"error df3 excel: {e}")
            # else:
            #     df3 = pd.DataFrame(None)    
                                    
            if uploaded_file4 is not None:
                if uploaded_file4.name.lower().endswith(('.csv')):
                    df4 = load_csv_df(uploaded_file4)
                    
                    # try:               
                    #     df4_head = pd.read_csv(uploaded_file4, encoding_errors='ignore', error_bad_lines=False, header=None)
                    #     head4_len = len(df4_head)                                                 
                    # except Exception as e:
                    #     print(e)  
                    #     st.info(f"error df4 head: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                    # if head4_len >= 40:
                    #     # majority of file comes thru, header likely at top
                    #     try:
                    #         uploaded_file4.seek(0)
                    #         df4 = pd.read_csv(uploaded_file4, encoding_errors='ignore')
                    #     except Exception as e:
                    #         print(e)  
                    #         st.info(f"error df4 load: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")   
                                                
                    #         try:
                    #             uploaded_file4.seek(0)
                    #             df4 = pd.read_csv(uploaded_file4, encoding_errors='ignore', error_bad_lines=False)
                    #         except Exception as e:
                    #             print(e)  
                    #             st.info(f"error df4 errorload: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                    
                    # else: 
                    #     # Majority of file does not come thru, header likely under some extra meta lines
                    #     try:
                    #         uploaded_file4.seek(0)                      
                    #         df4 = pd.read_csv(uploaded_file4, skiprows=head4_len, encoding_errors='ignore', error_bad_lines=False)
                    #     except Exception as e:
                    #         print(e)  
                    #         st.info(f"error df4 skiprow: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                
                if uploaded_file4.name.lower().endswith(('.xls','.xlsx')):
                    df4 = load_xls_df(uploaded_file4)                    
                    # try:
                    #     df4 = pd.read_excel(uploaded_file4, engine= 'openpyxl')
                    # except Exception as e:
                    #     print(e) 
                    #     st.error(f"error df4 excel: {e}")                    
            # else:
            #     df4 = pd.DataFrame(None)     
                
            #if uploaded_lasfile is not None:
                #try:
                    #df5 = lasio.read(uploaded_lasfile, engine='normal').df()
                    #df5.reset_index(inplace=True)
                #except Exception as e:
                    #print(e)
                    
            if not df.empty:  # only add the dataframe to the list if data was successfully read
                dfs_to_compare.append(df)
                
                # # Send df to db
                # with st.spinner('Loading data...'):
                #     try:
                #         try:
                #             db.validate_collection(collectionname)  # Try to validate a collection
                #             st.sidebar.write(f"{collectionname} Exists.")
                #         except pymongo.errors.OperationFailure:  # If the collection doesn't exist
                #             # # print("This collection doesn't exist")
                #             db[collectionname].insert_many(df.to_dict('records'), ordered=False)
                #             # db[collectionname].bulk_write(df.to_dict('records'), ordered=False)
                #             st.sidebar.write(f"{collectionname} Uploaded.")
                            
                            
                #     except Exception as e:
                #         print(e)
                #         st.sidebar.warning(f"Data upload error: {e}")
                
                # #  Test Azure Blob
                # with st.spinner('Loading data...'):
                #     blob_upload(df, collectionname)
                #     st.sidebar.write(f"{st.secrets('container')} Uploaded.")
            
            #  Test Azure Blob
            with st.spinner('Loading data...'):
                try:
                    if not df.empty: 
                        dbstatus = blob_upload(df, uploaded_file.name)
                        dbmsg.info(dbstatus)                
                except Exception as e:
                    print(e)
                    st.error(f"error df1upload: {e}")
                try:
                    if not df2.empty: 
                        dbstatus = blob_upload(df2, uploaded_file2.name)
                        dbmsg.info(dbstatus)               
                except Exception as e:
                    print(e)
                    # st.error(f"error df2upload: {e}")
                try:
                    if not df3.empty: 
                        dbstatus = blob_upload(df3, uploaded_file3.name)
                        dbmsg.info(dbstatus)               
                except Exception as e:
                    print(e)
                    # st.error(f"error df3upload: {e}")
                try:
                    if not df4.empty: 
                        dbstatus = blob_upload(df4, uploaded_file4.name)
                        dbmsg.info(dbstatus)               
                except Exception as e:
                    print(e)
                    # st.error(f"error df4upload: {e}")
                                    
                st.sidebar.write(f"{st.secrets('container')} Uploaded.")
                
    except Exception as e:
        print(e)
        # st.error(f"error upload: {e}")

    # Check filename against records
    try:
        if uploaded_file.name:
            # # https://blog.streamlit.io/create-a-search-engine-with-streamlit-and-google-sheets/    
            dfBR = pd.read_csv('data/records_USonly.csv')
            # Use a text_input to get the keywords to filter the dataframe
            # text_search = st.text_input("Search videos by title or speaker", value="")
            # text_search = '|'.join(uploaded_file.name)
            text_search = uploaded_file.name
            # Filter the dataframe using masks
            m1 = dfBR["Well_Name_Full"].str.contains(text_search)
            m2 = dfBR["Type"].str.contains(text_search)
            df_search = dfBR[m1 | m2]
            # Another way to show the filtered results
            # Show the cards
            N_cards_per_row = 3
            if text_search:        
                with st.expander ('Records Found:'):
                    for n_row, row in df_search.reset_index().iterrows():
                        i = n_row%N_cards_per_row
                        if i==0:
                            st.write("---")
                            cols = st.columns(N_cards_per_row, gap="large")
                        # draw the card
                        with cols[n_row%N_cards_per_row]:
                            st.caption(f"{row['Latitude'].strip()} - {row['Longitude'].strip()} - {row['Hole Section'].strip()} ")
                            st.markdown(f"**{row['Well_Name_Full'].strip()}**")
                            st.markdown(f"*{row['BitMfgr'].strip()}*")
                            st.markdown(f"**{row['Size']}**")
                            st.markdown(f"**{row['Type']}**")   
                            
                    # Show the results, if you have a text_search            
                    st.write(df_search)                  
    except Exception as e:
        print(e)
    
    # capture user options
    # chart_select = st.sidebar.selectbox(label='Do you want to show the Graphs?',options=['Yes', 'No'])
    # chart_select2 = st.sidebar.selectbox(label='Do you want to show the lines for the second dataframe?',options=['Yes', 'No'])
    chart_select = 'Yes'
    chart_select2 = 'Yes'
    chart_selects = [chart_select, chart_select2]
    #well_names = ['well 1', 'well 2']
    
    # dbmsg.info('Initializing Charts')
    dbmsg.progress(10,'Initializing Charts')
    
    #  Create indicator banner
    figmetric1 = go.Figure()
    # uploaded_file
    figmetric1.update_layout(autosize=True, height=100, title={'text':f'<b>Well 1:</b>  {well1name}','font':{'color':'#008C9A'}}, grid = {'rows': 1, 'columns': 4, 'pattern': "independent"},paper_bgcolor=gbg, margin={"l": 0, "r": 0, "t": 50, "b": 0}) 
    # figmetric1.update_layout(autosize=True, height=100, title={'text':'Well 1:','font':{'color':'#008C9A'}}, grid = {'rows': 1, 'columns': 4, 'pattern': "independent"},paper_bgcolor=gbg, margin={"l": 0, "r": 0, "t": 50, "b": 0}) 
    figmetric2 = go.Figure()
    figmetric2.update_layout(autosize=True, height=100, title={'text':f'<b>Well 2:</b>  {well2name}','font':{'color':'#8ACFDD'}}, grid = {'rows': 1, 'columns': 4, 'pattern': "independent"},paper_bgcolor=gbg, margin={"l": 0, "r": 0, "t": 50, "b": 0}) 
    figmetric3 = go.Figure()
    figmetric3.update_layout(autosize=True, height=100, title={'text':f'<b>Well 3:</b>  {well3name}','font':{'color':'#7A7D81'}}, grid = {'rows': 1, 'columns': 4, 'pattern': "independent"},paper_bgcolor=gbg, margin={"l": 0, "r": 0, "t": 50, "b": 0}) 
    figmetric4 = go.Figure()
    figmetric4.update_layout(autosize=True, height=100, title={'text':f'<b>Well 4:</b>  {well4name}','font':{'color':'#98C007'}}, grid = {'rows': 1, 'columns': 4, 'pattern': "independent"},paper_bgcolor=gbg, margin={"l": 0, "r": 0, "t": 50, "b": 0}) 
    
    # create & prepare basic figure for data traces
    fig = make_subplots(rows = 1, cols = 10, shared_yaxes=True, subplot_titles=[' ', 'Rate Of Penetration', 'Weight on Bit', 'Rotary RPM', 'Differential Pressure', 'Torque','Flow Rate','Gamma Ray','RigState','MSE'], specs=[[{},{},{},{},{},{},{},{},{},{}]])
    # adjust title heights
    for annotation in fig['layout']['annotations']: 
        annotation['yanchor']='bottom'
        annotation['y']=1.05
        annotation['yref']='paper'
    fig.update_xaxes(title_text="Datetime", rangemode="tozero",title_standoff=1, row=1, col=1)
    fig.update_xaxes(title_text="Rate Of Penetration",rangemode="tozero",title_standoff=1, row=1, col=2)
    fig.update_xaxes(title_text="Weight on Bit",rangemode="tozero",title_standoff=1, row=1, col=3)
    fig.update_xaxes(title_text="Rotary RPM",rangemode="tozero",title_standoff=1, row=1, col=4)
    fig.update_xaxes(title_text='Differential Pressure',rangemode="tozero",title_standoff=1, row=1, col=5)
    fig.update_xaxes(title_text='Torque',rangemode="tozero",title_standoff=1, row=1, col=6)
    fig.update_xaxes(title_text='Flow Rate',rangemode="tozero",title_standoff=1, row=1, col=7)
    fig.update_xaxes(title_text='Gamma Ray',rangemode="tozero",title_standoff=1, row=1, col=8)
    fig.update_xaxes(title_text='RigState',title_standoff=1, row=1, col=9)    
    fig.update_layout(xaxis9=dict(type='linear',range=[0, 12])) 
    fig.update_xaxes(title_text='MSE',rangemode="tozero",title_standoff=1, row=1, col=10)
    fig.update_layout(xaxis10=dict(type='log',rangemode="tozero")) 
    
    fig.update_yaxes(title_text="Depth", row=1, col=1, ) # autorange='reversed',tickformat='digits'
    # Test for top margin to prevent expanding from chart height button
    fig.update_layout(margin=dict(l=60,r=10,t=125,b=1,autoexpand=False),hovermode = "y")
    # 1000 = 900
    # width= 1000,
    # fig.update_layout(autosize=True, height= 1000,  xaxis={'side':'top'}, xaxis2={'side':'top'}, xaxis3={'side':'top'}, xaxis4={'side':'top'}, xaxis5={'side':'top'})
    fig.update_layout(autosize=True, height= 1000,  xaxis={'side':'top'}, xaxis2={'side':'top'}, xaxis3={'side':'top'}, xaxis4={'side':'top'}, xaxis5={'side':'top'}, xaxis6={'side':'top'}, xaxis7={'side':'top'}, xaxis8={'side':'top'}, xaxis9={'side':'top'}, xaxis10={'side':'top'})
    # fig.update_layout(paper_bgcolor='rgb(10,10,10)')
    # fig.update_layout(yaxis_rangeslider_visible=True)   
    
    buttons1 = [dict(args=["height", "1000"],label="Height: 1000px",method="relayout"),
                dict(args=["height", "2000"],label="2000px",method="relayout"),
                dict(args=["height", "5000"],label="5000px",method="relayout")]
    
    # fig.update_layout(updatemenus=[dict(type = "buttons",buttons=list([
    #                                     dict(args=["height", "1000"],label="Height: 1000px",method="relayout"),
    #                                     dict(args=["height", "2000"],label="2000px",method="relayout"),
    #                                     dict(args=["height", "5000"],label="5000px",method="relayout")]),
    #                                     direction="down",pad={"r": 10, "t": 10},showactive=True,x=0.01,xanchor="left",y=1.15,yanchor="top"),
    #                                 ])
    # fig.update_layout(annotations=[dict(text="Chart Height:", showarrow=False,x=0, y=1.1, yref="paper", align="left")])       
    fig.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg)
    fig.update_layout(newshape=dict(line_color='rgba(192,192,192,1)',fillcolor='rgba(192,192,192,0.6)'))
    
    # Color traces by updatemenu button  https://community.plotly.com/t/change-a-specific-trace-data-with-button/71050/7
    # def create(figure, name_list, change_dict_list):
    #     info = len([*figure.select_traces()])
    #     indices = [idx for idx, trace in enumerate(fig.select_traces()) if trace.name in name_list]        
    #     # analyze change_dict_list
    #     new_colors = []
    #     for d in change_dict_list:
    #         new_colors.append(d.get('line').get('color'))        
    #     # analyze traces
    #     existing_colors = []
    #     for trace in figure.select_traces():
    #         existing_colors.append(trace.line.color)        
    #     # substitute colors
    #     for idx, new_color in zip(indices, new_colors):
    #         existing_colors[idx] = new_color                
    #     # prepare output
    #     return {'line': [{'color': color} for color in existing_colors]}    
    # button = [dict(method='update',args=[
    #         create(fig,  ['Secondtrace', 'Forthtrace'], [{'line':{'color':'yellow'}}, {'line':{'color':'red'}}])
    #     ],label='button')]        
    # fig.layout.updatemenus = [{'buttons': button}]
    
    
    # Add range slider https://plotly.com/python/range-slider/
    # fig.update_layout(xaxis=dict(rangeslider=dict(visible=True)))
    # Plotly.relayout(gd, 'yaxis.range', [0, 5]);
    # fig.update_layout(xaxis=dict(rangeselector=dict(buttons=list([
    #                 dict(count=1,
    #                     label="1m",
    #                     step="month",
    #                     stepmode="backward"),
    #                 dict(count=6,
    #                     label="6m",
    #                     step="month",
    #                     stepmode="backward"),
    #                 dict(count=1,
    #                     label="YTD",
    #                     step="year",
    #                     stepmode="todate"),
    #                 dict(count=1,
    #                     label="1y",
    #                     step="year",
    #                     stepmode="backward"),
    #                 dict(step="all")
    #             ])),rangeslider=dict(visible=True),type="date"))
    
    # create & prepare basic figure for data traces
    figadv = make_subplots(rows = 1, cols = 5, shared_yaxes=True)
    figadv.update_xaxes(title_text="Datetime", row=1, col=1)
    figadv.update_xaxes(title_text="Rate Of Penetration", row=1, col=2)
    figadv.update_xaxes(title_text="Weight on Bit", row=1, col=3)
    figadv.update_xaxes(title_text="Rotary RPM", row=1, col=4)
    figadv.update_xaxes(title_text='Differential Pressure', row=1, col=5)
    figadv.update_xaxes(title_text='Torque', row=1, col=6)
    figadv.update_xaxes(title_text='Flow Rate', row=1, col=7)
    figadv.update_xaxes(title_text='Gamma Ray', row=1, col=8)
    figadv.update_xaxes(title_text='RigState', row=1, col=9)
    figadv.update_yaxes(title_text="Depth", row=1, col=1, autorange='reversed',tickformat="digits")
    figadv.update_yaxes(title_text="Depth", row=1, col=1, autorange='reversed',tickformat="digits")
    figadv.update_layout(autosize=True, height= 1000,  xaxis={'side':'top'}, xaxis2={'side':'top'}, xaxis3={'side':'top'}, xaxis4={'side':'top'}, xaxis5={'side':'top'}, xaxis6={'side':'top'}, xaxis7={'side':'top'}, xaxis8={'side':'top'}, xaxis9={'side':'top'}) 
    figadv.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg, newshape=dict(line_color='rgba(192,192,192,1)',fillcolor='rgba(192,192,192,0.6)'))
    
         
    figDVD = go.Figure()
    figDVD.update_layout(autosize=True, height=800, showlegend=True, legend=dict(yanchor="top",y=0.99,xanchor="right",x=0.99), title='Depth vs. Time')
    figDVD.update_layout(xaxis={'side':'top'})
    # figDVD.update_layout(xaxis={'side':'top','type':'date','tickformat':'%H:%M',dtick=24})
    figDVD.update_layout(yaxis={'title':'Depth','autorange':'reversed'})
    figDVD.update_layout(xaxis = dict(dtick=24))
    figDVD.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg)
    figDVD.layout.template='seaborn'
    # try:
    #     figDVD.update_layout(updatemenus=[        
    #         dict(
    #             buttons=list([
    #                 dict(
    #                     args=[{"xaxis.tickformat": "x%p"}],
    #                     label="Hours",
    #                     method="relayout"
    #                 ),
    #                 dict(
    #                     args=[{"xaxis.tickformat": "x%d"}],
    #                     label="Days",
    #                     method="relayout"
    #                 )
    #             ]),
    #             type = "buttons",direction="right",pad={"r": 10, "t": 10},showactive=True,x=0.1,xanchor="left",y=1.1,yanchor="top"
    #         ),])
    # except Exception as e:
    #     print(e)
    #     st.error(f"error buttons1 excel: {e}")
    
    figDVDI = make_subplots(rows = 1, cols = 2, column_widths=[0.6, 0.4], shared_yaxes=True)
    figDVDI.update_layout(autosize=True, height=800, showlegend=True)
    figDVDI.update_layout(xaxis={'side':'top'},xaxis2={'side':'top'},yaxis={'title':'Depth','autorange':'reversed'},yaxis2={'title':'Depth','autorange':'reversed'})
    figDVDI.update_xaxes(title_text="Depth vs. Time", row=1, col=1)    
    figDVDI.update_layout(xaxis = dict(dtick = 24))
    figDVDI.update_xaxes(title_text="Inclination", row=1, col=2)
    # figDVDI.update_layout(xaxis2=dict(type='linear',range=[1, 100],ticksuffix='°')) 
    figDVDI.update_layout(xaxis2=dict(type='linear',ticksuffix='°')) 
    figDVDI.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg)
    figDVDI.layout.template='seaborn'
    # try:    
    #     figDVDI.update_layout(updatemenus=[        
    #         dict(
    #             buttons=list([
    #                 dict(
    #                     args=[{"xaxis.ticktext": "x"}],
    #                     label="Hours",
    #                     method="relayout"
    #                 ),
    #                 dict(
    #                     args=[{"xaxis.ticktext": "x/24"}],
    #                     label="Days",
    #                     method="relayout"
    #                 )
    #             ]),
    #             type = "buttons",direction="right",pad={"r": 10, "t": 10},showactive=True,x=0.1,xanchor="left",y=1.1,yanchor="top"
    #         ),])
    # except Exception as e:
    #     print(e)
    #     st.error(f"error buttons2 excel: {e}")
    
    # figbxrop = go.Figure()
    figbxrop = make_subplots(rows = 2, cols = 2, subplot_titles=("ROP", "WOB", "RPM", "Diff Pres"), shared_yaxes=False)
    figbxrop.update_layout(autosize=True, height=800, xaxis={'side':'bottom'}, xaxis2={'side':'bottom'}, xaxis3={'side':'bottom'}, xaxis4={'side':'bottom'}, showlegend=False)
    figbxrop.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg)
    # legend=dict(yanchor="top",y=0.99,xanchor="right",x=0.99)
    # figbxrop.update_layout(title_text="Box Plot Styling Outliers")

    
    fig3DMap = go.Figure()
    fig3DMap.update_layout(autosize=True, height=800, showlegend=True, legend=dict(yanchor="top",y=0.99,xanchor="right",x=0.99), title='3D Projected Well Path')
    fig3DMap.update_layout(scene = dict(xaxis_title='Easting',yaxis_title='Northing',zaxis_title='Depth',zaxis = dict(autorange='reversed')))
    fig3DMap.update_layout(margin=dict(l=10, r=10))
    fig3DMap.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg)
    # fig3DMap.update_layout(margin=dict(l=65, r=50, b=65, t=90))
    
    figSurf1 = go.Figure()
    figSurf1.update_layout(autosize=True, height=1000, showlegend=True, legend=dict(yanchor="top",y=0.99,xanchor="right",x=0.99), title='XY Performance')
    figSurf1.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg)
    
    # figRotSlidemetric = go.Figure()
    # figRotSlidemetric.update_layout(autosize=True, height=100, title={'text':f'<b>Rotate vs. Sliding ROP</b>','font':{'color':'#98C007'}}, grid = {'rows': 2, 'columns': 4, 'pattern': "independent"},paper_bgcolor=gbg, margin={"l": 0, "r": 0, "t": 50, "b": 0}) 
    
    figRotSlide = make_subplots(rows = 4, cols = 4, column_widths=[0.1, 0.2, 0.2, 0.5], specs=[[{},{'type':'domain'},{},{}],[{},{'type':'domain'},{},{}],[{},{'type':'domain'},{},{}],[{},{'type':'domain'},{},{}]], shared_yaxes=False)
    figRotSlide.update_layout(autosize=True, height=800, showlegend=False)
    figRotSlide.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg)
    
    # figSlide1Pie = go.Figure() 
    # figSlide1Pie.update_layout(autosize=True, height=200, title='Sliding %')
    # figSlide1Bar = go.Figure() 
    # figSlide1Bar.update_layout(autosize=True, height=200, showlegend=True, legend=dict(yanchor="top",y=0.99,xanchor="right",x=0.99), title='Sliding %')
    # figSlide1Hor = go.Figure() 
    # figSlide1Hor.update_layout(autosize=True, height=200, showlegend=True, legend=dict(yanchor="top",y=0.99,xanchor="right",x=0.99), title='Sliding %')
    # figSlide1Pie.add_trace(go.Pie(labels=labels, values=[16, 15, 12, 6, 5, 4, 42], name="GHG Emissions"))
    
    
    figAnalytic = make_subplots(rows = 4, cols = 2, shared_yaxes=False)
    figAnalytic.update_layout(autosize=True, showlegend=False)
    figAnalytic.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg)
    
        
    # Add range slider
    # fig.update_layout(yaxis=dict(rangeselector=dict(buttons=list([
    #     dict(count=1,
    #         label="1m",
    #         step="month",
    #         stepmode="backward"),
    #     dict(count=6,
    #         label="6m",
    #         step="month",
    #         stepmode="backward"),
    #     dict(count=1,
    #         label="YTD",
    #         step="year",
    #         stepmode="todate"),
    #     dict(count=1,
    #         label="1y",
    #         step="year",
    #         stepmode="backward"),
    #     dict(step="all")
    #     ])),rangeslider=dict(visible=True),type="date")
    # )

    # showlegend = False, 

    # 7.1.17 testing timedelta axis formatting
    # df[Hdr.dt] = df[Hdr.dt] + pd.to_datetime(df['YYYY/MM/DD']+df['HH:MM:SS'])            
    # ds1_dt_min = df[Hdr.dt].min
    # ds1_dt_max = df[Hdr.wob].max
    # df[Hdr.dt] = df[Hdr.dt] - df[hdr.dt] - ds1_dt_min

    def data_cleanup(def_df):
            # dbmsg.info('Performing Miracle 1')
            dbmsg.progress(15, 'Performing Miracle 1')
            # Set all cols to upper case
            def_df.columns = def_df.columns.str.upper()
            # Clean unnamed colms, negate that match unnamed
            def_df.loc[:, ~def_df.columns.str.match('Unnamed')]
            # Remove Duplicate columns
            # https://sparkbyexamples.com/pandas/pandas-remove-duplicate-columns-from-dataframe/
            if not len(def_df.columns) == len(set(def_df.columns)):                
                dbmsg.progress(16, 'Miracle 1: Clean Duplicates')
                # def_df.T.drop_duplicates().T
                def_df.loc[:,~def_df.columns.duplicated()]
            # df2 = df.loc[:,~df.T.duplicated(keep='first')]
            # def_df = def_df.loc[:,~def_df.columns.duplicated()]
            # Rename to Convention
            def_df.rename(columns=final_column_name_map, inplace=True)
            # df.columns = [col + ' = ' + str(newElements.pop(0)) if col.startswith(stringMatch) else col for col in df.columns]
            def_num = def_df._get_numeric_data()
            # Null negative values
            def_num[def_num < 0] = np.nan
            def_num = def_num.dropna(axis = 0, how = 'all')
            # need to add code to remove holedepth that reverts *FIX*
            try:
                def_df = def_df[def_df[Hdr.hole_depth] >= 0.01]
            except Exception as e:
                st.error(f"error holedepth filter: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                st.write(def_df)
            # def_df = def_df[def_df[Hdr.hole_depth] != 0]
            
            def_df.reset_index(inplace=True)   
            
            # with st.expander ('W0 d_cleanup'):
            #     st.write(def_df)
                        
            return def_df

    def Data_Capping(def_df): 
        
        #  Flooring & Capping https://www.pluralsight.com/guides/cleaning-up-data-from-outliers
        # df["Income"] = np.where(df["Income"] <2960.0, 2960.0,df['Income'])
        # df["Income"] = np.where(df["Income"] >12681.0, 12681.0,df['Income']) 
        # for cc in df:
        #     if cc == [Hdr.rop] | [Hdr.wob] | [Hdr.rpm] | [Hdr.diff] | [Hdr.trq] | [Hdr.flowrate] | [Hdr.gr]:                        
        #         df[cc] = np.where(df[cc] < df[cc].quantile(qu_lim), df[cc].quantile(qu_lim), df[cc])
        #         df[cc] = np.where(df[cc] > df[cc].quantile(100-qu_lim), df[cc].quantile(100-qu_lim), df[cc])
        #         # Interpolate after capping to smoothen   
        #         # df[cc] = (df[cc].rolling(ravg_sample).mean())
        # dbmsg.info('Data capping')
        dbmsg.progress(20, 'Data capping')
        
        qu_lim = 0.01      
        if Hdr.rop in def_df.columns:                 
            def_df[Hdr.rop] = np.where(def_df[Hdr.rop] < def_df[Hdr.rop].quantile(qu_lim), def_df[Hdr.rop].rolling(avg5p_sample).mean(), def_df[Hdr.rop])
            def_df[Hdr.rop] = np.where(def_df[Hdr.rop] > def_df[Hdr.rop].quantile(1-qu_lim), def_df[Hdr.rop].rolling(avg5p_sample).mean(), def_df[Hdr.rop])
            
        if Hdr.wob in def_df.columns:
            def_df[Hdr.wob] = np.where(def_df[Hdr.wob] < def_df[Hdr.wob].quantile(qu_lim), def_df[Hdr.wob].rolling(avg5p_sample).mean(), def_df[Hdr.wob])
            def_df[Hdr.wob] = np.where(def_df[Hdr.wob] > def_df[Hdr.wob].quantile(1-qu_lim), def_df[Hdr.wob].rolling(avg5p_sample).mean(), def_df[Hdr.wob])
            
        if Hdr.rpm in def_df.columns:
            def_df[Hdr.rpm] = np.where(def_df[Hdr.rpm] < def_df[Hdr.rpm].quantile(qu_lim), def_df[Hdr.rpm].rolling(avg5p_sample).mean(), def_df[Hdr.rpm])
            def_df[Hdr.rpm] = np.where(def_df[Hdr.rpm] > def_df[Hdr.rpm].quantile(1-qu_lim), def_df[Hdr.rpm].rolling(avg5p_sample).mean(), def_df[Hdr.rpm])
            
        if Hdr.diff in def_df.columns:
            def_df[Hdr.diff] = np.where(def_df[Hdr.diff] < def_df[Hdr.diff].quantile(qu_lim), def_df[Hdr.diff].rolling(avg5p_sample).mean(), def_df[Hdr.diff])
            def_df[Hdr.diff] = np.where(def_df[Hdr.diff] > def_df[Hdr.diff].quantile(1-qu_lim), def_df[Hdr.diff].rolling(avg5p_sample).mean(), def_df[Hdr.diff])
            
        if Hdr.trq in def_df.columns:
            def_df[Hdr.trq] = np.where(def_df[Hdr.trq] < def_df[Hdr.trq].quantile(qu_lim), def_df[Hdr.trq].rolling(avg5p_sample).mean(), def_df[Hdr.trq])
            def_df[Hdr.trq] = np.where(def_df[Hdr.trq] > def_df[Hdr.trq].quantile(1-qu_lim), def_df[Hdr.trq].rolling(avg5p_sample).mean(), def_df[Hdr.trq])
                
        if Hdr.flowrate in def_df.columns:
            def_df[Hdr.flowrate] = np.where(def_df[Hdr.flowrate] < def_df[Hdr.flowrate].quantile(qu_lim), def_df[Hdr.flowrate].rolling(avg5p_sample).mean(), def_df[Hdr.flowrate])
            def_df[Hdr.flowrate] = np.where(def_df[Hdr.flowrate] > def_df[Hdr.flowrate].quantile(1-qu_lim), def_df[Hdr.flowrate].rolling(avg5p_sample).mean(), def_df[Hdr.flowrate])
        
        if Hdr.gr in def_df.columns:
            def_df[Hdr.gr] = np.where(def_df[Hdr.gr] < def_df[Hdr.gr].quantile(qu_lim), def_df[Hdr.gr].rolling(avg5p_sample).mean(), def_df[Hdr.gr])
            def_df[Hdr.gr] = np.where(def_df[Hdr.gr] > def_df[Hdr.gr].quantile(1-qu_lim), def_df[Hdr.gr].rolling(avg5p_sample).mean(), def_df[Hdr.gr])
        
        return def_df
        
    # loop over user data to create chart traces
    # with stcol0.spinner('Loading data...'):
    try:
        for df, chart_select in zip(dfs_to_compare, chart_selects):
            #  test code            
            # try:
            #     new_names = [final_column_name_map.get(i) for i in df.columns.str.replace("[^ ]*_","")]
            #     with st.expander('Test Col Parse'):
            #         st.write(new_names)
            # except Exception as e:
            #     st.error(f"error new-names: {e}")  
            
            # Test code for Regex column matching 
            # https://stackoverflow.com/questions/46162202/pandas-series-replace-using-dictionary-with-regex-keys
            # df['Col_1'].replace(repl_dict, regex=True)

            # try:
            #     dftest2 = df.columns.to_series().replace(final_column_name_map, regex=True)         
            #     with st.expander('Test Regex2'):
            #         st.write(dftest2)
            # except Exception as e:
            #     st.error(f"error regex test2: {e}")
            # try:
            #     dftest3 = df.rename(columns = {f"{k}": f"{v}" for k,v in regex2_column_name_map.items()})           
            #     with st.expander('Test Regex3'):
            #         st.write(dftest3)
            # except Exception as e:
            #     st.error(f"error regex test3: {e}")
            
            # # Clean unnamed colms, negate that match unnamed
            # df.loc[:, ~df.columns.str.match('Unnamed')]
            # # Rename to Convention
            # df.rename(columns=final_column_name_map, inplace=True)
            # # df.columns = [col + ' = ' + str(newElements.pop(0)) if col.startswith(stringMatch) else col for col in df.columns]
            # num = df._get_numeric_data()
            # # Null negative values
            # num[num < 0] = np.nan
            # df[df[Hdr.hole_depth] > 0]
            # df.reset_index(inplace=True)
            
            df = data_cleanup(df)
            
            count = 1
            # with stcol0.spinner('Finding time domain...'):
            try:
                # with st.expander("1"):
                #     st.write(df.columns)
                # if not set([Hdr.dt]).issubset(df.columns): 
                if Hdr.dt not in df.columns:
                    # df[Hdr.dt] = pd.to_datetime(df['YYYY/MM/DD'] + ' ' + df['HH:MM:SS'], errors='coerce')
                    # dbmsg.info('Normalizing time domain')
                    dbmsg.progress(30, 'Normalizing time domain')
                    # Check for 'YYYY/MM/DD', if not 'YYYY/MM/DD' then change to 'YYYY/MM/DD'
                    # https://stackoverflow.com/questions/16870663/how-do-i-validate-a-date-string-format-in-python
                    def validate(date_text):
                        try:
                            if date_text != datetime.strptime(date_text, "%Y-%m-%d").strftime('%Y-%m-%d'):
                                raise ValueError
                            return True
                        except ValueError:
                            return False
                    
                    # Fix date time format consolidation for Aramco data set
                    # https://stackoverflow.com/questions/54247148/pandas-check-if-any-column-is-date-time-and-change-it-to-date-format-string-y
                    try:
                        df[Hdr.dt] = pd.to_datetime(df[Hdr.date] + ' ' + df[Hdr.time], errors='coerce')             
                    except Exception as e:
                        st.error(f"error date time1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")  
                    try:
                        df['dt2'] = pd.to_datetime(df[Hdr.date].astype(str) + ' ' + df[Hdr.time].astype(str), errors='coerce')                 
                    except Exception as e:
                        st.error(f"error date time2: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")  
                elif len(df[Hdr.dt].value_counts()) < 0:
                    # df[Hdr.dt] = pd.to_datetime(df['YYYY/MM/DD'] + ' ' + df['HH:MM:SS'], errors='coerce')
                    df[Hdr.dt] = pd.to_datetime(df[Hdr.date] + ' ' + df[Hdr.time], errors='coerce')
                    try:
                        df['dt2'] = pd.to_datetime(df[Hdr.date].astype(str) + ' ' + df[Hdr.time].astype(str), errors='coerce')                 
                    except Exception as e:
                        st.error(f"error date time3: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")  
                else:
                    # Make sure dt format is correct
                    df[Hdr.dt] = pd.to_datetime(df[Hdr.dt], errors='coerce')

                try: 
                    dfst = df[Hdr.dt].values[1]
                    # df.loc[df[Hdr.rop]!='0','A'].index[0]
                    df['dtst'] = (dfst)
                    df['dte'] = ((df[Hdr.dt] - df['dtst']) / np.timedelta64(1, "h"))
                except Exception as e:
                    print(e)
                    st.error(f"error calculating elapsed time1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                
                # Check for kilo ft_lb torque, correct
                if set([Hdr.trq]).issubset(df.columns):
                    df[Hdr.trq] = np.where(np.mean(df[Hdr.trq]) < 1, df[Hdr.trq]*1000, df[Hdr.trq])
                
                # Check for rigstates
                if set([Hdr.sprstate,Hdr.substate]).issubset(df.columns):
                    # st.warning(f"rigstate present") 
                    # dbmsg.info('Performing Miracle 3')
                    dbmsg.progress(50, 'Performing Miracle 3')
                    df = miracle_rigstate(df)
                else:
                    st.warning(f"rigstate not present")
                    with st.expander ('W1 pre-miracle rigstate'):
                        st.write(df)
                    # dbmsg.info('Performing Mirale 3')
                    dbmsg.progress(50, 'Performing Miracle 3')
                    df = miracle_rigstate(df)
                    with st.expander ('W1 post-miracle rigstate'):
                        st.write(df)
                    
                    #  Run Rigstate logic
                
                # Create Synthetic Survey from data   
                dfsvy = []             
                if set([Hdr.hole_depth, 'ads-spr', Hdr.inc, Hdr.az]).issubset(df.columns):
                    # dbmsg.info('Performing Miracle 2')
                    dbmsg.progress(40,'Performing Miracle 2')
                    dfsvy = miracle_survey(df)  
                
                            
            except Exception as e:
                print(e)
                st.error(f"error date time: {e}  //// {traceback.print_exc()} //// {traceback.format_exc()}")                
                contents = [f"error date time: {e}", 'user info:', userdata,
                            f'File name: {uploaded_file.name}',
                            f'//// {traceback.print_exc()} //// {traceback.format_exc()}',
                            "You can find an audio file attached.", '/local/path/to/song.mp3'
                ]
                yag.send('ccasad@ulterra.com', 'Adira: error date time', contents, attachments=uploaded_file)
                
            try:
                # df2.rename(columns=final_column_name_map, inplace=True)
                # num = df2._get_numeric_data()
                # num[num < 0] = np.nan
                # df2.reset_index(inplace=True)
                
                # Clean up data
                data_cleanup(df2)
                
                count = 2
                        
                try:
                    df2[Hdr.dt] = pd.to_datetime(df2['YYYY/MM/DD'] + ' ' + df2['HH:MM:SS'])
                    avg_rop2 = np.mean(df2[Hdr.rop])
                    avg_wob2 = np.mean(df2[Hdr.wob])   
                          
                    try: 
                        df2st = df2[Hdr.dt].values[1] 
                        df2['dtst'] = (df2st)        
                        df2['dte'] = ((df2[Hdr.dt] - df2['dtst']) / np.timedelta64(1, "h")) 
                    except Exception as e:
                        print(e)
                        st.error(f"error calculating elapsed time3: {e}")
                        
                    # Check for kilo ft_lb torque, correct
                    if set([Hdr.trq]).issubset(df2.columns):
                        df2[Hdr.trq] = np.where(np.mean(df2[Hdr.trq]) < 1, df2[Hdr.trq]*1000, df2[Hdr.trq])
                                   
                    # Check for rigstates
                    if set([Hdr.sprstate,Hdr.substate]).issubset(df2.columns):
                        # st.warning(f"rigstate present") 
                        df2 = miracle_rigstate(df2)
                    else:
                        st.warning(f"rigstate not present")
                        df2 = miracle_rigstate(df2)
                                     
                    # Create Synthetic Survey from data
                    df2svy = []             
                    if set([Hdr.hole_depth, 'ads-spr', Hdr.inc, Hdr.az]).issubset(df2.columns):
                        df2svy = miracle_survey(df2)
                       
                except Exception as e:
                    print(e)                
                # try:
                #     df['depth'] = df['depth'].merge(df2[Hdr.hole_depth])
                # except Exception as e:
                #     print(e)
                #     st.error(f"error calculating elapsed time4: {e}")
                    
                # Set Plot DF view
                bdf2 = df2
                
                try:
                    # df3.rename(columns=final_column_name_map, inplace=True)
                    # num = df3._get_numeric_data()
                    # num[num < 0] = np.nan
                    # df3.reset_index(inplace=True)
                    
                    # Clean up data
                    data_cleanup(df3)
                    count = 3
                    try:
                        df3[Hdr.dt] = pd.to_datetime(df3['YYYY/MM/DD'] + ' ' + df3['HH:MM:SS'])                               
                        avg_rop3 = np.mean(df3[Hdr.rop])  
                        avg_wob3 = np.mean(df3[Hdr.wob])  
                        
                        try: 
                            df3st = df3[Hdr.dt].values[1] 
                            df3['dtst'] = (df3st)        
                            df3['dte'] = ((df3[Hdr.dt] - df3['dtst']) / np.timedelta64(1, "h")) 
                        except Exception as e:
                            print(e)
                            st.error(f"error calculating elapsed time5: {e}")
                                           
                        # Check for kilo ft_lb torque, correct
                        if set([Hdr.trq]).issubset(df3.columns):
                            df3[Hdr.trq] = np.where(np.mean(df3[Hdr.trq]) < 1, df3[Hdr.trq]*1000, df3[Hdr.trq])
                
                        # Check for rigstates
                        if set([Hdr.sprstate,Hdr.substate]).issubset(df3.columns):
                            # st.warning(f"rigstate present") 
                            df3 = miracle_rigstate(df3)
                        else:
                            st.warning(f"rigstate not present")
                            df3 = miracle_rigstate(df3)
                            
                        # Create Synthetic Survey from data 
                        df3svy = []             
                        if set([Hdr.hole_depth, 'ads-spr', Hdr.inc, Hdr.az]).issubset(df3.columns):
                            df3svy = miracle_survey(df3)                        
                            
                    except Exception as e:
                        print(e)  
                                                
                    # Set Plot DF view
                    bdf3 = df3
                    
                    try:
                        # df4.rename(columns=final_column_name_map, inplace=True)
                        # num = df4._get_numeric_data()
                        # num[num < 0] = np.nan
                        # df4.reset_index(inplace=True)
                        # Clean up data
                        data_cleanup(df4)
                        count = 4
                        try:
                            df4[Hdr.dt] = pd.to_datetime(df4['YYYY/MM/DD'] + ' ' + df4['HH:MM:SS'])                               
                            avg_rop4 = np.mean(df4[Hdr.rop])
                            avg_wob4 = np.mean(df4[Hdr.wob])         
                            try: 
                                df4st = df4[Hdr.dt].values[1] 
                                df4['dtst'] = (df4st)        
                                df4['dte'] = ((df4[Hdr.dt] - df4['dtst']) / np.timedelta64(1, "h"))
                            except Exception as e:
                                print(e)
                                st.error(f"error calculating elapsed time7: {e}")                                
                            
                            # Check for kilo ft_lb torque, correct
                            if set([Hdr.trq]).issubset(df4.columns):
                                df4[Hdr.trq] = np.where(np.mean(df4[Hdr.trq]) <= 1, df4[Hdr.trq]*1000, df4[Hdr.trq])
                           
                            # Check for rigstates
                            if set([Hdr.sprstate,Hdr.substate]).issubset(df4.columns):
                                # st.warning(f"rigstate present") 
                                df4 = miracle_rigstate(df4)
                            else:
                                st.warning(f"rigstate not present")
                                df4 = miracle_rigstate(df4)
                                    
                            # Create Synthetic Survey from data 
                            df4svy = []             
                            if set([Hdr.hole_depth, 'ads-spr', Hdr.inc, Hdr.az]).issubset(df4.columns):
                                df4svy = miracle_survey(df4)  
                                
                        except Exception as e:
                            print(e)              
                        # try:
                        #     df['depth'] = df['depth'].merge(df4[Hdr.hole_depth])
                        # except Exception as e:
                        #     print(e)
                        #     st.error(f"error calculating elapsed time8: {e}")
                            
                        # Set Plot DF view
                        bdf4 = df4
                    except Exception as e:
                        print(e)
                except Exception as e:
                    print(e)            
            except Exception as e:
                print(e)  
            
            
            # st.info(f'count is {count}')
            
            # if chart_select != DropOptions.yes:
            #     st.info('Continue')
            #     continue 
                    
        # Plot dataset 1
        try:
            if count >= 1:
                # Set Metric header
                # fig_meter.add_trace(go.Indicator(title='ROP',mode="number+delta",value=round(np.mean(df[Hdr.rop]),1),number={'suffix': ' ft/hr'},delta={'reference': round(np.mean(df[Hdr.rop]), 1),'relative':True,'valueformat': '.1%'},number_font_color="#008C9A",domain={'row': 0, 'column': 1}))
                # fig_meter.add_trace(go.Indicator(title='WOB',mode="number+delta",value=round(np.mean(df[Hdr.wob]),1),number={'suffix': ' kbs'},delta={'reference': round(np.mean(df[Hdr.wob]), 1),'relative':True,'valueformat': '.1%'},number_font_color="#008C9A",domain={'row': 0, 'column': 2}))
                
                # Set metric columns
                with hed1:
                    hed1title1, hed1title2 = st.columns([1,5], gap='small')
                    with hed1title1:
                        color1 = st.color_picker(f'**Well1:** {well1name}  \n\n **ROP:** {round(np.mean(df[Hdr.rop]),1)} ft/hr  **WOB:** {round(np.mean(df[Hdr.wob]),1)} kbs', color1, label_visibility='collapsed')                  
                    with hed1title2:
                        st.markdown(f"<span style='color:{color1}'> {well1name}</span>",unsafe_allow_html=True)
                        
                    hed1sub0, hed1sub1, hed1sub2, hed1sub3 = st.columns([1,2,2,1], gap='small')
                    with hed1sub1:
                        st.metric('ROP', round(np.mean(df[Hdr.rop]),1), delta=round(np.mean(df[Hdr.rop]), 1), delta_color="normal", help=None, label_visibility="visible")
                    with hed1sub2:
                        st.metric('WOB', round(np.mean(df[Hdr.wob]),1), delta=round(np.mean(df[Hdr.wob]), 1), delta_color="normal", help=None, label_visibility="visible")
                
                
                    
                # channel = [Hdr.rop,Hdr.wob,Hdr.rpm,Hdr.diff]
                # channel[1] = well1.multiselect('Channel 1', list(df),Hdr.rop)
                # channel[2] = well2.multiselect('Channel 2', list(df),Hdr.wob)
                # channel[3] = well3.multiselect('Channel 3', list(df),Hdr.rpm)
                # channel[4] = well4.multiselect('Channel 4', list(df),Hdr.diff)
                
                # channel = [df[Hdr.rop].head(), df[Hdr.wob].head(), df[Hdr.rpm].head(), df[Hdr.diff].head()]
                # channel[1] = well1.multiselect('Channel 1', list(df.columns), channel[1])
                # channel[2] = well1.multiselect('Channel 2', list(df.columns), channel[2])
                # channel[3] = well1.multiselect('Channel 3', list(df.columns), channel[3])
                # channel[4] = well1.multiselect('Channel 4', list(df.columns), channel[4])
                
                # chnlcount = 4                 
                # channels = stcol2.multiselect('Select Channels', list(df))
                # chnlwarn = stcol2.empty()
                # if len(channels) <= chnlcount:
                #     chnlwarn.warning(f'Select at least {chnlcount} curves.')
                # else:
                #     chnlwarn.empty()
                    
                    # # fig3 = make_subplots(rows=1, cols= len(channels), subplot_titles=channels, shared_yaxes=True)                    
                    # # create & prepare basic figure for data traces
                    # fig = make_subplots(rows=1, cols=len(channels), subplot_titles=channels, shared_yaxes=True)
                    # fig.update_xaxes(title_text="Datetime", row=1, col=1)
                    # fig.update_xaxes(title_text="Rate Of Penetration", row=1, col=2)
                    # fig.update_xaxes(title_text="Weight on Bit", row=1, col=3)
                    # fig.update_xaxes(title_text="Rotary RPM", row=1, col=4)
                    # fig.update_xaxes(title_text='Differential Pressure', row=1, col=5)
                    # fig.update_yaxes(title_text="Depth", row=1, col=1, autorange='reversed')

                    # 1000 = 1000
                    # fig.update_layout(autosize=True, height= 1000, plot_bgcolor=gbg, xaxis={'side':'top'}, xaxis2={'side':'top'}, xaxis3={'side':'top'}, xaxis4={'side':'top'}, xaxis5={'side':'top'})                     
                    # fig.update_layout(newshape=dict(line_color='rgba(192,192,192,1)',fillcolor='rgba(192,192,192,0.6)'))
                    
                    # channel_index = 1
                    # for channel in channels:
                    #     fig3.add_trace(go.Scatter(x=well_data[channel], y=well_data['DEPTH']), row=1, col=channel_index)
                    #     channel_index+=1
                    
                    # fig3.update_layout(height=1000, showlegend=False, yaxis={'title':'DEPTH','autorange':'reversed'})
                    # fig3.layout.template='seaborn'
                    # st.plotly_chart(fig3, use_container_width=True)
                    
                avg5p_sample = round(len(df.index) * 0.005)
                # rad0 = st_toggle_switch(label="Data Control:",key="Key1",default_value=False,label_after = False,inactive_color = '#D3D3D3',active_color="#11567f",track_color="#29B5E8")
                selOD = float(stcol2.selectbox('Bit Size',
                                      ('6.0','6.25','6.75','7.875','8.5','8.75','9.875','11.0','12.25','13.5','17.5'), index=8))
                            
                
                if set([Hdr.trq,Hdr.rpm,Hdr.rop,Hdr.wob]).issubset(df.columns):
                    # df['MSE'] = ((df[Hdr.wob]/selOD)+((120*np.pi*df[Hdr.rpm]*T)/(selOD*df[Hdr.rop])))
                    df['MSE'] = (((480*df[Hdr.trq]*df[Hdr.rpm])/((selOD**2)*df[Hdr.rop]))+((4*df[Hdr.wob])/(np.pi*(selOD**2))))
                    
                    #CMSE = 4 * WeightOnBit / (np.pi * BitSize * BitSize)  +  480 * RPM * Torque / (RateOfPenetration * BitSize * BitSize)

                # If unit_system = "US" Then
                #     MSE = 480 * bit_rpm * diff * torque_slope / (ROP * bit_dia ^ 2) / 1000 + 1.273 * wob / bit_dia ^ 2
                # ElseIf unit_system = "CAN" Then
                #     MSE = (480000 * bit_rpm * diff * torque_slope / (ROP * bit_dia ^ 2) / 1000 + 1.273 * wob * 1000 / bit_dia ^ 2) * 0.14504
                # End If
                # rad5 = stcol2.radio('Channels:', options = ['Standard', 'Advanced'], index=1)
                rad1 = stcol2.radio('Data Control:', options = ['Cleaned', 'Raw'], index=0)
                # st.write(rad1)
                # rad2 = stcol2.radio('Chart Size:', options = ['Fit to Screen', 'Full Height', 'Poster'], index=0)
                rad3 = stcol2.radio('Data Smoothing:', options = [f'Interpolated ({avg5p_sample})', 'Real'], index=1)
                rad4 = stcol2.radio('Rig State:', options = ['All', 'Drilling'], index=1)  
                
                
                if rad1 == 'Cleaned': 
                                           
                    # qu_lim = 0.01                    
                    # df[Hdr.rop] = np.where(df[Hdr.rop] < df[Hdr.rop].quantile(qu_lim), df[Hdr.rop].rolling(avg5p_sample).mean(), df[Hdr.rop])
                    # df[Hdr.rop] = np.where(df[Hdr.rop] > df[Hdr.rop].quantile(1-qu_lim), df[Hdr.rop].rolling(avg5p_sample).mean(), df[Hdr.rop])
                    
                    # df[Hdr.wob] = np.where(df[Hdr.wob] < df[Hdr.wob].quantile(qu_lim), df[Hdr.wob].rolling(avg5p_sample).mean(), df[Hdr.wob])
                    # df[Hdr.wob] = np.where(df[Hdr.wob] > df[Hdr.wob].quantile(1-qu_lim), df[Hdr.wob].rolling(avg5p_sample).mean(), df[Hdr.wob])
                    
                    # df[Hdr.rpm] = np.where(df[Hdr.rpm] < df[Hdr.rpm].quantile(qu_lim), df[Hdr.rpm].rolling(avg5p_sample).mean(), df[Hdr.rpm])
                    # df[Hdr.rpm] = np.where(df[Hdr.rpm] > df[Hdr.rpm].quantile(1-qu_lim), df[Hdr.rpm].rolling(avg5p_sample).mean(), df[Hdr.rpm])
                    
                    # df[Hdr.diff] = np.where(df[Hdr.diff] < df[Hdr.diff].quantile(qu_lim), df[Hdr.diff].rolling(avg5p_sample).mean(), df[Hdr.diff])
                    # df[Hdr.diff] = np.where(df[Hdr.diff] > df[Hdr.diff].quantile(1-qu_lim), df[Hdr.diff].rolling(avg5p_sample).mean(), df[Hdr.diff])
                    
                    df = Data_Capping(df)
                                
                # if rad5 == 'Standard':
                #     chnl_num = 5
                # elif rad5 == 'Advanced':
                #     chnl_num = 9
                
                chnl_num = 9
                    
                # if rad2 == 'Fit to Screen':
                #     1000 = 900
                #     fig.update_layout(height= 1000, width= 1000)
                # elif rad2 == 'Full Height':
                #     1000 = 5000
                #     fig.update_layout(height= 1000, width= 1000)
                # elif rad2 == 'Poster':
                #     1000 = 10000
                #     fig.update_layout(height= 1000, width= 2000) 
                                
                if rad4 == 'Drilling':
                    try:
                        df = df[df['ads-spr'] == 1]
                    except Exception as e:
                        print(e)
                        st.error(f"error rig state filter: {e}")    
                else:                        
                    df = df  
                    
                if rad3 == f'Interpolated ({avg5p_sample})':
                    ravg_sample = avg5p_sample
                    # round(len(df.index) * 0.005)
                                    
                    # try:
                    #     df[Hdr.crop] = (df[Hdr.rop].rolling(ravg_sample).mean())
                    # except Exception as e:
                    #     print(e)
                    #     st.error(f"error calculating rolling avg: {e}")
                    
                    # try:
                    #     df2[Hdr.crop] = (df2[Hdr.rop].rolling(ravg_sample).mean())
                    # except Exception as e:
                    #     print(e)  
                        
                    # try:
                    #     df3[Hdr.crop] = (df3[Hdr.rop].rolling(ravg_sample).mean())
                    # except Exception as e:
                    #     print(e)  
                        
                    # try:
                    #     df4[Hdr.crop] = (df4[Hdr.rop].rolling(ravg_sample).mean())
                    # except Exception as e:
                    #     print(e)                          
                elif rad3 == 'Real':
                    # df[Hdr.crop] = df[Hdr.rop]
                    ravg_sample = 1
                else:                    
                    ravg_sample = 1
                                            
                # if rad1 == 'Cleaned': 
                #     #  Flooring & Capping https://www.pluralsight.com/guides/cleaning-up-data-from-outliers
                #     # df["Income"] = np.where(df["Income"] <2960.0, 2960.0,df['Income'])
                #     # df["Income"] = np.where(df["Income"] >12681.0, 12681.0,df['Income'])
                    
                #     qu_lim = 0.01                    
                #     ravg_sample = round(len(df.index) * 0.005)
                #     for cc in df:
                #         if cc == Hdr.rop | Hdr.wob | Hdr.rpm | Hdr.diff | Hdr.trq | Hdr.flowrate | Hdr.gr:                        
                #             df[cc] = np.where(df[cc] < df[cc].quantile(qu_lim), df[cc].quantile(qu_lim), df[cc])
                #             df[cc] = np.where(df[cc] > df[cc].quantile(100-qu_lim), df[cc].quantile(100-qu_lim), df[cc])
                #             # Interpolate after capping to smoothen   
                #             # df[cc] = (df[cc].rolling(ravg_sample).mean())
                        
                #     # df[Hdr.rop] = np.where(df[Hdr.rop] < df[Hdr.rop].quantile(qu_lim), df[Hdr.rop].quantile(qu_lim), df[Hdr.rop])
                #     # df[Hdr.rop] = np.where(df[Hdr.rop] > df[Hdr.rop].quantile(100-qu_lim), df[Hdr.rop].quantile(100-qu_lim), df[Hdr.rop])
                    
                #     # df[Hdr.wob] = np.where(df[Hdr.wob] < df[Hdr.wob].quantile(qu_lim), df[Hdr.wob].quantile(qu_lim), df[Hdr.wob])
                #     # df[Hdr.wob] = np.where(df[Hdr.wob] > df[Hdr.wob].quantile(100-qu_lim), df[Hdr.wob].quantile(100-qu_lim), df[Hdr.wob])
                    
                #     # df[Hdr.rpm] = np.where(df[Hdr.rpm] < df[Hdr.rpm].quantile(qu_lim), df[Hdr.rpm].quantile(qu_lim), df[Hdr.rpm])
                #     # df[Hdr.rpm] = np.where(df[Hdr.rpm] > df[Hdr.rpm].quantile(100-qu_lim), df[Hdr.rpm].quantile(100-qu_lim), df[Hdr.rpm])
                    
                #     # df[Hdr.diff] = np.where(df[Hdr.diff] < df[Hdr.diff].quantile(qu_lim), df[Hdr.diff].quantile(qu_lim), df[Hdr.diff])
                #     # df[Hdr.diff] = np.where(df[Hdr.diff] > df[Hdr.diff].quantile(100-qu_lim), df[Hdr.diff].quantile(100-qu_lim), df[Hdr.diff])
                    
                #     # with stcol0.spinner('Recalculating...'):
                #     # df = df[(df[Hdr.rop] < df[Hdr.rop].quantile(0.99)) & (df[Hdr.rop] > df[Hdr.rop].quantile(0.01))]
                #     # df = df[(df[Hdr.wob] < df[Hdr.wob].quantile(0.99)) & (df[Hdr.wob] > df[Hdr.wob].quantile(0.01))]
                #     # df = df[(df[Hdr.rpm] < df[Hdr.rpm].quantile(0.99)) & (df[Hdr.rpm] > df[Hdr.rpm].quantile(0.01))]
                #     # df = df[(df[Hdr.diff] < df[Hdr.diff].quantile(0.99)) & (df[Hdr.diff] > df[Hdr.diff].quantile(0.01))]
                    
                #     # df = df[(df[Hdr.trq] < df[Hdr.trq].quantile(0.99)) & (df[Hdr.trq] > df[Hdr.trq].quantile(0.01))]
                #     # df = df[(df[Hdr.flowrate] < df[Hdr.flowrate].quantile(0.99)) & (df[Hdr.flowrate] > df[Hdr.flowrate].quantile(0.01))]
                #     # df = df[(df[Hdr.gr] < df[Hdr.gr].quantile(0.99)) & (df[Hdr.gr] > df[Hdr.gr].quantile(0.01))]
                                
                # dbmsg.info('Loading Well 1.')
                dbmsg.progress(60,'Loading Well 1.')
                if uploaded_file2 is None:                    
                    
                    make_traces(0, df, df, dfsvy, color1, chnl_num, ravg_sample, figmetric1, fig, figDVD, figDVDI, figRotSlide, figbxrop, fig3DMap)
                    
                    # # Test plots
                    # try:
                    #     if dfsvy.head(1):
                    #         if set(['3dx','3dy', Hdr.rop]).issubset(dfsvy.columns):
                    #             figSurf1.add_trace(go.Surface(x=dfsvy['3dx'],y=dfsvy['3dy'],z=dfsvy[Hdr.rop]))
                            
                    #             fig9 = go.Figure(data=[go.Surface(x=dfsvy['3dx'],y=dfsvy['3dy'],z=dfsvy[Hdr.rop])])

                    #             fig9.update_layout(title='Mt Bruno Elevation', autosize=False,
                    #                         width=500, height=500,
                    #                         margin=dict(l=65, r=50, b=65, t=90))
                    # except Exception as e:
                    #     print(e)
                    #     st.error(f"error dfsvy empty: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                        
                else:
                    make_traces(1, df, df, dfsvy, color1, chnl_num, ravg_sample, figmetric1, fig, figDVD, figDVDI, figRotSlide, figbxrop, fig3DMap)
                    
                # dfrot = df.loc[df['ads-rot'] == 1]
                # dfslide = df.loc[df['ads-slide'] == 1]
                # # figRotSlidemetric.add_trace(go.Indicator(title='Rotate',mode="number",value=round(dfrot[Hdr.crop].mean(),1),number={'suffix': ' ft/hr'},number_font_color="#008C9A",domain={'row': 0, 'column': 1}), row=1, col=1)
                # figRotSlide.add_trace(go.Box(y=dfrot[Hdr.crop],name=f'Avg ROP, {round(dfrot[Hdr.crop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color=color1), row=1, col=1)
                # figRotSlide.add_annotation(xref="x domain",yref="y domain",x=0.5, y=1.2, showarrow=False,text="ROTATING", row=1, col=1)
                # # figRotSlidemetric.add_trace(go.Indicator(title='Slide',mode="number",value=round(dfslide[Hdr.crop].mean(),1),number={'suffix': ' ft/hr'},number_font_color="#008C9A",domain={'row': 0, 'column': 2}), row=1, col=2)
                # figRotSlide.add_trace(go.Box(y=dfslide[Hdr.crop],name=f'Avg ROP, {round(dfslide[Hdr.crop].mean(),1)}',jitter=0.2,pointpos=0,boxpoints='outliers',marker_size=1,marker_color='Black',line_color=color1), row=1, col=2)
                # figRotSlide.add_annotation(xref="x domain",yref="y domain",x=0.5, y=1.2, showarrow=False,text="SLIDING", row=1, col=2)      
                
                
                
                
                # figRotSlide.add_trace(go.Scatter(name="Rotate", x=dfrot[Hdr.hole_depth], y=dfrot[Hdr.crop], markers= {'color':color1, 'width':0.4}), row=1, col=3)
                # figRotSlide.add_trace(go.Scatter(name="Slide", x=dfslide[Hdr.hole_depth], y=dfslide[Hdr.crop], markers= {'color':color3, 'width':0.4}), row=1, col=3)
                
                # figmetricRSrop = px.scatter_polar(dfrot, r=Hdr.az, theta=Hdr.inc, title='Rotate Net DLS', row=1, col=4)
                # figmetricRSrop = px.scatter_polar(dfslide, r=Hdr.az, theta=Hdr.inc, title='Slide Net DLS', row=2, col=4)
                
                # Test Altair plotting
                # scale = alt.Scale(
                #     domain=["sun", "fog", "drizzle", "rain", "snow"],
                #     range=["#e7ba52", "#a7a7a7", "#aec7e8", "#1f77b4", "#9467bd"],
                # )
                # color = alt.Color("weather:N", scale=scale)

                # We create two selections:
                # - a brush that is active on the top panel
                # - a multi-click that is active on the bottom panel
                # brush = alt.selection_interval(encodings=["y"])
                # click = alt.selection_multi(encodings=["color"])

                # # Top panel is scatter plot of temperature vs time
                # points = (alt.Chart().mark_point().encode(
                #         alt.X("monthdate(date):T", title="Date"),
                #         alt.Y("temp_max:Q",title="Maximum Daily Temperature (C)",scale=alt.Scale(domain=[-5, 40]),),
                #         color=alt.condition(brush, color, alt.value("lightgray")),
                #         size=alt.Size("precipitation:Q", scale=alt.Scale(range=[5, 200])),
                #     )
                #     .properties(width=550, height=300)
                #     .add_selection(brush)
                #     .transform_filter(click)
                # )
                # # Bottom panel is a bar chart of weather type
                # bars = (alt.Chart().mark_bar().encode(
                #         x="count()",
                #         y="weather:N",
                #         color=alt.condition(click, color, alt.value("lightgray")),
                #     )
                #     .transform_filter(brush)
                #     .properties(
                #         width=550,
                #     )
                #     .add_selection(click)
                # )
                
                
                
                # brush = alt.selection_interval(encodings=["y"])
                # # Left panel is scatter plot of depth vs time
                # altdvd = (alt.Chart().mark_point().encode(
                #         alt.X('dte', title="Time"),
                #         alt.Y(Hdr.hole_depth,title="Depth",),
                #     )
                #     .properties(height=900)
                #     .add_selection(brush)
                # )
                
                # altrop = (alt.Chart().mark_point().encode(
                #         alt.X(Hdr.crop, title="ROP"),
                #         alt.Y(Hdr.hole_depth,title="Depth",),
                #     )
                #     .properties(height=900)
                #     .transform_filter(brush)
                # )                
                # altchart = alt.hconcat(altdvd, altrop, data=df, title="test EDR Channels")
                
                
                                
                # attempted code for live avg of selection
                # selected_points = plotly_events(fig, click_event=False, hover_event=True)                
                # if len(selected_points) == 0:
                #     st.stop()
                
                # selected_x_value = selected_points[0]["x"]
                # selected_y_value = selected_points[0]["y"]
                
                # df_selected = df[
                #     (df[Hdr.crop] == selected_x_value)
                #     & (df[Hdr.hole_depth] == selected_y_value)
                # ]
                # st.write("Callback Data for selected point:")
                # st.write(df_selected)
                    
                # plot dataset 2
                try:
                    if count >= 2:
                        # Set metric columns
                        with hed2:
                            hed2title1, hed2title2 = st.columns([1,5], gap='small')
                            with hed2title1:
                                color2 = st.color_picker('w2', color2, label_visibility='collapsed')                  
                            with hed2title2:
                                st.markdown(f"<span style='color:{color2}'> {well2name}</span>",unsafe_allow_html=True)
                                
                            # hed2sub1, hed2sub2 = st.columns(2, gap='small')
                            hed2sub0, hed2sub1, hed2sub2, hed2sub3 = st.columns([1,2,2,1], gap='small')
                            with hed2sub1:
                                st.metric('ROP', round(np.mean(df2[Hdr.rop]),1), delta=round(np.mean(df[Hdr.rop]), 1), delta_color="normal", help=None, label_visibility="visible")
                            with hed2sub2:
                                st.metric('WOB', round(np.mean(df2[Hdr.wob]),1), delta=round(np.mean(df[Hdr.wob]), 1), delta_color="normal", help=None, label_visibility="visible")
                    
                                
                        # if agree:  
                        
                        if set([Hdr.trq,Hdr.rpm,Hdr.rop,Hdr.wob]).issubset(df2.columns):                       
                            df2['MSE'] = (((480*df2[Hdr.trq]*df2[Hdr.rpm])/((selOD**2)*df2[Hdr.rop]))+((4*df2[Hdr.wob])/(np.pi*(selOD**2))))   
                                               
                        if rad4 == 'Drilling':
                            df2 = df2[df2['ads-spr'] == 1]    
                        else:                        
                            df2 = df2 
                                        
                        if rad1 == 'Cleaned':                          
                            df2 = Data_Capping(df2)
                        
                        # dbmsg.info('Loading Well 2..')
                        dbmsg.progress(65,'Loading Well 2..')
                        make_traces(2, df, df2, df2svy, color2, chnl_num, ravg_sample, figmetric2, fig, figDVD, figDVDI, figRotSlide, figbxrop, fig3DMap)
                        
                        
                        # plot dataset 3
                        try:                            
                            if count >= 3:
                                # Set metric columns
                                with hed3:
                                    hed3title1, hed3title2 = st.columns([1,5], gap='small')
                                    with hed3title1:
                                        color3 = st.color_picker('w3', color3, label_visibility='collapsed')                  
                                    with hed3title2:
                                        st.markdown(f"<span style='color:{color3}'> {well3name}</span>",unsafe_allow_html=True)
                                        
                                    # hed3sub1, hed3sub2 = st.columns(2, gap='small')
                                    hed3sub0, hed3sub1, hed3sub2, hed3sub3 = st.columns([1,2,2,1], gap='small')
                                    with hed3sub1:
                                        st.metric('ROP', round(np.mean(df3[Hdr.rop]),1), delta=round(np.mean(df[Hdr.rop]), 1), delta_color="normal", help=None, label_visibility="visible")
                                    with hed3sub2:
                                        st.metric('WOB', round(np.mean(df3[Hdr.wob]),1), delta=round(np.mean(df[Hdr.wob]), 1), delta_color="normal", help=None, label_visibility="visible")
                            
                                
                                # if agree:                                  
                                if set([Hdr.trq,Hdr.rpm,Hdr.rop,Hdr.wob]).issubset(df3.columns):                               
                                    df3['MSE'] = (((480*df3[Hdr.trq]*df3[Hdr.rpm])/((selOD**2)*df3[Hdr.rop]))+((4*df3[Hdr.wob])/(np.pi*(selOD**2))))
                                
                                if rad4 == 'Drilling':
                                    df3 = df3[df3['ads-spr'] == 1]     
                                else:                        
                                    df3 = df3  
                                    
                                if rad1 == 'Cleaned':                       
                                    df3 = Data_Capping(df3)
                                dbmsg.progress(70, 'Loading Well 3...')
                                make_traces(3, df, df3, df3svy, color3, chnl_num, ravg_sample, figmetric3, fig, figDVD, figDVDI, figRotSlide, figbxrop, fig3DMap)
                                
                                
                                # plot dataset 4
                                try:                                    
                                    if count >= 4:
                                        # Set metric columns
                                        with hed4:
                                            hed4title1, hed4title2 = st.columns([1,5], gap='small')
                                            with hed4title1:
                                                color4 = st.color_picker('w2', color4, label_visibility='collapsed')                  
                                            with hed4title2:
                                                st.markdown(f"<span style='color:{color4}'> {well4name}</span>",unsafe_allow_html=True)
                                                
                                            # hed4sub1, hed4sub2 = st.columns(2, gap='small')
                                            hed4sub0, hed4sub1, hed4sub2, hed4sub3 = st.columns([1,2,2,1], gap='small')
                                            with hed4sub1:
                                                st.metric('ROP', round(np.mean(df4[Hdr.rop]),1), delta=round(np.mean(df[Hdr.rop]), 1), delta_color="normal", help=None, label_visibility="visible")
                                            with hed4sub2:
                                                st.metric('WOB', round(np.mean(df4[Hdr.wob]),1), delta=round(np.mean(df[Hdr.wob]), 1), delta_color="normal", help=None, label_visibility="visible")
                                    
                                                
                                        # if agree: 
                                        if set([Hdr.trq,Hdr.rpm,Hdr.rop,Hdr.wob]).issubset(df4.columns):
                                            df4['MSE'] = (((480*df4[Hdr.trq]*df4[Hdr.rpm])/((selOD**2)*df4[Hdr.rop]))+((4*df4[Hdr.wob])/(np.pi*(selOD**2))))
                                        
                                        if rad4 == 'Drilling':
                                            df4 = df4[df4['ads-spr'] == 1]     
                                        else:                        
                                            df4 = df4  
                                            
                                        if rad1 == 'Cleaned':                       
                                            df4 = Data_Capping(df4)   
                                            # df4[Hdr.rop] = np.where(df4[Hdr.rop] < df4[Hdr.rop].quantile(qu_lim), df4[Hdr.rop].rolling(avg5p_sample).mean(), df4[Hdr.rop])
                                            # df4[Hdr.rop] = np.where(df4[Hdr.rop] > df4[Hdr.rop].quantile(1-qu_lim), df4[Hdr.rop].rolling(avg5p_sample).mean(), df4[Hdr.rop])
                                            
                                            # df4[Hdr.wob] = np.where(df4[Hdr.wob] < df4[Hdr.wob].quantile(qu_lim), df4[Hdr.wob].rolling(avg5p_sample).mean(), df4[Hdr.wob])
                                            # df4[Hdr.wob] = np.where(df4[Hdr.wob] > df4[Hdr.wob].quantile(1-qu_lim), df4[Hdr.wob].rolling(avg5p_sample).mean(), df4[Hdr.wob])
                                            
                                            # df4[Hdr.rpm] = np.where(df4[Hdr.rpm] < df4[Hdr.rpm].quantile(qu_lim), df4[Hdr.rpm].rolling(avg5p_sample).mean(), df4[Hdr.rpm])
                                            # df4[Hdr.rpm] = np.where(df4[Hdr.rpm] > df4[Hdr.rpm].quantile(1-qu_lim), df4[Hdr.rpm].rolling(avg5p_sample).mean(), df4[Hdr.rpm])
                                            
                                            # df4[Hdr.diff] = np.where(df4[Hdr.diff] < df4[Hdr.diff].quantile(qu_lim), df4[Hdr.diff].rolling(avg5p_sample).mean(), df4[Hdr.diff])
                                            # df4[Hdr.diff] = np.where(df4[Hdr.diff] > df4[Hdr.diff].quantile(1-qu_lim), df4[Hdr.diff].rolling(avg5p_sample).mean(), df4[Hdr.diff]) 
                                        # dbmsg.info('Loading Well 4....')                     
                                        dbmsg.progress(75,'Loading Well 4....')                     
                                        make_traces(4, df, df4, df4svy, color4, chnl_num, ravg_sample, figmetric4, fig, figDVD, figDVDI, figRotSlide, figbxrop, fig3DMap)
                                    
                                except Exception as e:
                                    print(e)
                                    # st.info(f"error dataframe4: {e}")
                        except Exception as e:
                            print(e) 
                            # st.info(f"error dataframe3: {e}")  
                except Exception as e:
                    print(e)
                    # st.info(f"error dataframe2: {e}")
                
                config0 = dict({'displayModeBar': False})
                # config = dict(
                #     {
                #         'scrollZoom': True,
                #         'displayModeBar': True,
                #         # 'editable'              : True,
                #         'modeBarButtonsToAdd': [
                #             'drawline',
                #             'drawopenpath',
                #             'drawclosedpath',
                #             'drawcircle',
                #             'drawrect',
                #             'eraseshape',
                #         ],
                #         'toImageButtonOptions': {'format': 'svg'},
                #         'modeBarButtonsToRemove': [
                #             'zoom2d', 'pan2d', 'zoomIn2d' ,'zoomOut2d'
                #         ]
                #     }
                # )
                config = dict(
                    {
                        'scrollZoom': True,
                        'responsive': True,
                        # 'displayModeBar': True,
                        # 'editable'              : True,
                        'modeBarButtonsToAdd': [
                            # 'drawline',
                            # 'drawopenpath',
                            # 'drawclosedpath',
                            # 'drawcircle',
                            'drawrect',
                            'eraseshape',
                        ],
                        'displaylogo': False,
                        # 'toImageButtonOptions': {'format': 'jpg'},
                        # 'modeBarButtonsToRemove': [
                        #     'autoscale','toImage', 'zoom2d', 'pan2d', 'zoomIn2d' ,'zoomOut2d'
                        # ]
                    }
                )
                
                # Update chart axes
                try:
                    if count == 2:
                        if Hdr.rop in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.rop]),1)}</span>' 
                                                        f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.rop]),1)}</span>'),title_standoff=1, row=1, col=2)
                        if Hdr.wob in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.wob]),1)}</span>'  
                                                        f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.wob]),1)}</span>'),title_standoff=1, row=1, col=3)
                        if Hdr.rpm in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.rpm]),1)}</span>'  
                                                        f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.rpm]),1)}</span>'),title_standoff=1, row=1, col=4)
                        if Hdr.diff in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.diff]),1)}</span>'  
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.diff]),1)}</span>'),title_standoff=1, row=1, col=5)
                        if Hdr.trq in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.trq]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.trq]),1)}</span>'),title_standoff=1, row=1, col=6)
                        if Hdr.flowrate in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.flowrate]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.flowrate]),1)}</span>'),title_standoff=1, row=1, col=7)
                        if Hdr.gr in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.gr]),1)}</span>'  
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.gr]),1)}</span>'),title_standoff=1, row=1, col=8)
                        if 'ads-spr' in df.columns:
                            if Hdr.substate in df.columns:
                                fig.update_xaxes(title_text=f'<span style="color:#98C007;font-size:8px;"> <b>Super</b></span>       <span style="color:#008C9A"> <b>Sub</b></span>',title_standoff=1, row=1, col=9)
                        if 'MSE' in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df["MSE"]),1)}</span>'  
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2["MSE"]),1)}</span>'),title_standoff=1, row=1, col=10)
                        
                    
                    if count == 3:
                        if Hdr.rop in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.rop]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.rop]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.rop]),1)}</span>'),title_standoff=1, row=1, col=2)
                        if Hdr.wob in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.wob]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.wob]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.wob]),1)}</span>'),title_standoff=1, row=1, col=3)
                        if Hdr.rpm in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.rpm]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.rpm]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.rpm]),1)}</span>'),title_standoff=1, row=1, col=4)
                        if Hdr.diff in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.diff]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.diff]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.diff]),1)}</span>'),title_standoff=1, row=1, col=5)
                        if Hdr.trq in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.trq]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.trq]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.trq]),1)}</span>'),title_standoff=1, row=1, col=6)
                        if Hdr.flowrate in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.flowrate]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.flowrate]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.flowrate]),1)}</span>'),title_standoff=1, row=1, col=7)
                        if Hdr.gr in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.gr]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.gr]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.gr]),1)}</span>'),title_standoff=1, row=1, col=8)
                        if 'ads-spr' in df.columns:
                            if 'ads-sub' in df.columns:
                                fig.update_xaxes(title_text=(f'<span style="color:#98C007;font-size:8px;"> <b>Super</b></span>       <span style="color:#008C9A"> <b>Sub</b></span>'),title_standoff=1, row=1, col=9)
                        if 'MSE' in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df["MSE"]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2["MSE"]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3["MSE"]),1)}</span>'),title_standoff=1, row=1, col=10)
                    
                    if count == 4:
                        if Hdr.rop in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.rop]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.rop]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.rop]),1)}</span>' 
                                            f'<span style="color:{color4};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df4[Hdr.rop]),1)}</span>'),title_standoff=1, row=1, col=2)
                        if Hdr.wob in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.wob]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.wob]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.wob]),1)}</span>' 
                                            f'<span style="color:{color4};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df4[Hdr.wob]),1)}</span>'),title_standoff=1, row=1, col=3)
                        if Hdr.rpm in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.rpm]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.rpm]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.rpm]),1)}</span>' 
                                            f'<span style="color:{color4};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df4[Hdr.rpm]),1)}</span>'),title_standoff=1, row=1, col=4)
                        if Hdr.diff in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.diff]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.diff]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.diff]),1)}</span>' 
                                            f'<span style="color:{color4};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df4[Hdr.diff]),1)}</span>'),title_standoff=1, row=1, col=5)
                        if Hdr.trq in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.trq]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.trq]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.trq]),1)}</span>' 
                                            f'<span style="color:{color4};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df4[Hdr.trq]),1)}</span>'),title_standoff=1, row=1, col=6)
                        if Hdr.flowrate in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.flowrate]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.flowrate]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.flowrate]),1)}</span>' 
                                            f'<span style="color:{color4};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df4[Hdr.flowrate]),1)}</span>'),title_standoff=1, row=1, col=7)
                        if Hdr.gr in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df[Hdr.gr]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2[Hdr.gr]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3[Hdr.gr]),1)}</span>' 
                                            f'<span style="color:{color4};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df4[Hdr.gr]),1)}</span>'),title_standoff=1, row=1, col=8)
                        if 'ads-spr' in df.columns:
                            if 'ads-sub' in df.columns:
                                fig.update_xaxes(title_text=(f'<span style="color:#98C007"> <b>Super</b></span>       <span style="color:#008C9A"> <b>Sub</b></span>'),title_standoff=1, row=1, col=9)
                        if 'MSE' in df.columns:
                            fig.update_xaxes(title_text=(f'<span style="color:{color1};font-size:8px;"> <b>AVG:</b> {round(np.mean(df["MSE"]),1)}</span>' 
                                            f'<span style="color:{color2};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df2["MSE"]),1)}</span> <br>' 
                                            f'<span style="color:{color3};font-size:8px;"> <b>AVG:</b> {round(np.mean(df3["MSE"]),1)}</span>' 
                                            f'<span style="color:{color4};font-size:8px;"> <b>  AVG:</b> {round(np.mean(df4["MSE"]),1)}</span>'),title_standoff=1, row=1, col=10)
                
                except Exception as e:                  
                    st.info(f"error multiwell axes avges: {e}")
                    
                # Update updatemenus buttons
                try: 
                    if count == 1:
                        buttons2 = [dict(args=[{"x": [df['dte'], ],
                                        "y": [df[Hdr.hole_depth],],
                                        # "y2": [df[Hdr.hole_depth], ],
                                        # "y3": [df[Hdr.hole_depth], ],
                                        # "y4": [df[Hdr.hole_depth], ],
                                        # "y5": [df[Hdr.hole_depth], ],
                                        # "y6": [df[Hdr.hole_depth], ],
                                        # "y7": [df[Hdr.hole_depth], ],
                                        # "y8": [df[Hdr.hole_depth], ],
                                        # "y9": [df[Hdr.hole_depth], ],
                                        }],label="By Depth",method="restyle"),
                                    dict(args=[{"x": [df[Hdr.hole_depth], ],
                                        "y": [df['dte'], ],
                                        # "y2": [df['dte'], ],
                                        # "y3": [df['dte'], ],
                                        # "y4": [df['dte'], ],
                                        # "y5": [df['dte'], ],
                                        # "y6": [df['dte'], ],
                                        # "y7": [df['dte'], ],
                                        # "y8": [df['dte'], ],
                                        # "y9": [df['dte'], ],
                                    }],label="By Time",method="restyle")]
                    if count == 2:
                        buttons2 = [dict(args=[{"x": [df['dte'],df2['dte'], ],
                                        "y": [df[Hdr.hole_depth],df2[Hdr.hole_depth], ],
                                        # "y2": [df[Hdr.hole_depth],df2[Hdr.hole_depth], ],
                                        # "y3": [df[Hdr.hole_depth],df2[Hdr.hole_depth], ],
                                        # "y4": [df[Hdr.hole_depth],df2[Hdr.hole_depth], ],
                                        # "y5": [df[Hdr.hole_depth],df2[Hdr.hole_depth], ],
                                        # "y6": [df[Hdr.hole_depth],df2[Hdr.hole_depth], ],
                                        # "y7": [df[Hdr.hole_depth],df2[Hdr.hole_depth], ],
                                        # "y8": [df[Hdr.hole_depth],df2[Hdr.hole_depth], ],
                                        # "y9": [df[Hdr.hole_depth],df2[Hdr.hole_depth], ],
                                        }],label="By Depth",method="restyle"),
                                    dict(args=[{"x": [df[Hdr.hole_depth],df2[Hdr.hole_depth], ],
                                        "y": [df['dte'],df2['dte'], ],
                                        # "y2": [df['dte'],df2['dte'], ],
                                        # "y3": [df['dte'],df2['dte'], ],
                                        # "y4": [df['dte'],df2['dte'], ],
                                        # "y5": [df['dte'],df2['dte'], ],
                                        # "y6": [df['dte'],df2['dte'], ],
                                        # "y7": [df['dte'],df2['dte'],],
                                        # "y8": [df['dte'],df2['dte'],],
                                        # "y9": [df['dte'],df2['dte'],],
                                    }],label="By Time",method="restyle")]
                    if count == 3:
                        buttons2 = [dict(args=[{"x": [df['dte'],df2['dte'],df3['dte'], ],
                                        "y": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth], ],
                                        # "y2": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth], ],
                                        # "y3": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth], ],
                                        # "y4": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth], ],
                                        # "y5": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth], ],
                                        # "y6": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth], ],
                                        # "y7": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth], ],
                                        # "y8": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth], ],
                                        # "y9": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth], ],
                                        }],label="By Depth",method="restyle"),
                                    dict(args=[{"x": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth], ],
                                        "y": [df['dte'],df2['dte'],df3['dte'], ],
                                        # "y2": [df['dte'],df2['dte'],df3['dte'], ],
                                        # "y3": [df['dte'],df2['dte'],df3['dte'], ],
                                        # "y4": [df['dte'],df2['dte'],df3['dte'], ],
                                        # "y5": [df['dte'],df2['dte'],df3['dte'], ],
                                        # "y6": [df['dte'],df2['dte'],df3['dte'], ],
                                        # "y7": [df['dte'],df2['dte'],df3['dte'], ],
                                        # "y8": [df['dte'],df2['dte'],df3['dte'], ],
                                        # "y9": [df['dte'],df2['dte'],df3['dte'], ],
                                    }],label="By Time",method="restyle")]
                    if count == 4:
                        buttons2 = [dict(args=[{"x": [df['dte'],df2['dte'],df3['dte'],df4['dte']],
                                        "y": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth],df4[Hdr.hole_depth]],
                                        # "y2": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth],df4[Hdr.hole_depth]],
                                        # "y3": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth],df4[Hdr.hole_depth]],
                                        # "y4": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth],df4[Hdr.hole_depth]],
                                        # "y5": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth],df4[Hdr.hole_depth]],
                                        # "y6": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth],df4[Hdr.hole_depth]],
                                        # "y7": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth],df4[Hdr.hole_depth]],
                                        # "y8": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth],df4[Hdr.hole_depth]],
                                        # "y9": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth],df4[Hdr.hole_depth]],
                                        }],label="By Depth",method="restyle"),
                                    dict(args=[{"x": [df[Hdr.hole_depth],df2[Hdr.hole_depth],df3[Hdr.hole_depth],df4[Hdr.hole_depth]],
                                        "y": [df['dte'],df2['dte'],df3['dte'],df4['dte']],
                                        # "y2": [df['dte'],df2['dte'],df3['dte'],df4['dte']],
                                        # "y3": [df['dte'],df2['dte'],df3['dte'],df4['dte']],
                                        # "y4": [df['dte'],df2['dte'],df3['dte'],df4['dte']],
                                        # "y5": [df['dte'],df2['dte'],df3['dte'],df4['dte']],
                                        # "y6": [df['dte'],df2['dte'],df3['dte'],df4['dte']],
                                        # "y7": [df['dte'],df2['dte'],df3['dte'],df4['dte']],
                                        # "y8": [df['dte'],df2['dte'],df3['dte'],df4['dte']],
                                        # "y9": [df['dte'],df2['dte'],df3['dte'],df4['dte']],
                                    }],label="By Time",method="restyle")]
                    
                    fig.update_layout(updatemenus=[dict(active=0,buttons=buttons1,direction="down", bgcolor='white',x=0.01,y=1.13,xanchor='left',yanchor='top'), ])   
                    
                    fig.update_yaxes(type='linear',range=[df[Hdr.hole_depth].max(), df[Hdr.hole_depth].min()], tickformat = '000',dtick=500) # matches='y',    tickformat = '000',constrain='domain'                      
                                # dict(buttons=buttons2,direction="down",bgcolor='white',x=0.1,y=1.15,xanchor='left',yanchor='top')])               
                except Exception as e:                  
                    st.info(f"error axis swap button: {e}")

                # fig.update_layout(updatemenus=[dict(type = "buttons",buttons=list([
                #                                     dict(args=["height", "1000"],label="Height: 1000px",method="relayout"),
                #                                     dict(args=["height", "2000"],label="2000px",method="relayout"),
                #                                     dict(args=["height", "5000"],label="5000px",method="relayout")]),
                #                                     direction="down",pad={"r": 10, "t": 10},showactive=True,x=0.01,xanchor="left",y=1.15,yanchor="top"),
                #                                 ])
    
                # dbmsg.progress(80,'Render Metrics')
                # try:
                #     if count >= 1:                                            
                #         hed1.plotly_chart(figmetric1, use_container_width=True, config=config) 
                #         # color1 = hed5.color_picker('Well_1 Color', '#008C9A')
                                            
                #     if count >= 2:
                #         hed2.plotly_chart(figmetric2, use_container_width=True, config=config) 
                #         # color2 = hed5.color_picker('Well_2 Color', '#8ACFDD')
                        
                #     if count >= 3:
                #         hed3.plotly_chart(figmetric3, use_container_width=True, config=config) 
                #         # color3 = hed5.color_picker('Well_3 Color', '#7A7D81')
                        
                #     if count >= 4:
                #         hed4.plotly_chart(figmetric4, use_container_width=True, config=config) 
                #         # color4 = hed5.color_picker('Well_4 Color', '#98C007')
                        
                # except Exception as e:                  
                #     st.info(f"error plot figmetric1: {e}")
                
                # Overview
                with tabs[0]:
                    # dbmsg.info('Render Overview')  
                    dbmsg.progress(83,'Render Overview')                  
                    try: 
                        ovcol0.plotly_chart(figDVD, use_container_width=True, config=config)
                    except Exception as e:
                        ovcol0.info(f"error plot figDVD: {e}")
                    
                    try: 
                        ovcol1.plotly_chart(figbxrop, use_container_width=True, config=config)
                    except Exception as e:
                        ovcol1.info(f"error plot figDVD: {e}")
                                            
                # Channels                
                with tabs[1]:  
                    # dbmsg.info('Render Channels') 
                    dbmsg.progress(87,'Render Channels')                  
                    try: 
                        stcol1.plotly_chart(fig, use_container_width=True, config=config)                                                    
                    except Exception as e:    
                        st.info(f"error plot fig: {e}")
                                                
                    # try: 
                    #     stcol0.plotly_chart(figDVD, use_container_width=True, config=config)
                    # except Exception as e:
                    #     st.info(f"error plot figDVD: {e}")
                                        
                    with stcol2:
                        with st.expander('Add Comment'):                                            
                            formcomment = st.form("comment_form")       
                            annotation_ftin = formcomment.text_input("Start Depth:")
                            annotation_ftout = formcomment.text_input("End Depth:")                     
                            annotation_comment = formcomment.text_input("Comment text:", key="annotation_1")
                            submit = formcomment.form_submit_button("Submit")
                            # name = formcomment.text_input("Your name")
                            # reporttype = foformcommentrmpdf.selectbox(
                            #     "Choose Type",
                            #     ["Full Report", "Poster"],
                            #     index=0,
                            # )
                            # grade = formpdf.slider("Grade", 1, 100, 60)
                            # submit = formpdf.form_submit_button("Generate PDF")
                                # with st.form("comment_form"):
                                #     annotation_2 = st.text_input("Annotation text:", key="annotation_1")
                                #     submit = st.form_submit_button("Submit")
                            # annotation_2 = st.text_input("Annotation text:", key="annotation_2")
                            # x = np.arange(10)
                            # fig_2 = go.Figure(data=go.Scatter(x=x, y=x ** 2))
                            # chart_2 = st.empty()
                            # chart_2.plotly_chart(fig)
                            
                    if annotation_comment:
                        fig.add_annotation(text=annotation_comment,xref="paper",yref="paper",x=0,y=1.1,showarrow=False,align="left",xanchor="left",font=dict(size=20, color="#242526"))
                        stcol1.plotly_chart(fig, use_container_width=True, config=config) 
                    
                # DvD tab
                with tabs[2]:  
                    # dbmsg.info('Render DvD') 
                    dbmsg.progress(90,'Render DvD')                                       
                    try: 
                        st.plotly_chart(figDVDI, use_container_width=True, config=config)
                    except Exception as e:                  
                        st.info(e)
                                                              
                    # try: 
                    #     # TEST PLOTS
                    #     st.plotly_chart(figSurf1, use_container_width=True, config=config)
                    #     st.plotly_chart(fig9, use_container_width=True, config=config)
                        
                    # except Exception as e:                  
                    #     st.info(e)
                    
                
                # Survey(3D) Map tab           
                with tabs[3]:   
                    # dbmsg.info('Render Surveys')  
                    dbmsg.progress(92,'Render Surveys')                                                     
                    try: 
                        st.plotly_chart(fig3DMap, use_container_width=True, config=config)
                    except Exception as e:                  
                        st.info(e)
                                                      
                    st.write(dfsvy)
                
                # Sliding Tab
                with tabs[4]:
                    # try: 
                    #     st.altair_chart(altchart, use_container_width=True)
                    # except Exception as e:
                    #     st.info(e)
                    # dbmsg.info('Render Sliding') 
                    dbmsg.progress(95,'Render Sliding') 
                    st.plotly_chart(figRotSlide, use_container_width=True, config=config)
                    # tab4col0, tab4col1 = st.columns([2,6])                    
                    # with tab4col0:
                    #     st.plotly_chart(figRotSlidemetric, use_container_width=True, config=config)
                    
                    # with tab4col0:
                    #     st.plotly_chart(figRotSlide, use_container_width=True, config=config)
                                  
                
                
                # Toolface Tab
                with tabs[5]: 
                    # dbmsg.info('Render Altair test') 
                    dbmsg.progress(97,'Render Altair test') 
                    # get_chart_47479(df, True)
                    
                # Analysis Tab
                with tabs[6]:                     
                    dbmsg.progress(95,'Compute Analysis')                                       
                    # if count >= 1:
                    #     figAnalytic.add_trace(go.Scatter(name="WOB RPM ROP", x=df[Hdr.wob], y=df[Hdr.rpm], marker=dict(size=8, color=df[Hdr.rop],colorbar=dict(title="ROP"),colorscale="Plasma",showscale=False),mode="markers"), row=1, col=1)
                    #     figAnalytic.add_trace(go.Scatter(name="WOB DP ROP", x=df[Hdr.wob], y=df[Hdr.diff], marker=dict(size=8, color=df[Hdr.rop],colorbar=dict(title="ROP",len=0.25, y=0.75),colorscale="Plasma",showscale=True),mode="markers"), row=1, col=2)
                    #     # fig.add_trace(go.Scatter(x=values,y=values,marker=dict(size=16,cmax=39,cmin=0,color=values,colorbar=dict(title="Colorbar"),colorscale="Viridis"),mode="markers"))
                    
                    # if count >= 2:
                    #     figAnalytic.add_trace(go.Scatter(name="WOB RPM ROP", x=df2[Hdr.wob], y=df2[Hdr.rpm], marker=dict(size=8, color=df2[Hdr.rop],colorbar=dict(title="ROP"),colorscale="Plasma",showscale=False),mode="markers"), row=2, col=1)
                    #     figAnalytic.add_trace(go.Scatter(name="WOB DP ROP", x=df2[Hdr.wob], y=df2[Hdr.diff], marker=dict(size=8, color=df2[Hdr.rop],colorbar=dict(title="ROP",len=0.25, y=0.5),colorscale="Plasma",showscale=True,),mode="markers"), row=2, col=2)
                    
                    # if count >= 3:
                    #     figAnalytic.add_trace(go.Scatter(name="WOB RPM ROP", x=df3[Hdr.wob], y=df3[Hdr.rpm], marker=dict(size=8, color=df3[Hdr.rop],colorbar=dict(title="ROP"),colorscale="Plasma",showscale=False),mode="markers"), row=3, col=1)
                    #     figAnalytic.add_trace(go.Scatter(name="WOB DP ROP", x=df3[Hdr.wob], y=df3[Hdr.diff], marker=dict(size=8, color=df3[Hdr.rop],colorbar=dict(title="ROP",len=0.25, y=0.25),colorscale="Plasma",showscale=True),mode="markers"), row=3, col=2)
                    
                    # if count >= 4:
                    #     figAnalytic.add_trace(go.Scatter(name="WOB RPM ROP", x=df4[Hdr.wob], y=df4[Hdr.rpm], marker=dict(size=8, color=df4[Hdr.rop],colorbar=dict(title="ROP"),colorscale="Plasma",showscale=False),mode="markers"), row=4, col=1)
                    #     figAnalytic.add_trace(go.Scatter(name="WOB DP ROP", x=df4[Hdr.wob], y=df4[Hdr.diff], marker=dict(size=8, color=df4[Hdr.rop],colorbar=dict(title="ROP",len=0.25, y=0.0),colorscale="Plasma",showscale=True),mode="markers"), row=4, col=2)
                    
                    
                    # dbmsg.progress(97,'Render Analysis')   
                    # # figAnalytic.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg)
                    # st.plotly_chart(figAnalytic, use_container_width=True, config=config)
                    
                    # Viridis
                    
                # Data Quality
                with tabs[7]:         
                    # dbmsg.info('Render Data') 
                    dbmsg.progress(97,'Render Data')                                      
                    if count >= 1:
                        with st.expander('Data Table1'):
                            # AgGrid(df)
                            # st.write(df.describe(include='all').fillna("").astype("str"))
                            dbmsg.progress(97,'Render Data1') 
                            st.write(df)                               
                            dbmsg.progress(97,'Package Data1')                          
                            csv = convert_df(df)
                            st.download_button("Download Data",csv,"file.csv","text/csv",key='download-csv')
                    
                    if count >= 2:
                        dbmsg.progress(97,'Render Data2') 
                        with st.expander('Data Table2'):
                            # st.write(df2.describe(include='all').fillna("").astype("str"))
                            st.write(df2)                            
                            csv2 = convert_df(df2)
                            st.download_button("Download Data2",csv2,"file2.csv","text/csv",key='download-csv2')
                    
                    if count >= 3:
                        dbmsg.progress(97,'Render Data3') 
                        with st.expander('Data Table3'):
                            # st.write(df3.describe(include='all').fillna("").astype("str"))
                            st.write(df3)                            
                            csv3 = convert_df(df3)
                            st.download_button("Download Data",csv3,"file3.csv","text/csv",key='download-csv3')
                                            
                    if count >= 4:
                        dbmsg.progress(97,'Render Data4')   
                        with st.expander('Data Table4'):
                            # st.write(df4.describe(include='all').fillna("").astype("str"))
                            st.write(df4)                            
                            csv4 = convert_df(df4)
                            st.download_button("Download Data",csv4,"file4.csv","text/csv",key='download-csv4')


                    # https://docs.streamlit.io/library/api-reference/data/st.dataframe
                    # st.dataframe(df,column_config={
                    #     "name": "App name",
                    #     "stars": st.column_config.NumberColumn(
                    #         "Github Stars",
                    #         help="Number of stars on GitHub",
                    #         format="%d ⭐",
                    #     ),
                    #     "url": st.column_config.LinkColumn("App URL"),
                    #     "views_history": st.column_config.LineChartColumn(
                    #         "Views (past 30 days)", y_min=0, y_max=5000
                    #     ),},hide_index=True,)
                    
                    # for i in range(count):
                    #     with st.expander(f'Well {i+1}: Data Quality'):  
                    #         st.write("""The following plot can be used to identify the depth range of each of the logging curves.
                    #         To zoom in, click and drag on one of the tracks with the left mouse button. 
                    #         To zoom back out double click on the plot.""")
                    #         if i == 0:
                    #             well_data = df
                    #         if i == 1:
                    #             well_data = df2
                    #         if i == 2:
                    #             well_data = df3
                    #         if i == 3:
                    #             well_data = df4
                    #         else:                                
                    #             well_data = df
                                
                    #         # if i < 4:
                    #         if i == 0:
                    #             try:
                    #                 data_nan = well_data.notnull().astype('int')
                    #                 # Need to setup an empty list for len check to work
                    #                 curves = []
                    #                 columns = list(well_data.columns)
                    #                 columns.pop(-1) #pop off depth

                    #                 # col1_md, col2_md= st.columns(2)
                    #                 # selection = col1_md.radio('Select all data or custom selection', ('All Data', 'Custom Selection'))
                    #                 # fill_color_md = col2_md.color_picker('Select Fill Colour', '#9D0000')
                    #                 # # top_depth = col3_md.number_input('Top Depth', step=50.0, value=min_depth, min_value=min_depth, max_value=max_depth)
                    #                 # # bottom_depth = col4_md.number_input('Bottom Depth', step=50.0, value=max_depth, min_value=min_depth, max_value=max_depth)
                    #                 curves = columns
                    #                 curve_index = 1
                    #                 fig_missing = make_subplots(rows=1, cols= len(curves), subplot_titles=curves, shared_yaxes=True, horizontal_spacing=0.01)

                    #                 for curve in curves:
                    #                     fig_missing.add_trace(go.Scatter(x=data_nan[curve], y=well_data[Hdr.hole_depth],fill='tozerox',line=dict(width=0), fillcolor='#9D0000'), row=1, col=curve_index)
                    #                     fig_missing.update_xaxes(range=[0, 1], visible=False)
                    #                     fig_missing.update_xaxes(range=[0, 1], visible=False)
                    #                     curve_index+=1
                                    
                    #                 fig_missing.update_layout(height=1000, showlegend=False, yaxis={'title':'DEPTH','autorange':'reversed'})
                    #                 # rotate all the subtitles of 90 degrees
                    #                 for annotation in fig_missing['layout']['annotations']: 
                    #                         annotation['textangle']=-90
                    #                 fig_missing.layout.template='seaborn'
                    #                 st.plotly_chart(fig_missing, use_container_width=True)
                    #             except Exception as e:
                    #                 print(e)
                    #                 st.info(f"error data quality{i+1}: {e}")
                                
                    # if count >= 2:
                    #     with st.expander('Well 2: Data Quality'):  
                    #         st.write("""The following plot can be used to identify the depth range of each of the logging curves.
                    #         To zoom in, click and drag on one of the tracks with the left mouse button. 
                    #         To zoom back out double click on the plot.""")
                    #         try:
                    #             well_data = df2
                    #             data_nan = well_data.notnull().astype('int')
                    #             # Need to setup an empty list for len check to work
                    #             curves = []
                    #             columns = list(well_data.columns)
                    #             columns.pop(-1) #pop off depth

                    #             # col1_md, col2_md= st.columns(2)
                    #             # selection = col1_md.radio('Select all data or custom selection', ('All Data', 'Custom Selection'))
                    #             # fill_color_md = col2_md.color_picker('Select Fill Colour', '#9D0000')
                    #             # # top_depth = col3_md.number_input('Top Depth', step=50.0, value=min_depth, min_value=min_depth, max_value=max_depth)
                    #             # # bottom_depth = col4_md.number_input('Bottom Depth', step=50.0, value=max_depth, min_value=min_depth, max_value=max_depth)
                    #             curves = columns
                    #             curve_index = 1
                    #             fig_missing = make_subplots(rows=1, cols= len(curves), subplot_titles=curves, shared_yaxes=True, horizontal_spacing=0.01)

                    #             for curve in curves:
                    #                 fig_missing.add_trace(go.Scatter(x=data_nan[curve], y=well_data['DEPTH'], 
                    #                     fill='tozerox',line=dict(width=0), fillcolor='#9D0000'), row=1, col=curve_index)
                    #                 fig_missing.update_xaxes(range=[0, 1], visible=False)
                    #                 fig_missing.update_xaxes(range=[0, 1], visible=False)
                    #                 curve_index+=1
                                
                    #             fig_missing.update_layout(height=1000, showlegend=False, yaxis={'title':'DEPTH','autorange':'reversed'})
                    #             # rotate all the subtitles of 90 degrees
                    #             for annotation in fig_missing['layout']['annotations']: 
                    #                     annotation['textangle']=-90
                    #             fig_missing.layout.template='seaborn'
                    #             st.plotly_chart(fig_missing, use_container_width=True)
                    #         except Exception as e:
                    #             print(e)
                    #             st.info(f"error data quality1: {e}")
                        
                                
        except Exception as e:
            print(e)
            st.info(f"error dataframe1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
            
            
    except Exception as e:
        print(e)
        st.info(f"error dataframe loop: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
        
    # # Plot demo data    
    # try:
    #     for dfd, chart_select in zip(dfs_to_compare, chart_selects):
    #         dfd.rename(columns=final_column_name_map, inplace=True)
    #         num = df._get_numeric_data()
    #         num[num < 0] = np.nan
    #         df.reset_index(inplace=True)
    #         count = 1
    #         # with stcol0.spinner('Finding time domain...'):
    #         try:
    #             df[Hdr.dt] = df['YYYY/MM/DD']+df['HH:MM:SS']
                            
    #             # df[Hdr.dtn] = df['YYYY/MM/DD']+df['HH:MM:SS']
    #             df[Hdr.dtn]  = pd.DataFrame({'time': [pd.to_datetime(df['YYYY/MM/DD']+df['HH:MM:SS']).astype(int)/ 10**9]})
    #             # df_unix_sec = pd.to_datetime(df['time']).astype(int)/ 10**9
                
    #             # df[Hdr.dtn] = df[Hdr.dtn] + pd.to_datetime(df['YYYY/MM/DD']+df['HH:MM:SS'])            
    #             # ds1_dtn_min = df[Hdr.dtn].min
    #             # ds1_dtn_max = df[Hdr.dtn].max
    #             # df[Hdr.dtn] = df[Hdr.dt] - df[hdr.dtn] - ds1_dtn_min
                
    #         except Exception as e:
    #             print(e)
    #         try:
    #             df2.rename(columns=final_column_name_map, inplace=True)
    #             num = df2._get_numeric_data()
    #             num[num < 0] = np.nan
    #             df2.reset_index(inplace=True)
    #             count = 2
    #             try:
    #                 df2[Hdr.dt] = df2['YYYY/MM/DD']+df2['HH:MM:SS']
    #             except Exception as e:
    #                 print(e)
    #             try:
    #                 df3.rename(columns=final_column_name_map, inplace=True)
    #                 num = df3._get_numeric_data()
    #                 num[num < 0] = np.nan
    #                 df3.reset_index(inplace=True)
    #                 count = 3
    #                 try:
    #                     df3[Hdr.dt] = df3['YYYY/MM/DD']+df3['HH:MM:SS']
    #                 except Exception as e:
    #                     print(e)
    #                 try:
    #                     df4.rename(columns=final_column_name_map, inplace=True)
    #                     num = df4._get_numeric_data()
    #                     num[num < 0] = np.nan
    #                     df4.reset_index(inplace=True)
    #                     count = 4
    #                     try:
    #                         df4[Hdr.dt] = df4['YYYY/MM/DD']+df4['HH:MM:SS']
    #                     except Exception as e:
    #                         print(e)
    #                 except Exception as e:
    #                     print(e)
    #             except Exception as e:
    #                 print(e)
    #         except Exception as e:
    #             print(e)  
    #         if chart_select != DropOptions.yes:
    #             continue 
            
    #     # Plot dataset 1
    #     try:
    #         if count >= 1:
    #             # agree = stcol2.checkbox('Show Outliers')
    #             rad1 = stcol2.radio('Data Outliers:', options = ['Show', 'Cleaned'])
    #             rad2 = stcol2.radio('Chart Size:', options = ['Fit to Screen', 'Full Height'])
    #             if rad2 == 2:
    #                 fig.update_layout(height= 5000)
                
    #             if rad1 == 1:
    #                 with stcol0.spinner('Recalculating...'):
    #                     q_low_rop = df[Hdr.rop].quantile(0.01)
    #                     q_high_rop = df[Hdr.rop].quantile(0.99)
    #                     avg_rop = df[Hdr.rop].mean 
    #                     df = df[(df[Hdr.rop] < q_high_rop) & (df[Hdr.rop] > q_low_rop)]
    #                     q_low_wob = df[Hdr.wob].quantile(0.01)
    #                     q_high_wob = df[Hdr.wob].quantile(0.99)
    #                     df = df[(df[Hdr.wob] < q_high_wob) & (df[Hdr.wob] > q_low_wob)]
    #                     q_low_rpm = df[Hdr.rpm].quantile(0.01)
    #                     q_high_rpm = df[Hdr.rpm].quantile(0.99)
    #                     df = df[(df[Hdr.rpm] < q_high_rpm) & (df[Hdr.rpm] > q_low_rpm)]
    #                     q_low_diff = df[Hdr.diff].quantile(0.01)
    #                     q_high_diff = df[Hdr.diff].quantile(0.99)
    #                     df = df[(df[Hdr.diff] < q_high_diff) & (df[Hdr.diff] > q_low_diff)]
            
    #             with st.spinner('Loading channels...'):
    #                 fig.add_trace(go.Scatter(name="first", x=df[Hdr.dt], y=df[Hdr.hole_depth], line= {'color':'#121212'}), row=1, col=1)
    #                 fig.add_trace(go.Scatter(name="first", x=df[Hdr.rop], y=df[Hdr.hole_depth], line= {'color':'#008C9A'}), row=1, col=2)
    #                 # fig.add_annotation(xref='x domain',yref='y domain',x=0.01,y=0.9,text=avg_rop,showarrow=False,row=1, col=2)
    #                 fig.add_trace(go.Scatter(name="first", x=df[Hdr.wob], y=df[Hdr.hole_depth],line= {'color':'#8ACFDD'}), row=1, col=3)
    #                 fig.add_trace(go.Scatter(name="first", x=df[Hdr.rpm], y=df[Hdr.hole_depth],line= {'color':'#7A7D81'}), row=1, col=4)
    #                 # fig.add_trace(go.Scatter(name="first", x=df[Hdr.diff], y=df[Hdr.hole_depth],line= {'color':g4}), row=1, col=5)
    #                 fig.add_trace(go.Scatter(name="first", x=df[Hdr.diff], y=df[Hdr.hole_depth],line= {'color':'#98C007'}), row=1, col=5)
            
    #         if count == 1:                    
    #             # stcol1.write(avg_rop)  
    #             with tabs[0]:
    #                 stcol1.write(fig) 
    #     except Exception as e:
    #         print(e)
    # except Exception as e:
    #     print(e)
                
                    
    # with tabs[0]:
    #     stcol1.write(fig)                      
    #     st.write(fig) 
        
    #try:
            #fig4 = make_subplots(rows= 1, cols = 3, shared_yaxes=True)
            #fig4.update_xaxes(title_text="Gamma Ray", row=1, col=1)
            #fig4.update_xaxes(title_text="DTC", row=1, col=2)
            #fig4.update_xaxes(title_text='DTS', row=1, col=3)

            #fig4.update_yaxes(title_text="Depth", row=1, col=1, autorange='reversed')

            #fig4.add_trace(go.Scatter(x=df5['GR'], y=df5['DEPT']),  row=1, col =1)
            #fig4.add_trace(go.Scatter(x=df5['DTC'], y=df5['DEPT']),  row=1,col = 2)
            #fig4.add_trace(go.Scatter(x=df5['DTS'], y=df5['DEPT']), row=1,col =3)
            #st.write(fig4)
    #except Exception as e:
        #print(e)
        
        
        
    # # Viridis
    # try:
    #     fig2 = px.scatter(df, x="Differential Pressure (psi)", y='Weight on Bit (klbs)', color='Rate Of Penetration (ft_per_hr)',
    #                     color_continuous_scale=px.colors.sequential.Plasma,labels={
    #                     "Differential Pressure (psi)": "DIFF",
    #                     'Weight on Bit (klbs)': "WOB",
    #                     'Rate Of Penetration (ft_per_hr)': "ROP"
    #                 })
        
    #     fig2.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg)
    #     with tabs[6]:
    #         # dbmsg.info('Render Analysis')
    #         # dbmsg.progress(99,'Render Analysis') 
    #         st.write(fig2)
    # except Exception as e:
    #     print(e)

    # try:
    #     fig3 = px.scatter(df, x='Rotary RPM (RPM)',  y='Weight on Bit (klbs)', color='Rate Of Penetration (ft_per_hr)',
    #                       color_continuous_scale=px.colors.sequential.Plasma,labels={
    #                     "Rotary RPM (RPM)": "RPM",
    #                     'Weight on Bit (klbs)': "WOB",
    #                     'Rate Of Penetration (ft_per_hr)': "ROP"
    #                 })
    #     fig3.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor=gbg)
    #     with tabs[6]:
    #         st.write(fig3)
    # except Exception as e:
    #     print(e)
    
    
    dbmsg.success('Ready') 
    
    st.sidebar.markdown('<a href="mailto:ccasad@ulterra.com?subject=ADIRA Help & Feedback&body=Hey Chris, ADIRA rocks!"><button style="color:#43c6db;background-color:white;text-decoration:none;border-radius:4px;border:#43c6db;padding:10px 24px;">Email for Help & Feedback</button></a>', unsafe_allow_html=True)
    
    # function to replace text in pptx first slide with selected filters
    def replace_text(replacements, shapes):
        """function to replace text on a PowerPoint slide. Takes dict of {match: replacement, ... } and replaces all matches"""
        for shape in shapes:
            for match, replacement in replacements.items():
                if shape.has_text_frame:
                    if (shape.text.find(match)) != -1:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            whole_text = "".join(run.text for run in paragraph.runs)
                            whole_text = whole_text.replace(str(match), str(replacement))
                            for idx, run in enumerate(paragraph.runs):
                                if idx != 0:
                                    p = paragraph._p
                                    p.remove(run._r)
                            if bool(paragraph.runs):
                                paragraph.runs[0].text = whole_text
                                
                                
    # https://pptx-generator.streamlit.app/
    def create_adira_pptx():
        
        pptx = 'data/template.pptx'
        prs = Presentation(pptx)
        
        # declare positional variables
        width = Inches(8)
        left = Inches(2.5)
        top = Inches(1)
        
         # get stock info
        name = 'Ulterra'
        # sector = ticker.summary_profile[user_input]['sector']
        # industry = ticker.summary_profile[user_input]['industry']
        # employees = ticker.summary_profile[user_input]['fullTimeEmployees']
        # country = ticker.summary_profile[user_input]['country']
        # city = ticker.summary_profile[user_input]['city']
        # website = ticker.summary_profile[user_input]['website']
        # summary = ticker.summary_profile[user_input]['longBusinessSummary']
        # logo_url = 'https://logo.clearbit.com/' + website

        
        # try: 
        #     # declare pptx variables
        #     first_slide = prs.slides[0]
        #     # second_slide = prs.slides[1]
        #     shapes_1 = []
        #     # shapes_2 = []
        #     index_to_drop = []

        #     # create lists with shape objects
        #     for shape in first_slide.shapes:
        #         shapes_1.append(shape)

        #     # for shape in second_slide.shapes:
        #     #     shapes_2.append(shape)

        #     # initiate a dictionary of placeholders and values to replace
        #     replaces_1 = {
        #         '{company}': name,
        #         # '{date}': today
        #     }

        #     replaces_2 = {
        #         '{c}': name,
        #         # '{s}': sector,
        #         # '{i}': industry,
        #         # '{co}': country,
        #         # '{ci}': city,
        #         # '{ee}': "{:,}".format(employees),
        #         # '{w}': website,
        #         # '{summary}': summary
        #     }

        #     # run the function to replace placeholders with values
        #     replace_text(replaces_1, shapes_1)
        #     # replace_text(replaces_2, shapes_2)
        
        # except Exception as e:
        #     print(e)
        #     st.info(f"error ppt0: {e}")
            
        lyt=prs.slide_masters[0].slide_layouts[6] # choosing a slide layout
        
        save_bar.progress(10, text='Creating template..')    
        dbmsg.info('Creating template..') 
        
        slide1=prs.slides[0]  
        # txBox = slide1.shapes.add_textbox(Inches(4), Inches(4.75), Inches(2), Inches(1))
        txBox = slide1.shapes.add_textbox(left=Inches(4.33), top=Inches(4.2), width=Inches(1.5),height=Inches(0.25))
        tf = txBox.text_frame
        tf.text = f'Prepared {datetime.date.today()}'
        tf.paragraphs[0].font.size = Pt(10)
            
        # try:
        #     slide1=prs.slides.add_slide(lyt) # adding a slide
        #     title=slide1.shapes.title # assigning a title
        #     subtitle=slide1.placeholders[1] # placeholder for subtitle
        #     title.text="Hey,This is a Slide! How exciting!" # title
        #     subtitle.text="ulterra" # subtitle    
        # except Exception as e:
        #     print(e)
        #     st.info(f"error ppt00: {e}")
        
        # txBox = slide.shapes.add_textbox(left, top, width, height)
        # tf = txBox.text_frame
        # # Title
        # p = tf.add_paragraph()
        # p.text = 'Analysis'
        # p.font.size = Pt(18)
        # # First bullet point
        # p = tf.add_paragraph()
        # p.text = '• {} have been prepared during the week'.format(total_lines)
        # p.level = 1
        # # Second bullet point
        # p = tf.add_paragraph()
        # p.text = '• {} has been the busiest day with {} prepared'.format(busy_day, max_lines)
        # p.level = 1
        
        # Pg2
        try:
            slide2=prs.slides.add_slide(lyt) # adding a slide 
            txBox = slide2.shapes.add_textbox(Inches(0.87), Inches(0.4), Inches(1), Inches(0.25))
            tf = txBox.text_frame
            tf.text = f'OVERVIEW'
            tf.paragraphs[0].font.size = Pt(8)
            # tf = txBox.text_frame
            # p = tf.add_paragraph()
            # p.text = f'DATA CHANNELS'
            # # p.font = 'Lato Light' 
            # p.font.size = Pt(8)  
            
            save_bar.progress(30, text='Writing Overview')
            dbmsg.info('Writing Overview') 
                
            # Metrics
            try:
                if count >= 1:
                    filename = "figmetric1.png"
                    # figmetric1.write_image(filename, engine='kaleido') 
                    pio.write_image(figmetric1, filename, scale=2, width=412, height=132) 
                    slide2.shapes.add_picture(filename, left=Inches(0.45), top=Inches(0.79), width=Inches(2.15))    
            except NameError:
                pass        
            try:    
                if count >= 2:
                    filename = "figmetric2.png"
                    pio.write_image(figmetric2, filename, scale=2, width=412, height=132) 
                    slide2.shapes.add_picture(filename, left=Inches(2.77), top=Inches(0.79), width=Inches(2.15)) 
            except NameError:
                pass      
            try:   
                if count >= 3:
                    filename = "figmetric3.png"
                    pio.write_image(figmetric3, filename, scale=2, width=412, height=132) 
                    slide2.shapes.add_picture(filename, left=Inches(5.08), top=Inches(0.79), width=Inches(2.15)) 
            except NameError:
                pass     
            try:   
                if count >= 4: 
                    filename = "figmetric4.png"
                    pio.write_image(figmetric4, filename, scale=2, width=412, height=132) 
                    slide2.shapes.add_picture(filename, left=Inches(7.4), top=Inches(0.79), width=Inches(2.15))
            except NameError:
                pass       
            
            #  Depth vs Time
            filename = "figdvd.png"
            pio.write_image(figDVD, filename, scale=2, width=920, height=1178)             
            slide2.shapes.add_picture(filename, left=Inches(0.1), top=Inches(1.36), width=Inches(4.8))        
            # Box whisker summary        
            filename = "figbxrop.png"
            pio.write_image(figbxrop, filename, scale=2, width=920, height=1178)              
            slide2.shapes.add_picture(filename, left=Inches(4.9), top=Inches(1.36), width=Inches(4.8))      
        except Exception as e:
            print(e)
            st.info(f"error ppt2: {e}")
            
        # Pg3
        try:
            slide3=prs.slides.add_slide(lyt) # adding a slide         
            txBox = slide3.shapes.add_textbox(Inches(0.87), Inches(0.4), Inches(1), Inches(0.25))
            tf = txBox.text_frame
            tf.text = f'DATA CHANNELS'
            tf.paragraphs[0].font.size = Pt(8)
            # p = tf.add_paragraph()
            # p.text = f'DATA CHANNELS'
            # p.font = 'Lato Light' 
            # p.font.size = Pt(8) 
            
            save_bar.progress(40, text='Loading Channels')
            dbmsg.info('Loading Channels') 
            
            filename = "fig.png"
            pio.write_image(fig, filename, scale=2, width=1920, height=1278)
            # pio.write_image(fig, filename, scale=1, width=WIDTH, height=HEIGHT)
            slide3.shapes.add_picture(filename, left=Inches(0), top=Inches(0.65), width=Inches(10))      
        except Exception as e:
            print(e)
            st.info(f"error ppt3: {e}") 
        
            
        # Pg4 Depth vs Inclination
        try:    
            slide4=prs.slides.add_slide(lyt) # adding a slide          
            txBox = slide4.shapes.add_textbox(Inches(0.87), Inches(0.4), Inches(1), Inches(0.25))            
            tf = txBox.text_frame
            tf.text = f'INCLINATION'
            tf.paragraphs[0].font.size = Pt(8)
            # tf = txBox.text_frame
            # p = tf.add_paragraph()
            # p.text = f'INCLINATION'
            # p.font = 'Lato Light' 
            # p.font.size = Pt(8)   
            
            save_bar.progress(50, text='Finding Inclination')
            dbmsg.info('Finding Inclination') 
             
            filename = "figdvdi.png" 
            pio.write_image(figDVDI, filename, scale=2, width=960, height=658)
            # pio.write_image(figDVDI, filename, scale=1, width=WIDTH, height=HEIGHT)  
            slide4.shapes.add_picture(filename, left=Inches(0), top=Inches(0.65), width=Inches(10))     
        except Exception as e:
            print(e)
            st.info(f"error ppt4: {e}") 
        
            
        # Pg5 
        try:  
            slide5=prs.slides.add_slide(lyt) # adding a slide          
            txBox = slide5.shapes.add_textbox(Inches(0.87), Inches(0.4), Inches(1), Inches(0.25)).text_frame
            txBox.text = f'SURVEY'
            txBox.paragraphs[0].font.size = Pt(8)
            # tf = txBox.text_frame
            # p = tf.add_paragraph()
            # p.text = f'SURVEY'
            # p.font = 'Lato Light' 
            # p.font.size = Pt(8)  
            
            save_bar.progress(60, text='Plotting Trajectory')
            dbmsg.info('Plotting Trajectory') 
            
            filename = "fig3DMap.png" 
            pio.write_image(fig3DMap, filename, scale=2, width=960, height=658)            
            # pio.write_image(fig3DMap, filename, scale=1, width=1000, height=800)  
            slide5.shapes.add_picture(filename, left=Inches(0), top=Inches(0.65), width=Inches(10))
                    
        except Exception as e:
            print(e)
            st.info(f"error ppt5: {e}") 
        
            
        # Pg6 ROTATION / SLIDING
        try:  
            slide6=prs.slides.add_slide(lyt) # adding a slide          
            txBox = slide6.shapes.add_textbox(Inches(0.87), Inches(0.4), Inches(1), Inches(0.25))           
            tf = txBox.text_frame
            tf.text = f'ROTATION / SLIDING'
            tf.paragraphs[0].font.size = Pt(8)
            # tf = txBox.text_frame
            # p = tf.add_paragraph()
            # p.text = f'ROTATION / SLIDING'
            # p.font = 'Lato Light' 
            # p.font.size = Pt(8)  
            
            save_bar.progress(70, text='Mathing Slides')
            dbmsg.info('Mathing Slides') 
            
            filename = "figRotSlide.png" 
            pio.write_image(figRotSlide, filename, scale=2, width=960, height=658)            
            # pio.write_image(fig3DMap, filename, scale=1, width=1000, height=800)  
            slide6.shapes.add_picture(filename, left=Inches(0), top=Inches(0.65), width=Inches(10))
                    
        except Exception as e:
            print(e)
            st.info(f"error ppt6: {e}") 
        
            
        # Pg7 ANALYSIS
        try:  
            slide7=prs.slides.add_slide(lyt) # adding a slide          
            txBox = slide7.shapes.add_textbox(Inches(0.87), Inches(0.4), Inches(1), Inches(0.25))           
            tf = txBox.text_frame
            tf.text = f'ANALYSIS'
            tf.paragraphs[0].font.size = Pt(8)
            # tf = txBox.text_frame
            # p = tf.add_paragraph()
            # p.text = f'ANALYSIS'
            # p.font = 'Lato Light' 
            # p.font.size = Pt(8) 
            
            save_bar.progress(82, text='Compiling..')
            dbmsg.info('Compiling...') 
             
            filename = "figAnalytic.png" 
            pio.write_image(figAnalytic, filename, scale=2, width=1000, height=800)        
            slide7.shapes.add_picture(filename, left=Inches(0.1), top=Inches(0.65), width=Inches(9.25)) 
            # filename = "fig3.png" 
            # pio.write_image(fig3, filename, scale=2, width=446, height=658)         
            # slide7.shapes.add_picture(filename, left=Inches(5), top=Inches(0.65), width=Inches(4.65))
                    
        except Exception as e:
            print(e)
            st.info(f"error ppt7: {e}") 
        
            
        # # Pg8 DATA QUALITY
        # try:  
        #     slide8=prs.slides.add_slide(lyt) # adding a slide          
        #     txBox = slide8.shapes.add_textbox(Inches(0.87), Inches(0.4), Inches(1), Inches(0.25))        
        #     tf = txBox.text_frame
        #     tf.text = f'DATA QUALITY'
        #     tf.paragraphs[0].font.size = Pt(8)
        #     # tf = txBox.text_frame
        #     # p = tf.add_paragraph()
        #     # p.text = f'ANALYSIS'
        #     # p.font = 'Lato Light' 
        #     # p.font.size = Pt(8)  
        #     save_bar.progress(90, text='Quality Metrics')
        #     dbmsg.info('Quality Metrics') 
            
        #     filename = "fig_missing.png" 
        #     pio.write_image(fig_missing, filename, scale=2, width=960, height=658)            
        #     # pio.write_image(fig3DMap, filename, scale=1, width=1000, height=800)  
        #     slide8.shapes.add_picture(filename, left=Inches(0), top=Inches(0.65), width=Inches(10)) 
                     
        # except Exception as e:
        #     print(e)
        #     st.info(f"error ppt8: {e}") 
        
        # Close out
        binary_output = BytesIO()
        prs.save(binary_output)
        return binary_output
            

    # @st.cache_data
    # @st.cache_data(hash_funcs={builtins.function: create_adira_report})
    def create_adira_report():
        WIDTH = 792
        HEIGHT = 612
        pagetitle_x = 75
        pagetitle_y = 41
        gaf_y = 50
        
        try:
            pdf = adiraPDF(orientation = 'L', unit = 'pt', format='Letter') # Letter (215.9 × 279.4 mm) 1pt is 1/72in        
            pdf.set_auto_page_break(False, 0)              
        except Exception as e:
            print(e)
            st.info(f"error pdf4: {e}")      
        # pdf.add_font('Roboto','','Roboto.php')
        
        try:
            # Pg1
            pdf.add_page()
            pdf.set_font('Arial', '', 8) 
            # pdf.ln(20)
            # pdf.cell(pagetitle_x, 30, f'REPORT',0,1) 
            pdf.text(pagetitle_x, pagetitle_y, f'REPORT')     
            
            # Metric Banners
            if count >= 1:
                filename = "figmetric1.jpg"
                # figmetric1.write_image(filename, engine='kaleido') 
                pio.write_image(figmetric1, filename)        
                pdf.image("figmetric1.jpg", 50, 60, w=150)        
            if count >= 2:
                filename = "figmetric2.jpg"
                # figmetric2.write_image(filename, engine='kaleido') 
                pio.write_image(figmetric2, filename)               
                pdf.image("figmetric2.jpg", 230, 60, w=150)            
            if count >= 3:
                filename = "figmetric3.jpg"
                # figmetric3.write_image(filename, engine='kaleido')    
                pio.write_image(figmetric3, filename)            
                pdf.image("figmetric3.jpg", 410, 60, w=150)            
            if count >= 4:
                filename = "figmetric4.jpg"
                # figmetric4.write_image(filename, engine='kaleido')
                pio.write_image(figmetric4, filename)                
                pdf.image("figmetric4.jpg", 590, 60, w=150)
            
            #  Depth vs Time
            filename = "figdvd.jpg"
            # figDVD.write_image(filename, engine='kaleido') 
            pio.write_image(figDVD, filename)               
            pdf.image("figdvd.jpg", 5, 85, h= HEIGHT-100)        
            # Box whisker summary        
            filename = "figbxrop.jpg"
            # figbxrop.write_image(filename, engine='kaleido') 
            pio.write_image(figbxrop, filename)               
            pdf.image("figbxrop.jpg", 300, 85, h= HEIGHT-100)
            
            
            # Pg2
            pdf.add_page()       
            pdf.set_font('Arial', '', 8) 
            # pdf.cell(pagetitle_x, 30, f'CHANNELS',0,1)   
            pdf.text(pagetitle_x, pagetitle_y, f'CHANNELS')       
            # Channels
            filename = "figdvd.jpg"
            figDVD.write_image(filename, engine='kaleido')            
            pdf.image("figdvd.jpg", 5, gaf_y, w= 250, h= HEIGHT-150)
                    
            filename = "fig.jpg"
            # fig.write_image(filename, engine='kaleido')        
            pio.write_image(fig, filename, scale=1, width=WIDTH, height=HEIGHT) 
            pdf.image("fig.jpg", 300, gaf_y, w= WIDTH-320, h= HEIGHT-150)     
        except Exception as e:
            print(e)
            st.info(f"error pdf3: {e}")   
        
        
        # pdf.cell(10, 600, f'Created by ADIRA v2.34')
        # pdf.image('Ulterra_teal_250px.jpg', WIDTH-60, HEIGHT-20, 70)
        try:
            # Pg3        
            pdf.add_page()       
            pdf.set_font('Arial', '', 8) 
            # pdf.cell(pagetitle_x, 30, f'INCLINATION',0,1) 
            pdf.text(pagetitle_x, pagetitle_y, f'INCLINATION')     
            # Depth vs Inclination
            filename = "figdvdi.jpg"
            # figDVDI.write_image(filename, engine='kaleido')  
            pio.write_image(figDVDI, filename, scale=1, width=WIDTH, height=HEIGHT)      
            pdf.image("figdvdi.jpg", 30, gaf_y, w= WIDTH-60, h= HEIGHT-100)        
            
            
            # Pg4     
            pdf.add_page()       
            pdf.set_font('Arial', '', 8) 
            # pdf.cell(pagetitle_x, 30, f'SURVEY',0,1) 
            pdf.text(pagetitle_x, pagetitle_y, f'SURVEY')     
            # Depth vs Inclination
            filename = "fig3DMap.jpg"
            # fig3DMap.write_image(filename, engine='kaleido')              
            pio.write_image(fig3DMap, filename, scale=1, width=1000, height=800)      
            pdf.image("fig3DMap.jpg", 5, gaf_y, h= HEIGHT-100)
                
                        
            # Pg5  
            pdf.add_page()       
            pdf.set_font('Arial', '', 8) 
            # pdf.cell(pagetitle_x, 30, f'ROTATION / SLIDING',0,1) 
            pdf.text(pagetitle_x, pagetitle_y, f'ROTATION / SLIDING')     
            # Depth vs Inclination
            filename = "figRotSlide.jpg"
            # figRotSlide.write_image(filename, engine='kaleido')  
            pio.write_image(figRotSlide, filename, scale=1, width=1000, height=800)       
            pdf.image("figRotSlide.jpg", 5, gaf_y, h= HEIGHT-100)    

            
            # Pg6
            pdf.add_page()
            pdf.set_font('Arial', '', 8)  
            # pdf.cell(150, 30, f'ANALYSIS') 
            pdf.text(pagetitle_x, pagetitle_y, f'ANALYSIS')     
            
            # filename = "fig2.jpg"
            # fig2.write_image(filename, engine='kaleido')        
            # pdf.image("fig2.jpg", 5, gaf_y, h= 350)            
            filename = "figAnalytic.jpg"
            figAnalytic.write_image(filename, engine='kaleido')        
            pdf.image("figAnalytic.jpg", 400, gaf_y, h= 350)
        
        except Exception as e:
            print(e)
            st.info(f"error pdf2: {e}")
        # pdf.cell(5, 205, f'Created by ADIRA v2.34')
        
        # Close out
        pdf.output("output1.pdf", 'f')
        with open("output1.pdf", "rb") as pdf_file:
            PDFbyte = pdf_file.read()
        return PDFbyte
    
    # @st.cache_data(hash_funcs={builtins.function: create_adira_poster})
    def create_adira_poster():        
        WIDTH = 1684
        HEIGHT = 2384
        
        pdf = adiraPDF('P', 'pt', (1684, 2384)) # A1 (594 × 841mm or 23.39 × 33.11 inch) 1pt is 1/72in     
        pdf.set_auto_page_break(False)
        
        pdf.add_page()
        pdf.set_font('Arial', '', 24)  
        pdf.image('adira-logo-name.jpg', 5, 5, 200)
        pdf.cell(5, 205, f'REPORT')        
        
        filename = "fig.jpg"
        fig.write_image(filename, engine='kaleido')        
        pdf.image("fig.jpg", 5, 230, WIDTH-20, HEIGHT-400)
        
        pdf.cell(5, 205, f'Created by ADIRA v2.34')
        pdf.image('Ulterra_teal_250px.jpg', WIDTH-260, HEIGHT-80)
        
        #  Close out
        pdf.output("ouptut2.pdf", 'f')
        with open("ouptut2.pdf", "rb") as pdf_file:
            PDFbyte = pdf_file.read()
        return PDFbyte
    

    if count >= 1:
        # # Create an in-memory buffer
        # buffer = io.BytesIO()
        # # Save the figure as a pdf to the buffer
        # figDVD.write_image(file=buffer, format="pdf")
        # # Download the pdf from the buffer
        # hed6.download_button(
        #     label="Generate Report",
        #     data=buffer,
        #     file_name="figure.pdf",
        #     mime="application/pdf",
        # )          
        
        # # Test code for creating a report
        # if hed6.button('Generate Report2'):
        #     figures = [figDVD]                    
        #     images_html = figure_to_base64(figures)
        #     report_html = create_html_report("template.html", images_html)
        #     convert_html_to_pdf(report_html, "adira_report.pdf")
        
        # if hed6.button('View Report2'):
        #     combine_plotly_figs_to_html((figDVD,fig,fig2,fig3), "adira_report2.pdf", separator=',', auto_open=True)
                
        
        # btndl = hed6.button('Prepare Reports')
        
        # if btndl:
        #     try:
        #         hed6.download_button(
        #             label="Create Report",
        #             data=reportpdf,
        #             file_name="Adira Report.pdf",
        #             mime="application/octet-stream",
        #         )
        #     except Exception as e:
        #         print(e)
        #         st.info(f"error pdf: {e}")
                
        #     try:
        #         hed6.download_button(
        #             label="Create Poster",
        #             data=posterpdf,
        #             file_name="Adira Poster.pdf",
        #             mime="application/octet-stream",
        #         )
        #     except Exception as e:
        #         print(e)
        #         st.info(f"error pdf: {e}")
        
        # Reports tab
        with tabs[8]:
            formpdf = st.form("template_form")
            # name = formpdf.text_input("Your name")
            reporttype = formpdf.selectbox(
                "Choose Type",
                ["PowerPoint", "Full Report", "Poster" ],
                index=0,
            )
            # grade = formpdf.slider("Grade", 1, 100, 60)
            
            with st.spinner('Request...'): 
                submit = formpdf.form_submit_button("Request Report")               
            
            if submit:                
                save_bar = st.progress(0, text='Initializing..')
                with st.spinner('Constructing...'):                    
                    # dbmsg.info('Constructing...') 
                    dbmsg.progress(0, text='Constructing..')
                    if reporttype == 'PowerPoint':
                        try:
                            reportppt = create_adira_pptx()
                            try:
                                dbmsg.info('Preparing Download')                                 
                                save_bar.progress(95, text='Preparing Download')
                                st.download_button(
                                    label="Download PowerPoint",
                                    data=reportppt,
                                    file_name=f"Adira_PPTX_{datetime.datetime.now()}.pptx",
                                    mime="application/octet-stream",
                                ) 
                                save_bar.progress(100, text='Ready')
                                dbmsg.success('Ready') 
                            except Exception as e:
                                print(e)
                                st.info(f"error pptx dl1: {e}")                             
                                             
                        except Exception as e:
                            print(e)
                            st.info(f"error pptx: {e}")
                            
                    elif reporttype == 'Full Report':
                        try:                            
                            reportpdf = create_adira_report()
                            st.download_button(
                                label="Download Report",
                                data=reportpdf,
                                file_name=f"Adira_Report_{datetime.datetime.now()}.pdf",
                                mime="application/octet-stream",
                            )
                        except Exception as e:
                            print(e)
                            st.info(f"error pdf: {e}")
                    
                    elif reporttype == 'Poster':
                        try:
                            posterpdf = create_adira_poster()
                            st.download_button(
                                label="Download Poster",
                                data=posterpdf,
                                file_name=f"Adira_Poster_{datetime.datetime.now()}.pdf",
                                mime="application/octet-stream",
                            )
                        except Exception as e:
                            print(e)
                            st.info(f"error pdf: {e}")
                            

        
    #  Download sample data
    if uploaded_file is None:
        st.sidebar.write('Download Sample Data:')
        stsbc1, stsbc2 = st.sidebar.columns(2)
        with stsbc1:
            with open("data/WellDemo 24-23.csv", "rb") as file:
                stsbc1.download_button(
                    label="Sample1",
                    data=file,
                    file_name='adira_sample_data1.csv',
                    mime='text/csv')
        with stsbc2:
            with open("data/WellDemo 27-34.csv", "rb") as file:
                stsbc2.download_button(
                    label="Sample2",
                    data=file,
                    file_name='adira_sample_data2.csv',
                    mime='text/csv')
    

        
        # with open("ouptut1.pdf", "rb") as pdf_file:
        #     PDFbyte = pdf_file.read()

        # hed6.download_button(label="Export_Report5",
        #                     data=PDFbyte,
        #                     file_name="test.pdf",
        #                     mime='application/octet-stream')
        
        # #  Pylatex example
        # def gen_pdf():
        #     # initialize a Document            
        #     geometry_options = {"right": "2cm", "left": "2cm"}
        #     doc = Document('tmppdf', geometry_options=geometry_options)

        #     # this is a sample of a document, you could add more sections
        #     with doc.create(MiniPage(align='c')):
        #         doc.append(LargeText("Title"))
            
        #     # plot saved as png
        #     image_filename = 'path/to/image.png'
        #     with doc.create(Section('Section Title')):
        #         with doc.create(Figure(position='h!')) as fig:
        #             fig.add_image(image_filename,
        #                         width='15cm')
        #             fig.add_caption('Caption for the figure')
            
        #     # generate the pdf file
        #     # this file will be generated at the Streamlit server side under the name *tmppdf.pdf*
        #     doc.generate_pdf("tmppdf", clean_tex=False, compiler='pdfLaTex')
            
        #     # Open the file and read it as bytes
        #     with open("tmppdf.pdf", "rb") as pdf_file:
        #         PDFbyte = pdf_file.read()
        
        #     #  return the bytes object created *PDFbyte*since the data argument in the download button must be string or bytes or file
        #     return PDFbyte

        # # the download button will get the generated file stored in Streamlit server side, and download it at the user's side
        # hed6.download_button(label="Download PDF Report",
        #                 key='download_pdf_btn',
        #                 data=gen_pdf(),
        #                 file_name='name_of_your_file.pdf', # this might be changed from browser after pressing on the download button
        #                 mime='application/octet-stream',)

    # st.success('Done!')
    
    
    
    # # # # Create PDF
    
        
    # def create_title(day, pdf):
    #     # Unicode is not yet supported in the py3k version; use windows-1252 standard font
    #     pdf.set_font('Arial', '', 24)  
    #     pdf.ln(60)
    #     pdf.write(5, f"Test Adira Report")
    #     pdf.ln(10)
    #     pdf.set_font('Arial', '', 16)
    #     pdf.write(4, f'{day}')
    #     pdf.ln(5)

    # def create_analytics_report(day=TEST_DATE, filename="report.pdf"):
    #     pdf = FPDF() # A4 (210 by 297 mm)

    #     states = ['Massachusetts', 'New Hampshire']

    #     ''' First Page '''
    #     pdf.add_page()
    #     pdf.image("./resources/letterhead_cropped.png", 0, 0, WIDTH)
    #     create_title(day, pdf)

    #     plot_usa_case_map("./tmp/usa_cases.png", day=day)
    #     prev_days = 250
    #     plot_states(states, days=prev_days, filename="./tmp/cases.png", end_date=day)
    #     plot_states(states, days=prev_days, mode=Mode.DEATHS, filename="./tmp/deaths.png", end_date=day)

    #     pdf.image("./tmp/usa_cases.png", 5, 90, WIDTH-20)
    #     pdf.image("./tmp/cases.png", 5, 200, WIDTH/2-10)
    #     pdf.image("./tmp/deaths.png", WIDTH/2, 200, WIDTH/2-10)

    #     # ''' Second Page '''
    #     # pdf.add_page()

    #     # plot_daily_count_states(states, day=day, filename="./tmp/cases_day.png")
    #     # plot_daily_count_states(states, day=day, mode=Mode.DEATHS, filename="./tmp/deaths_day.png")
    #     # pdf.image("./tmp/cases_day.png", 5, 20, WIDTH/2-10)
    #     # pdf.image("./tmp/deaths_day.png", WIDTH/2, 20, WIDTH/2-10)

    #     # prev_days = 7
    #     # plot_states(states, days=prev_days, filename="./tmp/cases2.png", end_date=day)
    #     # plot_states(states, days=prev_days, mode=Mode.DEATHS, filename="./tmp/deaths2.png", end_date=day)
    #     # pdf.image("./tmp/cases2.png", 5, 110, WIDTH/2-10)
    #     # pdf.image("./tmp/deaths2.png", WIDTH/2, 110, WIDTH/2-10)

    #     # prev_days = 30
    #     # plot_states(states, days=prev_days, filename="./tmp/cases3.png", end_date=day)
    #     # plot_states(states, days=prev_days, mode=Mode.DEATHS, filename="./tmp/deaths3.png", end_date=day)
    #     # pdf.image("./tmp/cases3.png", 5, 200, WIDTH/2-10)
    #     # pdf.image("./tmp/deaths3.png", WIDTH/2, 200, WIDTH/2-10)

    #     # ''' Third Page '''
    #     # pdf.add_page()

    #     # plot_global_case_map("./tmp/global_cases.png", day=day)

    #     # countries = ['US', 'India', 'Brazil']
    #     # prev_days = 7
    #     # plot_countries(countries, days=prev_days, filename="./tmp/cases4.png", end_date=day)
    #     # plot_countries(countries, days=prev_days, mode=Mode.DEATHS, filename="./tmp/deaths4.png", end_date=day)

    #     # pdf.image("./tmp/global_cases.png", 5, 20, WIDTH-20)
    #     # pdf.image("./tmp/cases4.png", 5, 130, WIDTH/2-10)
    #     # pdf.image("./tmp/deaths4.png", WIDTH/2, 130, WIDTH/2-10)

    #     pdf.output(filename, 'F')


    # if __name__ == '__main__':
    #     yesterday = (datetime.today() - timedelta(days=1)).strftime("%m/%d/%y").replace("/0","/").lstrip("0")
    #     yesterday = "10/10/20" # Uncomment line for testing
    
    #     create_analytics_report(yesterday)
    
    
    # def plot_usa_case_map(filename=None, day=None):
    #     df = load_relevant_data()
    #     dates = list(df.columns)
    #     df = df.groupby('Province_State')[dates].agg('sum')
    #     create_usa_figure(df, filename, day)

    # def plot_global_case_map(filename=None, day=None):
    #     df = load_relevant_data(us_data=False)
    #     dates = list(df.columns)
    #     df = df.groupby('Country/Region')[dates].agg('sum')
    #     create_global_figure(df, filename, day)

    # def create_usa_figure(df, filename, day):
    #     day = day if day else yesterday # default to yesterday's date if not provided

    #     df['Cases'] = df.diff(axis=1)[day]
    #     df['state'] = [us_state_abbrev.get(x, None) for x in list(df.index)]
        

    #     fig = px.choropleth(df,
    #                     locations="state",
    #                     locationmode="USA-states",
    #                     scope="usa",
    #                     color="Cases",
    #                     hover_name="state",
    #                     color_continuous_scale='Peach',
    #                     title=f"US Daily Cases, {day}",
    #                     width=1000,
    #                     #height=500,
    #                     range_color=[0,3000])

    #     fig.update_layout(margin=dict(l=0, r=0, t=70, b=0), title={"font": {"size": 20}, "x":0.5},)
    #     filename = filename if filename else "usa_chart.png"
    #     fig.write_image(filename, engine='kaleido')

    # def create_global_figure(df, filename, day):
    #     day = day if day else yesterday # default to yesterday's date if not provided

    #     df['Cases'] = df.diff(axis=1)[day]
    #     df['Country'] = df.index

    #     fig = px.choropleth(df,
    #                     locations="Country",
    #                     locationmode="country names",
    #                     scope="world", # Try 'europe', 'africa', 'asia', 'south america', 'north america'
    #                     color="Cases",
    #                     hover_name="Country",
    #                     #projection="miller",
    #                     color_continuous_scale='Peach',
    #                     title=f"Global Daily Cases, {day}",
    #                     width=1000,
    #                     #height=500,
    #                     range_color=[0,50000])

    #     fig.update_layout(margin=dict(l=0, r=0, t=70, b=20), title={"font": {"size": 20}, "x":0.5},)
    #     filename = filename if filename else "global_chart.png"
    #     fig.write_image(filename, engine='kaleido')

    # if __name__ == '__main__':
    #     yesterday = (datetime.today() - timedelta(days=1)).strftime("%m/%d/%y")
    #     # Uncomment below line for testing
    #     yesterday = "10/10/20"

    #     plot_usa_case_map(day=yesterday) # saves as usa_chart.png by default
    #     plot_global_case_map(day=yesterday) # saves as global_chart.png by default

# asyncio.run(do_all_asynchronously(delay=delay))
# st.sidebar.markdown('<a href="mailto:ccasad@ulterra.com">Help & Feedback</a>', unsafe_allow_html=True)

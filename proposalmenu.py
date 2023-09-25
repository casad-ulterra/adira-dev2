from pathlib import Path
import streamlit as st
import pandas as pd
import numpy as np
from enum import Enum
from pandas.api.types import (
    is_categorical_dtype,
    is_datetime64_any_dtype,
    is_numeric_dtype,
    is_object_dtype,
)

# Plotly imports
from plotly.subplots import make_subplots
import plotly.graph_objects as go
import plotly.figure_factory as ff
import plotly.express as px
import plotly.io as pio
import pip
from PIL import Image
import datetime
# from datetime import datetime, timedelta
import io
import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import traceback
from streamlit_tags import st_tags, st_tags_sidebar
from geopy.distance import geodesic
from geopy.geocoders import Nominatim

import requests
import zipfile
import json
from pdf2image import convert_from_path
import yagmail
import openpyxl
from fuzzywuzzy import process
from fuzzywuzzy import fuzz
from enverus_developer_api import DeveloperAPIv3
import csv
from tempfile import mkdtemp
from streamlit_datalist import stDatalist

#  Azure blob setup
from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions

# from st_oauth import st_oauth
# # https://medium.com/streamlit/oauth-component-for-streamlit-e05f00874fbc
# st.markdown("## This (and above) is always seen")
# id = st_oauth('myoauth', 'Click to login via OAuth')
# # id = st_oauth(‘myoauth’)
# st.markdown("## This (and below) is only seen after authentication")


# initialize Nominatim API
geolocator = Nominatim(user_agent="adirageolocation")
# initialize yagmail
yag = yagmail.SMTP(st.secrets['yagmailusername'], st.secrets['yagmailpassword'])
# initialize azure blob storage
connection_string = 'DefaultEndpointsProtocol=https;AccountName=' + st.secrets['blobname'] + ';AccountKey=' + st.secrets['blobkey'] + ';EndpointSuffix=core.windows.net'
blob_service = BlobServiceClient.from_connection_string(connection_string)

def blob_upload(df, file_name, bupdate):
    blob_client = blob_service.get_blob_client(container=st.secrets['container2'], blob=file_name)
    try:
        output = df.to_csv(index=False, encoding="utf-8")        
        blob_client.upload_blob(output, overwrite=bupdate)
        status = 'Saved'
    except Exception as e:
        print(e)
        st.error(f"error blob df-to-csv: {e}")    
    # try:
    #     blob_exists = blob_client.exists()
    #     if not blob_exists:
    #         blob_client.upload_blob(output, overwrite=False)
    #         status = 'Loaded'
    #     else:            
    #         status = 'Exists'
    #     # , blob_type="BlockBlob")
    # except Exception as e:
    #     print(e)
    #     st.error(f"error blob upload: {e}")
    return status 

@st.cache_resource
def blob_download(filename):
    # blob_client = blob_service.get_blob_client(container=st.secrets['container'], blob=df_name)    
    # #get a list of all blob files in the container
    # blob_list = []
    # for blob_i in blob_client.list_blobs():
    #     blob_list.append(blob_i.name)
    
    # for blob_i in blob_list:
    #     #generate a shared access signature for each blob file
    #     sas_i = generate_blob_sas(account_name = account_name,
    #                                 container_name = container_name,
    #                                 blob_name = blob_i,
    #                                 account_key=account_key,
    #                                 permission=BlobSasPermissions(read=True),
    #                                 expiry=datetime.utcnow() + timedelta(hours=1))
    
    #generate a shared access signature for each blob file
    sas_i = generate_blob_sas(account_name = st.secrets['blobname'],
                                container_name = st.secrets['container2'],
                                blob_name = filename,
                                account_key=st.secrets['blobkey'],
                                permission=BlobSasPermissions(read=True),
                                expiry=datetime.datetime.utcnow() + datetime.timedelta(hours=1))
    sas_url = 'https://' + st.secrets['blobname'] + '.blob.core.windows.net/' + st.secrets['container2'] + '/' + filename + '?' + sas_i 
    
    df = pd.read_csv(sas_url)        
    # try:
    #     output = df.to_csv(index=False, encoding="utf-8")        
    #     blob_client.upload_blob(output, overwrite=False)
    #     status = 'Saved'
    # except Exception as e:
    #     print(e)
    #     st.error(f"error blob df-to-csv: {e}")  
    return df 

def runningToggle(bool):
  if bool == True:
    hide_streamlit_style = """
              <style>

              div[class='css-4z1n4l ehezqtx5']{
                background: rgba(0, 0, 0, 0.3);
                color: #fff;
                border-radius: 15px;
                height: 40px;
                max-width: 160px;


                position: fixed;
                top: 50%;
                left: 50%;
                transform: translate(-50%, -50%);
                width: 50%;
              }

              </style>
              """
    st.markdown(hide_streamlit_style, unsafe_allow_html=True) 
    
# Standard Color Palette
class upal(str, Enum):
    teal1 = '05929f'
    teal2 = '3F7077'
    gray1 = 'CBC9C9'
    gray2 = '9B9DA0'
    gray3 = '7A7D81'
    green1 = '98C21F'

# https://tools.nov.com/PowerSections/
pdmotors = ['',
          '1.688" 5:6 8.31 Rev/Gal 2.3',
          '1.688" 5:6 15.00 Rev/Gal 5.0', 
          '2.125" 5:6 12.20 Rev/Gal 6.0',
          '2.875" 5:6 3.38 Rev/Gal 3.5',    
          '2.875" 5:6 3.56 Rev/Gal 4.7',
          '2.875" 7:8 2.92 Rev/Gal 3.6',     
          '2.875" 5:6 5.52 Rev/Gal 7.0',
          '2.875" 4:5 2.85 Rev/Gal 4.0',     
          '2.875" 7:8 3.39 Rev/Gal 3.7',
          '3.125" 7:8 1.74 Rev/Gal 3.0',     
          '3.125" 5:6 2.15 Rev/Gal 3.5',
          '3.125" 5:6 2.64 Rev/Gal 5.0',     
          '3.125" 7:8 1.23 Rev/Gal 2.5',
          '3.375" 7:8 1.52 Rev/Gal 3.0',     
          '3.375" 5:6 2.09 Rev/Gal 3.5',
          '3.5" 5:6 2.66 Rev/Gal 3.0',     
          '3.75" 7:8 1.62 Rev/Gal 6.7',
          '4.125" 7:8 0.99 Rev/Gal 5.0',    
          '4.75" 4:5 0.96 Rev/Gal 6.3',
          '4.75" 5:6 0.99 Rev/Gal 8.3',    
          '4.75" 7:8 0.49 Rev/Gal 2.2',
          '4.75" 7:8 0.50 Rev/Gal 3.8',    
          '4.75" 9:10 0.54 Rev/Gal 4.0',
          '4 5" 6:7 0.79 Rev/Gal 6.4', 
          '5" 6:7 0.78 Rev/Gal 8.0',
          '5" 7:8 0.24 Rev/Gal 2.6',   
          '5" 6:7 0.83 Rev/Gal 6.0',
          '5" 6:7 0.81 Rev/Gal 8.0', 
          '5" 7:8 0.27 Rev/Gal 2.6',
          '5" 7:8 0.36 Rev/Gal 3.7', 
          '5" 7:8 0.51 Rev/Gal 3.8',
          '5" 7:8 0.25 Rev/Gal 2.6',   
          '5.125" 6:7 0.58 Rev/Gal 7.8',
          '5.125" 5:6 0.80 Rev/Gal 9.5',    
          '5.125" 9:10 0.46 Rev/Gal 6.1', 
          '5.500" 6:7 0.58 Rev/Gal 7.8',     
          '6.25" 7:8 0.34 Rev/Gal 4.8',
          '6.5" 7:8 0.27 Rev/Gal 6.7',     
          '6.565" 7:8 0.28 Rev/Gal 5.0',
          '6.75" 4:5 0.47 Rev/Gal 7.0', 
          '6.75" 7:8 0.15 Rev/Gal 2.9',
          '6.75" 7:8 0.27 Rev/Gal 5.0',    
          '6.75" 4:5 0.50 Rev/Gal 7.0',
          '6.75" 6:7 0.30 Rev/Gal 5.0',     
          '6.75" 6:7 0.29 Rev/Gal 6.0',
          '6.75" 7:8 0.14 Rev/Gal 3.3',     
          '6.75" 7:8 0.28 Rev/Gal 5.0',
          '6.75" 7:8 0.23 Rev/Gal 5.7',     
          '6.75" 7:8 0.17 Rev/Gal 2.9',
          '6.75" 7:8 0.30 Rev/Gal 6.4', 
          '6.75" 7:8 0.28 Rev/Gal 6.4',
          '6.75" 7:8 0.26 Rev/Gal 6.0',    
          '7" 7:8 0.31 Rev/Gal 6.0',
          '7" 7:8 0.31 Rev/Gal 7.5',     
          '7" 8:9 0.09 Rev/Gal 2.1',
          '7" 7:8 0.23 Rev/Gal 5.7',    
          '7" 5:6 0.41 Rev/Gal 8.2',
          '7" 5:6 0.45 Rev/Gal 11.1',   
          '7" 6:7 0.32 Rev/Gal 7.8',
          '7" 5:6 0.41 Rev/Gal 8.2',   
          '7.25" 7:8 0.26 Rev/Gal 7.2',
          '8" 4:5 0.25 Rev/Gal 5.3',     
          '8" 6:7 0.18 Rev/Gal 5.0',
          '8" 7:8 0.16 Rev/Gal 4.0 3',
          '8" 9:10 0.07 Rev/Gal 2.7',
          '8" 6:7 0.17 Rev/Gal 4.0',
          '8" 7:8 0.15 Rev/Gal 5.9 4',
          '8.25" 7:8 0.16 Rev/Gal 4.0 3',
          '8.25" 9:10 0.11 Rev/Gal 3.9',
          '8.500 8:9 0.17 Rev/Gal 6.7',
          '9.625" 6:7 0.10 Rev/Gal 3.5 6',
          '9.625" 7:8 0.11 Rev/Gal 4.8 6', 
          '9.625" 3:4 0.23 Rev/Gal 6.0 6', 
          '9.625" 6:7 0.13 Rev/Gal 5.0 7', 
          '9.625" 6:7 0.13 Rev/Gal 6.0 6', 
          '9.625" 7:8 0.11 Rev/Gal 4.8 6',    
          '11.25" 7:8 0.11 Rev/Gal 4.8 6', 
          '11.25" 7:8 0.11 Rev/Gal 4.8 6',
          'Other',     
]

# data = [    
#     'region':['Western US','Western US','Western US','Eastern US','Eastern US','Eastern US','Eastern US'],
#     'district':['Permian','Rockies','Williston','ArkLaTex','South Texas','MidCon','NEUS'],
#     'description':["In the Midland Basin, we're continuing to maintain high ROP throughout operations in stacked plays and the wide variety of interbedded formations that can lead to serious problems for operators. Ulterra also works alongside operators in the Delaware Basin to analyze extreme drilling parameters and introduce customized drill bits that can handle the toughest conditions with peak performance.",
#         "Our team of exs have experience across all basins in the Rocky Mountains and parts of California, such as the Denver Basin, Green River Basin, Powder River Basin, and others. With industry-leading experience for solving drilling challenges, including hard carbonates, gas, water flows, and slower drilling, Ulterra works to provide the right PDC bit to meet any requirement.",
#         "Ulterra's NODAK district focuses on helping operators drill through complex geology across the basin. We provide ex customer service for all areas spanning from portions of North Dakota, South Dakota, and Montana. This region covers the Williston Basin, including formations such as the Three Forks and Bakken.",
#         "Across the region including the Cotton Valley, and Haynesville Shale, operators encounter various drilling difficulties that require new bit designs and proper bit selection. Ulterra works alongside operators to solve problems that are faced in these difficult areas, and the challenges faced while drilling through the Travis Peak and Cotton Valley sands.",
#         "Ulterra helps operators from the Austin Chalk and Eagle Ford Shales to areas of the Gulf Coast achieve impressive drilling results. Our successes in South Texas have led to more application and performance improvements across other regions. In fact, our exs pioneered the use of Ulterra's CounterForce, SplitBlade, and FastBack bits throughout the South Texas oilfields.",
#         "Ulterra's involvement in the MidCon region includes plays in all of Oklahoma, Kansas, Arkansas, and the Texas Panhandle. The primary formations in this region include the Woodford, Meramec, Mississippian, Osage, and Springer. When operators face challenges in this region, like drilling through cy formations, hard rock, and interbedded formations, Ulterra is the supplier to count on to efficiently meet their drilling goals.",
#         "Ulterra provides operators in the Marcellus and Utica shales with the latest PDC bit developments that can tackle everything from curve and lateral runs to complex air drilling applications. One of Ulterra's latest innovations—AirRaid™—is the first PDC bit specifVertically designed for air drilling with outstanding durability and field-proven performance."
#     ]
# }
data = [['Permian','Western US',"In the Midland Basin, we're continuing to maintain high ROP throughout operations in stacked plays and the wide variety of interbedded formations that can lead to serious problems for operators. Ulterra also works alongside operators in the Delaware Basin to analyze extreme drilling parameters and introduce customized drill bits that can handle the toughest conditions with peak performance."],
        ['Rockies','Western US',"Our team of exs have experience across all basins in the Rocky Mountains and parts of California, such as the Denver Basin, Green River Basin, Powder River Basin, and others. With industry-leading experience for solving drilling challenges, including hard carbonates, gas, water flows, and slower drilling, Ulterra works to provide the right PDC bit to meet any requirement."],
        ['Williston','Western US',"Ulterra's NODAK district focuses on helping operators drill through complex geology across the basin. We provide ex customer service for all areas spanning from portions of North Dakota, South Dakota, and Montana. This region covers the Williston Basin, including formations such as the Three Forks and Bakken."],
        ['ArkLaTex','Eastern US',"Across the region including the Cotton Valley, and Haynesville Shale, operators encounter various drilling difficulties that require new bit designs and proper bit selection. Ulterra works alongside operators to solve problems that are faced in these difficult areas, and the challenges faced while drilling through the Travis Peak and Cotton Valley sands."],
        ['South Texas','Eastern US',"Ulterra helps operators from the Austin Chalk and Eagle Ford Shales to areas of the Gulf Coast achieve impressive drilling results. Our successes in South Texas have led to more application and performance improvements across other regions. In fact, our exs pioneered the use of Ulterra's CounterForce, SplitBlade, and FastBack bits throughout the South Texas oilfields."],    
        ['MidCon','Eastern US',"Ulterra's involvement in the MidCon region includes plays in all of Oklahoma, Kansas, Arkansas, and the Texas Panhandle. The primary formations in this region include the Woodford, Meramec, Mississippian, Osage, and Springer. When operators face challenges in this region, like drilling through cy formations, hard rock, and interbedded formations, Ulterra is the supplier to count on to efficiently meet their drilling goals."],
        ['NEUS''Eastern US',"Ulterra provides operators in the Marcellus and Utica shales with the latest PDC bit developments that can tackle everything from curve and lateral runs to complex air drilling applications. One of Ulterra's latest innovations—AirRaid™—is the first PDC bit specifVertically designed for air drilling with outstanding durability and field-proven performance."],
        ['Canada','North America',"Ulterra can draw it's roots in the Canadian oilfield back all the way to United Diamond in the mid-90s. With reach into all Canadian plays across Saskaquewon, Ala, and into the Montney in British Colombia. When operators face challenges in this region, like drilling through cy formations, hard rock, and interbedded formations, Ulterra is the supplier to count on to efficiently meet their drilling goals."],
        
    ]
districtnames = ['Permian','ArkLaTex','South Texas','Rockies','Williston','NEUS','Canada']
sectionlist = {1:'Surface',2:'Vertical',3:'Drill Out',4:'Intermediate',5:'Curve',6:'Lateral',7:'Extended Lateral',8:'Extended Lateral',9:'Extended Lateral',10:'Extended Lateral'}
sectionlist3 = {1:'Vertical',2:'Curve',3:'Lateral',}
sectionlist4 = {1:'Surface',2:'Vertical',3:'Curve',4:'Lateral',}
sectionsize_index = {'Surface':9,'Vertical':23,'Drill Out':23,'Intermediate':30,'Vertical/Curve/Lateral':34,'Vertical/Curve':34,'Curve':34,'Curve/Lateral':34,'Lateral':34,'Extended Lateral':47}
sectionsize = {'Surface':17.5,'Vertical':12.25,'Drill Out':12.25,'Intermediate':9.875,'Vertical/Curve/Lateral':8.75,'Vertical/Curve':8.75,'Curve':8.75,'Curve/Lateral':6.75,'Lateral':6.75,'Extended Lateral':6.75}
sectionrop = {'Surface':200,'Vertical':180,'Drill Out':170,'Intermediate':150,'Vertical/Curve/Lateral':130,'Vertical/Curve':140,'Curve':90,'Curve/Lateral':80,'Lateral':75,'Extended Lateral':60}
customernames = ['Aethon','Comstock','ConocoPhillips','Ovintiv','EOG','Oxy','Devon','Chevron','Other']
# find logos here https://companieslogo.com/eog-resources/logo/#google_vignette
# alternate logo lookup option logo_url = 'https://logo.clearbit.com/' + website
basins = pd.DataFrame(data, columns=['district','region','description'])

def proposalmaker():
    
    # padding = 0
    # st.markdown(f''' <style>
    #             .reportview-container .sidebar-content {{
    #                 padding-top: {padding}rem;
    #             }}
    #             .appview-container .main .block-container{{
    #                 padding-top: {padding}rem;    
    #                 padding-right: {padding}rem;
    #                 padding-left: {padding}rem;
    #                 padding-bottom: {padding}rem;
    #             }}
    #             .reportview-container .main .block-container {{
    #                 padding-top: {padding}rem;
    #                 padding-right: {padding}rem;
    #                 padding-left: {padding}rem;
    #                 padding-bottom: {padding}rem;
    #             }}
    #         </style> ''',
    #         unsafe_allow_html=True,
    #     )
    
    # st.markdown(f""" <style>
    #     .reportview-container .main .block-container{{
    #         padding-top: {padding}rem;
    #         padding-right: {padding}rem;
    #         padding-left: {padding}rem;
    #         padding-bottom: {padding}rem;
    #     }} </style> """, unsafe_allow_html=True)
    
    
    hide_img_fs = '''
        <style>
        button[title="View fullscreen"]{
            visibility: hidden;}
        </style>
        '''
    st.markdown(hide_img_fs, unsafe_allow_html=True)
    
    
    # # st.markdown( """<style>
    # #     div[class*="stcaption"] body {font-size:  2rem !important;color: black;}
    # #     div[class*="stTextArea"] label {font-size:  2rem !important;color: black;}
    # #     div[class*="stTextInput"] label {font-size:  2rem !important;color: black;}
    # #     div[class*="stSelectBox"] label {font-size:  2rem !important;}
    # #     div[class*="stMultiSelect"] label {font-size:  2rem !important;color: black;}
    # #     div[class*="stNumberInput"] label {font-size:  2rem !important;color: black;}
    # #     .big-font {font-size: 1.5rem !important;color: #43c6db;}
    # #     </style>""", unsafe_allow_html=True)
    
    st.markdown( """<style>
        div[class*="stcaption"] body {font-size:  2rem;color: black;}
        div[class*="stTextArea"] label {font-size:  2rem;color: black;}
        div[class*="stTextInput"] label {font-size:  2rem;color: black;}
        div[class*="stSelectBox"] label {font-size:  2rem;}
        div[class*="stMultiSelect"] label {font-size:  2rem;color: black;}
        div[class*="stNumberInput"] label {font-size:  2rem;color: black;}
        .big-font {font-size: 1.5rem;color: #43c6db;}
        </style>""", unsafe_allow_html=True)
    
    
    st.markdown(
        '''
        <style>
        .streamlit-expanderHeader {
            background-color: #37a8b2;
            color: white; # Adjust this for expander header color
        }
        </style>
        ''',
        unsafe_allow_html=True
    )    
    
        # .streamlit-expanderContent {
        #     background-color: #f7f7f7;
        #     color: black; # Expander content color
        # }
    
    # Moving triathalon running animation
    # runningToggle(True)
    
    # st.markdown(".stTextInput > label {font-size:105%; font-weight:bold;} ",unsafe_allow_html=True) #for all text-input label sections
    # st.markdown(".stSelectBox > label {font-size:105%; font-weight:bold;} ",unsafe_allow_html=True) #for all multi-select label sections
    # st.markdown(".stMultiSelect > label {font-size:105%; font-weight:bold;} ",unsafe_allow_html=True) #for all multi-select label sections
    
    uploaded_bpm = st.sidebar.file_uploader(label='Upload BPM file here', type= ['xlsx'], accept_multiple_files=False) 
    
    st.sidebar.markdown('<a href="mailto:ccasad@ulterra.com?subject=ADIRA Proposals Help & Feedback&body=Hey Chris, ADIRA rocks!"><button style="color:#43c6db;background-color:white;text-decoration:none;border-radius:4px;border:#43c6db;padding:10px 24px;">Email for Help & Feedback</button></a>', unsafe_allow_html=True)
    
    if not uploaded_bpm:
        # if "showpartnumber" not in st.session_state:
            # st.session_state.sizecount = 1
            # st.session_state.sectioncount = 3
            # st.session_state.wsectioncount = 1
            # st.session_state.showpartnumber = False
            # # st.session_state.showgage = True
            # # st.session_state.showpartnumber2 = False
            # # st.session_state.showgage2 = True
            # st.session_state.disabled = False
            
        if 'wellname' not in st.session_state:
            st.session_state.bpm = False
            st.session_state.bpmoverwrite = False  
            st.session_state.csv = False
            st.session_state.csvoverwrite = False 
            st.session_state.district = "Permian"
            st.session_state.custoemp = "Aethon"
            st.session_state.preparedby = " "
            st.session_state.preparedfor = " "
            st.session_state.wellname = " "
            st.session_state.rigname = " "
            st.session_state.lat = ""
            st.session_state.long = ""
            # st.session_state.opprofile = ''
            # st.session_state.sectioncount = 3
            st.session_state.spreadrate = '100,000'
            st.session_state.triprate = '1000'
            
        
    st.title('ADIRA Bit Menu')
    st.caption('Custom bit availability menu builder.')    
    
    # with st.container():                  
    #     col1, col2, col3 = st.columns(3)
    #     with col1:
    #         district = st.selectbox('Select district:',['Permian','ArkLaTex','South Texas','Rockies','Williston','NEUS'],)
    #         customer = st.selectbox('Select Customer:',['ConocoPhillips','Ovintiv','EOG','TapRock','DiamondBack','Chevron'],)
    #         # sections = st.multiselect('Which Sections?',['Surface','Vertical','Drill Out','Intermediate','Vertical/Curve','Vertical/Curve/Lateral','Curve','Curve/Lateral','Lateral','Extended Lateral'],)
        
    #     with col2:
    #         rigname = st.text_input("Enter Rig Name/Number:",)
    #         wellname = st.text_input("Enter Well Name:",)
    #     with col3:
    #         st.write('Display Options:')
    #         showpartnumber = st.checkbox("Show Part Number",value=False, key="showpartnumber")
    #         showgage = st.checkbox("Show Gage Type",value=True, key="showgage")        
        
    #     # Add css to make text bigger
    #     # st.markdown(
    #     #     """
    #     #     <style>
    #     #     textarea {
    #     #         font-size: 3rem !important;
    #     #     }
    #     #     input {
    #     #         font-size: 3rem !important;
    #     #     }
    #     #     </style>
    #     #     """,
    #     #     unsafe_allow_html=True,
    #     # )

            
    # st.write(
    #     """This app accomodates the blog [here](https://blog.streamlit.io/auto-generate-a-dataframe-filtering-ui-in-streamlit-with-filter_dataframe/)
    #     and walks you through one example of how the Streamlit
    #     Data Science Team builds add-on functions to Streamlit.
    #     """
    # )


    def convert_png_transparent(src_file, dst_file, bg_color=(255,255,255)):
        image = Image.open(src_file).convert("RGBA")
        array = np.array(image, dtype=np.ubyte)
        mask = (array[:,:,:3] == bg_color).all(axis=2)
        alpha = np.where(mask, 0, 255)
        array[:,:,-1] = alpha
        Image.fromarray(np.ubyte(array)).save(dst_file, "png")
        
    def filter_dataframe(dffltr: pd.DataFrame) -> pd.DataFrame:
        """
        Adds a UI on top of a dataframe to let viewers filter columns

        Args:
            df (pd.DataFrame): Original dataframe

        Returns:
            pd.DataFrame: Filtered dataframe
        """
        modify = True
        # modify = st.checkbox("Add filters")

        if not modify:
            return df

        dffltr = dffltr.copy()

        # Try to con datetimes into a standard format (datetime, no timezone)
        for col in dffltr.columns:
            if is_object_dtype(dffltr[col]):
                try:
                    dffltr[col] = pd.to_datetime(dffltr[col])
                except Exception:
                    pass

            if is_datetime64_any_dtype(dffltr[col]):
                dffltr[col] = dffltr[col].dt.tz_localize(None)

        modification_container = st.container()

        with modification_container:
            to_filter_columns = st.multiselect("Filter dataframe on", dffltr.columns)
            for column in to_filter_columns:
                left, right = st.columns((1, 20))
                left.write("↳")
                # Treat columns with < 10 unique values as categorVertical
                if is_categorical_dtype(dffltr[column]) or dffltr[column].nunique() < 10:
                    user_cat_input = right.multiselect(
                        f"Values for {column}",
                        dffltr[column].unique(),
                        default=list(dffltr[column].unique()),
                    )
                    df = dffltr[dffltr[column].isin(user_cat_input)]
                elif is_numeric_dtype(dffltr[column]):
                    _min = float(dffltr[column].min())
                    _max = float(dffltr[column].max())
                    step = (_max - _min) / 100
                    user_num_input = right.slider(
                        f"Values for {column}",
                        _min,
                        _max,
                        (_min, _max),
                        step=step,
                    )
                    df = dffltr[dffltr[column].between(*user_num_input)]
                elif is_datetime64_any_dtype(dffltr[column]):
                    user_date_input = right.date_input(
                        f"Values for {column}",
                        value=(
                            dffltr[column].min(),
                            dffltr[column].max(),
                        ),
                    )
                    if len(user_date_input) == 2:
                        user_date_input = tuple(map(pd.to_datetime, user_date_input))
                        start_date, end_date = user_date_input
                        dffltr = dffltr.loc[dffltr[column].between(start_date, end_date)]
                else:
                    user_text_input = right.text_input(
                        f"Substring or regex in {column}",
                    )
                    if user_text_input:
                        df = dffltr[dffltr[column].str.contains(user_text_input)]

        return df
    
    # uploaded_file = st.sidebar.file_uploader('Upload CSV or Excel File here', type=['csv', 'xlsx'])    
    # if uploaded_file is not None:
        # df = pd.read_csv(uploaded_file, encoding_errors='ignore')
    
    # old hardcoded DDA sample
    # df = pd.read_csv('data/adira book3.csv', encoding_errors='ignore')
    
    @st.cache_data
    def get_bit_records(csvdir):
        # dfrec = pd.read_csv('data/records_USonly.csv') 
        dfrec = pd.read_csv(csvdir) 
        #  Clean up data        
        dfrec['Latitude'] = dfrec['Latitude'].clip(upper=90)
        dfrec['Longitude'] = dfrec['Longitude'].clip(upper=180)
        dfrec['ROP'] = dfrec['ROP'].clip(upper=500)
        dfrec['ROP'] = dfrec['ROP'].replace(500, np.nan)
        # https://codereview.stackexchange.com/questions/185389/dropping-rows-from-a-pandas-dataframe-where-some-of-the-columns-have-value-0
        dropcolumns = ['DepthOut', 'ROP', 'Latitude', 'Longitude'] 
        dfrec = dfrec.replace(0, np.nan).dropna(axis=0, how='any', subset=dropcolumns).fillna(0)
        return dfrec 
    
    @st.cache_resource
    def get_bit_list(url,headers):
        response = requests.get(url=url, headers=headers)
        df_list = pd.DataFrame(json.loads(response.text))
        return df_list
    
    # API DDA
    url = 'https://bithub.ulterra.com/adira_dda_fetch'
    headers = {'Adira_API_Key': st.secrets['Adira_API_Key']}
    # response = requests.get(url=url, headers=headers)
    # df = pd.DataFrame(json.loads(response.text))
    df = get_bit_list(url, headers)                
    
    
    @st.cache_resource
    def get_bit_files(url,headers):
        response = requests.get(url=url, headers=headers)
        zipper = zipfile.ZipFile(io.BytesIO(response.content))
        return zipper
    
    @st.cache_data
    def get_geolocation(lat, long):
        glocation = geolocator.reverse(lat+","+long)
        return glocation
        
    
    if df is not None:         
        
        # with tab2:
        with st.container():
            def BitMenu():
                #  Version 1 of work flow             
                col1, col2 = st.columns([3,1])
                with col1:
                    district = st.selectbox('Select district?',districtnames,)
                    customer = st.selectbox('Select Customer?',customernames,)
                    sections = st.multiselect('Which Sections?',['Surface','Vertical','Drill Out','Intermediate','Vertical/Curve','Vertical/Curve/Lateral','Curve','Curve/Lateral','Lateral','Extended Lateral'],)
                
                with col2:
                    st.write('Display Options:')
                    showpartnumber = st.checkbox("Show Part Number",value=False, key="showpartnumber")
                    showgage = st.checkbox("Show Gage Type",value=True, key="showgage")
                                
                st.divider()       
                
                def definesections(name,df_bits,sizeindex): 
                    with st.container():                    
                        st.caption(name)
                        # df.sort_values(by="Company").Company.unique() 
                        scol1,scol2 = st.columns([1,5])
                        with scol1:
                            defsize = st.selectbox('Select Surface Sizes:',options=df_bits['BitSize'].sort_values(ascending=False).unique(),index=sizeindex,key=f'{name}size')
                        with scol2:
                            defbits = st.multiselect('Select Bits:',options=df_bits['PartNumber'].sort_values(ascending=False).loc[df_bits['BitSize'] == defsize], key=f'{name}bits')                        
                            # localbit = st.multiselect('Select **Bit** (Primary, Backup):',options=df['PartNumber'].sort_values(ascending=False), key=f'{i}bit', help='First selection is Primary, second is backup.')
                            
                        try:
                            if defbits is not None:                        
                                df_section = df_bits.loc[df_bits['PartNumber'].isin(defbits)]
                                df_section['Price'] = None
                                df_section['Backup'] = None
                                scols = st.columns(len(defbits))
                                for count, bit in enumerate(defbits): 
                                    with scols[count]:
                                        st.caption(bit)
                                                                                             
                                        if os.path.exists(f'{bit}-1.jpg'):
                                            image = f'{bit}-1.jpg'       
                                        elif os.path.exists(f'data/Bit Pictures/{bit}-1.jpg'):                                
                                            image = f'data/Bit Pictures/{bit}-1.jpg' 
                                        elif os.path.exists(f'data/Bit Pictures/{bit}-1.JPG'):                                
                                            image = f'data/Bit Pictures/{bit}-1.JPG'
                                        else:
                                            image = f'data/Bit Pictures/blank.jpg'     
                                            
                                        # imagepng = f'{bit}-1.png'
                                        # convert_png_transparent(image, imagepng)
                                        st.image(image,caption=f"{df_bits.loc[df_bits[df_bits['PartNumber'] == bit].index, 'BitSize'].values[0]}  - {df_bits.loc[df_bits[df_bits['PartNumber'] == bit].index, 'BitType'].values[0]}",width=250)

                                        # # st.image('https://static.streamlit.io/examples/cat.jpg',use_column_width='auto')
                                        # if os.path.exists(f'data/Bit Pictures/{bit}-1.jpg'):                                
                                        #     # image = Image.open(f'data/Bit Pictures/{bit}-1.jpg') 
                                        #     st.image(f'data/Bit Pictures/{bit}-1.jpg',width=250)
                                        # elif os.path.exists(f'data/Bit Pictures/{bit}-1.JPG'):     
                                        #     st.image(f'data/Bit Pictures/{bit}-1.JPG',width=250)
                                        # else:
                                        #     st.image(f'data/Bit Pictures/blank.jpg',width=250)                            
                                        # st.image(image,use_column_width='auto')
                                        
                                        df_section.loc[df_section[df_section['PartNumber'] == bit].index,'Price'] = st.text_input('Price:', '$', key=f'{bit}price')
                                        
                                        # df_section.loc[df_section[df_section['PartNumber'] == bit].index,'Backup'] = st.radio("Pickup:",['Primary','Secondary'],horizontal=True, key=f'{bit}backup')
                                        df_section.loc[df_section[df_section['PartNumber'] == bit].index,'Backup'] = st.checkbox("Backup",value=False, key=f'{bit}backup')
                                        
                                        # st.session_state.df_bits = pd.concat([st.session_state.df_bits, pd.DataFrame({'nr': input_values})],ignore_index=True)
                        except Exception as e:
                            print(e)
                            st.info(f"error priceloop: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                        
                    st.divider()
                    return defsize,defbits,df_section        
                
                # Section first, then size
                if 'Surface' in sections: 
                    bitsize_s,bits_s,df_s = definesections('Surface',df,9)            
                if 'Vertical' in sections:              
                    bitsize_v,bits_v,df_v = definesections('Vertical',df,23)            
                if 'Drill Out' in sections:  
                    bitsize_do,bits_do,df_do = definesections('Drill Out',df,23)            
                if 'Intermediate' in sections:   
                    bitsize_i,bits_i,df_i = definesections('Intermediate',df,30)            
                if 'Vertical/Curve/Lateral' in sections:  
                    bitsize_vcl,bits_vcl,df_vcl = definesections('Vertical/Curve/Lateral',df,34)            
                if 'Vertical/Curve' in sections:           
                    bitsize_vc,bits_vc,df_vc = definesections('Vertical/Curve',df,34)            
                if 'Curve' in sections:            
                    bitsize_c,bits_c,df_c = definesections('Curve',df,34)            
                if 'Curve/Lateral' in sections:          
                    bitsize_cl,bits_cl,df_cl = definesections('Curve/Lateral',df,34)
                if 'Lateral' in sections:       
                    bitsize_l,bits_l,df_l = definesections('Lateral',df,34)
                if 'Extended Lateral' in sections: 
                    bitsize_el,bits_el,df_el = definesections('Extended Lateral',df,47)
                
                def create_pptx_multipage(pagefootage):                
                    
                    save_bar.progress(1, text='Initializing..')             
                    pptx = 'data/Adira Proposal Template.pptx'
                    prs = Presentation(pptx)
                    
                    # declare positional variables
                    WIDTH = Inches(10)
                    HEIGHT = Inches(7.5)
                    # left = Inches(2.5)
                    # top = Inches(1)
                    
                    # get stock info
                    name = 'Ulterra'
                    
                    try:   
                        
                        # def resize_image(url):
                        #     """function to resize logos while keeping aspect ratio. Accepts URL as an argument and return an image object"""

                        #     # Open the image file
                        #     image = Image.open(requests.get(url, stream=True).raw)

                        #     # if a logo is too high or too wide then make the background container twice as big
                        #     if image.height > 140:
                        #         container_width = 220 * 2
                        #         container_height = 140 * 2

                        #     elif image.width > 220:
                        #         container_width = 220 * 2
                        #         container_height = 140 * 2
                        #     else:
                        #         container_width = 220
                        #         container_height = 140

                        #     # Create a new image with the same aspect ratio as the original image
                        #     new_image = Image.new('RGBA', (container_width, container_height))

                        #     # Calculate the position to paste the image so that it is centered
                        #     x = (container_width - image.width) // 2
                        #     y = (container_height - image.height) // 2

                        #     # Paste the image onto the new image
                        #     new_image.paste(image, (x, y))

                        #     return new_image
                        
                        def add_image(slide, image, left, top, width):
                            """function to add an image to the PowerPoint slide and specify its position and width"""
                            slide.shapes.add_picture(image, left=left, top=top, width=width)

                        # function to replace text in pptx first slide with selected filters
                        def replace_text(replacements, shapes):
                            """function to replace text on a PowerPoint slide. Takes dict of {match: replacement, ... } and replaces all matches"""
                            for shape in shapes:
                                for match, replacement in replacements.items():
                                    if shape.has_text_frame:
                                        if (shape.text.find(match)) != -1:
                                            text_frame = shape.text_frame
                                            for paragraph in text_frame.paragraphs:
                                                whole_text = "".join(run.text for run in paragraph.runs)
                                                whole_text = whole_text.replace(str(match), str(replacement))
                                                for idx, run in enumerate(paragraph.runs):
                                                    if idx != 0:
                                                        p = paragraph._p
                                                        p.remove(run._r)
                                                if bool(paragraph.runs):
                                                    paragraph.runs[0].text = whole_text

        
                        save_bar.progress(5, text='Creating template..') 
                        lyt=prs.slide_masters[0].slide_layouts[2] # choosing a slide layout
                                        
                        # declare pptx variables
                        first_slide = prs.slides[0]
                        second_slide = prs.slides[2]
                        # fourth_slide = prs.slides[3]
                        shapes_1 = []
                        shapes_2 = []
                        shapes_4 = []
                        index_to_drop = []

                        # create lists with shape objects
                        for shape in first_slide.shapes:
                            shapes_1.append(shape)

                        for shape in second_slide.shapes:
                            shapes_2.append(shape)
                            
                        # for shape in fourth_slide.shapes:
                        #     shapes_4.append(shape)

                        # initiate a dictionary of placeholders and values to replace
                        replaces_cover = {
                            '{company}': customer,
                            '{document}':'BIT AVAILABILITY & PRICING',
                            '{basin}': district,
                            '{year}': '2023',
                            '{wellname}': '',
                            '{rigname}': '',
                            }

                        replaces_basin = {
                            '{region}': basins.loc[basins[basins['district'] == district].index,'region'].values[0],
                            '{district}': basins.loc[basins[basins['district'] == district].index,'district'].values[0],
                            '{description}': basins.loc[basins[basins['district'] == district].index,'description'].values[0],
                        }
                        
                        replaces_county = {
                            '{county}': district,
                        }
                        

                        # run the function to replace placeholders with values
                        replace_text(replaces_cover, shapes_1)
                        replace_text(replaces_basin, shapes_2)
                        replace_text(replaces_county, shapes_4)
                        
                        
                        prs.slides[2].shapes.add_picture(f'data/maps/map-{district}.png', left=Inches(-0.02), top=Inches(-0.04)) 
                        
                        #  Pg1
                        # slide1=prs.slides[0]
                        # txBox = slide1.shapes.add_textbox(left=Inches(0.7), top=Inches(0.44), width=Inches(3), height=Inches(0.3))
                        # tf = txBox.text_frame
                        # tf.text = f'ROCK REPORT'
                        # tf.paragraphs[0].font.size = Pt(10)
                        # tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  
                        # # tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  
                        
                        # table of info
                        # rows = 4		# number of rows in table
                        # cols = 6		# number of columns in table
                        # # Add Shape object (table)
                        # # Arguments: number of rows, number of columns, x-coordinate of top-left corner, x-coordinate of top-left corner, width,  height
                        # table_shape = slide1.shapes.add_table(rows, cols, left=Inches(0.5), top=Inches(1), width=Inches(6.5), height=Inches(2))
                        # table = table_shape.table	# Create Table object

                        # tbl =  table_shape._element.graphic.graphicData.tbl
                        # # List of Style IDs: https://github.com/scanny/python-pptx/issues/27#issuecomment-263076372
                        # style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
                        # tbl[0][-1].text = style_id
                        
                    except Exception as e:
                        print(e)
                        st.info(f"error ppt intro: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                    # try:
                    #     https://stackoverflow.com/questions/65330597/exporting-dataframe-to-pptx-using-python
                    #     save_bar.progress(15, text='Page 1: Header..')                     
                    #     def set_cell_info(tablecell, text, aligned):                        
                    #         try:
                    #             tablecell.text =  text
                    #             tablecell.text_frame.paragraphs[0].alignment = aligned
                    #             tablecell.text_frame.paragraphs[0].font.size = Pt(8)  
                    #             tablecell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)                        
                    #         except Exception as e:
                    #             print(e)
                    #             st.info(f"error ppt tablecell: {tablecell}: {e}")                         
                    #         return tablecell
                        
                    #     set_cell_info(table.cell(0, 0), str(las_file.well.COMP.descr.capitalize()), PP_ALIGN.RIGHT)
                    #     set_cell_info(table.cell(1, 0), str(las_file.well.WELL.descr.capitalize()), PP_ALIGN.RIGHT)
                    #     set_cell_info(table.cell(2, 0), str(las_file.well.UWI.descr.capitalize()), PP_ALIGN.RIGHT)
                    #     set_cell_info(table.cell(3, 0), str(las_file.well.LOC.descr.capitalize()), PP_ALIGN.RIGHT)                
                    #     set_cell_info(table.cell(0, 1), str(las_file.well.COMP.value), PP_ALIGN.LEFT)
                    #     set_cell_info(table.cell(1, 1), str(las_file.well.WELL.value), PP_ALIGN.LEFT)
                    #     set_cell_info(table.cell(2, 1), str(las_file.well.UWI.value), PP_ALIGN.LEFT)
                    #     set_cell_info(table.cell(3, 1), str(las_file.well.LOC.value), PP_ALIGN.LEFT)                
                    #     set_cell_info(table.cell(0, 2), str(las_file.well.FLD.descr.capitalize()), PP_ALIGN.RIGHT)
                    #     set_cell_info(table.cell(1, 2), str(las_file.well.CNTY.descr.capitalize()), PP_ALIGN.RIGHT)
                    #     set_cell_info(table.cell(2, 2), str(las_file.well.STAT.descr.capitalize()), PP_ALIGN.RIGHT)
                    #     set_cell_info(table.cell(3, 2), str(las_file.well.CTRY.descr.capitalize()), PP_ALIGN.RIGHT)                
                    #     set_cell_info(table.cell(0, 3), str(las_file.well.FLD.value), PP_ALIGN.LEFT)
                    #     set_cell_info(table.cell(1, 3), str(las_file.well.CNTY.value), PP_ALIGN.LEFT)
                    #     set_cell_info(table.cell(2, 3), str(las_file.well.STAT.value), PP_ALIGN.LEFT)
                    #     set_cell_info(table.cell(3, 3), str(las_file.well.CTRY.value), PP_ALIGN.LEFT)                
                    #     set_cell_info(table.cell(0, 4), str(las_file.well.STRT.descr.capitalize()), PP_ALIGN.RIGHT)
                    #     set_cell_info(table.cell(1, 4), str(las_file.well.STOP.descr.capitalize()), PP_ALIGN.RIGHT)
                    #     set_cell_info(table.cell(2, 4), str(las_file.well.STEP.descr.capitalize()), PP_ALIGN.RIGHT)                
                    #     set_cell_info(table.cell(0, 5), str(las_file.well.STRT.value), PP_ALIGN.LEFT)
                    #     set_cell_info(table.cell(1, 5), str(las_file.well.STOP.value), PP_ALIGN.LEFT)
                    #     set_cell_info(table.cell(2, 5), str(las_file.well.STEP.value), PP_ALIGN.LEFT)
                        
                    # except Exception as e:
                    #     print(e)
                    #     st.info(f"File missing header info: {e}")                     

                    try:   
                        # slidelast=prs.slides.add_slide(lyt) # adding a slide   
                        def createbitcard(aslide, df_section, pn, lpos, tpos):
                            shapes2 = aslide.shapes
                            bshape = shapes2.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left=Inches(lpos), top=Inches(tpos), width=Inches(2), height=Inches(4.5))
                            try:
                                # https://python-pptx.readthedocs.io/en/latest/user/autoshapes.html#adjusting-an-autoshape
                                adjs = bshape.adjustments
                                
                                # # Size & Type
                                # txBox = aslide.shapes.add_textbox(left=Inches(0), top=Inches(0), width=Inches(2), height=Inches(0.22))            
                                # tf = txBox.text_frame
                                # tf.text = f"{bshape.adjustments}"
                                # tf.paragraphs[0].font.size = Pt(18)
                                
                                adjs[1].effective_value = 0.1                                       
                            except Exception as e:
                                print(e)
                                # st.info(f"error shaperound: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                    
                            bshape.fill.solid()
                            # bshape.fill.fore_color.rgb = RGBColor(212, 230, 232)
                            bshape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                            bshape.line.fill.background()
                            # D4E6E8
                            # set lighter or darker, -0.2 is 20% darker, 0.4 is 40% lighter
                            # shape.fill.fore_color.brightness = 0.4
                            # set fill to transparent (no fill)
                            # shape.fill.background()
                                                                            
                            if os.path.exists(f'{pn}-1.jpg'):
                                image = f'{pn}-1.jpg'       
                            elif os.path.exists(f'data/Bit Pictures/{pn}-1.jpg'):                                
                                image = f'data/Bit Pictures/{pn}-1.jpg' 
                            elif os.path.exists(f'data/Bit Pictures/{pn}-1.JPG'):                                
                                image = f'data/Bit Pictures/{pn}-1.JPG'
                            else:
                                image = f'data/Bit Pictures/blank.jpg'     
                                
                            imagepng = f'{pn}-1.png'
                            convert_png_transparent(image, imagepng)
                            aslide.shapes.add_picture(imagepng, left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2)) 
                            
                            # if os.path.exists(f'data/Bit Pictures/{pn}-1.jpg'):
                            #     aslide.shapes.add_picture(f'data/Bit Pictures/{pn}-1.jpg', left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2))                            
                            #     # image = f'data/Bit Pictures/{pn}-1.jpg'
                            # elif os.path.exists(f'data/Bit Pictures/{pn}-1.JPG'):
                            #     aslide.shapes.add_picture(f'data/Bit Pictures/{pn}-1.JPG', left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2))
                            #     # image = f'data/Bit Pictures/{pn}-1.JPG'
                            # else:                           
                            #     # image = f'data/Bit Pictures/blank.jpg'
                            #     aslide.shapes.add_picture(f'data/Bit Pictures/blank.jpg', left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2)) 
                            # # aslide.shapes.add_picture(image, left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2)) 
                            
                            # Size & Type
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos), top=Inches(tpos+2.15), width=Inches(2), height=Inches(0.4))            
                            tf = txBox.text_frame
                            # bitsize = df_section.loc[df_section['PartNumber'] == pn,'BitSize']
                            tf.text = f"{df_section.loc[df_section['PartNumber'] == pn,'BitSize'].values[0]} {df_section.loc[df_section['PartNumber'] == pn,'BitType'].values[0]}"
                            tf.paragraphs[0].font.size = Pt(18)
                            txBox.fill.solid()
                            txBox.fill.fore_color.rgb = RGBColor(0, 147, 159)                    
                            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  
                            
                            #  Features
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+0.1), top=Inches(tpos+2.55), width=Inches(1.9), height=Inches(1))            
                            tf = txBox.text_frame
                            feat1 = df_section.loc[df_section['PartNumber'] == pn,'PartNumber'].values[0] if showpartnumber else ''
                            feat2 = f"{df_section.loc[df_section['PartNumber'] == pn,'GaugeLength'].values[0]} Gage Length" if showgage else ''
                            feat3 = df_section.loc[df_section['PartNumber'] == pn,'GaugeType'].values[0] if showgage else ''
                            feat4 = df_section.loc[df_section['PartNumber'] == pn,'PerformancePackage'].values[0]
                            tf.text = f"{feat1} \v {feat2} \v {feat3} \v {feat4}"
                            tf.paragraphs[0].font.size = Pt(12)
                            
                            # Price
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos), top=Inches(tpos+3.65), width=Inches(2), height=Inches(0.4))            
                            tf = txBox.text_frame
                            # bitsize = df_section.loc[df_section['PartNumber'] == pn,'BitSize'].values[0]
                            tf.text = f"{df_section.loc[df_section['PartNumber'] == pn,'Price'].values[0]}"
                            tf.paragraphs[0].font.size = Pt(18)
                            txBox.fill.solid()
                            txBox.fill.fore_color.rgb = RGBColor(242, 242, 242)
                            
                            # # Status / Price Notes
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+0.1), top=Inches(tpos+4.06), width=Inches(1.9), height=Inches(0.4))            
                            tf = txBox.text_frame
                            tf.text = f"{'Backup' if df_section.loc[df_section['PartNumber'] == pn,'Backup'].values[0] == 'True' else ''}"
                            tf.paragraphs[0].font.size = Pt(10)
                            
                            
                        
                        def createsectionpage(section,df_section,bits):     
                            newslide=prs.slides.add_slide(lyt) # adding a slide              
                            # Section Name
                            txBox = newslide.shapes.add_textbox(left=Inches(1.1), top=Inches(0.24), width=Inches(3), height=Inches(0.33))            
                            tf = txBox.text_frame
                            tf.text = f'{section}'
                            tf.paragraphs[0].font.size = Pt(24)
                            tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                            # subtitle             
                            txBox = newslide.shapes.add_textbox(left=Inches(1.1), top=Inches(0.67), width=Inches(5), height=Inches(0.33))            
                            tf = txBox.text_frame
                            tf.text = f'BIT AVAILABILITY AND PRICING'
                            tf.paragraphs[0].font.size = Pt(22)
                            # description             
                            txBox = newslide.shapes.add_textbox(left=Inches(1.1), top=Inches(1.1), width=Inches(5), height=Inches(1.35))            
                            tf = txBox.text_frame
                            tf.text = f"We combine our company's unique product design with a rapid \v prototyping manufacturing process, leading to reliable and \v high-end performance drill bits customers can count on in \v the {district}. The proof is in the numbers."
                            tf.paragraphs[0].font.size = Pt(12)
                            
                            bitcount = 0
                            for bit in bits:
                                bitcount += 1
                                # txBox = newslide.shapes.add_textbox(left=Inches(0.5), top=Inches(1+bitcount), width=Inches(1), height=Inches(2))            
                                # tf = txBox.text_frame
                                # tf.text = f'data/Bit Pictures/{bit}-1.jpg'
                                # tf.paragraphs[0].font.size = Pt(8)
                                # if os.path.exists(f'data/Bit Pictures/{bit}-1.jpg'):
                                #     newslide.shapes.add_picture(f'data/Bit Pictures/{bit}-1.jpg', left=Inches(0.25), top=Inches(4.75), width=Inches(2))  
                                # elif os.path.exists(f'data/Bit Pictures/{bit}-1.JPG'):
                                #     newslide.shapes.add_picture(f'data/Bit Pictures/{bit}-1.JPG', left=Inches(0.25), top=Inches(4.75), width=Inches(2)) 
                                # else:
                                #     newslide.shapes.add_picture(f'data/Bit Pictures/blank.jpg', left=Inches(0.25), top=Inches(4.75), width=Inches(2)) 
                                
                                if len(bits) <= 2:
                                    start_h = 1.5
                                    start_v = 2.1
                                    createbitcard(newslide,df_section,bit,(start_h + (2 * (bitcount-1)) + (0.25 * (bitcount-1))),start_v)
                                else:
                                    start_h = 1.5
                                    start_v = 2.1
                                    createbitcard(newslide,df_section,bit,(start_h + (2 * (bitcount-1))+ (0.25 * (bitcount-1))),start_v)
                                
                        if 'Surface' in sections:  
                            createsectionpage('Surface',df_s,bits_s)  
                            
                            # Verticalslide=prs.slides.add_slide(lyt) # adding a slide                  
                            # txBox = Verticalslide.shapes.add_textbox(left=Inches(0.5), top=Inches(1), width=Inches(1), height=Inches(2))            
                            # tf = txBox.text_frame
                            # tf.text = f'{bits_s}'
                            # tf.paragraphs[0].font.size = Pt(8)
                            # bitcount = 0
                            # for bit in bits_s:
                            #     bitcount += 1                
                            #     txBox = Verticalslide.shapes.add_textbox(left=Inches(0.5), top=Inches(1+bitcount), width=Inches(1), height=Inches(2))            
                            #     tf = txBox.text_frame
                            #     tf.text = f'data/Bit Pictures/{bit}-1.jpg'
                            #     tf.paragraphs[0].font.size = Pt(8)
                            #     if os.path.exists(f'data/Bit Pictures/{bit}-1.jpg'):
                            #         Verticalslide.shapes.add_picture(f'data/Bit Pictures/{bit}-1.jpg', left=Inches(0.25), top=Inches(4.75), width=Inches(2))  
                            #     elif os.path.exists(f'data/Bit Pictures/{bit}-1.JPG'):
                            #         Verticalslide.shapes.add_picture(f'data/Bit Pictures/{bit}-1.JPG', left=Inches(0.25), top=Inches(4.75), width=Inches(2)) 
                            #     else:
                            #         Verticalslide.shapes.add_picture(f'data/Bit Pictures/blank.jpg', left=Inches(0.25), top=Inches(4.75), width=Inches(2)) 
                            
                        # if 'Vertical' in sections:    
                        #     createsectionpage('Vertical',df_v,bits_v)   
                            
                        # if 'Drill Out' in sections: 
                        #     createsectionpage('Drill Out',df_do,bits_do)    
                            
                        # if 'Intermediate' in sections:    
                        #     createsectionpage('Intermediate',df_i,bits_i)  
                            
                        # if 'Curve' in sections: 
                        #     createsectionpage('Curve',df_c,bits_c)    
                            
                            
                        if 'Vertical' in sections:              
                            createsectionpage('Vertical',df_v,bits_v)      
                            
                        if 'Drill Out' in sections:  
                            createsectionpage('Drill Out',df_do,bits_do)  
                            
                        if 'Intermediate' in sections:   
                            createsectionpage('Intermediate',df_i,bits_i)  
                            
                        if 'Vertical/Curve/Lateral' in sections:  
                            createsectionpage('Vertical/Curve/Lateral',df_vcl,bits_vcl)  
                            
                        if 'Vertical/Curve' in sections:           
                            createsectionpage('Vertical/Curve',df_vc,bits_vc)  
                            
                        if 'Curve' in sections:            
                            createsectionpage('Curve',df_c,bits_c)  
                            
                        if 'Curve/Lateral' in sections:          
                            createsectionpage('Curve/Lateral',df_cl,bits_cl)   
                            
                        if 'Lateral' in sections:       
                            createsectionpage('Lateral',df_l,bits_l)     
                            
                        if 'Extended Lateral' in sections: 
                            createsectionpage('Extended Lateral',df_el,bits_el)      
                            
                            # filename = "xplot.png"
                            # pio.write_image(xplot, filename, scale=2, width=873, height=450)
                            # slidelast.shapes.add_picture(filename, left=Inches(0.25), top=Inches(4.75), width=Inches(7))                         
                            
                    except Exception as e:
                        print(e)
                        st.info(f"error ppt-slides: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                            
                        
                    # # Pg2
                    # try:
                    #     slide2=prs.slides.add_slide(lyt) # adding a slide 
                    #     txBox = slide2.shapes.add_textbox(Inches(0.87), Inches(0.4), Inches(1), Inches(0.25))
                    #     tf = txBox.text_frame
                    #     tf.text = f'OVERVIEW'
                    #     tf.paragraphs[0].font.size = Pt(8)  
                    #     #  Depth vs Time
                    #     filename = "fig1.jpg"
                    #     pio.write_image(fig1, filename, scale=2, width=(72*4.8))             
                    #     slide2.shapes.add_picture(filename, left=Inches(0.1), top=Inches(1.36), width=Inches(4.8))        
                            
                    # except Exception as e:
                    #     print(e)
                    #     st.info(f"error ppt2: {e}")
                    
                    # Close out
                    # save_bar.progress(92, text='Saving..') 
                    binary_output = BytesIO()
                    prs.save(binary_output)
                    return binary_output
            
            
                with st.form(key='template_form'):
                    fc1, fc2 = st.columns(2)
                    with fc1:
                        reporttype = st.selectbox('Choose Format',['PowerPoint',],index=0)
                    with fc2:
                        reportsize = st.selectbox('Choose Type',['Multipage',],index=0)
                        # pagefootage = st.selectbox('Choose Ratio',[2000,1000,5000],index=0)
                        pagefootage = 5
                    with st.spinner('Request...'): 
                        submit = st.form_submit_button('Request Menu')
                
                if submit:                    
                    save_bar = st.progress(0, text='Initializing..')
                    with st.spinner('Constructing...'):    
                        if reporttype == 'PowerPoint':
                            if reportsize == 'Multipage':
                                try:
                                    reportppt = create_pptx_multipage(pagefootage)
                                    try:                
                                        save_bar.progress(95, text='Preparing Download')
                                        st.download_button(
                                            label="Download PowerPoint",
                                            data=reportppt,
                                            file_name=f"Adira_BitMenu_{customer}_{datetime.datetime.now()}.pptx",
                                            mime="application/octet-stream",
                                        ) 
                                        save_bar.progress(100, text='Ready')
                                    except Exception as e:
                                        st.info(f"error pptx dl1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                             
                                except Exception as e:
                                    st.info(f"error pptx: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
            
            # Call app
            BitMenu()


# import pandas as pd
# from geopy.distance import geodesic

# def calculate_distance(source_coords, target_coords):
#     return geodesic(source_coords, target_coords).miles

# def calculate_ranking_score(record, target_record, vehicle_class):
#     distance_score = calculate_distance((record['latitude'], record['longitude']), 
#                                         (target_record['latitude'], target_record['longitude']))
#     size_score = abs(record['Size'] - target_record['Size'])
#     speed_score = abs(record['speed'] - target_record['speed'])

#     # Check if the vehicle class is available in the record, otherwise consider it as a separate class
#     if 'class' in record:
#         class_score = 0 if record['class'] == vehicle_class else 1
#     else:
#         class_score = 1

#     # Assuming the target_record includes the start and end mileage points
#     start_mileage = target_record['start_mileage']
#     end_mileage = target_record['end_mileage']
#     segment_start = record['start_mileage']
#     segment_end = record['end_mileage']

#     # Calculate the absolute difference between the middle segment start and end mileage
#     segment_mileage_diff = abs((segment_end - segment_start) - (end_mileage - start_mileage))

#     # You can adjust the weights below based on how you want to prioritize the factors.
#     # For example, if distance is more important, increase its weight, and vice versa.
#     distance_weight = 0.25
#     size_weight = 0.15
#     speed_weight = 0.15
#     class_weight = 0.2
#     segment_weight = 0.25

#     total_score = (distance_weight * distance_score +
#                    size_weight * size_score +
#                    speed_weight * speed_score +
#                    class_weight * class_score +
#                    segment_weight * segment_mileage_diff)

#     return total_score

# def rank_records(csv_file, target_record, vehicle_class):
#     df = pd.read_csv(csv_file)

#     # Calculate ranking score for each record in the dataframe
#     df['ranking_score'] = df.apply(lambda row: calculate_ranking_score(row, target_record, vehicle_class), axis=1)

#     # Sort the dataframe based on ranking score in ascending order
#     ranked_df = df.sort_values(by='ranking_score')

#     return ranked_df

# def find_high_ranking_matches(csv_file, target_record, vehicle_class, num_matches=3):
#     ranked_records = rank_records(csv_file, target_record, vehicle_class)
#     high_ranking_matches = ranked_records.head(num_matches)
    
#     while len(high_ranking_matches) < num_matches:
#         # Gradually broaden the class pool and re-rank
#         vehicle_class = 'Any'
#         ranked_records = rank_records(csv_file, target_record, vehicle_class)
#         high_ranking_matches = ranked_records.head(num_matches)

#     return high_ranking_matches

# def find_high_ranking_matches2(csv_file, target_record, vehicle_class, num_matches=3, max_distance=100):
#     ranked_records = rank_records(csv_file, target_record, vehicle_class)
#     high_ranking_matches = ranked_records.head(num_matches)
#     current_distance = 0

#     while len(high_ranking_matches) < num_matches and current_distance <= max_distance:
#         current_distance += 10  # Increase the distance search by 10 miles
#         # Re-rank the records based on the increased distance
#         ranked_records = rank_records(csv_file, target_record, vehicle_class)

#         # Filter the records based on the current distance and vehicle class
#         filtered_records = ranked_records[
#             (ranked_records['ranking_score'] <= current_distance) & 
#             ((ranked_records['class'] == vehicle_class) | (vehicle_class == 'Any'))
#         ]

#         high_ranking_matches = filtered_records.head(num_matches)

#     if len(high_ranking_matches) < num_matches:
#         # If still not enough high ranking matches, expand the search to other class sizes
#         vehicle_class = 'Any'
#         ranked_records = rank_records(csv_file, target_record, vehicle_class)
#         high_ranking_matches = ranked_records.head(num_matches)

#     return high_ranking_matches

# def find_high_ranking_matches3(csv_file, target_record, vehicle_class, num_matches=3, max_distance=100):
#     df = pd.read_csv(csv_file)
#     key_column = 'matching_key'  # Name of the column to store the matching key or tag

#     # Check if there are already records with the same matching key
#     previous_matches = df[df[key_column].notnull()]

#     if len(previous_matches) >= num_matches:
#         # If we have enough previous matches, use them to expedite future matching
#         filtered_records = df[df[key_column].isin(previous_matches[key_column])]
#     else:
#         # Otherwise, perform a new search using the previous logic
#         ranked_records = rank_records(csv_file, target_record, vehicle_class)
#         high_ranking_matches = ranked_records.head(num_matches)
#         current_distance = 0

#         while len(high_ranking_matches) < num_matches and current_distance <= max_distance:
#             current_distance += 10  # Increase the distance search by 10 miles
#             # Re-rank the records based on the increased distance
#             ranked_records = rank_records(csv_file, target_record, vehicle_class)

#             # Filter the records based on the current distance and vehicle class
#             filtered_records = ranked_records[
#                 (ranked_records['ranking_score'] <= current_distance) & 
#                 ((ranked_records['class'] == vehicle_class) | (vehicle_class == 'Any'))
#             ]

#             high_ranking_matches = filtered_records.head(num_matches)

#         if len(high_ranking_matches) < num_matches:
#             # If still not enough high ranking matches, expand the search to other class sizes
#             vehicle_class = 'Any'
#             ranked_records = rank_records(csv_file, target_record, vehicle_class)
#             high_ranking_matches = ranked_records.head(num_matches)

#         # Store the matching key or tag in the DataFrame for future matching
#         matching_key = 'key_{}'.format(len(previous_matches) + 1)
#         df.at[high_ranking_matches.index, key_column] = matching_key
#         filtered_records = df[df[key_column].isin([matching_key])]

#     return filtered_records

# # Sample usage
# if __name__ == "__main__":
#     # Replace 'target_record' with the record you want to use as the target for matching
#     target_record = {
#         'latitude': 37.7749,
#         'longitude': -122.4194,
#         'Size': 10,
#         'speed': 50,
#         'start_mileage': 0,
#         'end_mileage': 1000
#     }

#     csv_file = 'path/to/your/csv/file.csv'
#     vehicle_class = 'SUV'  # Replace this with the desired vehicle class
#     num_matches = 3  # The number of high ranking matches to find

#     high_ranking_matches = find_high_ranking_matches(csv_file, target_record, vehicle_class, num_matches)
#     print(high_ranking_matches)





    # # import pandas as pd
    # # from geopy.distance import geodesic

    # def calculate_distance(source_coords, target_coords):
    #     return geodesic(source_coords, target_coords).miles

    # def calculate_ranking_score(record, target_record, vehicle_class):
    #     distance_score = calculate_distance((record['Latitude'], record['Longitude']), 
    #                                         (target_record['Latitude'], target_record['Longitude']))
    #     size_score = abs(record['Size'] - target_record['Size'])
    #     speed_score = abs(record['Speed'] - target_record['Speed'])

    #     # Check if the vehicle class is available in the record, otherwise consider it as a separate class
    #     if 'class' in record:
    #         class_score = 0 if record['class'] == vehicle_class else 1
    #     else:
    #         class_score = 1

    #     # Assuming the target_record includes the start and end mileage points
    #     start_mileage = target_record['start_mileage']
    #     end_mileage = target_record['end_mileage']
    #     segment_start = record['start_mileage']
    #     segment_end = record['end_mileage']

    #     # Calculate the absolute difference between the middle segment start and end mileage
    #     segment_mileage_diff = abs((segment_end - segment_start) - (end_mileage - start_mileage))

    #     # You can adjust the weights below based on how you want to prioritize the factors.
    #     # For example, if distance is more important, increase its weight, and vice versa.
    #     distance_weight = 0.25
    #     size_weight = 0.15
    #     speed_weight = 0.15
    #     class_weight = 0.2
    #     segment_weight = 0.25

    #     total_score = (distance_weight * distance_score +
    #                 size_weight * size_score +
    #                 speed_weight * speed_score +
    #                 class_weight * class_score +
    #                 segment_weight * segment_mileage_diff)

    #     return total_score

    # def rank_records(csv_file, target_record, vehicle_class):
    #     dfrr = pd.read_csv(csv_file)

    #     # Calculate ranking score for each record in the dataframe
    #     dfrr['ranking_score'] = dfrr.apply(lambda row: calculate_ranking_score(row, target_record, vehicle_class), axis=1)

    #     # Sort the dataframe based on ranking score in ascending order
    #     ranked_df = dfrr.sort_values(by='ranking_score')

    #     return ranked_df

    # def find_high_ranking_matches(csv_file, target_record, vehicle_class, num_matches=3, max_distance=100):
    #     dfhrr = pd.read_csv(csv_file)
    #     key_column = 'matching_key'  # Name of the column to store the matching key or tag

    #     # Check if there are already records with the same matching key
    #     previous_matches = dfhrr[dfhrr[key_column].notnull()]

    #     if len(previous_matches) >= num_matches:
    #         # If we have enough previous matches, use them to expedite future matching
    #         filtered_records = dfhrr[dfhrr[key_column].isin(previous_matches[key_column])]
    #     else:
    #         # Otherwise, perform a new search using the previous logic
    #         ranked_records = rank_records(csv_file, target_record, vehicle_class)
    #         high_ranking_matches = ranked_records.head(num_matches)
    #         current_distance = 0

    #         while len(high_ranking_matches) < num_matches and current_distance <= max_distance:
    #             current_distance += 10  # Increase the distance search by 10 miles
    #             # Re-rank the records based on the increased distance
    #             ranked_records = rank_records(csv_file, target_record, vehicle_class)

    #             # Filter the records based on the current distance and vehicle class
    #             filtered_records = ranked_records[
    #                 (ranked_records['ranking_score'] <= current_distance) & 
    #                 ((ranked_records['class'] == vehicle_class) | (vehicle_class == 'Any'))
    #             ]

    #             high_ranking_matches = filtered_records.head(num_matches)

    #         if len(high_ranking_matches) < num_matches:
    #             # If still not enough high ranking matches, expand the search to other class sizes
    #             vehicle_class = 'Any'
    #             ranked_records = rank_records(csv_file, target_record, vehicle_class)
    #             high_ranking_matches = ranked_records.head(num_matches)

    #         # Store the matching key or tag in the DataFrame for future matching
    #         matching_key = 'key_{}'.format(len(previous_matches) + 1)
    #         df.at[high_ranking_matches.index, key_column] = matching_key
    #         filtered_records = dfhrr[dfhrr[key_column].isin([matching_key])]

    #     return filtered_records

    # # Sample usage
    # if __name__ == "__main__":
    #     # Replace 'target_record' with the record you want to use as the target for matching
    #     # target_record = {
    #     #     'latitude': 37.7749,
    #     #     'longitude': -122.4194,
    #     #     'Size': 10,
    #     #     'speed': 50,
    #     #     'start_mileage': 0,
    #     #     'end_mileage': 1000
    #     # }
        
    #     target_lat = lat # st.number_input("Latitude:")
    #     target_lon = long # st.number_input("Longitude:")
    #     target_size = dfchart.loc[1,'Size'] # st.number_input("Size:")
    #     target_speed = dfchart.loc[1,'ROP'] # st.number_input("Speed:")
    #     target_start_mileage = dfchart.loc[1,'Din'] # st.number_input("Start Mileage:")
    #     target_end_mileage = dfchart.loc[1,'Dout'] # st.number_input("End Mileage:")
    #     target_class = dfchart.loc[1,'Size'] # st.text_input("Vehicle Class:")

    #     target_record = {
    #         'Latitude': target_lat,
    #         'Longitude': target_lon,
    #         'Size': target_size,
    #         'ROP': target_speed,
    #         'DepthIn': target_start_mileage,
    #         'DepthOut': target_end_mileage,
    #     }

    #     csv_file = 'records_USonly.csv'
    #     vehicle_class = 'SUV'  # Replace this with the desired vehicle class
    #     num_matches = 3  # The number of high ranking matches to find
    #     max_distance = 100  # The maximum distance for initial GPS location search

    #     high_ranking_matches = find_high_ranking_matches(csv_file, target_record, vehicle_class, num_matches, max_distance)
    #     print(high_ranking_matches)


    # def main():
    #     st.title("Matchmaking Web App")
    #     st.write("Upload a CSV data file and test the matchmaking algorithm.")

    #     # uploaded_file = st.file_uploader("Choose a CSV file", type=["csv"])
    #     uploaded_file = 'records_USonly.csv'
        
    #     if uploaded_file:
    #         df = pd.read_csv(uploaded_file)

    #         # Create input fields for the target record
    #         st.caption("Enter Target Record Information:")
    #         target_lat = lat # st.number_input("Latitude:")
    #         target_lon = long # st.number_input("Longitude:")
    #         target_size = dfchart.loc[1,'Size'] # st.number_input("Size:")
    #         target_speed = dfchart.loc[1,'ROP'] # st.number_input("Speed:")
    #         target_start_mileage = dfchart.loc[1,'Din'] # st.number_input("Start Mileage:")
    #         target_end_mileage = dfchart.loc[1,'Dout'] # st.number_input("End Mileage:")
    #         target_class = dfchart.loc[1,'Size'] # st.text_input("Vehicle Class:")

    #         target_record = {
    #             'Latitude': target_lat,
    #             'Longitude': target_lon,
    #             'Size': target_size,
    #             'ROP': target_speed,
    #             'DepthIn': target_start_mileage,
    #             'DepthOut': target_end_mileage,
    #         }

    #         if target_class:
    #             target_record['Size'] = target_class

    #         # Perform the matchmaking algorithm and display results
    #         st.caption("Matchmaking Results:")
    #         matched_records = find_high_ranking_matches(df, target_record, vehicle_class=target_class)
    #         st.write(matched_records)

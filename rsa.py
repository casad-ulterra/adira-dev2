from pathlib import Path
import streamlit as st
import pandas as pd
import numpy as np
from enum import Enum
import altair as alt

# Plotly imports
from plotly.subplots import make_subplots
import plotly.graph_objects as go
import plotly.figure_factory as ff
import plotly.express as px
import plotly.io as pio
import pip
from PIL import Image
from fpdf import FPDF
import datetime
import io
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import traceback

# Standard Color Palette
class upal(str, Enum):
    teal1 = '05929f'
    teal2 = '3F7077'
    gray1 = 'CBC9C9'
    gray2 = '9B9DA0'
    gray3 = '7A7D81'
    green1 = '98C21F'

class adiraPDF(FPDF):
    def header(self):
        # # Select Arial bold 15
        # self.set_font('Arial', 'B', 15)       
        self.set_font('Arial', '', 12)  
        self.image('images/adira-logo-name-lrg.png', 10, 10, 100)
        # self.cell(30, 30, f'header-REPORT',0,1)
        
    def footer(self):
        # Go to 1.5 cm from bottom
        self.set_y(-15)
        # Select Arial italic 8
        self.set_font('Arial', 'I', 8)
        # Print centered page number
        self.cell(0, 10, 'Page %s' % self.page_no(), 0, 0, 'C')
        
        self.image('Ulterra_teal_250px.png', self.set_x(-20), self.set_y(-35), 140)
        self.text(40, 2360, f'Created by ADIRA v2.34')


colh1, colh2 = st.columns(2)
# colh1.header('Options')

def convert_df(df):
   return df.to_csv(index=False).encode('utf-8')

def header(las_file):
    if not las_file:
        st.info('No file has been uploaded')
    
    else:
        def make_grid(cols,rows):
            grid = [0]*cols
            for i in range(cols):
                with st.container():
                    grid[i] = st.columns(rows)
            return grid

        mygrid = make_grid(4,4)      
        if hasattr(las_file.well, 'COMP'):
            mygrid[0][0].write(f'<b>{las_file.well.COMP.descr.capitalize()}:</b> {las_file.well.COMP.value}', unsafe_allow_html=True)
        if hasattr(las_file.well, 'WELL'):
            mygrid[1][0].write(f'<b>{las_file.well.WELL.descr.capitalize()}:</b> {las_file.well.WELL.value}', unsafe_allow_html=True)
        if hasattr(las_file.well, 'UWI'):
            mygrid[2][0].write(f'<b>{las_file.well.UWI.descr.capitalize()}:</b> {las_file.well.UWI.value}', unsafe_allow_html=True)
        if hasattr(las_file.well, 'LOC'):
            mygrid[3][0].write(f'<b>{las_file.well.LOC.descr.capitalize()}:</b> {las_file.well.LOC.value}', unsafe_allow_html=True)
        
        if hasattr(las_file.well, 'FLD'):      
            mygrid[0][1].write(f'<b>{las_file.well.FLD.descr.capitalize()}:</b> {las_file.well.FLD.value}', unsafe_allow_html=True)
        if hasattr(las_file.well, 'CNTY'):
            mygrid[1][1].write(f'<b>{las_file.well.CNTY.descr.capitalize()}:</b> {las_file.well.CNTY.value}', unsafe_allow_html=True)
        if hasattr(las_file.well, 'STAT'):
            mygrid[2][1].write(f'<b>{las_file.well.STAT.descr.capitalize()}:</b> {las_file.well.STAT.value}', unsafe_allow_html=True)
        if hasattr(las_file.well, 'CTRY'):
            mygrid[3][1].write(f'<b>{las_file.well.CTRY.descr.capitalize()}:</b> {las_file.well.CTRY.value}', unsafe_allow_html=True)
            
        if hasattr(las_file.well, 'STRT'):            
            mygrid[0][2].write(f'<b>{las_file.well.STRT.descr.capitalize()}:</b> {las_file.well.STRT.value}', unsafe_allow_html=True)
        if hasattr(las_file.well, 'STOP'):
            mygrid[1][2].write(f'<b>{las_file.well.STOP.descr.capitalize()}:</b> {las_file.well.STOP.value}', unsafe_allow_html=True)
        if hasattr(las_file.well, 'STEP'):
            mygrid[2][2].write(f'<b>{las_file.well.STEP.descr.capitalize()}:</b> {las_file.well.STEP.value}', unsafe_allow_html=True)
        
        # with st.expander(f'{las_file.well.WELL.descr.capitalize()}: {las_file.well.WELL.value}', expanded=True): 
        #     # for item in las_file.well:
        #     #     st.write(f"<b>{item.descr.capitalize()} ({item.mnemonic}):</b> {item.value}", unsafe_allow_html=True)
                
        #     mygrid = make_grid(4,4)            
        #     mygrid[0][0].write(f'<b>{las_file.well.COMP.descr.capitalize()}:</b> {las_file.well.COMP.value}', unsafe_allow_html=True)
        #     mygrid[1][0].write(f'<b>{las_file.well.WELL.descr.capitalize()}:</b> {las_file.well.WELL.value}', unsafe_allow_html=True)
        #     mygrid[2][0].write(f'<b>{las_file.well.UWI.descr.capitalize()}:</b> {las_file.well.UWI.value}', unsafe_allow_html=True)
        #     mygrid[3][0].write(f'<b>{las_file.well.LOC.descr.capitalize()}:</b> {las_file.well.LOC.value}', unsafe_allow_html=True)
            
        #     mygrid[0][1].write(f'<b>{las_file.well.FLD.descr.capitalize()}:</b> {las_file.well.FLD.value}', unsafe_allow_html=True)
        #     mygrid[1][1].write(f'<b>{las_file.well.CNTY.descr.capitalize()}:</b> {las_file.well.CNTY.value}', unsafe_allow_html=True)
        #     mygrid[2][1].write(f'<b>{las_file.well.STAT.descr.capitalize()}:</b> {las_file.well.STAT.value}', unsafe_allow_html=True)
        #     mygrid[3][1].write(f'<b>{las_file.well.CTRY.descr.capitalize()}:</b> {las_file.well.CTRY.value}', unsafe_allow_html=True)
            
        #     mygrid[0][2].write(f'<b>{las_file.well.STRT.descr.capitalize()}:</b> {las_file.well.STRT.value}', unsafe_allow_html=True)
        #     mygrid[1][2].write(f'<b>{las_file.well.STOP.descr.capitalize()}:</b> {las_file.well.STOP.value}', unsafe_allow_html=True)
        #     mygrid[2][2].write(f'<b>{las_file.well.STEP.descr.capitalize()}:</b> {las_file.well.STEP.value}', unsafe_allow_html=True)
            
        
def plot(las_file, well_data):
    # st.title('LAS File Visualisation')
    
    st.sidebar.markdown('<a href="mailto:ccasad@ulterra.com?subject=ADIRA Help & Feedback&body=Hey Chris, ADIRA rocks!"><button style="color:#43c6db;background-color:white;text-decoration:none;border-radius:4px;border:#43c6db;padding:10px 24px;">Email for Help & Feedback</button></a>', unsafe_allow_html=True)
        
    if not las_file:
        st.info(' ')

    else:
        columns = list(well_data.columns)
        # st.write('Expand one of the following to visualise your well data.')
        # st.write("""Each plot can be interacted with. To change the scales of a plot/track, click on the left hand or right hand side of the scale and change the value as required.""")
                
        # with st.expander('Rock Strength Analysis', expanded=True):
        if las_file:
            # col1_h, col2_h = st.columns([1,4])
            # col1_h.header('Options')

            # hist_curve = col1_h.selectbox('Select a Curve', columns)
            # log_option = col1_h.radio('Select Linear or Logarithmic Scale', ('Linear', 'Logarithmic'))
            # hist_col = col1_h.color_picker('Select Histogram Colour')
            # st.write('Color is'+hist_col)
            
            # if log_option == 'Linear':
            #     log_bool = False
            # elif log_option == 'Logarithmic':
            #     log_bool = True
        
            # histogram = px.histogram(well_data, x=hist_curve, log_x=log_bool)
            # histogram.update_traces(marker_color=hist_col)
            # histogram.layout.template='seaborn'
            # col2_h.plotly_chart(histogram, use_container_width=True)
            
            # uRSA = px.area(well_data, x="DEPTH", y=["SS","SH"], color=well_data.columns, pattern_shape=["SS","SH"], pattern_shape_sequence=[".", "x", "+"])
            # col2_h.plotly_chart(uRSA, use_container_width=True)
            
            # try 1
            
            # uRSA = px.area(well_data, x="DEPTH", y=["SS","SH"])
            # st.plotly_chart(uRSA, use_container_width=True)
                                    
            # # The 'shape' property is an enumeration that may be specified as: - One of the following enumeration values: ['', '/', '\\', 'x', '-', '|', '+', '.']
            # fig9 = go.Figure() 
            # fig9.add_trace(go.Scatter(x=well_data['SH'], y=well_data['DEPTH'], fill='tozerox', name='Shale', fillcolor ='#bebebe', line= {'color':'#bebebe', 'width':0}, stackgroup='one', fillpattern=dict(bgcolor='#bebebe', shape='-'),orientation='h'))
            # fig9.add_trace(go.Scatter(x=well_data['SS'], y=well_data['DEPTH'], fill='tonextx', name='Sandstone', fillcolor ='#ffff00', line= {'color':'#ffff00', 'width':0}, stackgroup='one', fillpattern=dict(bgcolor='#ffff00', shape='.'),orientation='h'))   
            # fig9.add_trace(go.Scatter(x=well_data['LS'], y=well_data['DEPTH'], fill='tonextx', name='Limestone', fillcolor ='#80ffff', line= {'color':'#80ffff', 'width':0}, stackgroup='one', fillpattern=dict(bgcolor='#80ffff', shape='+'),orientation='h'))      
            # fig9.add_trace(go.Scatter(x=well_data['SI'], y=well_data['DEPTH'], fill='tonextx', name='Siltstone', fillcolor ='#7cfc00', line= {'color':'#7cfc00', 'width':0}, stackgroup='one', fillpattern=dict(bgcolor='#7cfc00', shape='|'),orientation='h'))      
            # fig9.add_trace(go.Scatter(x=well_data['DO'], y=well_data['DEPTH'], fill='tonextx', name='Dolomite', fillcolor ='#8080ff', line= {'color':'#8080ff', 'width':0}, stackgroup='one', fillpattern=dict(bgcolor='#8080ff', shape='/'),orientation='h'))      
            # fig9.add_trace(go.Scatter(x=well_data['AN'], y=well_data['DEPTH'], fill='tonextx', name='Anhydrite', fillcolor ='#ff80ff', line= {'color':'#ff80ff', 'width':0}, stackgroup='one', fillpattern=dict(bgcolor='#ff80ff', shape='\\'),orientation='h'))      
            # fig9.add_trace(go.Scatter(x=well_data['SL'], y=well_data['DEPTH'], fill='tonextx', name='Salt', fillcolor ='#7ddfbe', line= {'color':'#7ddfbe', 'width':0}, stackgroup='one', fillpattern=dict(bgcolor='#7ddfbe', shape='x'),orientation='h'))      
            # fig9.add_trace(go.Scatter(x=well_data['CO'], y=well_data['DEPTH'], fill='tonextx', name='Coal', fillcolor ='black', line= {'color':'black', 'width':0}, stackgroup='one', fillpattern=dict(bgcolor='black', shape=''),orientation='h')) 
            # fig9.update_layout(height=1000, showlegend=True, yaxis={'title':'DEPTH','autorange':'reversed'}, title='LithologyA', xaxis=dict(type='linear',range=[1, 100],ticksuffix='%'))
            # fig9.update_layout(xaxis={'side':'top'}, legend=dict(yanchor="top", y=1.05, xanchor="left", x=0.3,orientation="h"))
            
            # fig9.layout.template='plotly_dark'
            # fig9.layout.template='seaborn'
            # fig9.update_layout(paper_bgcolor="rgb(10,10,10)", plot_bgcolor="rgb(10,10,10)")
            # col1_h.plotly_chart(fig9, use_container_width=True)
            
            # fig9b = go.Figure()
            # fig9b.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['SH'], fill='tozeroy', name='Shale', stackgroup='one', groupnorm='percent',orientation='h'))
            # fig9b.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['SS'], fill='tonexty', name='Sandstone', stackgroup='one',orientation='h')) 
            # fig9b.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['LS'], fill='tonexty', name='Limestone', stackgroup='one',orientation='h')) 
            # fig9b.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['SI'], fill='tonexty', name='Siltstone', stackgroup='one',orientation='h')) 
            # fig9b.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['DO'], fill='tonexty', name='Dolomite', stackgroup='one',orientation='h')) 
            # fig9b.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['AN'], fill='tonexty', name='Anhydrite', stackgroup='one',orientation='h')) 
            # fig9b.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['SL'], fill='tonexty', name='SL', stackgroup='one',orientation='h')) 
            # fig9b.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['CO'], fill='tonexty', name='CO', stackgroup='one',orientation='h')) 
            # fig9b.update_layout(height=350, showlegend=True, xaxis={'title':'DEPTH','autorange':'reversed'}, title='LithologyB')
            # # yaxis=dict(type='linear',range=[1, 100],ticksuffix='%')
            # fig9b.layout.template='seaborn'
            # # fig9b.update_layout(paper_bgcolor="rgb(10,10,10)", plot_bgcolor="rgb(10,10,10)")
            # st.plotly_chart(fig9b, use_container_width=True)
            
            # # Streamlit API call code    https://discuss.streamlit.io/t/using-an-api-from-inside-streamlit-app/37477
            # import requests
            # import streamlit as st

            # def call_api(filename,bdata=None,overlay=False,api_key='helloworld',language='eng'):
            #     payload = {'isOverlayRequired': overlay,'apikey': api_key,'language': language}
            #     if bdata is not None:
            #         r = requests.post('https://api.ocr.space/parse/image',files={filename: bdata},data=payload)
            #     else:
            #         with open(filename, 'rb') as f:
            #             r = requests.post('https://api.ocr.space/parse/image',files={filename: f},data=payload)

            #     return r.content.decode()

            # if __name__ == '__main__':
            #     uploaded_file = st.file_uploader("Choose file")
            #     is_submit = st.button('Process File')                
            #     if is_submit and uploaded_file is not None:
            #         bytes_data = uploaded_file.getvalue()
            #         ret_json = call_api(uploaded_file.name,bdata=bytes_data,api_key=st.secrets['ocr_api_key']['mykey']  # 'helloworld'
            #         )

            #         st.json(ret_json)
            
            
            # https://stackoverflow.com/questions/65056691/plotly-how-to-show-more-than-2-x-axes-titles-ranges-on-the-same-subplot
            # # initial subplot with two traces
            # figa = make_subplots(rows=1, cols=6,shared_yaxes=True)
            # figa.add_trace(go.Scatter(x=[1, 2, 3], y=[4, 5, 6]),row=1, col=1)
            # figa.add_trace(go.Scatter(x=[20, 30, 40], y=[50, 60, 70]),row=1, col=2)
            # figa.add_trace(go.Scatter(x=[20, 30, 40], y=[50, 60, 70]),row=1, col=3)
            # figa.add_trace(go.Scatter(x=[20, 30, 40], y=[50, 60, 70]),row=1, col=4)
            # figa.add_trace(go.Scatter(x=[20, 30, 40], y=[50, 60, 70]),row=1, col=5)
            # figa.add_trace(go.Scatter(x=[20, 30, 40], y=[50, 60, 70]),row=1, col=6)
            # # figa.add_trace(go.Scatter(x=[50, 60, 70], y=[500, 600, 700]),row=1, col=3)
            # figa.update_layout(height=600, width=800,title_text="Subplots with shared x-axes")

            # # extra data where xaxis3 is shared with subplot 1
            # figa.add_trace(go.Scatter(x=[11, 12, 13],y=[6, 5, 4],name="xaxis4 data",xaxis="x7"))
            # # some adjustmentns for xaxis3
            # figa.update_layout(xaxis7=dict(title="xaxis4 title",titlefont=dict(color="#9467bd"),tickfont=dict(color="#9467bd"),anchor="free",overlaying="x1",side="right",position=1.0))
            # # extra data where xaxis4 is shared with subplot 2
            # figa.add_trace(go.Scatter(x=[50, 60, 70],y=[60, 60, 60],name="xaxis5 data",xaxis="x8",yaxis = 'y2'))
            # # some adjustments for xaxis4
            # figa.update_layout(xaxis8=dict(title="xaxis5 title",titlefont=dict(color="#9467bd"),tickfont=dict(color="#9467bd"),anchor="free",overlaying="x2",side="right",position=0.0))
            
            # figa.add_trace(go.Scatter(x=[50, 60, 70],y=[60, 60, 60],name="xaxis6 data",xaxis="x9",yaxis = 'y3'))
            # # some adjustments for xaxis4
            # figa.update_layout(xaxis9=dict(title="xaxis6 title",titlefont=dict(color="#9467bd"),tickfont=dict(color="#9467bd"),anchor="free",overlaying="x3",side="right",position=0.0))
            
            # figa.add_trace(go.Scatter(x=[50, 60, 70],y=[60, 60, 60],name="xaxis6 data",xaxis="x12",yaxis = 'y6'))
            # # some adjustments for xaxis4
            # figa.update_layout(xaxis12=dict(title="xaxis12 title",titlefont=dict(color="#9467bd"),tickfont=dict(color="#9467bd"),anchor="free",overlaying="x6",side="right",position=0.0))
            
            
            # # # extra data where xaxis4 is shared with subplot 3
            # # figa.add_trace(go.Scatter(x=[50, 60, 70],y=[60, 60, 60],name="xaxis4 data",xaxis="x5",yaxis = 'y2'))
            # # # some adjustments for xaxis4
            # # figa.update_layout(xaxis5=dict(title="xaxis5 title",titlefont=dict(color="#9467bd"),tickfont=dict(color="#9467bd"),anchor="free",overlaying="x3",side="right",position=0.0))

            # # # make room to display double x-axes
            # # figa.update_layout(yaxis1=dict(domain=[0.1, 1]),yaxis2=dict(domain=[0.1, 1]),)
            # # # not critical, but just to put a little air in there
            # # figa.update_layout(xaxis1=dict(domain=[0.0, 0.4]),xaxis2=dict(domain=[0.6, 1]),)
            
            # # # # make room to display double x-axes
            # # figa.update_layout(yaxis1=dict(domain=[0.1, 1]),yaxis2=dict(domain=[0.1, 1]),yaxis3=dict(domain=[0.1, 1]),)
            # # # # not critical, but just to put a little air in there
            # # figa.update_layout(xaxis1=dict(domain=[0.0, 0.3]),xaxis2=dict(domain=[0.31, 0.6]),xaxis3=dict(domain=[0.61, 1]),)

            # figa.update_layout(title={'text': "Plot Title",'y':0.88,'x':0.42,'xanchor': 'left','yanchor': 'top'})

            # # fig.show()
            # with st.expander ('Test plots'):
            #     st.plotly_chart(figa, use_container_width=True)
            
            
            
            
            
            # # https://plotly.com/python/range-slider/
            # try:
            #     # Load data
            #     df = pd.read_csv(
            #         "https://raw.githubusercontent.com/plotly/datasets/master/finance-charts-apple.csv")
            #     df.columns = [col.replace("AAPL.", "") for col in df.columns]
            #     # Create figure
            #     fig3 = go.Figure()
            #     fig3.add_trace(go.Scatter(x=list(df.Date), y=list(df.High),orientation='v'))

            #     # Set title
            #     fig3.update_layout(title_text="Time series with range slider and selectors")

            #     # Add range slider
            #     fig3.update_layout(xaxis=dict(rangeselector=dict(
            #                 buttons=list([
            #                     dict(count=1,
            #                         label="1m",
            #                         step="month",
            #                         stepmode="backward"),
            #                     dict(count=6,
            #                         label="6m",
            #                         step="month",
            #                         stepmode="backward"),
            #                     dict(count=1,
            #                         label="YTD",
            #                         step="year",
            #                         stepmode="todate"),
            #                     dict(count=1,
            #                         label="1y",
            #                         step="year",
            #                         stepmode="backward"),
            #                     dict(step="all")
            #                 ])),rangeslider=dict(visible=True),type="date"))
            #     fig3.show()
            
            # except Exception as e:
            #     print(e)
            #     st.info(f'error vertical slider test: {e}')
            
            # set color vars
            clrGR = '#121212'
            clrPE = '#8ACFDD'
            clrNPHI = '#8ACFDD'
            clrRHOB = '#98C007'
            clrDTS = '#7A7D81'
            clrDTC = '#8ACFDD'
            clrUCS = '#008C9A'
            clrUCS1 = '#3F7077'
            clrUCS2 = '#05929f'
            clrUCS3 = '#8ACFDD'
            clrUCS4 = '#98C21F'
            clrAbr = '#3F94A9'
            clrTA = '#2C4D6E'
            clrImp = '#F0676E'
            clrTI = '#8C0E14'
            
            
            @st.cache_resource()
            def create_rsa_chart():
                # fig1 = make_subplots(rows=1, cols= 4, subplot_titles=('GR','Lithology','Sonic','CUCS'),shared_yaxes=True) 
                # fig1 = make_subplots(rows=1, cols= 6, subplot_titles=['Lithology', 'Gamma Ray', 'Sonic', 'Rock Strength', 'Impact', 'Abrasion'], specs = [[{},{},{},{},{},{}]], horizontal_spacing = 0.01, shared_yaxes=True)
                fig1 = make_subplots(rows=1, cols= 7, specs = [[{},{},{},{},{},{},{}]], horizontal_spacing = 0.02, shared_yaxes=True)
                # fig1.update_layout(xaxis1=dict(domain=[0.0, 0.16]),xaxis2=dict(domain=[0.17, 0.32]),xaxis3=dict(domain=[0.33, 0.48]),xaxis4=dict(domain=[0.49, 0.64]),xaxis5=dict(domain=[0.65, 0.81]),xaxis6=dict(domain=[0.62, 1]),)
                # fig1 = make_subplots(rows=1, cols= 6, shared_yaxes=True)
                # adjust title heights
                # for annotation in fig1['layout']['annotations']: 
                #     annotation['yanchor']='bottom'
                #     annotation['y']=1.1
                #     annotation['yref']='paper'
                
                fig1.update_xaxes(showgrid=True, gridwidth=1,)
                # fig1.update_yaxes(showgrid=True, gridwidth=500,range=[well_data['DEPTH'].min().min(), well_data['DEPTH'].max().max()])
                # pattern fill options  ['', '/', '\\', 'x', '-', '|', '+', '.']
                # old shale gray #bebebe, old sandstone yellow #ffff00
                fig1.add_trace(go.Scatter(x=well_data['SH'], y=well_data['DEPTH'], name='Shale', line= {'color':'#e1e3db', 'width':0}, stackgroup='one',orientation='h', fillpattern=dict(bgcolor='#e1e3db',size=6,shape='-'), hoveron='points',text="Shale",hoverinfo='text+x+y',legendgroup='2'), row=1, col=1)
                fig1.add_trace(go.Scatter(x=well_data['SS'], y=well_data['DEPTH'], name='Sandstone', line={'color':'#f2c949', 'width':0}, stackgroup='one',orientation='h', fillpattern=dict(bgcolor='#f2c949',size=6,shape='.'),hoveron='points',text="Sandstone",hoverinfo='text+x+y',legendgroup='2'), row=1, col=1)
                fig1.add_trace(go.Scatter(x=well_data['LS'], y=well_data['DEPTH'], name='Limestone', line= {'color':'#bebd8f', 'width':0}, stackgroup='one',orientation='h', fillpattern=dict(bgcolor='#bebd8f',size=8,shape='+'),hoveron='points',text="Limestone",hoverinfo='text+x+y',legendgroup='2'), row=1, col=1)
                fig1.add_trace(go.Scatter(x=well_data['SI'], y=well_data['DEPTH'], name='Siltstone', line= {'color':'#6a6239', 'width':0}, stackgroup='one',orientation='h', fillpattern=dict(bgcolor='#6a6239',size=8,shape='-'),hoveron = 'points',text="Siltstone",hoverinfo='text+x+y',legendgroup='2'), row=1, col=1)
                fig1.add_trace(go.Scatter(x=well_data['DO'], y=well_data['DEPTH'], name='Dolomite', line= {'color':'#097384', 'width':0}, stackgroup='one',orientation='h', fillpattern=dict(bgcolor='#097384',size=8,shape='/'),hoveron = 'points',text="Dolomite",hoverinfo='text+x+y',legendgroup='2'), row=1, col=1)
                fig1.add_trace(go.Scatter(x=well_data['AN'], y=well_data['DEPTH'], name='Anhydrite', line= {'color':'#ff80ff', 'width':0}, stackgroup='one',orientation='h', fillpattern=dict(bgcolor='#ff80ff',size=6,shape='x'),hoveron = 'points',text="Anhydrite",hoverinfo='text+x+y',legendgroup='2'), row=1, col=1)
                fig1.add_trace(go.Scatter(x=well_data['SL'], y=well_data['DEPTH'], name='Salt', line= {'color':'#edf7f6', 'width':0}, stackgroup='one',orientation='h', fillpattern=dict(bgcolor='#edf7f6',size=6,shape='+'),hoveron = 'points',text="Salt",hoverinfo='text+x+y',legendgroup='2'), row=1, col=1)
                fig1.add_trace(go.Scatter(x=well_data['CO'], y=well_data['DEPTH'], name='Coal', line= {'color':'#36454f', 'width':0}, stackgroup='one',orientation='h', fillpattern=dict(bgcolor='#36454f',size=8,shape=''),hoveron = 'points',text="Coal",hoverinfo='text+x+y',legendgroup='2'), row=1, col=1)
                fig1.update_xaxes(title_text='Lithology', row=1, col=1)               
                if 'CALI' in well_data.columns:
                    fig1.update_xaxes(title_standoff=1, row=1, col=1)
                # fig1.update_layout(xaxis1=dict(title='Lithology', side='top', position=0.0))
                
                if 'PE' in well_data.columns:
                    fig1.add_trace(go.Scatter(name='PE', x=well_data['PE'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrPE, 'width':0.5}), row=1, col=2)
                    fig1.update_xaxes(title_text='PE', titlefont=dict(color=clrPE),tickfont=dict(color=clrPE), rangemode='tozero', title_standoff=1, zeroline=True, row=1, col=2) 
                
                if 'RHOB' in well_data.columns:
                    fig1.add_trace(go.Scatter(name='RHOB', x=well_data['RHOB'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrRHOB, 'width':0.5}), row=1, col=3)
                    fig1.update_xaxes(title_text='RHOB', titlefont=dict(color=clrRHOB),tickfont=dict(color=clrRHOB), title_standoff=1, rangemode="tozero",zeroline=True,row=1, col=3)
                            
                if 'DTC' in well_data.columns:
                    fig1.add_trace(go.Scatter(name='DTC', x=well_data['DTC'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrDTC, 'width':0.5}), row=1, col=4)
                    fig1.update_xaxes(title_text='DTC', titlefont=dict(color=clrDTC),tickfont=dict(color=clrDTC), title_standoff=1, rangemode="tozero",zeroline=True,row=1, col=4)
                    # fig1.update_layout(xaxis3=dict(type='linear',range=[140, 40])) 
                    
                if 'CUCS' in well_data.columns:                
                    try:
                        fig1.add_trace(go.Scatter(name='UCS1', x=well_data['CUCS'], y=well_data['DEPTH'],showlegend=False, fill='tozerox',line= {'color':well_data['CUCS'], 'width':0.5}), row=1, col=5)
                    except Exception as e:
                        print(e)
                        try:
                            fig1.add_trace(go.Scatter(name='UCS2', x=well_data['CUCS'], y=well_data['DEPTH'],showlegend=False, fill='tozerox',marker=dict(size=2,color=well_data['CUCS'],colorbar=dict(title="Colorbar"),colorscale="Plasma",showscale=False),mode="markers"), row=1, col=5)
                        except Exception as e:
                            print(e)
                            fig1.add_trace(go.Scatter(name='UCS3', x=well_data['CUCS'], y=well_data['DEPTH'],showlegend=False, fill='tozerox',line= {'color':clrUCS, 'width':0.5}), row=1, col=5)
                    
                    # !!!!!!!!!!!!!!!!!!!!!color by value technique https://stackoverflow.com/questions/73597115/customize-color-of-bar-chart-based-on-value-in-plotly
                    # fig = go.Figure(data=[go.Bar(x=x, y=y,marker=dict(color=np.cos(x / 3),colorscale="Bluered_r",showscale=True,colorbar=dict(title="value"),),)])
                                
                    # # fig1.add_trace(go.Scatter(name='UCS2', x=well_data['CUCS'], y=well_data['DEPTH'],showlegend=False,marker=dict(size=2,color=well_data['CUCS'],colorbar=dict(title="Colorbar"),colorscale="Plasma",showscale=False),mode="markers"), row=1, col=5)   
                    # fig1.add_trace(go.Scatter(name='UCS2', x=well_data['CUCS'], y=well_data['DEPTH'],showlegend=False,line= {'color':clrUCS1, 'width':0.5}), row=1, col=5)   
                    # fig1.add_trace(go.Scatter(name='UCS2', x=well_data['CUCS'].where(well_data['CUCS'] <= np.percentile(well_data['CUCS'], 25)), y=well_data['DEPTH'],showlegend=False,hoverinfo='skip',line= {'color':clrUCS1, 'width':0.5}), row=1, col=5)   
                    # fig1.add_trace(go.Scatter(name='UCS2', x=well_data['CUCS'].where(well_data['CUCS'] >= np.percentile(well_data['CUCS'], 75)), y=well_data['DEPTH'],showlegend=False,hoverinfo='skip',line= {'color':clrUCS4, 'width':0.5}), row=1, col=5)   
                    
                    fig1.update_xaxes(title_text='UCS (Calculated)',titlefont=dict(color=clrUCS),tickfont=dict(color=clrUCS), title_standoff=1, rangemode="tozero",zeroline=True,row=1, col=5)
                # fig.add_trace(go.Scatter(x=values,y=values,marker=dict(size=16,cmax=39,cmin=0,color=values,colorbar=dict(title="Colorbar"),colorscale="Viridis"),mode="markers"))
                        
                # Impact Channel            
                # x5 = go.XAxis(title= 'Instant', type='log', range=[0,3], side='top', position=1, titlefont=go.Font(color='SteelBlue'))
                # x15 = go.XAxis(title= 'Total', side='top', position=0.5, titlefont=go.Font(color='DarkOrange'))
                
                # example of filled line chart by colorscale https://community.plotly.com/t/plotly-graph-object-scatter-plot-fill-color-map/71479/9
                # try:
                #     fig1.add_trace(go.Contour(z=well_data['Impact'], colorscale='Inferno', contours_coloring='heatmap', line_width=0), row=1, col=5)
                # except Exception as e:
                #     print(e)
                #     st.info(f'error colorscale fill: {e}')
                
                # fig1.update_layout(xaxis5=dict(domain=[0.0, 0.4]))
                # fig1.update_xaxes(title_text='Impact', type='log', range=[0,3], row=1, col=5)
                # fig1.update_layout(xaxis2=dict(domain=[0.6, 0.9]),
                #     yaxis2=dict(title=“yaxis1 title”,titlefont=dict(color=“red”), tickfont=dict(color=“red”)),
                #     yaxis3=dict(title=“yaxis2 title”,titlefont=dict(color=“orange”),tickfont=dict(color=“orange”), anchor=“free”, overlaying=“y2”, side=“left”, position=0.5),
                #     yaxis4=dict(title=“yaxis3 title”,titlefont=dict(color=“pink”),tickfont=dict(color=“pink”), anchor=“x2”,overlaying=“y2”, side=“right”),
                #     yaxis5=dict(title=“yaxis4 title”,titlefont=dict(color=“cyan”),tickfont=dict(color=“cyan”), anchor=“free”,overlaying=“y2”, side=“right”,position=1))
                # example of secondary x axis 
                # https://stackoverflow.com/questions/61803438/python-plotly-figure-with-secondary-x-axis-linked-to-primary
                
                
                # Abrasion Channel    
                
                if 'Total Abrasion' in well_data.columns:      
                    # x6 = go.XAxis(title= 'Instant', type='log', range=[0,2], side='top', position=1, titlefont=go.Font(color='SteelBlue'))
                    # x16 = go.XAxis(title= 'Total', side='top', position=0.5,  titlefont=go.Font(color='DarkOrange'))
                    fig1.add_trace(go.Scatter(name='Total Abrasion',x=well_data['Total Abrasion'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrTA, 'width':0.5}), row=1, col=6) 
                    # fig1.update_layout(xaxis6=dict(title='Abrasion',titlefont=dict(color='#121212'),tickfont=dict(color='#121212'), type='log', range=[0,2], side='top',position=0.0))
                    # fig1.add_trace(go.Scatter(x=well_data['Abrasion'], y=well_data['DEPTH'],showlegend=False, line= {'color':'#121212', 'width':0.4}, marker=dict(cmax=well_data['Abrasion'].astype(int).max(),cmin=0,color=list(range(well_data['Abrasion'].astype(int).max())),colorbar=dict(title='Colorbar'),colorscale='Viridis'),), row=1, col=6)       
                    fig1.update_xaxes(title_text='Total Abrasion',titlefont=dict(color=clrTA),tickfont=dict(color=clrTA), title_standoff=1, rangemode="tozero",zeroline=True,row=1, col=6)
                
                if 'Total Impact' in well_data.columns:
                    fig1.add_trace(go.Scatter(name='Total Impact',x=well_data['Total Impact'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrTI, 'width':0.5}), row=1, col=7) 
                    
                    # fig1.update_layout(xaxis5=dict(title='Impact',titlefont=dict(color='#121212'),tickfont=dict(color='#121212'), type='log', range=[0,3],side='top',position=0.0))          
                    # fig1.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.7,y=1.045,text=(f'Instant'),textangle=0,showarrow=False,font=dict(color='#121212',size=12))
                    fig1.update_xaxes(title_text='Total Impact',titlefont=dict(color=clrTI),tickfont=dict(color=clrTI), title_standoff=1, rangemode="tozero",zeroline=True,row=1, col=7) # type='log',
                
                #  ADD SECONDARY TRACES AND X AXIS            
                # if set(['GR']).issubset(well_data.columns):    
                if 'CALI' in well_data.columns:
                    fig1.add_trace(go.Scatter(name='CALI', x=well_data['CALI'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrGR, 'width':0.4}, xaxis='x14',yaxis = 'y1')) 
                    fig1.update_layout(xaxis14=dict(title='CALI',titlefont=dict(color=clrGR),tickfont=dict(color=clrGR), rangemode="tozero",anchor='free',overlaying='x1',side='top',title_standoff=5,position=1.0))
                
                if 'GR' in well_data.columns:
                    fig1.add_trace(go.Scatter(name='GR', x=well_data['GR'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrGR, 'width':0.4}, xaxis='x8',yaxis = 'y2')) 
                    # fig1.add_trace(go.Scatter(name='GR', x=well_data['GR'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrGR, 'width':0.4}, xaxis='x8',yaxis = 'y2')) 
                    well_data['GR'].where(well_data['GR'] < 40)
                    fig1.update_layout(xaxis8=dict(title='Gamma Ray',titlefont=dict(color=clrGR),tickfont=dict(color=clrGR), rangemode="tozero",anchor='free',overlaying='x2',side='top',title_standoff=5,position=1.0))
                
                if 'NPHI' in well_data.columns:
                    fig1.add_trace(go.Scatter(name='NPHI', x=well_data['NPHI'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrNPHI, 'width':0.4}, xaxis='x9',yaxis = 'y3'))
                    fig1.update_layout(xaxis9=dict(title='NPHI',titlefont=dict(color=clrNPHI),tickfont=dict(color=clrNPHI), rangemode="tozero",anchor='free',overlaying='x3',side='top',title_standoff=5,position=1.0))
                    
                if 'DTS' in well_data.columns:
                    fig1.add_trace(go.Scatter(name='DTS', x=well_data['DTS3'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrDTS, 'width':0.4}, xaxis='x10',yaxis = 'y4'))
                    fig1.update_layout(xaxis10=dict(title='DTS',titlefont=dict(color=clrDTS),tickfont=dict(color=clrDTS), rangemode="tozero",anchor='free',overlaying='x4',side='top',title_standoff=5,position=1.0)) 
                            
                if 'Abrasion' in well_data.columns:
                    fig1.add_trace(go.Scatter(name='Abrasion',x=well_data['Abrasion'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrAbr, 'width':1}, xaxis='x12',yaxis = 'y6'))    
                    fig1.update_layout(xaxis12=dict(title='Abrasion',titlefont=dict(color=clrAbr),tickfont=dict(color=clrAbr), anchor='free',overlaying='x6',rangemode="tozero",zeroline=True,side='top',title_standoff=5,position=1.0))
                                            
                if 'Impact' in well_data.columns:
                    fig1.add_trace(go.Scatter(name='Impact',x=well_data['Impact'], y=well_data['DEPTH'],showlegend=False, line= {'color':clrImp, 'width':1}, xaxis='x13',yaxis = 'y7'))   
                    fig1.update_layout(xaxis13=dict(title='Impact',titlefont=dict(color=clrImp),tickfont=dict(color=clrImp), anchor='free',overlaying='x7',rangemode="tozero",zeroline=True,side='top',title_standoff=5,position=1.0))
                                    
                # fig1.add_annotation(xref='paper',yref='paper',xanchor='left',x=0.77,y=1.045,text=(f'Total'),textangle=0,showarrow=False,font=dict(color='#008C9A',size=12))  
                
                # fig1.update_layout(xaxis6=dict(domain=[0.0, 0.4]))
                # https://stackoverflow.com/questions/71185900/python-plotly-apply-log-scale-on-a-specific-axis-by-index
                
                # if 'PE' in well_data.columns:
                #     fig1.update_layout(xaxis8=dict(title='PE',titlefont=dict(color='#8ACFDD'),tickfont=dict(color='#8ACFDD'), rangemode="tozero",anchor='free',overlaying='x2',side='top',title_standoff=5,position=1.0))
                
                # if 'DTS' in well_data.columns:
                #     fig1.update_layout(xaxis10=dict(title='DTS',titlefont=dict(color='#8ACFDD'),tickfont=dict(color='#8ACFDD'), rangemode="tozero",anchor='free',overlaying='x4',side='top',title_standoff=5,position=1.0)) 
                
                # if 'Total Abrasion' in well_data.columns:
                #     fig1.update_layout(xaxis12=dict(title='Total Abrasion',titlefont=dict(color='#008C9A'),tickfont=dict(color='#008C9A'), anchor='free',overlaying='x6',rangemode="tozero",zeroline=True,side='top',title_standoff=5,position=1.0))
                        
                # if 'Total Impact' in well_data.columns:
                #     fig1.update_layout(xaxis13=dict(title='Total Impact',titlefont=dict(color='#008C9A'),tickfont=dict(color='#008C9A'), anchor='free',overlaying='x7',rangemode="tozero",zeroline=True,side='top',title_standoff=5,position=1.0))
                    
                # fig1.update_layout(yaxes=dict(domain=[0, 0.95]))
                fig1.update_layout(title={'text': "Rock Strength Analysis",'y':1,'x':0.42,'xanchor': 'left','yanchor': 'top'})
                fig1.update_layout(legend=dict(orientation="v",yanchor="auto",y=1,xanchor="right",  x=-0.03))
                # hoverstate fix for filled area charts https://community.plotly.com/t/values-0-on-stacked-area-chart/28301
                # autorange='reversed' creates the 5% extra chart space on top and bottom per https://community.plotly.com/t/reversed-axis-with-range-specified/3806/3 
                fig1.update_layout(autosize=True, height=1000,showlegend=True,legend_tracegroupgap = 180, hovermode = "y", yaxis={'title':'DEPTH'},) # ,'autorange':'reversed'
                fig1.update_yaxes(type='linear',range=[well_data['DEPTH'].max(),well_data['DEPTH'].min()], tickformat = '000',constrain='domain', nticks=20 ) # matches='y', dtick=500
                fig1.update_layout(
                    yaxis_tickformatstops = [
                        dict(dtickrange=[None, 1000], value="%H:%M:%S.%L ms"),
                        dict(dtickrange=[1000, 60000], value="%H:%M:%S s"),
                        dict(dtickrange=[60000, 3600000], value="%H:%M m"),
                        dict(dtickrange=[3600000, 86400000], value="%H:%M h"),
                        dict(dtickrange=[86400000, 604800000], value="%e. %b d"),
                        dict(dtickrange=[604800000, "M1"], value="%e. %b w"),
                        dict(dtickrange=["M1", "M12"], value="%b '%y M"),
                        dict(dtickrange=["M12", None], value="%Y Y")
                    ]
                )

                # fig1.update_layout(margin=dict(l=140,r=10,t=100,b=1,autoexpand=False))
                # # fig1.update_layout(margin=dict(r=10,t=100,b=1,autoexpand=False))
                fig1.update_layout(xaxis={'side':'top'}, xaxis2={'side':'top'}, xaxis3={'side':'top'}, xaxis4={'side':'top'}, xaxis5={'side':'top'}, xaxis6={'side':'top'},xaxis7={'side':'top'})
                fig1.update_layout(yaxis=dict(domain=[0, 0.95]),yaxis2=dict(domain=[0, 0.95]),yaxis3=dict(domain=[0, 0.95]),yaxis4=dict(domain=[0, 0.95]),yaxis5=dict(domain=[0, 0.95]),yaxis6=dict(domain=[0, 0.95]),yaxis7=dict(domain=[0, 0.95]))
                # fig1.update_layout(yaxis=dict(domain=[0, 0.95],range=[well_data['DEPTH'].min(),well_data['DEPTH'].max()]),
                #                    yaxis2=dict(domain=[0, 0.95],range=[well_data['DEPTH'].min(),well_data['DEPTH'].max()]),
                #                    yaxis3=dict(domain=[0, 0.95],range=[well_data['DEPTH'].min(),well_data['DEPTH'].max()]),
                #                    yaxis4=dict(domain=[0, 0.95],range=[well_data['DEPTH'].min(),well_data['DEPTH'].max()]),
                #                    yaxis5=dict(domain=[0, 0.95],range=[well_data['DEPTH'].min(),well_data['DEPTH'].max()]),
                #                    yaxis6=dict(domain=[0, 0.95],range=[well_data['DEPTH'].min(),well_data['DEPTH'].max()]),
                #                    yaxis7=dict(domain=[0, 0.95],range=[well_data['DEPTH'].min(),well_data['DEPTH'].max()])
                #                    )
                fig1.update_layout(xaxis=dict(type='linear',range=[0, 100],ticksuffix='%')) 
                # fig1.update_layout(yaxis=dict(type='linear',range=[well_data['DEPTH'].min(),well_data['DEPTH'].max()], tickformat = '000'),dtick=500) # ,dtick=500
                
                
                fig1.update_xaxes(ticks="inside")
                
                # minDep = min(well_data['DEPTH'].notnull(), key=float)
                # maxDep = max(well_data['DEPTH'].notnull(), key=float)
                # ymin = float(minDep)
                # ymax = float(maxDep)
                # fig1.update_layout(yaxis=dict(range=[ymin, ymax]))
                
                # attempt to keep margin the same when chart gets taller
                # fig1.update_yaxes(automargin=True)
                # fig1.update_layout(yaxis={range: [well_data['DEPTH'].min(), well_data['DEPTH'].max()]})
                # fig1.update_layout(xaxis5=dict(type='linear',range=[1, 6])) 
                fig1.update_layout(barmode='group')
                # fig1.update_bars(bargap=0.30,bargroupgap=0.0)
                # xaxis=dict(type='linear',range=[1, 100],ticksuffix='%')
                fig1.layout.template='seaborn'
                fig1.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(245,245,245,0.5)')
                fig1.update_layout(margin=dict(l=140,r=10,t=100,b=1,autoexpand=False))
                fig1.update_layout(updatemenus=[dict(type = "dropdown",direction = "down",buttons=list([dict(args=["height", "1000"],label="Height: 1000px",method="relayout"),dict(args=["height", "2000"],label="2000px",method="relayout"),dict(args=["height", "5000"],label="5000px",method="relayout")]),
                                                    pad={"r": 10, "t": 10},showactive=True,bgcolor='White',x=0.02,xanchor="left",y=1.1,yanchor="top"),])
                # ,{'xaxis12': {'title_standoff': '5'}},
                #,{'xaxis12':{'domain':{'0', '0.95'}}},
                # fig1.update_layout(annotations=[dict(text="Chart Height:", showarrow=False,x=0, y=1.08, yref="paper", align="left")])
                # Menu button color code https://stackoverflow.com/questions/72532725/change-background-color-of-menu-button-in-plotly
                
                # col2_h.plotly_chart(fig1, use_container_width=True)
                st.plotly_chart(fig1, use_container_width=True)
                # 'xaxis12': {'title_standoff': '5'},
                # xaxis12=dict(title_standoff=5,position=1.0)
            
            # run chart
            create_rsa_chart()
            
            #  Crossplots 
            col1, col2 = st.columns(2)
            
            
            if 'CUCS' in well_data.columns:
                # col1.write('Strength Histogram')
                histogram = px.histogram(well_data, x='CUCS', log_x=False, title='Rock Strength Histogram')
                histogram.update_traces(marker_color='black')            
                # # The two histograms are drawn on top of another https://plotly.com/python/histograms/
                # fig.add_trace(go.Histogram(x=x1))
                # fig.update_layout(barmode='stack')
                
                # histogram.layout.template='seaborn'
                col1.plotly_chart(histogram, use_container_width=True)
            
            
            if set(['NPHI','RHOB', 'DEPTH']).issubset(well_data.columns):
                # col2.write('Porosity / Density Crossplot') 
                # https://plotly.com/python/builtin-colorscales/           
                xplot = px.scatter(well_data, x='NPHI', y='RHOB', color='DEPTH', color_continuous_scale='Plasma', log_x=True, log_y=True,title='Porosity / Density Crossplot')
                # xplot.layout.template='seaborn'
                col2.plotly_chart(xplot, use_container_width=True)
            
            # fig9c = go.Figure()
            # fig9c.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['SH'], mode='lines', line=dict(width=0.5, color='gray'), hoveron = 'points+fills', text="SH", hoverinfo='text+y', name='Shale', stackgroup='one', groupnorm='percent'))
            # fig9c.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['SS'], mode='lines', line=dict(width=0.5, color='#fefaea'), hoveron = 'points+fills', text="SS", hoverinfo='text+y', name='Sandstone', stackgroup='one')) 
            # fig9c.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['LS'], mode='lines', line=dict(width=0.5, color='#fefaea'), hoveron = 'points+fills', text="LS", hoverinfo='text+y', name='Limestone', stackgroup='one')) 
            # fig9c.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['SI'], mode='lines', line=dict(width=0.5, color='#fefaea'), hoveron = 'points+fills', text="SI", hoverinfo='text+y', name='Siltstone', stackgroup='one')) 
            # fig9c.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['DO'], mode='lines', line=dict(width=0.5, color='#fefaea'), hoveron = 'points+fills', text="DO", hoverinfo='text+y', name='Dolomite', stackgroup='one')) 
            # fig9c.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['AN'], mode='lines', line=dict(width=0.5, color='#fefaea'), hoveron = 'points+fills', text="AN", hoverinfo='text+y', name='Anhydrite', stackgroup='one')) 
            # fig9c.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['SL'], mode='lines', line=dict(width=0.5, color='#fefaea'), hoveron = 'points+fills', text="SL", hoverinfo='text+y', name='SL', stackgroup='one')) 
            # fig9c.add_trace(go.Scatter(x=well_data['DEPTH'], y=well_data['CO'], mode='lines', line=dict(width=0.5, color='#fefaea'), hoveron = 'points+fills', text="CO", hoverinfo='text+y', name='CO', stackgroup='one'))
            # # fig9b.update_traces(orientation='h')
            # fig9c.update_layout(height=350, showlegend=True, xaxis={'title':'DEPTH','autorange':'reversed'}, title='LithologyC',
            #     yaxis=dict(
            #         type='linear',
            #         range=[1, 100],
            #         ticksuffix='%'))
            # fig9c.layout.template='seaborn'
            # # fig9b.update_layout(paper_bgcolor="rgb(10,10,10)", plot_bgcolor="rgb(10,10,10)")
            # st.plotly_chart(fig9c, use_container_width=True)
            
            
            # # Altair chart text code 
            # df3 = pd.DataFrame(
            #     np.random.rand(10, 4),
            #     columns= ["NO2","C2H5CH","VOC","CO"])
            # # generate a date range to be used as the x axis
            # df3['date'] =  pd.date_range("2014-01-01", periods=10, freq="m")
            # df_melted = pd.melt(df3,id_vars=['date'],var_name='parameter', value_name='value')
            # c = alt.Chart(df_melted, title='measure of different elements over time').mark_line().encode(
            #     x='date', y='value', color='parameter')
            # st.altair_chart(c, use_container_width=True)


            # x=['Winter', 'Spring', 'Summer', 'Fall']
            # fig2 = go.Figure()
            # fig2.add_trace(go.Scatter(
            #     x=x, y=[40, 20, 30, 40],
            #     mode='lines',
            #     line=dict(width=0.5, color='rgb(184, 247, 212)'),
            #     stackgroup='one',
            #     groupnorm='percent' # sets the normalization for the sum of the stackgroup
            # ))
            # fig2.add_trace(go.Scatter(
            #     x=x, y=[50, 70, 40, 60],
            #     mode='lines',
            #     line=dict(width=0.5, color='rgb(111, 231, 219)'),
            #     stackgroup='one'
            # ))
            # fig2.add_trace(go.Scatter(
            #     x=x, y=[70, 80, 60, 70],
            #     mode='lines',
            #     line=dict(width=0.5, color='rgb(127, 166, 238)'),
            #     stackgroup='one'
            # ))
            # fig2.add_trace(go.Scatter(
            #     x=x, y=[100, 100, 100, 100],
            #     mode='lines',
            #     line=dict(width=0.5, color='rgb(131, 90, 241)'),
            #     stackgroup='one'
            # ))
            # fig2.update_traces(orientation='h')            
            # fig2.update_layout(
            #     showlegend=True,
            #     xaxis_type='category',
            #     yaxis=dict(
            #         type='linear',
            #         range=[1, 100],
            #         ticksuffix='%'))                        
            # col2_h.plotly_chart(fig2, use_container_width=True)

            
            # uRSA = px.area(well_data, x="DEPTH", y=["SS","SH"], color=well_data.columns, pattern_shape=["SS","SH"], pattern_shape_sequence=[".", "x", "+"])
            # col2_h.plotly_chart(uRSA, use_container_width=True)
            
            # try 2
            # x=well_data["DEPTH"]

            # fig = go.Figure()
            # fig.add_trace(go.Scatter(
            #     x=x, y=well_data["SS"],
            #     hoverinfo='x+y',
            #     mode='lines',
            #     line=dict(width=0.0, color='rgb(131, 90, 241)'),
            #     stackgroup='one' # define stack group
            # ))
            # fig.add_trace(go.Scatter(
            #     x=x, y=well_data["SS"],
            #     hoverinfo='x+y',
            #     mode='lines',
            #     line=dict(width=0.0, color='rgb(111, 231, 219)'),
            #     stackgroup='one'
            # ))

            # fig.update_layout(yaxis_range=(0, 100))            
            # col2_h.write(fig)
            
            # st.table(data=well_data)
            
            # @st.cache 
            # @st.cache(hash_funcs={builtins.function: create_adira_poster})
            
            # https://pptx-generator.streamlit.app/
            def create_pptx_poster():                
                                
                save_bar.progress(1, text='Initializing..') 
                pptx = 'data/template_poster.pptx'
                prs = Presentation(pptx)
                
                # declare positional variables
                WIDTH = Inches(24)
                HEIGHT = Inches(36)
                left = Inches(2.5)
                top = Inches(1)
                
                # get stock info
                name = 'Ulterra'
                
                try:    
                    save_bar.progress(5, text='Creating template..') 
                    lyt=prs.slide_masters[0].slide_layouts[1] # choosing a slide layout
                    
                    #  Pg1
                    save_bar.progress(10, text='Poster: Setup..') 
                    slide1=prs.slides[0]
                    txBox = slide1.shapes.add_textbox(Inches(2.14), Inches(1.63), Inches(2), Inches(1))
                    tf = txBox.text_frame
                    tf.text = f'REPORT'
                    tf.paragraphs[0].font.size = Pt(24)
                    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  
                    
                    rows = 4		# number of rows in table
                    cols = 6		# number of columns in table
                    # Add Shape object (table)
                    # Arguments: number of rows, number of columns, x-coordinate of top-left corner, x-coordinate of top-left corner, width,  height
                    table_shape = slide1.shapes.add_table(rows, cols, left=Inches(7), top=Inches(0.42), width=Inches(14), height=Inches(2))
                    table = table_shape.table	# Create Table object

                    tbl =  table_shape._element.graphic.graphicData.tbl
                    # List of Style IDs: https://github.com/scanny/python-pptx/issues/27#issuecomment-263076372
                    style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
                    tbl[0][-1].text = style_id
                    # table.first_row = False 
                    # table.first_col = False                    
                    # table.horz_banding = False
                    # table.vert_banding = False
                    
                except Exception as e:
                    print(e)
                    st.info(f"error ppt intro: {e}") 
                
                # # Column Heading Text Settings
                # category = ['Category_1', 'Category_2',  'Category_3']
                # for i in range(len(category)):
                #     cell = table.cell(i+1, 0)	 # Getting a Cell object
                #     cell.text = category[i]	   # Set the value with the text property
                
                try:
                    # table.cell(0, 0).text = str(las_file.well.COMP.descr.capitalize())
                    # table.cell(0, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                    # # table.cell(0, 0).fill.background()     
                    # try:               
                    #     table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(18)  
                    #     table.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    # except Exception as e:
                    #     print(e)
                    #     st.info(f"error ppt font color: {e}") 
                    
                    save_bar.progress(15, text='Poster: Header..') 
                    def set_cell_info(tablecell, text, aligned):                        
                        try:
                            tablecell.text =  text
                            tablecell.text_frame.paragraphs[0].alignment = aligned
                            tablecell.text_frame.paragraphs[0].font.size = Pt(16)  
                            tablecell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)                        
                        except Exception as e:
                            print(e)
                            st.info(f"error ppt tablecell: {tablecell}: {e}") 
                            
                        return tablecell
                    
                    # table.cell.fill.solid()
                    # table.cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    set_cell_info(table.cell(0, 0), str(las_file.well.COMP.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(1, 0), str(las_file.well.WELL.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(2, 0), str(las_file.well.UWI.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(3, 0), str(las_file.well.LOC.descr.capitalize()), PP_ALIGN.RIGHT)                    
                                        
                    set_cell_info(table.cell(0, 1), str(las_file.well.COMP.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(1, 1), str(las_file.well.WELL.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(2, 1), str(las_file.well.UWI.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(3, 1), str(las_file.well.LOC.value), PP_ALIGN.LEFT)                    
                    
                    set_cell_info(table.cell(0, 2), str(las_file.well.FLD.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(1, 2), str(las_file.well.CNTY.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(2, 2), str(las_file.well.STAT.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(3, 2), str(las_file.well.CTRY.descr.capitalize()), PP_ALIGN.RIGHT)                   
                    
                    set_cell_info(table.cell(0, 3), str(las_file.well.FLD.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(1, 3), str(las_file.well.CNTY.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(2, 3), str(las_file.well.STAT.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(3, 3), str(las_file.well.CTRY.value), PP_ALIGN.LEFT)
                                    
                    set_cell_info(table.cell(0, 4), str(las_file.well.STRT.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(1, 4), str(las_file.well.STOP.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(2, 4), str(las_file.well.STEP.descr.capitalize()), PP_ALIGN.RIGHT)
                                        
                    set_cell_info(table.cell(0, 5), str(las_file.well.STRT.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(1, 5), str(las_file.well.STOP.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(2, 5), str(las_file.well.STEP.value), PP_ALIGN.LEFT)
                    
                except Exception as e:
                    print(e)
                    # st.info(f"File missing header info: {e}") 
                                    
                # # Set values in other cells
                # value1 = [[1, 2, 3],[4, 5, 6],[7, 8, 9]]
                # for i in range(len(value1)):
                #     for k in range(len(value1[i])):
                #         cell = table.cell(k+1, i+1)	   # Getting a cell object
                #         cell.text = str(value1[i][k])	 # Set the value with the text property
                
                # # Column Heading Format
                # for i in range(len(table.columns)):	        # Getting the number of table columns 
                #     cell = table.cell(0, i)
                #     pg = cell.text_frame.paragraphs[0]	    # Getting paragraph object
                #     pg.font.size = Pt(15)		                # Setting the font size of paragraphs
                #     pg.font.color.rgb = RGBColor(255, 0, 0) # Setting the font color for paragraphs
                #     pg.aligment = PP_ALIGN.CENTER	          # Setting paragraph string positioning (centered)

                # #---------------------------------------------------------------------------------------
                # # Row Heading Format
                # for i in range(len(table.rows)):   # Getting the number of table rows
                #     cell = table.cell(i, 0)
                #     cell.fill.solid()			         # fill a cell
                #     cell.fill.fore_color.rgb = RGBColor(100, 200, 50)   # Specify cell fill color
                
                
                # chartbk = slide1.shapes.add_shape(autoshape_type_id=MSO_SHAPE.RECTANGLE,left=Inches(1), top=Inches(3), width=Inches(23), height=Inches(31.3))
                # chartbk.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                try:
                    save_bar.progress(20, text='Poster: Background..') 
                    # https://stackoverflow.com/questions/38202582/how-do-i-add-transparency-to-shape-in-python-pptx
                    def SubElement(parent, tagname, **kwargs):
                        element = OxmlElement(tagname)
                        element.attrib.update(kwargs)
                        parent.append(element)
                        return element

                    def _set_shape_transparency(shape, alpha):
                        """ Set the transparency (alpha) of a shape"""
                        ts = shape.fill._xPr.solidFill
                        sF = ts.get_or_change_to_srgbClr()
                        sE = SubElement(sF, 'a:alpha', val=str(alpha))

                    ##Add a  box to the slide                
                    chartbk = slide1.shapes.add_shape(autoshape_type_id=MSO_SHAPE.RECTANGLE,left=Inches(1), top=Inches(3), width=Inches(23), height=Inches(30.86))
                    ## Make the box 
                    chartbkFill = chartbk.fill
                    chartbkFill.solid()
                    chartbkFillColour = chartbkFill.fore_color
                    chartbkFillColour.rgb = RGBColor(255, 255, 255)
                    ## Set the transparency to 30%
                    _set_shape_transparency(chartbk,70000)
                    # https://python-pptx.readthedocs.io/en/latest/user/autoshapes.html#line
                    chartbk.line.width = Pt(0)
                    chartbk.line.color.rgb = RGBColor(255, 255, 255)
                except Exception as e:
                    print(e)
                    st.info(f"error ppt backshape: {e}") 

                try:
                    save_bar.progress(30, text='Poster: Channels..') 
                    filename = "figure1.png"
                    # pio.write_image(fig1, filename, scale=2, width=2208, height=3030)                    
                    pio.write_image(fig1, filename, scale=2, width=2182, height=2675)
                    slide1.shapes.add_picture(filename, left=Inches(0.38), top=Inches(3), width=Inches(22.72)) 
                    
                    save_bar.progress(50, text='Poster: Histogram..') 
                    filename = "histogram.png"
                    pio.write_image(histogram, filename, scale=2, width=873, height=450)
                    slide1.shapes.add_picture(filename, left=Inches(2.91), top=Inches(29.61), width=Inches(9.09)) 
                    
                    
                    save_bar.progress(75, text='Poster: Crossplot..') 
                    filename = "xplot.png"
                    pio.write_image(xplot, filename, scale=2, width=873, height=450)
                    slide1.shapes.add_picture(filename, left=Inches(14.01), top=Inches(29.61), width=Inches(9.09)) 
                    
                except Exception as e:
                    print(e)
                    st.info(f"error ppt pic: {e}") 
                    
                # # Pg2
                # try:
                #     slide2=prs.slides.add_slide(lyt) # adding a slide 
                #     txBox = slide2.shapes.add_textbox(Inches(0.87), Inches(0.4), Inches(1), Inches(0.25))
                #     tf = txBox.text_frame
                #     tf.text = f'OVERVIEW'
                #     tf.paragraphs[0].font.size = Pt(8)  
                #     #  Depth vs Time
                #     filename = "fig1.jpg"
                #     pio.write_image(fig1, filename, scale=2, width=(72*4.8))             
                #     slide2.shapes.add_picture(filename, left=Inches(0.1), top=Inches(1.36), width=Inches(4.8))        
                        
                # except Exception as e:
                #     print(e)
                #     st.info(f"error ppt2: {e}")
                
                # Close out
                save_bar.progress(80, text='Saving..') 
                binary_output = BytesIO()
                prs.save(binary_output)
                return binary_output
            
            
            def create_pptx_multipage(pagefootage):                
                    
                save_bar.progress(1, text='Initializing..')             
                pptx = 'data/template_portrait.pptx'
                prs = Presentation(pptx)
                
                # declare positional variables
                WIDTH = Inches(7.5)
                HEIGHT = Inches(10)
                # left = Inches(2.5)
                # top = Inches(1)
                
                # get stock info
                name = 'Ulterra'
                
                try:    
                    save_bar.progress(5, text='Creating template..') 
                    lyt=prs.slide_masters[0].slide_layouts[1] # choosing a slide layout
                    
                    #  Pg1
                    slide1=prs.slides[0]
                    txBox = slide1.shapes.add_textbox(left=Inches(0.7), top=Inches(0.44), width=Inches(3), height=Inches(0.3))
                    tf = txBox.text_frame
                    tf.text = f'ROCK REPORT'
                    tf.paragraphs[0].font.size = Pt(10)
                    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  
                    # tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  
                    
                    rows = 4		# number of rows in table
                    cols = 6		# number of columns in table
                    # Add Shape object (table)
                    # Arguments: number of rows, number of columns, x-coordinate of top-left corner, x-coordinate of top-left corner, width,  height
                    table_shape = slide1.shapes.add_table(rows, cols, left=Inches(0.5), top=Inches(1), width=Inches(6.5), height=Inches(2))
                    table = table_shape.table	# Create Table object

                    tbl =  table_shape._element.graphic.graphicData.tbl
                    # List of Style IDs: https://github.com/scanny/python-pptx/issues/27#issuecomment-263076372
                    style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
                    tbl[0][-1].text = style_id
                    
                except Exception as e:
                    print(e)
                    st.info(f"error ppt intro: {e}") 
                
                try:
                    save_bar.progress(15, text='Page 1: Header..')                     
                    def set_cell_info(tablecell, text, aligned):                        
                        try:
                            tablecell.text =  text
                            tablecell.text_frame.paragraphs[0].alignment = aligned
                            tablecell.text_frame.paragraphs[0].font.size = Pt(8)  
                            tablecell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)                        
                        except Exception as e:
                            print(e)
                            st.info(f"error ppt tablecell: {tablecell}: {e}") 
                            
                        return tablecell
                    
                    set_cell_info(table.cell(0, 0), str(las_file.well.COMP.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(1, 0), str(las_file.well.WELL.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(2, 0), str(las_file.well.UWI.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(3, 0), str(las_file.well.LOC.descr.capitalize()), PP_ALIGN.RIGHT)
                    
                    set_cell_info(table.cell(0, 1), str(las_file.well.COMP.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(1, 1), str(las_file.well.WELL.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(2, 1), str(las_file.well.UWI.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(3, 1), str(las_file.well.LOC.value), PP_ALIGN.LEFT)
                    
                    set_cell_info(table.cell(0, 2), str(las_file.well.FLD.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(1, 2), str(las_file.well.CNTY.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(2, 2), str(las_file.well.STAT.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(3, 2), str(las_file.well.CTRY.descr.capitalize()), PP_ALIGN.RIGHT)
                    
                    set_cell_info(table.cell(0, 3), str(las_file.well.FLD.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(1, 3), str(las_file.well.CNTY.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(2, 3), str(las_file.well.STAT.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(3, 3), str(las_file.well.CTRY.value), PP_ALIGN.LEFT)
                    
                    set_cell_info(table.cell(0, 4), str(las_file.well.STRT.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(1, 4), str(las_file.well.STOP.descr.capitalize()), PP_ALIGN.RIGHT)
                    set_cell_info(table.cell(2, 4), str(las_file.well.STEP.descr.capitalize()), PP_ALIGN.RIGHT)
                    
                    set_cell_info(table.cell(0, 5), str(las_file.well.STRT.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(1, 5), str(las_file.well.STOP.value), PP_ALIGN.LEFT)
                    set_cell_info(table.cell(2, 5), str(las_file.well.STEP.value), PP_ALIGN.LEFT)
                    
                except Exception as e:
                    print(e)
                    st.info(f"File missing header info: {e}") 
                      

                try:
                    # https://stackoverflow.com/questions/61724509/how-to-crop-an-image-by-using-python-pptx-after-it-inserted-into-slide
                    # Find total length of data and divid by pagefootage
                    try:
                        pages = int(np.ceil(int(well_data.shape[0]) / int(pagefootage)))
                    except Exception as e:
                        st.info(f"error pagefooter: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                        pages = 1 
                    
                    # st.write(f'pages: {pages}')
                    
                    save_bar.progress(15, text='Creating header..') 
                    
                    filename2 = "figurehead.png"
                    pio.write_image(fig1, filename2, scale=2, width=1746, height=1000) 
                    
                    save_bar.progress(16, text='Creating graphics..') 
                    filename = "figure1.png"
                    # pio.write_image(fig1, filename, scale=2, width=2208, height=3030)                    
                    # pio.write_image(fig1, filename, scale=1, width=2182, height=pages*2675)              
                    pio.write_image(fig1, filename, scale=2, width=672, height=pages*809)
                                        
                    # placeholder1 = slide1.shapes.add_picture(filename, left=Inches(0.25), top=Inches(1), width=Inches(7), height=Inches(8.58)) 
                    # placeholder1.crop_top = (1-1)/9
                    # placeholder1.crop_bottom = 8/9
                    
                    save_bar.progress(17, text='Writing page..') 
                    
                    def add_slide(prs, layout, title, number):
                        """Return slide newly added to `prs` using `layout` and having `title`."""
                        slide = prs.slides.add_slide(layout) # adding a slide                                 
                        txBox = slide.shapes.add_textbox(left=Inches(0.87), top=Inches(0.4), width=Inches(1), height=Inches(0.25))           
                        tf = txBox.text_frame
                        tf.text = f'{title}'
                        tf.paragraphs[0].font.size = Pt(8)
                        # placeholder = slide.placeholders[1]
                        # placeholder = placeholder.insert_picture(filename)
                        
                        rangeinst = (90 / pages)
                        rangetop = 5 + (number * rangeinst)
                        rangebot = rangetop + rangeinst
                        
                        # # fig1.update_yaxes(range=[rangetop, rangebot])
                        # fig1.update_layout(yaxis_range=[rangetop,rangebot])
                        # filename2 = "figure1.png"
                        # # pio.write_image(fig1, filename, scale=2, width=2208, height=3030)                    
                        # pio.write_image(fig1, filename2, scale=1, width=2182, height=2675)
                        # placeholder2 = slide.shapes.add_picture(filename2, left=Inches(0.25), top=Inches(1), width=Inches(7), height=Inches(8.58))                         
                        # txBox = slide.shapes.add_textbox(left=Inches(0.87), top=Inches(1), width=Inches(1), height=Inches(0.25)) 
                        # tf = txBox.text_frame
                        # tf.text = f'numbers {number} step {rangeinst} top {rangetop} bot {rangebot}'
                        
                        # https://pythonprogramming.altervista.org/inserting-an-image-in-powerpoint-with-python/
                        
                        # filename = "figure1.png"
                        # pio.write_image(fig1, filename, scale=2, width=2208, height=3030)                    
                        # pio.write_image(fig1, filename, scale=2, width=2182, height=int(np.ceil(int(well_data.shape[0]))*(pagefootage/2675)))
                        placeholder = slide.shapes.add_picture(filename, left=Inches(0.25), top=Inches(1.2), width=Inches(7), height=Inches(8.43)) 
                        # placeholder.crop_top = (number-1) * (placeholder.height.inches / pages)
                        # placeholder.crop_bottom = (pages-number) * (placeholder.height.inches / pages)
                        
                        # txBox = slide.shapes.add_textbox(left=Inches(0.87), top=Inches(1), width=Inches(1), height=Inches(0.25)) 
                        # tf = txBox.text_frame
                        # tf.text = f'numbers {(number-1)} and {(number-1) / pages} and {float(round((number-1) / pages,1))}'
                        placeholder.crop_top = (number - 1) / pages
                        # txBox = slide.shapes.add_textbox(left=Inches(0.87), top=Inches(2), width=Inches(1), height=Inches(0.25)) 
                        # tf = txBox.text_frame
                        # tf.text = f'numbers {(pages-number)} and {(pages-number) / pages} and {float(round((pages-number) / pages,1))}'
                        placeholder.crop_bottom = (pages - number) / pages
                        
                        # Set chart header
                        placeholder = slide.shapes.add_picture(filename2, left=Inches(0.25), top=Inches(0.85), width=Inches(7), height=Inches(0.32)) 
                        placeholder.crop_top = 0.08
                        placeholder.crop_bottom = 0.86
                        placeholder.crop_left = 0.06
                        
                        return slide
                                        
                    for i in range(1,pages):
                        prog = int(20+i*(80/pages))
                        save_bar.progress(prog, text=f'Page {1+i}: Creating..') 
                        add_slide(prs, lyt, f'Depth{(well_data["DEPTH"]/i)*pages}', i)
                    
                    # Add Histograms to last page
                    try:    
                        save_bar.progress(85, text='Creating analysis page..') 
                        slidelast=prs.slides.add_slide(lyt) # adding a slide          
                        txBox = slidelast.shapes.add_textbox(Inches(0.87), Inches(0.4), Inches(1), Inches(0.25))            
                        tf = txBox.text_frame
                        tf.text = f'Analysis'
                        tf.paragraphs[0].font.size = Pt(8)
                                      
                        save_bar.progress(87, text='Analysis: Histogram..') 
                        filename = "histogram.png"
                        pio.write_image(histogram, filename, scale=2, width=873, height=450)
                        slidelast.shapes.add_picture(filename, left=Inches(0.25), top=Inches(1), width=Inches(7)) 
                        
                        save_bar.progress(90, text='Analysis: Crossplot..') 
                        filename = "xplot.png"
                        pio.write_image(xplot, filename, scale=2, width=873, height=450)
                        slidelast.shapes.add_picture(filename, left=Inches(0.25), top=Inches(4.75), width=Inches(7))                         
                          
                    except Exception as e:
                        print(e)
                        st.info(f"error ppt-slidelast: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                        
                except Exception as e:
                    print(e)
                    st.info(f"error ppt pic: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                # # Pg2
                # try:
                #     slide2=prs.slides.add_slide(lyt) # adding a slide 
                #     txBox = slide2.shapes.add_textbox(Inches(0.87), Inches(0.4), Inches(1), Inches(0.25))
                #     tf = txBox.text_frame
                #     tf.text = f'OVERVIEW'
                #     tf.paragraphs[0].font.size = Pt(8)  
                #     #  Depth vs Time
                #     filename = "fig1.jpg"
                #     pio.write_image(fig1, filename, scale=2, width=(72*4.8))             
                #     slide2.shapes.add_picture(filename, left=Inches(0.1), top=Inches(1.36), width=Inches(4.8))        
                        
                # except Exception as e:
                #     print(e)
                #     st.info(f"error ppt2: {e}")
                
                # Close out
                save_bar.progress(92, text='Saving..') 
                binary_output = BytesIO()
                prs.save(binary_output)
                return binary_output
            
            
            def create_pdf_poster():
                WIDTH = 1684
                HEIGHT = 2384
                ltr_WIDTH = 792
                ltr_HEIGHT = 612
                pagetitle_x = 75
                pagetitle_y = 41
                gaf_y = 50
                
                # pdf = adiraPDF(orientation = 'L', unit = 'pt', format='Letter') # Letter (215.9 × 279.4 mm) 1pt is 1/72in
                pdf = adiraPDF('P', 'pt', (1684, 2384)) # A1 (594 × 841mm or 23.39 × 33.11 inch) 1pt is 1/72in
                pdf.set_auto_page_break(False, 0)
                
                # Pg1
                pdf.add_page()
                pdf.set_font('Arial', '', 8)
                pdf.text(pagetitle_x, pagetitle_y, f'ROCK STRENGTH REPORT')
                                
                filename = "fig1.jpg"
                # fig1.write_image(filename, engine='kaleido')
                # pdf.image("fig1.jpg", 5, 230, WIDTH-20, HEIGHT-400)
                pio.write_image(fig1, filename, scale=1, width=WIDTH, height=HEIGHT)
                pdf.image("fig1.jpg", 5, gaf_y, w=WIDTH-60, h= HEIGHT-100)
                
                # Add header info, after image for layering to prevent the image covering it
                pdf.text(500, 120, f'<b>{las_file.well.COMP.descr.capitalize()}:</b> {las_file.well.COMP.value}')
                pdf.text(500, 130, f'<b>{las_file.well.WELL.descr.capitalize()}:</b> {las_file.well.WELL.value}')
                pdf.text(500, 140, f'<b>{las_file.well.UWI.descr.capitalize()}:</b> {las_file.well.UWI.value}')
                pdf.text(500, 150, f'<b>{las_file.well.LOC.descr.capitalize()}:</b> {las_file.well.LOC.value}')
                
                pdf.text(650, 120, f'<b>{las_file.well.FLD.descr.capitalize()}:</b> {las_file.well.FLD.value}')
                pdf.text(650, 130, f'<b>{las_file.well.CNTY.descr.capitalize()}:</b> {las_file.well.CNTY.value}')
                pdf.text(650, 140, f'<b>{las_file.well.STAT.descr.capitalize()}:</b> {las_file.well.STAT.value}')
                pdf.text(650, 150, f'<b>{las_file.well.CTRY.descr.capitalize()}:</b> {las_file.well.CTRY.value}')
                
                pdf.text(750, 120, f'<b>{las_file.well.STRT.descr.capitalize()}:</b> {las_file.well.STRT.value}')
                pdf.text(750, 130, f'<b>{las_file.well.STOP.descr.capitalize()}:</b> {las_file.well.STOP.value}')
                pdf.text(750, 140, f'<b>{las_file.well.STEP.descr.capitalize()}:</b> {las_file.well.STEP.value}')
                
                
                #  Close out
                pdf.output("ouptut2.pdf", 'f')
                with open("ouptut2.pdf", "rb") as pdf_file:
                    PDFbyte = pdf_file.read()
                return PDFbyte
            
            # formpdf = st.form("template_form")
            # # name = formpdf.text_input("Your name")
            # reporttype = formpdf.selectbox(
            #     'Choose Type',
            #     ['PowerPoint', 'PDF'],
            #     index=0,
            # )
            # with st.spinner('Request...'): 
            #     submit = formpdf.form_submit_button('Request Report')
            
            # grade = formpdf.slider("Grade", 1, 100, 60)
            
            with st.form(key='template_form'):
                fc1, fc2 = st.columns(2)
                with fc1:
                    reporttype = st.selectbox('Choose Format',['PowerPoint', 'PDF'],index=0)
                with fc2:
                    reportsize = st.selectbox('Choose Type',['Poster','Multipage'],index=0)
                    pagefootage = st.selectbox('Choose Ratio',[2000,1000,5000],index=0)
                with st.spinner('Request...'): 
                    submit = st.form_submit_button('Request Report')
            
            if submit:                    
                save_bar = st.progress(0, text='Initializing..')
                with st.spinner('Constructing...'):    
                    if reporttype == 'PowerPoint':
                        if reportsize == 'Poster':
                            try:
                                reportppt = create_pptx_poster()
                                try:                
                                    save_bar.progress(95, text='Preparing Download')
                                    st.download_button(
                                        label="Download PowerPoint",
                                        data=reportppt,
                                        file_name=f"Adira_Rock_Report_{datetime.datetime.now()}.pptx",
                                        mime="application/octet-stream",
                                    ) 
                                    save_bar.progress(100, text='Ready')
                                except Exception as e:
                                    st.info(f"error pptx dl1: {e}")                                             
                            except Exception as e:
                                st.info(f"error pptx: {e}")
                                
                        if reportsize == 'Multipage':
                            try:
                                reportppt = create_pptx_multipage(pagefootage)
                                try:                
                                    save_bar.progress(95, text='Preparing Download')
                                    st.download_button(
                                        label="Download PowerPoint",
                                        data=reportppt,
                                        file_name=f"Adira_Rock_Report_{datetime.datetime.now()}.pptx",
                                        mime="application/octet-stream",
                                    ) 
                                    save_bar.progress(100, text=f'Ready ({int(np.ceil(int(well_data.shape[0]) / int(pagefootage)))} pages)')
                                except Exception as e:
                                    st.info(f"error pptx dl1: {e}")                                             
                            except Exception as e:
                                st.info(f"error pptx: {e}")
                            
                    elif reporttype == 'PDF':
                        try:                            
                            reportpdf = create_pdf_poster()                
                            save_bar.progress(95, text='Preparing Download')
                            st.download_button(
                                label='Download Report',
                                data=reportpdf,
                                file_name=f"Adira_Rock_Report_{datetime.datetime.now()}.pdf",
                                mime="application/octet-stream",
                            )
                            save_bar.progress(100, text='Ready')
                        except Exception as e:
                            print(e)
                            st.info(f"error pdf: {e}")
            
        with st.expander('Log Plot'):    
            curves = st.multiselect('Select Curves To Plot', columns)
            if len(curves) <= 1:
                st.warning('Please select at least 2 curves.')
            else:
                curve_index = 1
                fig3 = make_subplots(rows=1, cols= len(curves), subplot_titles=curves, shared_yaxes=True)

                for curve in curves:
                    fig3.add_trace(go.Scatter(x=well_data[curve], y=well_data['DEPTH']), row=1, col=curve_index)
                    curve_index+=1
                
                fig3.update_layout(height=1000, showlegend=False, yaxis={'title':'DEPTH','autorange':'reversed'})
                fig3.layout.template='seaborn'
                st.plotly_chart(fig3, use_container_width=True)

        with st.expander('Histograms'):
            col1_h, col2_h = st.columns(2)
            col1_h.header('Options')

            hist_curve = col1_h.selectbox('Select a Curve', columns)
            log_option = col1_h.radio('Select Linear or Logarithmic Scale', ('Linear', 'Logarithmic'))
            hist_col = col1_h.color_picker('Select Histogram Colour')
            st.write('Color is'+hist_col)
            
            if log_option == 'Linear':
                log_bool = False
            elif log_option == 'Logarithmic':
                log_bool = True
        

            histogram = px.histogram(well_data, x=hist_curve, log_x=log_bool)
            histogram.update_traces(marker_color=hist_col)
            histogram.layout.template='seaborn'
            col2_h.plotly_chart(histogram, use_container_width=True)

        
        with st.expander('Crossplot'):
            
            col1, col2 = st.columns(2)
            with col1:
                with st.form(key='template_form_xplot'):
                    st.write('Options')

                    xplot_x = st.selectbox('X-Axis', columns)
                    xplot_y = st.selectbox('Y-Axis', columns)
                    xplot_col = st.selectbox('Colour By', columns)
                    
                    subcol1, subcol2 = st.columns(2)
                    with subcol1:
                        xplot_x_log = st.radio('X Axis - Linear or Logarithmic', ('Linear', 'Logarithmic'))
                        xplot_y_log = st.radio('Y Axis - Linear or Logarithmic', ('Linear', 'Logarithmic'))
                        xplot_y_rev = st.checkbox('Y Axis - Reverse', False)

                    if xplot_x_log == 'Linear':
                        xplot_x_bool = False
                    elif xplot_x_log == 'Logarithmic':
                        xplot_x_bool = True
                    
                    if xplot_y_log == 'Linear':
                        xplot_y_bool = False
                    elif xplot_y_log == 'Logarithmic':
                        xplot_y_bool = True

                    with st.spinner('Request...'): 
                        submit_xplot = st.form_submit_button('Update Crossplot')
                
            if submit_xplot:   
                with col2:     
                    st.write('Crossplot')           
                    xplot = px.scatter(well_data, x=xplot_x, y=xplot_y, color=xplot_col, log_x=xplot_x_bool, log_y=xplot_y_bool)
                    if xplot_y_rev is True:                
                        xplot.update_layout(yaxis={'autorange':'reversed'})
                    xplot.layout.template='plotly_white'
                    st.plotly_chart(xplot, use_container_width=True)

            
    
def raw_data(las_file, well_data):
    # st.title('LAS File Data Info')
    
    with st.expander('Data Info'): 
        if not las_file:            
            st.info(' ')
        else:
            st.write('**Curve Information**')
            for count, curve in enumerate(las_file.curves):
                # st.write(f"<b>Curve:</b> {curve.mnemonic}, <b>Units: </b>{curve.unit}, <b>Description:</b> {curve.descr}", unsafe_allow_html=True)
                st.write(f"   {curve.mnemonic} ({curve.unit}): {curve.descr}", unsafe_allow_html=True)
            st.write(f"<b>There are a total of: {count+1} curves present within this file</b>", unsafe_allow_html=True)
            
            st.write('<b>Curve Statistics</b>', unsafe_allow_html=True)
            st.write(well_data.describe()) 
            
            st.write('<b>Raw Data Values</b>', unsafe_allow_html=True)
            st.dataframe(data=well_data)
            
            csv = convert_df(well_data)
            st.download_button(
                "Download Data",
                csv,
                "file.csv",
                "text/csv",
                key='download-csv'
                )
            
#  RDA Formulas
#  ESTIMATES SHALE %
        # Present Columns:
        # [SpontaneousPotential]
        # Absent Columns:
        # [GammaRay]

        # Calculate
        # Condition:true
        # Affected Column:  [C%SH]
        # Affected Column Expression:  ([SpontaneousPotential] - [SpontaneousPotentialMin])/([SpontaneousPotentialMax] - [SpontaneousPotentialMin])

        # Present Columns:
        # [GammaRay]

        # Calculate
        # Condition:[GammaRay] >= [GammaRayMax]
        # Affected Column:  [C%SH]
        # Affected Column Expression:  1


        # Calculate
        # Condition:[GammaRay] <= [GammaRayMin]
        # Affected Column:  [C%SH]
        # Affected Column Expression:  0


        # Calculate
        # Condition:[GammaRay] < [GammaRayMax] && [GammaRay] > [GammaRayMin]
        # Affected Column:  [C%SH]
        # Affected Column Expression:  ([GammaRay] - [GammaRayMin])/([GammaRayMax] - [GammaRayMin])
        
# NeutronPorosity
        # Present Columns:
        # [NeutronPorosity]

        # Calculate
        # Condition:[NeutronPorosity] > 1
        # Affected Column:  [NeutronPorosity]
        # Affected Column Expression:  [NeutronPorosity] /100

# Estimated Percentages
        # Present Columns:
        # [SonicP]

        # Calculate
        # Condition:[SonicP] < 55
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  1 - [C%SH]
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]
        # Absent Columns:
        # [SonicP]

        # Calculate
        # Condition:[BulkDensity] > 2.65
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  1 - [C%SH]
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]
        # [SonicP]
        # Absent Columns:
        # [PhotoElectric]
        # [NeutronPorosity]

        # Calculate
        # Condition:[RHOLS2] <=0
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  1 - [C%SH]
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]
        # [SonicP]
        # Absent Columns:
        # [PhotoElectric]
        # [NeutronPorosity]

        # Calculate
        # Condition:([RHOLS2] >= -.04) && ([RHOLS2] <= 0.04)
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  1 - [C%SH]
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]
        # [PhotoElectric]

        # Calculate
        # Condition:([PEDO] <= 0) && ([PELS] >0) && ([PEDO] / ([PEDO] - [PELS]) > .5) && ([BulkDensity] >= 1.9)
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  (1 - [C%SH])
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]
        # [PhotoElectric]

        # Calculate
        # Condition:([PELS] <= 0) && ([BulkDensity] >= 1.9)
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  1 - [C%SH]
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]

        # Calculate
        # Condition:[BulkDensity] > 2.71 
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  1 - [C%SH]
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]
        # [SonicP]
        # Absent Columns:
        # [PhotoElectric]
        # [NeutronPorosity]

        # Calculate
        # Condition:([RHODO2] >= -.02) && ([RHODO2] <= 0.02)
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  1 - [C%SH]
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]
        # [NeutronPorosity]
        # Absent Columns:
        # [PhotoElectric]

        # Calculate
        # Condition:([RHODO] >= -.01) && ([RHODO] <= .01)
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  1 - [C%SH]
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]
        # [PhotoElectric]

        # Calculate
        # Condition:([PEDO] >= -.1) && ([PEDO] <= .1) && ([BulkDensity] >= 1.9)
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  (1 - [C%SH])
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0.0002

        # Present Columns:
        # [SonicP]

        # Calculate
        # Condition:[SonicP] < 47
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  1 - [C%SH]
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]
        # [PhotoElectric]

        # Calculate
        # Condition:([PhotoElectric] > 4.5) && ([BulkDensity] >= 2.85)
        # Affected Column:  [C%SH]
        # Affected Column Expression:  0
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  1
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]

        # Calculate
        # Condition:[BulkDensity] > 2.87 
        # Affected Column:  [C%SH]
        # Affected Column Expression:  0
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  1
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]

        # Calculate
        # Condition:[BulkDensity] < 1.9 
        # Affected Column:  [C%SH]
        # Affected Column Expression:  0
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  1
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]
        # [SonicP]

        # Calculate
        # Condition:([BulkDensity] <= 2.05) && ([BulkDensity] > 1.7) && ([SonicP] < 75)
        # Affected Column:  [C%SH]
        # Affected Column Expression:  0
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  1
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]
        # [PhotoElectric]

        # Calculate
        # Condition:([PhotoElectric] > 4.4) && ([BulkDensity] <= 2.05)
        # Affected Column:  [C%SH]
        # Affected Column Expression:  0
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  1
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Present Columns:
        # [BulkDensity]

        # Calculate
        # Condition:[BulkDensity] < 1.5 
        # Affected Column:  [C%SH]
        # Affected Column Expression:  0
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  1

        # Present Columns:
        # [PhotoElectric]

        # Calculate
        # Condition:[PhotoElectric] < 0.5 
        # Affected Column:  [C%SH]
        # Affected Column Expression:  0
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  1


        # Calculate
        # Condition:true
        # Affected Column:  [C%SS]
        # Affected Column Expression:  1 - [C%SH]
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

        # Absent Columns:
        # [SpontaneousPotential]
        # [GammaRay]

        # Calculate
        # Condition:true
        # Affected Column:  [C%SH]
        # Affected Column Expression:  0
        # Affected Column:  [C%SS]
        # Affected Column Expression:  0
        # Affected Column:  [C%LS]
        # Affected Column Expression:  0
        # Affected Column:  [C%DO]
        # Affected Column Expression:  0
        # Affected Column:  [C%AN]
        # Affected Column Expression:  0
        # Affected Column:  [C%SL]
        # Affected Column Expression:  0
        # Affected Column:  [C%CO]
        # Affected Column Expression:  0

#  DTS
        # Present Columns:
        # [SonicP]

        # Calculate
        # Condition:[SonicS] = [NULL]
        # Affected Column:  [DTS]
        # Affected Column Expression:  ([C%SH] * [SH_Ratio] + [C%SS] * [SS_Ratio] + [C%LS] * [LS_Ratio] + [C%DO] * [DO_Ratio] + [C%AN] * [AN_Ratio] + [C%SL] * [SL_Ratio] + [C%CO] * [CO_Ratio] )  * [SonicP]

#  CUCS
        # Present Columns:
        # [SonicS]

        # Calculate
        # Condition:true
        # Affected Column:  [CUCS]
        # Affected Column Expression:   2670000 * Pow([SonicS] -58, -1.36)

#  CCCS
        # Present Columns:
        # [MudWeight]
        # [CIFA]

        # Calculate
        # Condition:true
        # Affected Column:  [CCCS]
        # Affected Column Expression:    [CUCS] + [MudWeight]*[DEPT]*0.052 * (1+Sin ( [CIFA] * [PI]/180] )) / (1 - Sin ([CIFA] * [PI]/180] )) 

#  Constants
        # PI 3.14159265358979
        #  CO_Ratio 1.76
        #  SL_Ratio 2.15
        #  AN_ratio 2.45
        #  DO_Ratio 1.8
        # LS_Ratio 2.1
        # SS_Ratio 1.6
        # SH_Ratio 1.725
        # GammaRayMax 140
        # KellyBushingHeight 0
    
# Categroies
#  Depth
    #  DEPT
# Mudweight
    #  MW
# PhotoElectric
    # PE
    # PEF
    # PEFA
    # PhotoElectric
# NuetronPorosity
    #  CN
    # FCNL
    # NPHI
    # NPOR
    # NuetronPorosity'
# BulkDensity
    # DEN
    # DENS
    # RHO
    # RHOB
    # ZDEN
#  SonicS
    # [DTS]
    # [DTSM]
    # [S_DT]
    # [SDTS]
    # [SonicS]
# SonicP
    # [AC]
    # [ACC]
    # [DT]
    # [DTC]
    # [DTCO]
    # [P_DT]
    # [PSON]
    # [SON]
    # [SonicP]
# SpontaneousPotential
    # [SP]
#  GammaRay
    # [CGR]
    # [GR]
    # [GRD]
    # [GammaRay]
    # [GAMMARAY]


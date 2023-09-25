from pathlib import Path
import streamlit as st
import pandas as pd
import numpy as np
from enum import Enum
from pandas.api.types import (
    is_categorical_dtype,
    is_datetime64_any_dtype,
    is_numeric_dtype,
    is_object_dtype,
)

# Plotly imports
from plotly.subplots import make_subplots
import plotly.graph_objects as go
import plotly.figure_factory as ff
import plotly.express as px
import plotly.io as pio
import pip
from PIL import Image
import datetime
# from datetime import datetime, timedelta
import io
import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import traceback
from streamlit_tags import st_tags, st_tags_sidebar
from geopy.distance import geodesic
from geopy.geocoders import Nominatim

import requests
import zipfile
import json
from pdf2image import convert_from_path
import yagmail
import openpyxl
from fuzzywuzzy import process
from fuzzywuzzy import fuzz
from enverus_developer_api import DeveloperAPIv3
import csv
from tempfile import mkdtemp
from streamlit_datalist import stDatalist

#  Azure blob setup
from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions

# from st_oauth import st_oauth
# # https://medium.com/streamlit/oauth-component-for-streamlit-e05f00874fbc
# st.markdown("## This (and above) is always seen")
# id = st_oauth('myoauth', 'Click to login via OAuth')
# # id = st_oauth(‘myoauth’)
# st.markdown("## This (and below) is only seen after authentication")


# initialize Nominatim API
geolocator = Nominatim(user_agent="adirageolocation")
# initialize yagmail
yag = yagmail.SMTP(st.secrets['yagmailusername'], st.secrets['yagmailpassword'])
# initialize azure blob storage
connection_string = 'DefaultEndpointsProtocol=https;AccountName=' + st.secrets['blobname'] + ';AccountKey=' + st.secrets['blobkey'] + ';EndpointSuffix=core.windows.net'
blob_service = BlobServiceClient.from_connection_string(connection_string)

def blob_upload(df, file_name, bupdate):
    blob_client = blob_service.get_blob_client(container=st.secrets['container2'], blob=file_name)
    try:
        output = df.to_csv(index=False, encoding="utf-8")        
        blob_client.upload_blob(output, overwrite=bupdate)
        status = 'Saved'
    except Exception as e:
        print(e)
        st.error(f"error blob df-to-csv: {e}")    
    # try:
    #     blob_exists = blob_client.exists()
    #     if not blob_exists:
    #         blob_client.upload_blob(output, overwrite=False)
    #         status = 'Loaded'
    #     else:            
    #         status = 'Exists'
    #     # , blob_type="BlockBlob")
    # except Exception as e:
    #     print(e)
    #     st.error(f"error blob upload: {e}")
    return status 

@st.cache_resource
def blob_download(filename):
    # blob_client = blob_service.get_blob_client(container=st.secrets['container'], blob=df_name)    
    # #get a list of all blob files in the container
    # blob_list = []
    # for blob_i in blob_client.list_blobs():
    #     blob_list.append(blob_i.name)
    
    # for blob_i in blob_list:
    #     #generate a shared access signature for each blob file
    #     sas_i = generate_blob_sas(account_name = account_name,
    #                                 container_name = container_name,
    #                                 blob_name = blob_i,
    #                                 account_key=account_key,
    #                                 permission=BlobSasPermissions(read=True),
    #                                 expiry=datetime.utcnow() + timedelta(hours=1))
    
    #generate a shared access signature for each blob file
    sas_i = generate_blob_sas(account_name = st.secrets['blobname'],
                                container_name = st.secrets['container2'],
                                blob_name = filename,
                                account_key=st.secrets['blobkey'],
                                permission=BlobSasPermissions(read=True),
                                expiry=datetime.datetime.utcnow() + datetime.timedelta(hours=1))
    sas_url = 'https://' + st.secrets['blobname'] + '.blob.core.windows.net/' + st.secrets['container2'] + '/' + filename + '?' + sas_i 
    
    df = pd.read_csv(sas_url)        
    # try:
    #     output = df.to_csv(index=False, encoding="utf-8")        
    #     blob_client.upload_blob(output, overwrite=False)
    #     status = 'Saved'
    # except Exception as e:
    #     print(e)
    #     st.error(f"error blob df-to-csv: {e}")  
    return df 
    
@st.cache_resource
def enverus_df():    
    # https://github.com/enverus-ea/enverus-developer-api
    try:
        v3 = DeveloperAPIv3(secret_key=st.secrets['enveruskey'],retries=5,backoff_factor=1)
    except Exception as e:
        print(e)
        st.error(f"error enverus1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
        # generate new key https://app.enverus.com/provisioning/directaccess
        
    df = pd.DataFrame()
    # df2 = pd.DataFrame()
    # df3 = pd.DataFrame()
    
    # # Get well records updated after 2018-08-01 and without deleted dates
    # for row in v3.query('wells', updateddate='gt(2018-08-01)', deleteddate='null'):
    #     print(row)
    # # Get permit records with approved dates between 2018-03-01 and 2018-06-01
    # for row in v3.query('rigs', spuddate='btw(2018-03-01,2018-06-01)'):
    #     print(row)
    
    # test that worked
    options = dict(ENVOperator="AETHON ENERGY",updateddate='gt(2023-01-01)',pagesize=10000,deleteddate='null')
    # v3.query(dataset, **options)
    
    with st.spinner('formation'):
        count = v3.count("formation-tops", updateddate='gt(2023-01-01)', deleteddate="null")
        with st.expander(f'enverus tops {count}'):
            try:
                docs = v3.docs("formation-tops") # pagesize=10000, updateddate='gt(2023-01-01)', deleteddate='null'
                df = pd.DataFrame(docs)
                st.write(df)
            except Exception as e:
                print(e)
                st.error(f"error enverussample: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
    
    with st.spinner('permits'):
        count = v3.count("permits", updateddate='gt(2023-01-01)', deleteddate="null")
        with st.expander(f'enverus permits {count}'):
            try:
                # docs = v3.docs("permits") # pagesize=10000,updateddate='gt(2023-01-01)',  deleteddate='null'
                # df = pd.DataFrame(docs) 
                # st.write(df)               
                records = v3.query("permits",ENVOperator="AETHON ENERGY",updateddate='gt(2023-01-01)',pagesize=10000,deleteddate='null') # pagesize=10000,updateddate='gt(2023-01-01)',  deleteddate='null'
                df = pd.DataFrame(records)
                st.write(df)
            except Exception as e:
                print(e)
                st.error(f"error enverussample: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
            
    # tempdir = mkdtemp()
    # path = os.path.join(tempdir, "rigs.csv")

    # dataset = "rigs"
    # options = dict(pagesize=10000, deleteddate="null")

    # # v3.query(dataset, **options)
    # for row in v3.query("rigs", pagesize=1000, deleteddate="null"):
    #     # print(row)
    #     st.write(row)
    # df = v3.to_csv(query, path, log_progress=True, delimiter=",", quoting=csv.QUOTE_MINIMAL)

    # # Example 1, Create a pandas dataframe from a dataset query
    # with st.spinner('rigs'):
    #     count = v3.count("rigs", updateddate='gt(2023-01-01)', deleteddate="null")
    #     with st.expander(f'enverus rigs {count}'):
    #         try:
    #             # df = v3.to_dataframe('rigs', pagesize=10000, deleteddate='null') # spuddate='btw(2018-03-01,2018-06-01)'
    #             docs = v3.docs("rigs",  ) # pagesize=10000, updateddate='gt(2023-01-01)', deleteddate='null'
    #             df = pd.DataFrame(docs)
    #             st.write(df)
    #         except Exception as e:
    #             print(e)
    #             st.error(f"error enverus1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
    # # Example 2, Create a Texas rigs dataframe, replacing the state abbreviation with the complete name and removing commas from Operator names
    # try:
    #     df2 = v3.to_dataframe(
    #         dataset="rigs",
    #         deleteddate="null",
    #         pagesize=100000,
    #         stateprovince="TX",
    #         converters={
    #             "StateProvince": lambda x: "TEXAS",
    #             "ENVOperator": lambda x: x.replace(",", "")
    #         }
    #         )
    #     st.write(df2)
    # except Exception as e:
    #     print(e)
    #     st.error(f"error enverus3: {e}")          
    # try:
    #     df = v3.to_dataframe('rigs',pagesize=10000, deleteddate='null')
    #     st.write(df)
    # except Exception as e:
    #     print(e)
    #     st.error(f"error enverus2: {e}")
    
    # # Example 3, to csv
    # tempdir = mkdtemp()
    # path = os.path.join(tempdir, "rigs.csv")
    
    # dataset = "rigs"
    # options = dict(pagesize=10000, deleteddate="null")

    # v3.query(dataset, **options)
    # v3.to_csv(query, "rigs.csv", log_progress=True, delimiter=",", quoting=csv.QUOTE_MINIMAL)
    
    
    return df


@st.cache_resource
def enverus_rig_list(**options):    
    # https://github.com/enverus-ea/enverus-developer-api
    try:
        v3 = DeveloperAPIv3(secret_key=st.secrets['enveruskey'],retries=5,backoff_factor=1)
    except Exception as e:
        print(e)
        st.error(f"error enverus1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
        # generate new key https://app.enverus.com/provisioning/directaccess
    df = pd.DataFrame()        
    try:            
        # count = v3.count("rigs", updateddate='gt(2023-01-01)', deleteddate="null")
        # records = v3.query("rigs",updateddate='gt(2023-01-01)',pagesize=10000,deleteddate='null') # pagesize=10000,updateddate='gt(2023-01-01)',  deleteddate='null'        
        # options = dict(updateddate='gt(2023-01-01)',pagesize=10000,deleteddate='null')
        records = v3.query("rigs", **options)  
        df = pd.DataFrame(records)    
    except Exception as e:
        print(e)
        st.error(f"error enverusrigs: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")    
    return df

@st.cache_resource
def enverus_permit_list(**options):    
    # https://github.com/enverus-ea/enverus-developer-api
    try:
        v3 = DeveloperAPIv3(secret_key=st.secrets['enveruskey'],retries=5,backoff_factor=1)
        # generate new key https://app.enverus.com/provisioning/directaccess
        df = pd.DataFrame() 
        records = v3.query("permits", **options)  
        df = pd.DataFrame(records)    
    except Exception as e:
        print(e)
        st.error(f"error enveruspermits: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")    
    return df

@st.cache_resource
def enverus_ftop_list(**options):    
    # https://github.com/enverus-ea/enverus-developer-api
    try:
        v3 = DeveloperAPIv3(secret_key=st.secrets['enveruskey'],retries=5,backoff_factor=1)
        # generate new key https://app.enverus.com/provisioning/directaccess
        
        docs = v3.docs("formation-tops") # pagesize=10000, updateddate='gt(2023-01-01)', deleteddate='null'
        df = pd.DataFrame(docs)
        # df = pd.DataFrame() 
        # records = v3.query("formation-tops", **options)  
        # df = pd.DataFrame(records)    
    except Exception as e:
        print(e)
        st.error(f"error enverusftops: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")    
    return df

def runningToggle(bool):
  if bool == True:
    hide_streamlit_style = """
              <style>

              div[class='css-4z1n4l ehezqtx5']{
                background: rgba(0, 0, 0, 0.3);
                color: #fff;
                border-radius: 15px;
                height: 40px;
                max-width: 160px;


                position: fixed;
                top: 50%;
                left: 50%;
                transform: translate(-50%, -50%);
                width: 50%;
              }

              </style>
              """
    st.markdown(hide_streamlit_style, unsafe_allow_html=True) 
    
# Standard Color Palette
class upal(str, Enum):
    teal1 = '05929f'
    teal2 = '3F7077'
    gray1 = 'CBC9C9'
    gray2 = '9B9DA0'
    gray3 = '7A7D81'
    green1 = '98C21F'

# https://tools.nov.com/PowerSections/
pdmotors = ['',
          '1.688" 5:6 8.31 Rev/Gal 2.3',
          '1.688" 5:6 15.00 Rev/Gal 5.0', 
          '2.125" 5:6 12.20 Rev/Gal 6.0',
          '2.875" 5:6 3.38 Rev/Gal 3.5',    
          '2.875" 5:6 3.56 Rev/Gal 4.7',
          '2.875" 7:8 2.92 Rev/Gal 3.6',     
          '2.875" 5:6 5.52 Rev/Gal 7.0',
          '2.875" 4:5 2.85 Rev/Gal 4.0',     
          '2.875" 7:8 3.39 Rev/Gal 3.7',
          '3.125" 7:8 1.74 Rev/Gal 3.0',     
          '3.125" 5:6 2.15 Rev/Gal 3.5',
          '3.125" 5:6 2.64 Rev/Gal 5.0',     
          '3.125" 7:8 1.23 Rev/Gal 2.5',
          '3.375" 7:8 1.52 Rev/Gal 3.0',     
          '3.375" 5:6 2.09 Rev/Gal 3.5',
          '3.5" 5:6 2.66 Rev/Gal 3.0',     
          '3.75" 7:8 1.62 Rev/Gal 6.7',
          '4.125" 7:8 0.99 Rev/Gal 5.0',    
          '4.75" 4:5 0.96 Rev/Gal 6.3',
          '4.75" 5:6 0.99 Rev/Gal 8.3',    
          '4.75" 7:8 0.49 Rev/Gal 2.2',
          '4.75" 7:8 0.50 Rev/Gal 3.8',    
          '4.75" 9:10 0.54 Rev/Gal 4.0',
          '4 5" 6:7 0.79 Rev/Gal 6.4', 
          '5" 6:7 0.78 Rev/Gal 8.0',
          '5" 7:8 0.24 Rev/Gal 2.6',   
          '5" 6:7 0.83 Rev/Gal 6.0',
          '5" 6:7 0.81 Rev/Gal 8.0', 
          '5" 7:8 0.27 Rev/Gal 2.6',
          '5" 7:8 0.36 Rev/Gal 3.7', 
          '5" 7:8 0.51 Rev/Gal 3.8',
          '5" 7:8 0.25 Rev/Gal 2.6',   
          '5.125" 6:7 0.58 Rev/Gal 7.8',
          '5.125" 5:6 0.80 Rev/Gal 9.5',    
          '5.125" 9:10 0.46 Rev/Gal 6.1', 
          '5.500" 6:7 0.58 Rev/Gal 7.8',     
          '6.25" 7:8 0.34 Rev/Gal 4.8',
          '6.5" 7:8 0.27 Rev/Gal 6.7',     
          '6.565" 7:8 0.28 Rev/Gal 5.0',
          '6.75" 4:5 0.47 Rev/Gal 7.0', 
          '6.75" 7:8 0.15 Rev/Gal 2.9',
          '6.75" 7:8 0.27 Rev/Gal 5.0',    
          '6.75" 4:5 0.50 Rev/Gal 7.0',
          '6.75" 6:7 0.30 Rev/Gal 5.0',     
          '6.75" 6:7 0.29 Rev/Gal 6.0',
          '6.75" 7:8 0.14 Rev/Gal 3.3',     
          '6.75" 7:8 0.28 Rev/Gal 5.0',
          '6.75" 7:8 0.23 Rev/Gal 5.7',     
          '6.75" 7:8 0.17 Rev/Gal 2.9',
          '6.75" 7:8 0.30 Rev/Gal 6.4', 
          '6.75" 7:8 0.28 Rev/Gal 6.4',
          '6.75" 7:8 0.26 Rev/Gal 6.0',    
          '7" 7:8 0.31 Rev/Gal 6.0',
          '7" 7:8 0.31 Rev/Gal 7.5',     
          '7" 8:9 0.09 Rev/Gal 2.1',
          '7" 7:8 0.23 Rev/Gal 5.7',    
          '7" 5:6 0.41 Rev/Gal 8.2',
          '7" 5:6 0.45 Rev/Gal 11.1',   
          '7" 6:7 0.32 Rev/Gal 7.8',
          '7" 5:6 0.41 Rev/Gal 8.2',   
          '7.25" 7:8 0.26 Rev/Gal 7.2',
          '8" 4:5 0.25 Rev/Gal 5.3',     
          '8" 6:7 0.18 Rev/Gal 5.0',
          '8" 7:8 0.16 Rev/Gal 4.0 3',
          '8" 9:10 0.07 Rev/Gal 2.7',
          '8" 6:7 0.17 Rev/Gal 4.0',
          '8" 7:8 0.15 Rev/Gal 5.9 4',
          '8.25" 7:8 0.16 Rev/Gal 4.0 3',
          '8.25" 9:10 0.11 Rev/Gal 3.9',
          '8.500 8:9 0.17 Rev/Gal 6.7',
          '9.625" 6:7 0.10 Rev/Gal 3.5 6',
          '9.625" 7:8 0.11 Rev/Gal 4.8 6', 
          '9.625" 3:4 0.23 Rev/Gal 6.0 6', 
          '9.625" 6:7 0.13 Rev/Gal 5.0 7', 
          '9.625" 6:7 0.13 Rev/Gal 6.0 6', 
          '9.625" 7:8 0.11 Rev/Gal 4.8 6',    
          '11.25" 7:8 0.11 Rev/Gal 4.8 6', 
          '11.25" 7:8 0.11 Rev/Gal 4.8 6',
          'Other',     
]

# data = [    
#     'region':['Western US','Western US','Western US','Eastern US','Eastern US','Eastern US','Eastern US'],
#     'district':['Permian','Rockies','Williston','ArkLaTex','South Texas','MidCon','NEUS'],
#     'description':["In the Midland Basin, we're continuing to maintain high ROP throughout operations in stacked plays and the wide variety of interbedded formations that can lead to serious problems for operators. Ulterra also works alongside operators in the Delaware Basin to analyze extreme drilling parameters and introduce customized drill bits that can handle the toughest conditions with peak performance.",
#         "Our team of exs have experience across all basins in the Rocky Mountains and parts of California, such as the Denver Basin, Green River Basin, Powder River Basin, and others. With industry-leading experience for solving drilling challenges, including hard carbonates, gas, water flows, and slower drilling, Ulterra works to provide the right PDC bit to meet any requirement.",
#         "Ulterra's NODAK district focuses on helping operators drill through complex geology across the basin. We provide ex customer service for all areas spanning from portions of North Dakota, South Dakota, and Montana. This region covers the Williston Basin, including formations such as the Three Forks and Bakken.",
#         "Across the region including the Cotton Valley, and Haynesville Shale, operators encounter various drilling difficulties that require new bit designs and proper bit selection. Ulterra works alongside operators to solve problems that are faced in these difficult areas, and the challenges faced while drilling through the Travis Peak and Cotton Valley sands.",
#         "Ulterra helps operators from the Austin Chalk and Eagle Ford Shales to areas of the Gulf Coast achieve impressive drilling results. Our successes in South Texas have led to more application and performance improvements across other regions. In fact, our exs pioneered the use of Ulterra's CounterForce, SplitBlade, and FastBack bits throughout the South Texas oilfields.",
#         "Ulterra's involvement in the MidCon region includes plays in all of Oklahoma, Kansas, Arkansas, and the Texas Panhandle. The primary formations in this region include the Woodford, Meramec, Mississippian, Osage, and Springer. When operators face challenges in this region, like drilling through cy formations, hard rock, and interbedded formations, Ulterra is the supplier to count on to efficiently meet their drilling goals.",
#         "Ulterra provides operators in the Marcellus and Utica shales with the latest PDC bit developments that can tackle everything from curve and lateral runs to complex air drilling applications. One of Ulterra's latest innovations—AirRaid™—is the first PDC bit specifVertically designed for air drilling with outstanding durability and field-proven performance."
#     ]
# }
data = [['Permian','Western US',"In the Midland Basin, we're continuing to maintain high ROP throughout operations in stacked plays and the wide variety of interbedded formations that can lead to serious problems for operators. Ulterra also works alongside operators in the Delaware Basin to analyze extreme drilling parameters and introduce customized drill bits that can handle the toughest conditions with peak performance."],
        ['Rockies','Western US',"Our team of exs have experience across all basins in the Rocky Mountains and parts of California, such as the Denver Basin, Green River Basin, Powder River Basin, and others. With industry-leading experience for solving drilling challenges, including hard carbonates, gas, water flows, and slower drilling, Ulterra works to provide the right PDC bit to meet any requirement."],
        ['Williston','Western US',"Ulterra's NODAK district focuses on helping operators drill through complex geology across the basin. We provide ex customer service for all areas spanning from portions of North Dakota, South Dakota, and Montana. This region covers the Williston Basin, including formations such as the Three Forks and Bakken."],
        ['ArkLaTex','Eastern US',"Across the region including the Cotton Valley, and Haynesville Shale, operators encounter various drilling difficulties that require new bit designs and proper bit selection. Ulterra works alongside operators to solve problems that are faced in these difficult areas, and the challenges faced while drilling through the Travis Peak and Cotton Valley sands."],
        ['South Texas','Eastern US',"Ulterra helps operators from the Austin Chalk and Eagle Ford Shales to areas of the Gulf Coast achieve impressive drilling results. Our successes in South Texas have led to more application and performance improvements across other regions. In fact, our reps pioneered the use of Ulterra's CounterForce, SplitBlade, and FastBack bits throughout the South Texas oilfields."],    
        ['North Texas','Eastern US',"Ulterra helps operators achieve impressive drilling results. Our successes in North Texas have led to more application and performance improvements across other regions."],    
        ['MidCon','Eastern US',"Ulterra's involvement in the MidCon region includes plays in all of Oklahoma, Kansas, Arkansas, and the Texas Panhandle. The primary formations in this region include the Woodford, Meramec, Mississippian, Osage, and Springer. When operators face challenges in this region, like drilling through cy formations, hard rock, and interbedded formations, Ulterra is the supplier to count on to efficiently meet their drilling goals."],
        ['NEUS''Eastern US',"Ulterra provides operators in the Marcellus and Utica shales with the latest PDC bit developments that can tackle everything from curve and lateral runs to complex air drilling applications. One of Ulterra's latest innovations—AirRaid™—is the first PDC bit specifVertically designed for air drilling with outstanding durability and field-proven performance."],
        ['Canada','North America',"Ulterra can draw it's roots in the Canadian oilfield back all the way to United Diamond in the mid-90s. With reach into all Canadian plays across Saskaquewon, Ala, and into the Montney in British Colombia. When operators face challenges in this region, like drilling through cy formations, hard rock, and interbedded formations, Ulterra is the supplier to count on to efficiently meet their drilling goals."],
        
    ]
districtnames = ['Permian','ArkLaTex','South Texas','MidCon','Rockies','Williston','NEUS','North Texas','Canada']
sectionlist = {1:'Surface',2:'Vertical',3:'Drill Out',4:'Intermediate',5:'Curve',6:'Lateral',7:'Extended Lateral',8:'Extended Lateral',9:'Extended Lateral',10:'Extended Lateral'}
sectionlist3 = {1:'Vertical',2:'Curve',3:'Lateral',}
sectionlist4 = {1:'Surface',2:'Vertical',3:'Curve',4:'Lateral',}
sectionsize_index = {'Surface':9,'Vertical':23,'Drill Out':23,'Intermediate':30,'Vertical/Curve/Lateral':34,'Vertical/Curve':34,'Curve':34,'Curve/Lateral':34,'Lateral':34,'Extended Lateral':47}
sectionsize = {'Surface':17.5,'Vertical':12.25,'Drill Out':12.25,'Intermediate':9.875,'Vertical/Curve/Lateral':8.75,'Vertical/Curve':8.75,'Curve':8.75,'Curve/Lateral':6.75,'Lateral':6.75,'Extended Lateral':6.75}
sectionrop = {'Surface':200,'Vertical':180,'Drill Out':170,'Intermediate':150,'Vertical/Curve/Lateral':130,'Vertical/Curve':140,'Curve':90,'Curve/Lateral':80,'Lateral':75,'Extended Lateral':60}
customernames = ['Aethon','Comstock','ConocoPhillips','Ovintiv','EOG','Oxy','Devon','Chevron','Other']
# find logos here https://companieslogo.com/eog-resources/logo/#google_vignette
# alternate logo lookup option logo_url = 'https://logo.clearbit.com/' + website
basins = pd.DataFrame(data, columns=['district','region','description'])

def proposalmaker():
    # df_env = pd.DataFrame()
    # try:
    #     df_env = enverus_df()
    # except Exception as e:
    #     print(e)
    #     with st.expander('Enverus Read Error'):
    #         st.info(f"error enverusdf: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")    
    
    # rigoptions = dict(updateddate='gt(2023-01-01)',pagesize=10000,deleteddate='null')
    # df_env_rigs = enverus_rig_list(**rigoptions)
    # # with st.expander(f'enverus rigs {df_env_rigs.shape[0]}'):
    # #     st.write(df_env_rigs)
    # #     csv = df_env_rigs.to_csv(index=False).encode('utf-8')
    # #     st.download_button("Download Data",csv,"Enverus_Rigs.csv","text/csv",key='download-csv')
    
    # permitoptions = dict(updateddate='gt(2023-01-01)',fields='ENVOperator, WellName,WellID, API_UWI, Formation, Latitude, Longitude,',PermitStatus='ACTIVE',pagesize=10000,deleteddate='null')     # ENVOperator=customer,
    # df_env_permits = enverus_permit_list(**permitoptions)
    # with st.expander(f'enverus permits {df_env_permits.shape[0]}'):
    #     st.write(df_env_permits)
    #     csv = df_env_permits.to_csv(index=False).encode('utf-8')
    #     st.download_button("Download Data",csv,"Enverus_Permits.csv","text/csv",key='download-csv')
    
    # ftopoptions = dict(updateddate='gt(2023-07-01)',pagesize=10000,deleteddate='null')
    # df_env_ftops = enverus_ftop_list(**ftopoptions)
    # with st.expander(f'enverus rigs {df_env_ftops.shape[0]}'):
    #     st.write(df_env_ftops)
    #     csv = df_env_ftops.to_csv(index=False).encode('utf-8')
    #     st.download_button("Download Data",csv,"Enverus_Ftops.csv","text/csv",key='download-csv')
    
    
    
    # padding = 0
    # st.markdown(f''' <style>
    #             .reportview-container .sidebar-content {{
    #                 padding-top: {padding}rem;
    #             }}
    #             .appview-container .main .block-container{{
    #                 padding-top: {padding}rem;    
    #                 padding-right: {padding}rem;
    #                 padding-left: {padding}rem;
    #                 padding-bottom: {padding}rem;
    #             }}
    #             .reportview-container .main .block-container {{
    #                 padding-top: {padding}rem;
    #                 padding-right: {padding}rem;
    #                 padding-left: {padding}rem;
    #                 padding-bottom: {padding}rem;
    #             }}
    #         </style> ''',
    #         unsafe_allow_html=True,
    #     )
    
    # st.markdown(f""" <style>
    #     .reportview-container .main .block-container{{
    #         padding-top: {padding}rem;
    #         padding-right: {padding}rem;
    #         padding-left: {padding}rem;
    #         padding-bottom: {padding}rem;
    #     }} </style> """, unsafe_allow_html=True)
    
    
    hide_img_fs = '''
        <style>
        button[title="View fullscreen"]{
            visibility: hidden;}
        </style>
        '''
    st.markdown(hide_img_fs, unsafe_allow_html=True)
    
    
    # # st.markdown( """<style>
    # #     div[class*="stcaption"] body {font-size:  2rem !important;color: black;}
    # #     div[class*="stTextArea"] label {font-size:  2rem !important;color: black;}
    # #     div[class*="stTextInput"] label {font-size:  2rem !important;color: black;}
    # #     div[class*="stSelectBox"] label {font-size:  2rem !important;}
    # #     div[class*="stMultiSelect"] label {font-size:  2rem !important;color: black;}
    # #     div[class*="stNumberInput"] label {font-size:  2rem !important;color: black;}
    # #     .big-font {font-size: 1.5rem !important;color: #43c6db;}
    # #     </style>""", unsafe_allow_html=True)
    
    st.markdown( """<style>
        div[class*="stcaption"] body {font-size:  2rem;color: black;}
        div[class*="stTextArea"] label {font-size:  2rem;color: black;}
        div[class*="stTextInput"] label {font-size:  2rem;color: black;}
        div[class*="stSelectBox"] label {font-size:  2rem;}
        div[class*="stMultiSelect"] label {font-size:  2rem;color: black;}
        div[class*="stNumberInput"] label {font-size:  2rem;color: black;}
        .big-font {font-size: 1.5rem;color: #43c6db;}
        </style>""", unsafe_allow_html=True)
    
    
    st.markdown(
        '''
        <style>
        .streamlit-expanderHeader {
            background-color: #37a8b2;
            color: white; # Adjust this for expander header color
        }
        </style>
        ''',
        unsafe_allow_html=True
    )    
    
        # .streamlit-expanderContent {
        #     background-color: #f7f7f7;
        #     color: black; # Expander content color
        # }
    
    # Moving triathalon running animation
    # runningToggle(True)
    
    # st.markdown(".stTextInput > label {font-size:105%; font-weight:bold;} ",unsafe_allow_html=True) #for all text-input label sections
    # st.markdown(".stSelectBox > label {font-size:105%; font-weight:bold;} ",unsafe_allow_html=True) #for all multi-select label sections
    # st.markdown(".stMultiSelect > label {font-size:105%; font-weight:bold;} ",unsafe_allow_html=True) #for all multi-select label sections
    
    uploaded_bpm = st.sidebar.file_uploader(label='Upload BPM file here', type= ['xlsx'], accept_multiple_files=False) 
    
    st.sidebar.markdown('<a href="mailto:ccasad@ulterra.com?subject=ADIRA Proposals Help & Feedback&body=Hey Chris, ADIRA rocks!"><button style="color:#43c6db;background-color:white;text-decoration:none;border-radius:4px;border:#43c6db;padding:10px 24px;">Email for Help & Feedback</button></a>', unsafe_allow_html=True)
        
    if not uploaded_bpm:
        # if "showpartnumber" not in st.session_state:
            # st.session_state.sizecount = 1
            # st.session_state.sectioncount = 3
            # st.session_state.wsectioncount = 1
            # st.session_state.showpartnumber = False
            # # st.session_state.showgage = True
            # # st.session_state.showpartnumber2 = False
            # # st.session_state.showgage2 = True
            # st.session_state.disabled = False
            
        if 'wellname' not in st.session_state:
            st.session_state.bpm = False
            st.session_state.bpmoverwrite = False  
            st.session_state.csv = False
            st.session_state.csvoverwrite = False 
            st.session_state.district = "Permian"
            st.session_state.custoemp = "Aethon"
            st.session_state.preparedby = " "
            st.session_state.preparedfor = " "
            st.session_state.wellname = " "
            st.session_state.rigname = " "
            st.session_state.lat = ""
            st.session_state.long = ""
            # st.session_state.opprofile = ''
            # st.session_state.sectioncount = 3
            st.session_state.spreadrate = '100,000'
            st.session_state.triprate = '1000'
            
        
    st.header('ADIRA Proposals')
    st.caption('Select **Bit Menu** or **Well Proposal** to start...')    
    
    # with st.container():                  
    #     col1, col2, col3 = st.columns(3)
    #     with col1:
    #         district = st.selectbox('Select district:',['Permian','ArkLaTex','South Texas','Rockies','Williston','NEUS'],)
    #         customer = st.selectbox('Select Customer:',['ConocoPhillips','Ovintiv','EOG','TapRock','DiamondBack','Chevron'],)
    #         # sections = st.multiselect('Which Sections?',['Surface','Vertical','Drill Out','Intermediate','Vertical/Curve','Vertical/Curve/Lateral','Curve','Curve/Lateral','Lateral','Extended Lateral'],)
        
    #     with col2:
    #         rigname = st.text_input("Enter Rig Name/Number:",)
    #         wellname = st.text_input("Enter Well Name:",)
    #     with col3:
    #         st.write('Display Options:')
    #         showpartnumber = st.checkbox("Show Part Number",value=False, key="showpartnumber")
    #         showgage = st.checkbox("Show Gage Type",value=True, key="showgage")        
        
    #     # Add css to make text bigger
    #     # st.markdown(
    #     #     """
    #     #     <style>
    #     #     textarea {
    #     #         font-size: 3rem !important;
    #     #     }
    #     #     input {
    #     #         font-size: 3rem !important;
    #     #     }
    #     #     </style>
    #     #     """,
    #     #     unsafe_allow_html=True,
    #     # )

            
    # st.write(
    #     """This app accomodates the blog [here](https://blog.streamlit.io/auto-generate-a-dataframe-filtering-ui-in-streamlit-with-filter_dataframe/)
    #     and walks you through one example of how the Streamlit
    #     Data Science Team builds add-on functions to Streamlit.
    #     """
    # )


    def convert_png_transparent(src_file, dst_file, bg_color=(255,255,255)):
        image = Image.open(src_file).convert("RGBA")
        array = np.array(image, dtype=np.ubyte)
        mask = (array[:,:,:3] == bg_color).all(axis=2)
        alpha = np.where(mask, 0, 255)
        array[:,:,-1] = alpha
        Image.fromarray(np.ubyte(array)).save(dst_file, "png")
        
    def filter_dataframe(dffltr: pd.DataFrame) -> pd.DataFrame:
        """
        Adds a UI on top of a dataframe to let viewers filter columns

        Args:
            df (pd.DataFrame): Original dataframe

        Returns:
            pd.DataFrame: Filtered dataframe
        """
        modify = True
        # modify = st.checkbox("Add filters")

        if not modify:
            return df

        dffltr = dffltr.copy()

        # Try to con datetimes into a standard format (datetime, no timezone)
        for col in dffltr.columns:
            if is_object_dtype(dffltr[col]):
                try:
                    dffltr[col] = pd.to_datetime(dffltr[col])
                except Exception:
                    pass

            if is_datetime64_any_dtype(dffltr[col]):
                dffltr[col] = dffltr[col].dt.tz_localize(None)

        modification_container = st.container()

        with modification_container:
            to_filter_columns = st.multiselect("Filter dataframe on", dffltr.columns)
            for column in to_filter_columns:
                left, right = st.columns((1, 20))
                left.write("↳")
                # Treat columns with < 10 unique values as categorVertical
                if is_categorical_dtype(dffltr[column]) or dffltr[column].nunique() < 10:
                    user_cat_input = right.multiselect(
                        f"Values for {column}",
                        dffltr[column].unique(),
                        default=list(dffltr[column].unique()),
                    )
                    df = dffltr[dffltr[column].isin(user_cat_input)]
                elif is_numeric_dtype(dffltr[column]):
                    _min = float(dffltr[column].min())
                    _max = float(dffltr[column].max())
                    step = (_max - _min) / 100
                    user_num_input = right.slider(
                        f"Values for {column}",
                        _min,
                        _max,
                        (_min, _max),
                        step=step,
                    )
                    df = dffltr[dffltr[column].between(*user_num_input)]
                elif is_datetime64_any_dtype(dffltr[column]):
                    user_date_input = right.date_input(
                        f"Values for {column}",
                        value=(
                            dffltr[column].min(),
                            dffltr[column].max(),
                        ),
                    )
                    if len(user_date_input) == 2:
                        user_date_input = tuple(map(pd.to_datetime, user_date_input))
                        start_date, end_date = user_date_input
                        dffltr = dffltr.loc[dffltr[column].between(start_date, end_date)]
                else:
                    user_text_input = right.text_input(
                        f"Substring or regex in {column}",
                    )
                    if user_text_input:
                        df = dffltr[dffltr[column].str.contains(user_text_input)]

        return df
    
    # uploaded_file = st.sidebar.file_uploader('Upload CSV or Excel File here', type=['csv', 'xlsx'])    
    # if uploaded_file is not None:
        # df = pd.read_csv(uploaded_file, encoding_errors='ignore')
    
    # old hardcoded DDA sample
    # df = pd.read_csv('data/adira book3.csv', encoding_errors='ignore')
    
    @st.cache_data
    def get_bit_records(csvdir):
        # dfrec = pd.read_csv('data/records_USonly.csv') 
        dfrec = pd.read_csv(csvdir) 
        #  Clean up data        
        dfrec['Latitude'] = dfrec['Latitude'].clip(upper=90)
        dfrec['Longitude'] = dfrec['Longitude'].clip(upper=180)
        dfrec['ROP'] = dfrec['ROP'].clip(upper=500)
        dfrec['ROP'] = dfrec['ROP'].replace(500, np.nan)
        # https://codereview.stackexchange.com/questions/185389/dropping-rows-from-a-pandas-dataframe-where-some-of-the-columns-have-value-0
        dropcolumns = ['DepthOut', 'ROP', 'Latitude', 'Longitude'] 
        dfrec = dfrec.replace(0, np.nan).dropna(axis=0, how='any', subset=dropcolumns).fillna(0)
        return dfrec 
    
    @st.cache_resource
    def get_bit_list(url,headers):
        response = requests.get(url=url, headers=headers)
        df_list = pd.DataFrame(json.loads(response.text))
        return df_list
        
    @st.cache_resource
    def get_bit_files(url,headers):
        response = requests.get(url=url, headers=headers)
        zipper = zipfile.ZipFile(io.BytesIO(response.content))  
        zipper.extractall('.')  
        # return zipper
    
    @st.cache_data
    def get_geolocation(lat, long):
        glocation = geolocator.reverse(lat+","+long)
        return glocation
        
    # API DDA
    url = 'https://bithub.ulterra.com/adira_dda_fetch'
    headers = {'Adira_API_Key': st.secrets['Adira_API_Key']}
    # response = requests.get(url=url, headers=headers)
    # df = pd.DataFrame(json.loads(response.text))
    df = get_bit_list(url, headers) 
    
    if df is not None:          
        
        # tab1 = st.tabs(["      Well Proposal      "])         
        
    # with tab1:  
        def WellProposal():
            # Version 2 of work flow            
        # with st.container():
            # load proposal history
            df_history = blob_download('Proposal_history.csv')
            # st.write(df_history)
                                
            # Load Bit Record DF
            df_br = get_bit_records('data/records_USonly.csv')                    
            df_br['Drilled'] = (df_br['DepthOut'] - df_br['DepthIn'])
            bpm_port = {}
            # bpm_table = pd.DataFrame()
            
            if uploaded_bpm:
                st.session_state.bpm = True
            else:                        
                st.session_state.bpm = False
                
            best_match = None            
            riglist = [' ']
            riglist.extend(df_br['Rig_Name_Full'].sort_values(ascending=False).unique().tolist())
            # riglist.extend(df_env_rigs['RigName_Number'].sort_values(ascending=False).unique().tolist())
            
            if st.session_state.bpm is True:
                # Read and fill page
                # try:
                #     df_bpm = pd.read_excel(uploaded_bpm, engine= 'openpyxl')                                                  
                # except Exception as e:
                #     print(e)  
                #     df_bpm = None
                #     st.info(f"error df_bpm: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")      
                # https://discuss.streamlit.io/t/read-particular-cells-values-from-an-excel-file/18590               
                wb = openpyxl.load_workbook(uploaded_bpm, data_only=True,)
                sheet = wb['ProgramDashboard']
                # B4 = sheet['B4']
                # C1 = sheet['C1'] # read direct value in cell C1
                # C1 = sheet.cell(row=1,column=3) # or this has the same effect
                # print(C1.value)   

                for i in range(1,14):
                    bpm_port[i] = sheet[f'B{i}'].value
                
                bpm_port[14] = sheet[f'G2'].value # depthin
                
                try:
                    bpm_table=pd.read_excel(uploaded_bpm,usecols='I:Y',sheet_name='ProgramDashboard',header=0)[0:27]
                    # st.write(bpm_table)
                except Exception as e:
                    st.info(f"error bpmsrc: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    # try:
                    #     bpmsrc=pd.read_excel(rf'{wb.name}',usecols='I:Y',sheet_name='ProgramDashboard',header=0)[0:27]
                    #     st.write(bpmsrc)
                    # except Exception as e:
                    #     st.info(f"error bpmsrc2: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                        
                
                # if 'preparedby' not in st.session_state: 
                # If uploaded file, inject header data
                if st.session_state.bpm is True and st.session_state.bpmoverwrite is False:   
                    # # st.session_state.district = " "
                    # st.session_state.custoemp = bpm_port[2]
                    st.session_state.preparedby = bpm_port[11]
                    st.session_state.preparedfor = bpm_port[6]
                    st.session_state.wellname = bpm_port[4]
                    # st.session_state.rigname = bpm_port[3]
                    # fuzzy match https://typesense.org/learn/fuzzy-string-matching-python/
                    # best_match = process.extractOne(bpm_port[3], df_br['Rig_Name_Full'].sort_values(ascending=False).unique().tolist(), scorer=fuzz.token_sort_ratio)
                    best_match = process.extractOne(bpm_port[3], riglist, scorer=fuzz.token_sort_ratio)
                    # riglist.extend(df_env_rigs['RigName_Number'].sort_values(ascending=False).unique().tolist())
                    # st.write(best_match)                        
                    st.session_state.rigname = str(list(best_match)[0])
                    if str(list(best_match)[0]) not in riglist:
                        riglist.extend(str(list(best_match)[0]))
                    bpm_port[10] = bpm_port[10].replace(" ", "")
                    st.session_state.lat = bpm_port[10].split(',', 1)[0] 
                    st.session_state.long = bpm_port[10].rsplit(',', 1)[-1] 
                    # st.session_state.sectioncount = " "
                    st.session_state.spreadrate = str(float(bpm_port[12])*24)
                    st.session_state.triprate = str(bpm_port[13])
                    
                    
                    
                    # Set BPM Table to df                        
                    st.session_state.sectioncount = int(bpm_table['Run'].max())
                    
                # for i in range(1,26):
                #     if len(str(sheet[f'J{i}'].value)) > 1: 
                #         st.session_state.sectioncount += 1
                #         for j in range(8,24):                                    
                #             bpm_table[i,j-8] = sheet[f'{chr(65+j)}{i}'].value
                #             st.write(f'{j}-{chr(65+j)}{i}')
                #             st.write(sheet[f"{chr(65+j)}{i}"].value)
                #             st.write(bpm_table[i,j-8])
                #     else:
                #         break

                # df["balance"] = df.name.apply(lambda x: wb[x][f'G{wb[x].max_row}'].value)

                # st.write(bpm_table)
            else:
                for i in range(1,14):
                    bpm_port[i] = ' '
                    
            if st.session_state.csv is True and 'opprofile' in st.session_state:
                # Read and fill page
                csv_table = df_history.loc[df_history['wellname']==st.session_state.opprofile]
                csv_table = csv_table.sort_values(['Run'], ascending=True).reset_index(drop=True)
                csv_port = csv_table.iloc[0]
                # st.write(csv_port)
                
                # If uploaded file, inject header data
                if st.session_state.csv is True and st.session_state.csvoverwrite is False:
                    try:
                        if not pd.isna(csv_port['district']):
                            st.session_state.district = str(csv_port['district'])
                        if not pd.isna(csv_port['operator']):
                            st.session_state.customer = str(csv_port['operator'])
                        if not pd.isna(csv_port['preparedby']):
                            st.session_state.preparedby = str(csv_port['preparedby'])
                        if not pd.isna(csv_port['preparedfor']):
                            st.session_state.preparedfor = str(csv_port['preparedfor'])                                
                        if not pd.isna(csv_port['wellname']):
                            st.session_state.wellname = csv_port['wellname']
                        if not pd.isna(csv_port['targetformation']):
                            st.session_state.targetformation = str(csv_port['targetformation'])
                        if not pd.isna(csv_port['rigname']):
                            st.session_state.rigname = str(csv_port['rigname'])
                        # fuzzy match https://typesense.org/learn/fuzzy-string-matching-python/
                        # best_match = process.extractOne(csv_port['rigname'], df_br['Rig_Name_Full'].sort_values(ascending=False).unique().tolist(), scorer=fuzz.token_sort_ratio)
                        # st.session_state.rigname = str(list(best_match)[0])
                        if not pd.isna(csv_port['lat']):
                            st.session_state.lat = f"{csv_port['lat']}"
                        if not pd.isna(csv_port['long']):
                            st.session_state.long = f"{csv_port['long']}"
                        st.session_state.sectioncount = int(csv_port['sectioncount'])
                        if not pd.isna(csv_port['spreadrate']):
                            st.session_state.spreadrate = str(csv_port['spreadrate']).replace(",", "")
                        if not pd.isna(csv_port['triprate']):
                            st.session_state.triprate = str(csv_port['triprate']).replace(",", "")
                        
                        # # Set BPM Table to df                        
                        # st.session_state.sectioncount = int(csv_port['Run'].max())
                    except Exception as e:
                        st.info(f"error csvport: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
            # else:
            #     for i in range(1,14):
            #         bpm_port[i] = ' '
            
            col1, col2, col3 = st.columns(3)
            with col1:
                district = st.selectbox('Select **District**:',districtnames,key='district')
                custoemp = st.selectbox('Select **Customer**:',customernames,key='custoemp')
                uploaded_custlogo = None
                if custoemp == 'Other': 
                    custcol1, custcol2 = st.columns(2)
                    with custcol1:
                        customer = st.text_input("Customer Name:",)
                    with custcol2:
                        uploaded_custlogo = st.file_uploader('Upload logo here', type=['png']) 
                else:
                    customer = custoemp 
                    st.session_state.customer1 = customer
                
                # permitoptions = dict(ENVOperator=st.session_state.customer1,PermitStatus='ACTIVE',updateddate='gt(2023-01-01)',pagesize=10000,deleteddate='null')
                # df_env_permits = enverus_permit_list(**permitoptions)
                # df_env_permits['WellName'] # WellID API_UWI Formation Latitude Longitude
                if uploaded_custlogo is not None:
                    # customerlogo = uploaded_custlogo.getvalue()
                    customerlogo = uploaded_custlogo
                else:
                    customerlogo = f'data/customer logos/{customer}.png'
                
                preparedfor = st.text_input("Prepared for:",key='preparedfor')
                preparedby = st.text_input("Ulterra Rep:",key='preparedby')
                
                # https://github.com/gagan3012/streamlit-tags
                # sectionkeys = st_tags(label='Type Section Names:',text='Press enter to add more',value=['Surface', 'Vertical'],suggestions=['Surface', 'Vertical', 'Drill Out', 'Intermediate', 'Curve', 'Lateral', 'Vertical/Curve', 'Vertical/Curve/Lateral', 'Curve/Lateral','Extended Lateral'],maxtags = 10,key='1')
                # sectioncount = len(sectionkeys)
            with col2:
                # wellname = st.text_input("Enter Well Name:",key='wellname')
                try:
                    wellname = st.text_input("Enter Well Name:",key='wellname')
                except Exception as e:
                    st.text_input("Enter Well Name2:",key='wellname')
                    st.info(f"error wellname: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                
                # rigname = st.text_input("Enter Rig Name:",)
                # list(','.join('%0.3f' %x for x in y) for y in df.values)
                fuzzp = None
                if best_match is not None:
                    fuzzp = str(list(best_match)[0])
                rigname = st.selectbox(f"Select Rig Name: {fuzzp}",riglist, key='rigname')
                # https://github.com/RobertoFN/streamlit-datalist
                # rigname = stDatalist(f"Select Rig Name: {st.session_state.rigname}", riglist, key='rigname')
                targetformation = st.text_input("Target Formation:", key='targetformation')
                
                gpscol1, gpscol2 = st.columns(2,gap="small")
                with gpscol1:
                    lat = st.text_input("Enter Latitude:",key='lat')
                with gpscol2:
                    long = st.text_input("Enter Longitude:",key='long')
                
                if lat and long:
                    # glocation = geolocator.reverse(lat+","+long)
                    try:
                        glocation = get_geolocation(lat, long)
                        gaddress = glocation.raw['address']                        
                        # gcity = gaddress.get('city', '')
                        gstate = gaddress.get('state', '')
                        gcounty = gaddress.get('county', '')
                        # gcountry = gaddress.get('country', '')
                        # gcode = gaddress.get('country_code')
                        # gzipcode = gaddress.get('postcode')
                    except Exception as e:
                        # st.text_input("Enter Well Name2:",key='wellname')                                
                        gstate = ''
                        gcounty = ''
                        st.info(f"error glocation: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                else:
                    gstate = ''
                    gcounty = ''
                    
                # sections = st.multiselect('Which Sections?',['Surface','Vertical','Drill Out','Intermediate','Vertical/Curve','Vertical/Curve/Lateral','Curve','Curve/Lateral','Lateral','Extended Lateral'],)
                # sectioncount = st.number_input('Number of Sections:', min_value=1, max_value=10, value=1, step=1,help='Select the number of sections in your well.')                                                                            
                sectioncount = st.slider('Number of Bit Runs:',min_value=0,max_value=15,step=1,help='Select the number of sections in your well. (Selecting 8 or more may result in formatting errors.)',key='sectioncount')                      
                
            with col3:
                st.write('Assumptions:')
                ratecol1, ratecol2 = st.columns(2,gap="small")
                with ratecol1:
                    spreadrate = st.text_input('Spread Rate/day:',disabled=False, key=f'spreadrate')
                    spreadrate.replace(',', '') 
                with ratecol2:
                    triprate = st.text_input('Trip Rate/hr:',disabled=False, key=f'triprate')
                    triprate.replace(',', '') 
                optctrl1, optctrl2 = st.columns(2,gap="small")
                with optctrl1:
                    st.write('Display Options:')
                    # Smartfill will attempt to pre-load as many data powered defaults as possible
                    smartfill = st.checkbox("AutoFill",value=True, key="smartfill",disabled=True) 
                    # add smarts to filter history selection by operator name                            
                    # Load Previous Form
                    def loadrecord():
                        st.session_state.csv = True
                        st.session_state.csvoverwrite = False 
                        # Fetch records from history into DF
                        
                        # st.session_state.label_list.append(st.session_state.label_input)
                        # st.session_state.count += 1
                        
                    with st.form(key='template_form_loadprev'):
                        profilelist = [' ']
                        profilelist.extend(df_history['wellname'].loc[(df_history['operator'] == customer)].unique().tolist()) #  .sort_values(ascending=False).unique().tolist())
                        opprofile = st.selectbox('Select Profile:', options=profilelist, key='opprofile', label_visibility='collapsed',) # disabled=True
                        with st.spinner('Loading...'): 
                            loadrecord = st.form_submit_button('Load', on_click=loadrecord)
                        
                        
                    
                with optctrl2:
                    st.write('Smart Options:')
                    smartoffset = st.checkbox("C-Rank Offsets",value=True, key="smartoffset")  
                    bitstats = st.checkbox("Calculate Bit Stats",value=True, key="bitstats")  
                    graycomps = st.checkbox("Gray Competitors",value=True, key="graycomps")  
                    rsstoggle = st.checkbox("RSS",value=True, key="rsstoggle", disabled=True)  
                    # showperfchart = st.checkbox("Performance Chart",value=True, key="showperfchart")  
                    # showdvdchart = st.checkbox("Time/Depth Chart",value=True, key="showdvdchart")    
                    # showcostchart = st.checkbox("Value Chart",value=True, key="showcostchart")           
                
            
            st.divider()



            
            # def calculate_distance(source_coords, target_coords):
            #     try:
            #         miles = geodesic(source_coords, target_coords).miles
            #     except Exception as e:
            #         miles = 9999
            #         st.write(f'source {source_coords} || target {target_coords}')
            #         st.info(f"error calculate_distance: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")   
                
            #     # return geodesic(source_coords, target_coords).miles
            #     return miles

            # def calculate_ranking_score(record, target_record, vehicle_class):
            #     # if -90 <= record['Latitude'] <= 90 and -180 <= record['Longitude'] <= 180:
            #     record.clip(upper=pd.Series({'Latitude': 90}), axis=1)
            #     record.clip(upper=pd.Series({'Longitude': 180}), axis=1)
            #     distance_score = calculate_distance((record['Latitude'], record['Longitude']), 
            #                                         (float(target_record['Latitude']), float(target_record['Longitude'])))
            #     # else:
            #     #     distance_score = 9999
            #     size_score = abs(record['Size'] - target_record['Size'])
            #     speed_score = abs(record['ROP'] - target_record['ROP'])

            #     # Check if the vehicle class is available in the record, otherwise consider it as a separate class
            #     if 'class' in record:
            #         class_score = 0 if record['Size'] == vehicle_class else 1
            #     else:
            #         class_score = 1

            #     # Assuming the target_record includes the start and end mileage points
            #     start_mileage = float(target_record['DepthIn'])
            #     end_mileage = float(target_record['DepthOut'])
            #     segment_start = record['DepthIn']
            #     segment_end = record['DepthOut']
            #     # Calculate the absolute difference between the middle segment start and end mileage
            #     segment_mileage_diff = abs((segment_end - segment_start) - (end_mileage - start_mileage))
                
            #     # You can adjust the weights below based on how you want to prioritize the factors.
            #     # For example, if distance is more important, increase its weight, and vice versa.
            #     distance_weight = 0.25
            #     size_weight = 0.15
            #     speed_weight = 0.15
            #     class_weight = 0.2
            #     segment_weight = 0.25

            #     total_score = (distance_weight * (1/distance_score) +
            #                 size_weight * size_score +
            #                 speed_weight * speed_score +
            #                 class_weight * class_score +
            #                 segment_weight * segment_mileage_diff)

            #     return total_score
            
            
            # def calculate_ranking_score2(record, target_record, vehicle_class):
            #     # if -90 <= record['Latitude'] <= 90 and -180 <= record['Longitude'] <= 180:
            #     record.clip(upper=pd.Series({'Latitude': 90}), axis=1)
            #     record.clip(upper=pd.Series({'Longitude': 180}), axis=1)
            #     record['distance_score'] = calculate_distance((record['Latitude'], record['Longitude']), 
            #                                         (float(target_record['Latitude']), float(target_record['Longitude'])))
            #     # else:
            #     #     distance_score = 9999
            #     record['size_score'] = abs(record['Size'] - target_record['Size'])
            #     record['speed_score'] = abs(record['ROP'] - target_record['ROP'])

            #     # Check if the vehicle class is available in the record, otherwise consider it as a separate class
            #     if 'class' in record:
            #         record.class_score = 0 if record['Size'] == vehicle_class else 1
            #     else:
            #         record.class_score = 1

            #     # Assuming the target_record includes the start and end mileage points
            #     start_mileage = float(target_record['DepthIn'])
            #     end_mileage = float(target_record['DepthOut'])
            #     segment_start = record['DepthIn']
            #     segment_end = record['DepthOut']
            #     # Calculate the absolute difference between the middle segment start and end mileage
            #     record['segment_mileage_diff'] = abs((segment_end - segment_start) - (end_mileage - start_mileage))
                
            #     # You can adjust the weights below based on how you want to prioritize the factors.
            #     # For example, if distance is more important, increase its weight, and vice versa.
            #     distance_weight = 0.25
            #     size_weight = 0.15
            #     speed_weight = 0.15
            #     class_weight = 0.2
            #     segment_weight = 0.25

            #     record['total_score'] = (distance_weight * (1/record.distance_score) +
            #                 size_weight * record.size_score +
            #                 speed_weight * record.speed_score +
            #                 class_weight * record.class_score +
            #                 segment_weight * record.segment_mileage_diff)

            #     return record
            
            def calc_distance(row, site_coords):
                station_coords = (row['Latitude'], row['Longitude'])                    
                try:
                    miles = geodesic(site_coords, station_coords).miles
                except Exception as e:
                    miles = 9999                        
                    st.write(f'site {site_coords} station {station_coords}')
                    # st.write(f'source {source_coords} || target {target_coords}')
                    st.info(f"error vector calc_distance: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                return(miles)
            
            @st.cache_data
            def vector_ranking_score(df_ref, target_record):
                # st.write('vector target')
                # st.write(target_record)
                # st.write(df_ref)
                criteria_tag = ''
                record = df_ref.loc[(df_ref['Size'] == target_record['Size'])]
                criteria_tag += f"{target_record['Size']}"
                record['Drilled'] = (record['DepthOut'] - record['DepthIn'])
                record['Hrs'] = (record['Drilled'] / record['ROP'])
                # record = record.loc[(record['Drilled'] >= 100 )]                    
                record = record.loc[(record['DepthIn'].between((0.8 * float(target_record['DepthIn'])),(1.2 * float(target_record['DepthIn']))))] 
                criteria_tag += f', 20% +/- Depthin'
                # record = record.clip(upper=pd.Series({'Latitude': 90}), axis=1)
                # record = record.clip(upper=pd.Series({'Longitude': 180}), axis=1)
                # record['Latitude'] = record['Latitude'].clip(upper=90)
                # record['Longitude'] = record['Longitude'].clip(upper=180)
                # dropcolumns = ['DepthOut', 'ROP', 'Latitude', 'Longitude'] # https://codereview.stackexchange.com/questions/185389/dropping-rows-from-a-pandas-dataframe-where-some-of-the-columns-have-value-0
                # record = record.replace(0, pd.np.nan).dropna(axis=0, how='any', subset=dropcolumns).fillna(0)
                record['distance_score'] = record.apply(calc_distance, site_coords=(float(target_record['Latitude']),  float(target_record['Longitude'])), axis=1)
                # Logic to return at least 30 results
                results_threshold = 20
                offset_rad = 1
                if len(record.loc[(record['distance_score'] <= 1)] ) > results_threshold:
                    record = record.loc[(record['distance_score'] <= 10)]  
                    offset_rad = 1
                elif len(record.loc[(record['distance_score'] <= 3)] ) > results_threshold:
                    record = record.loc[(record['distance_score'] <= 10)]  
                    offset_rad = 3
                elif len(record.loc[(record['distance_score'] <= 5)] ) > results_threshold:
                    record = record.loc[(record['distance_score'] <= 10)] 
                    offset_rad = 5 
                elif len(record.loc[(record['distance_score'] <= 10)] ) > results_threshold:
                    record = record.loc[(record['distance_score'] <= 10)]  
                    offset_rad = 10
                elif len(record.loc[(record['distance_score'] <= 40)] ) > results_threshold:
                    record = record.loc[(record['distance_score'] <= 40)]  
                    offset_rad = 40
                elif len(record.loc[(record['distance_score'] <= 100)] ) > results_threshold:
                    record = record.loc[(record['distance_score'] <= 100)]  
                    offset_rad = 100
                elif len(record.loc[(record['distance_score'] <= 500)] ) > results_threshold:
                    record = record.loc[(record['distance_score'] <= 500)]  
                    offset_rad = 500
                
                criteria_tag += f', {offset_rad}-mile'
                dfchecktd = record.loc[(record['DG_RP'].isin(['TD','BHA']))]
                if len(dfchecktd) > 5 and dfchecktd['BitMfgr'].value_counts()['ULT'] > 0:
                    record = record.loc[(record['DG_RP'] == 'TD')]
                    criteria_tag += f', TD/BHA only'
                
                # record['distance_score'] = calculate_distance((record['Latitude'], record['Longitude']), (target_record['Latitude'], target_record['Longitude']))                    
                # study_lat = float(target_record['Latitude'])
                # study_lon = float(target_record['Longitude'])
                # record['distance_score'] = geodesic((record['Latitude'], record['Longitude']), (study_lat, study_lon)).miles
                # size_score = abs(record['Size'] - target_record['Size'])
                record['speed_score']  = abs(record['ROP']) # - target_record['ROP'])

                # # Check if the vehicle class is available in the record, otherwise consider it as a separate class
                # if 'class' in record:
                #     class_score = 0 if record['Size'] == vehicle_class else 1
                # else:
                #     class_score = 1
                
                # Assuming the target_record includes the start and end mileage points
                # start_mileage = target_record['DepthIn']
                # end_mileage = target_record['DepthOut']
                # segment_start = record['DepthIn']
                # segment_end = record['DepthOut']
                
                record['depthin_score'] = abs(record['DepthIn'] - float(target_record['DepthIn']))
                # Calculate the absolute difference between the middle segment start and end mileage
                record['segment_depth_diff']  = abs(((record['DepthOut'] - record['DepthIn'])) / ((float(target_record['DepthOut']) - float(target_record['DepthIn']))))
                # filter out ftg drilled that are less than 75% or more than 300% target 
                record = record.loc[(record['segment_depth_diff'].between(0.75,3))] 
                criteria_tag += f', 75%-300% Drilled Range'
                record['segment_depth_diff'] = record['segment_depth_diff'].clip(upper=1) 
                criteria_tag += f', No 1%'

                # You can adjust the weights below based on how you want to prioritize the factors.
                # For example, if distance is more important, increase its weight, and vice versa.
                rig_match = 0.4
                distance_weight = 0.5
                speed_weight = 0.1
                depth_weight = 0.1
                # size_weight = 0.15
                # class_weight = 0.2

                # record['total_score']  = (distance_weight * (1/record['distance_score']) +
                #             size_weight * 1 +
                #             speed_weight * record['speed_score'] +
                #             class_weight * 1 +
                #             depth_weight * record['segment_depth_diff'])
                record['speed_score'] = (speed_weight * record['speed_score'])
                record['depth_score'] = (depth_weight * record['segment_depth_diff'])
                record['total_score']  = (distance_weight * (1/record['distance_score']) +
                                        speed_weight * record['speed_score'] +
                                        depth_weight * record['segment_depth_diff'])
                
                
                record2 = record.sort_values(by=['total_score'], ascending=False,ignore_index=True) 
                return record2, offset_rad, criteria_tag
            
            # def rank_records(csv_file, target_record, vehicle_class):
            #     dfrr = pd.read_csv(csv_file)

            #     # Calculate ranking score for each record in the dataframe
            #     dfrr['ranking_score'] = dfrr.apply(lambda row: calculate_ranking_score(row, target_record, vehicle_class), axis=1)

            #     # Sort the dataframe based on ranking score in ascending order
            #     ranked_df = dfrr.sort_values(by='ranking_score')
            #     st.write(ranked_df)
            #     return ranked_df

            # def find_high_ranking_matches(csv_file, target_record, vehicle_class, num_matches=3, max_distance=100):
            #     dfhrr = pd.read_csv(csv_file)
            #     key_column = 'comp_key'  # Name of the column to store the matching key or tag

            #     # Check if there are already records with the same matching key
            #     previous_matches = dfhrr[dfhrr[key_column].notnull()]

            #     if len(previous_matches) >= num_matches:
            #         # If we have enough previous matches, use them to expedite future matching
            #         filtered_records = dfhrr[dfhrr[key_column].isin(previous_matches[key_column])]
            #     else:
            #         # Otherwise, perform a new search using the previous logic
            #         ranked_records = rank_records(csv_file, target_record, vehicle_class)
            #         high_ranking_matches = ranked_records.head(num_matches)
            #         current_distance = 0

            #         while len(high_ranking_matches) < num_matches and current_distance <= max_distance:
            #             current_distance += 10  # Increase the distance search by 10 miles
            #             # Re-rank the records based on the increased distance
            #             ranked_records = rank_records(csv_file, target_record, vehicle_class)

            #             # Filter the records based on the current distance and vehicle class
            #             filtered_records = ranked_records[
            #                 (ranked_records['ranking_score'] <= current_distance) & 
            #                 ((ranked_records['class'] == vehicle_class) | (vehicle_class == 'Any'))
            #             ]

            #             high_ranking_matches = filtered_records.head(num_matches)

            #         if len(high_ranking_matches) < num_matches:
            #             # If still not enough high ranking matches, expand the search to other class sizes
            #             vehicle_class = 'Any'
            #             ranked_records = rank_records(csv_file, target_record, vehicle_class)
            #             high_ranking_matches = ranked_records.head(num_matches)

            #         # Store the matching key or tag in the DataFrame for future matching
            #         matching_key = 'key_{}'.format(len(previous_matches) + 1)
            #         df.at[high_ranking_matches.index, key_column] = matching_key
            #         filtered_records = dfhrr[dfhrr[key_column].isin([matching_key])]

            #     return filtered_records
                
            
            
            
            # https://developers.google.com/maps/documentation/geocoding/requests-reverse-geocoding
            # https://maps.googleapis.com/maps/api/geocode/json?latlng=40.714224,-73.961452&key=YOUR_API_KEY

            # def matchmain(lat,long,size,speed,runentry,runexit,classtype):
            #     # uploaded_file = st.file_uploader("Choose a CSV file", type=["csv"])
            #     uploaded_file = 'data/records_USonly.csv'                    
            #     # if uploaded_file:
            #     #     df = pd.read_csv(uploaded_file)
                
            #     # Create input fields for the target record
            #     st.write(f"Target Record Information: {lat}|{long}|{size}|{speed}|{runentry}|{runexit}|{classtype}")
            #     target_lat = lat # st.number_input("Latitude:")
            #     target_lon = long # st.number_input("Longitude:")
            #     target_size = size # st.number_input("Size:")
            #     target_speed = speed # st.number_input("Speed:")
            #     target_start_mileage = runentry # st.number_input("Start Mileage:")
            #     target_end_mileage = runexit # st.number_input("End Mileage:")
            #     target_class = classtype # st.text_input("Vehicle Class:")
                
            #     if len(str(target_lat)) > 0 and len(str(target_lon)) > 0:                        
            #         if -90 <= float(target_lat) <= 90 and -180 <= float(target_lon) <= 180:
            #             target_record = {
            #                 'Latitude': target_lat,
            #                 'Longitude': target_lon,
            #                 'Size': target_size,
            #                 'ROP': target_speed,
            #                 'DepthIn': target_start_mileage,
            #                 'DepthOut': target_end_mileage,
            #             }

            #         if target_class:
            #             target_record['Size'] = target_class

            #         # Perform the matchmaking algorithm and display results
            #         st.write("Matchmaking Results:")
            #         matched_records = find_high_ranking_matches(uploaded_file, target_record, vehicle_class=target_class)
                    
            #         # matched_records = calculate_ranking_score(df_br, target_record, target_class)
                                        
            #         # st.dataframe(matched_records)
            #         return matched_records
                
            # def matchmain2(lat,long,size,speed,runentry,runexit,classtype):
            #     # uploaded_file = st.file_uploader("Choose a CSV file", type=["csv"])
            #     uploaded_file = 'data/records_USonly.csv'                    
            #     # if uploaded_file:
            #     #     df = pd.read_csv(uploaded_file)
                
            #     # df_br = pd.read_csv('data/records_USonly.csv')
            #     # st.write(df_br)
            #     # Create input fields for the target record
            #     st.write(f"Target Record Information: {lat}|{long}|{size}|{speed}|{runentry}|{runexit}|{classtype}")
            #     target_lat = float(lat) if lat is not None else 9999 # st.number_input("Latitude:")
            #     target_lon = float(long) if long is not None else 9999 # st.number_input("Longitude:")
            #     target_size = float(size) if size is not None else 0 # st.number_input("Size:")
            #     target_speed = float(speed) if speed is not None else 0 # st.number_input("Speed:")
            #     target_start_mileage = float(runentry) if runentry is not None else 0 # st.number_input("Start Mileage:")
            #     target_end_mileage = float(runexit) if runexit is not None else 0 # st.number_input("End Mileage:")
            #     target_class = classtype # st.text_input("Vehicle Class:")
                
            #     if len(str(target_lat)) > 0 and len(str(target_lon)) > 0:
            #         if -90 <= target_lat <= 90 and -180 <= target_lon <= 180:
            #             target_record = {
            #                 'Latitude': target_lat,
            #                 'Longitude': target_lon,
            #                 'Size': target_size,
            #                 'ROP': target_speed,
            #                 'DepthIn': target_start_mileage,
            #                 'DepthOut': target_end_mileage,
            #             }

            #         if target_class:
            #             target_record['Size'] = target_class

            #         # Perform the matchmaking algorithm and display results
            #         # st.write("Matchmaking Results:")
            #         # matched_records = find_high_ranking_matches(uploaded_file, target_record, vehicle_class=target_class)
                    
            #         matched_records = calculate_ranking_score2(df_br, target_record, target_class)
                                        
            #         # st.dataframe(matched_records)
            #         return matched_records
                    
            
            welldata = []        
            dfwell = pd.DataFrame(welldata, columns=['Run','Section','Size','Bit','Type','Price','Priceft','Pricetotal','DBR','Backup','Backuptype','Din','Dout','WOBin','WOBout','ROP','RPMin','RPMax','Motorspecs','Flowrate','Comment'])
            # sections = []
            # sheader = []
            for i in  range(0, sectioncount + 1): 
                dfwell.loc[i] = [0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0]                

            # Section fig loop container
            figsec = {}
            figbit = {}
            figbitrop = {}
            figbitftg = {}
            df_bitrop = {}
            df_bitftg = {}  
            offset_rad = {}
            offset_logic = {}
            vectored_records2 = {}   
            vectored_groups = {}  
            vectored_group_count = {}
            df_br_rig = {}  
            figrigdate = {}        
            figcompavg = {}             
            figcompoffset = {}
            
            # @st.cache_data(experimental_allow_widgets=True)
            def definesections(sindex,df_bits):                     
                section,defsize,defbit,defbittype,defbitbackup,defbackuptype,defsectiondepthin,defsectiondepthout,defsectionwob,wobin,wobout,defsectionrop,defsectionrpm,rpmin,rpmax,defsectionmotorspeed,defsectionflowrate,price,priceft,pricetotal,dbr,comment = [],[],[],[],[],[],[],[],[],[],[],[],[],[],[],[],[],[],[],[],[],[]
                # section,defsize,defbit,defbittype,defbitbackup,defbackuptype,defsectiondepthin,defsectiondepthout,defsectionwob,wobin,wobout,defsectionrop,defsectionrpm,rpmin,rpmax,defsectionmotorspeed,defsectionflowrate,price,priceft,pricetotal,dbr,comment = 0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0
                
                for i in range(1, sindex + 1): 
                    
                    # @st.cache_resource(experimental_allow_widgets=True)
                    # def sectioncard():
                    try:                            
                        if uploaded_bpm:
                            # # load defaults if file uploaded
                            # st.session_state[f'{i}size'] = bpm_table[1,i] # 1run
                            # # st.session_state[f'{i}size'] = bpm_table[2,i] # 2 size
                            # # st.session_state[f'{i}type'] = bpm_table[3,i] # 3 type
                            # st.session_state[f'{i}bit'] = bpm_table[4,i] # 4 part number                        
                            # if i <= 1: # depthin
                            #     st.session_state[f'{i}depthin'] = bpm_port[14]  # 5 out
                            # st.session_state[f'{i}depthout'] = bpm_table[5,i] # 5 out
                            # st.session_state[f'{i}rop'] = bpm_table[10,i] # 10 rop
                            # # st.session_state[f'{i}hrs'] = bpm_table[6,i] # 6 hrs
                            # # st.session_state[f'{i}ftg'] = bpm_table[9,i] # 9 ftg
                            # # Costs / Comments
                            # st.session_state[f'{i}price'] = bpm_table[12,i] # 12 rotary$
                            # st.session_state[f'{i}priceft'] = bpm_table[13,i] # 13 dolFt
                            # st.session_state[f'{i}dbr'] = bpm_table[14,i] # 14 dbr
                            # st.session_state[f'{i}totalprice'] = bpm_table[15,i] # 15 total cost
                            # st.session_state[f'{i}comment'] = bpm_table[16,i] # 16 comments
                            # # st.session_state[f'{i}pricemin'] = bpm_table[17,i] # 17 minimum$
                            
                            # # Parameters
                            # if len(str(bpm_table[7,i])) > 0:
                            #     st.session_state[f'{i}setparams'] = True
                            #     st.session_state[f'{i}wobrange'] = bpm_table[7,i].split('/', 1)[0]+','+bpm_table[7,i].split('/', 1)[-1] # 7 wob 10/35 bpm_port[10].split(',', 1)[0]                         
                            #     st.session_state[f'{i}rpmrange'] = bpm_table[8,i].split('/', 1)[0]+','+bpm_table[8,i].split('/', 1)[-1] # 8 rpm 160/200
                            #     st.session_state[f'{i}flowrate'] = bpm_table[11,i] # 11 gpm
                            
                            # load defaults if file uploaded
                            if st.session_state.bpm is True and st.session_state.bpmoverwrite is False:               
                                # if f'{i}size' not in st.session_state:
                                # st.session_state[f'{i}sectionpick'] = bpm_table.loc[i-1,'Run'] # 1run
                                st.session_state[f'{i}size'] =bpm_table.loc[i-1,'Size'] # 2 size
                                # st.session_state[f'{i}type'] = bbpm_table.loc[i,'Type'] # 3 type
                                if not pd.isnull(bpm_table.loc[i-1,'Part #']):
                                    st.session_state[f'{i}bit'] = [str(bpm_table.loc[i-1,'Part #'])] # 4 part number 
                                else:
                                    if i >= 1:
                                        if not pd.isnull(st.session_state[f'{i-1}bit']):  
                                            st.session_state[f'{i}bit'] = st.session_state[f'{i-1}bit']
                                        else:                                
                                            st.session_state[f'{i}bit'] = ''
                                    else:                                
                                        st.session_state[f'{i}bit'] = ''
                                
                                # if st.session_state[f'{i}bit'] == 'NaN':
                                #     st.session_state[f'{i}bit'] = st.session_state[f'{i-1}bit']                                    
                                #     if st.session_state[f'{i}bit'] == 'NaN':
                                #         st.session_state[f'{i}bit'] = ''
                                # if st.session_state[f'{i}bit'] is None:
                                #     st.session_state[f'{i}bit'] = ''
                                        
                                if i <= 1: # depthin
                                    st.session_state[f'{i}depthin'] = str(bpm_port[14])  # 5 out
                                st.session_state[f'{i}depthout'] = str(bpm_table.loc[i-1,'Out']) # 5 out
                                st.session_state[f'{i}rop'] = str(bpm_table.loc[i-1,'ROP']) # 10 rop
                                # st.session_state[f'{i}hrs'] = bpm_table.loc[i,'HRS'] # 6 hrs
                                # st.session_state[f'{i}ftg'] = bpm_table.loc[i,'FTG'] # 9 ftg
                                # Costs / Comments
                                if not pd.isnull(bpm_table.loc[i-1,'Rotary($)']):
                                    st.session_state[f'{i}price'] = str(bpm_table.loc[i-1,'Rotary($)']) # 12 rotary$
                                else:
                                    st.session_state[f'{i}price'] = str(bpm_table.loc[i-1,'Total Cost']) # 12 rotary$
                                
                                if not pd.isnull(bpm_table.loc[i-1,'Rotary($)']):
                                    st.session_state[f'{i}priceft'] = str(bpm_table.loc[i-1,'DolFt']) # 13 dolFt
                                else:
                                    st.session_state[f'{i}priceft'] = str(0) # 13 dolFt
                                    
                                st.session_state[f'{i}dbr'] = str(bpm_table.loc[i-1,'DBR($)']) # 14 dbr
                                # st.session_state[f'{i}totalprice'] = str(bpm_table.loc[i-1,'Total Cost']) # 15 total cost
                                st.session_state[f'{i}comment'] = str(bpm_table.loc[i-1,'Comments']) # 16 comments
                                # st.session_state[f'{i}pricemin'] = bpm_table.loc[i,'Minimum($)'] # 17 minimum$
                                
                                # Parameters
                                if len(str(bpm_table.loc[i-1,'WOB'])) > 1:
                                    st.session_state[f'{i}setparams'] = True
                                    st.session_state[f'{i}wobrange'] = (int(bpm_table.loc[i-1,'WOB'].split('/', 1)[0]),int(bpm_table.loc[i-1,'WOB'].split('/', 1)[-1])) # 7 wob 10/35 bpm_port[10].split(',', 1)[0]                         
                                    st.session_state[f'{i}rpmrange'] = (int(bpm_table.loc[i-1,'RPM'].split('/', 1)[0]),int(bpm_table.loc[i-1,'RPM'].split('/', 1)[-1])) # 8 rpm 160/200
                                    st.session_state[f'{i}flowrate'] = bpm_table.loc[i-1,'GPM'] # 11 gpm

                                if i == sindex:
                                    # end BPMoverwrite command
                                    st.session_state.bpmoverwrite = True
                                    
                    except Exception as e:
                        print(e)
                        st.info(f"error bpm_table: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                    try:
                        if st.session_state.csv is True and st.session_state.csvoverwrite is False:               
                            # if f'{i}size' not in st.session_state:
                            # st.write(csv_table)
                            
                            if not pd.isna(csv_table.loc[i-1,'Section']) and len(str(csv_table.loc[i-1,'Section'])):
                                st.session_state[f'{i}sectionpick'] = list(csv_table.loc[i-1,'Section'].split("/"))
                            if not pd.isna(csv_table.loc[i-1,'Size']) and len(str(csv_table.loc[i-1,'Size'])):
                                st.session_state[f'{i}size'] = csv_table.loc[i-1,'Size']
                            if not pd.isna(csv_table.loc[i-1,'Bit']) and len(str(csv_table.loc[i-1,'Bit'])):
                                st.session_state[f'{i}bit'] = [str(csv_table.loc[i-1,'Bit'])]
                                if not pd.isna(csv_table.loc[i-1,'Backup']) and len(str(csv_table.loc[i-1,'Backup'])):
                                    st.session_state[f'{i}bit'].append(str(csv_table.loc[i-1,'Backup']))                                        
                            if i <= 1: # depthin
                                if not pd.isna(csv_table.loc[i-1,'Din']) and len(str(csv_table.loc[i-1,'Din'])):
                                    st.session_state[f'{i}depthin'] = f"{csv_table.loc[i-1,'Din']}"
                            if not pd.isna(csv_table.loc[i-1,'Dout']) and len(str(csv_table.loc[i-1,'Dout'])):
                                st.session_state[f'{i}depthout'] = f"{csv_table.loc[i-1,'Dout']}"
                            if not pd.isna(csv_table.loc[i-1,'ROP']) and len(str(csv_table.loc[i-1,'ROP'])):
                                st.session_state[f'{i}rop'] = f"{csv_table.loc[i-1,'ROP']}"
                            if not pd.isna(csv_table.loc[i-1,'Price']) and len(str(csv_table.loc[i-1,'Price'])):
                                st.session_state[f'{i}price'] = f"{csv_table.loc[i-1,'Price']}"
                            if not pd.isna(csv_table.loc[i-1,'Priceft']) and len(str(csv_table.loc[i-1,'Priceft'])):
                                st.session_state[f'{i}priceft'] = f"{csv_table.loc[i-1,'Priceft']}"
                            if not pd.isna(csv_table.loc[i-1,'DBR']) and len(str(csv_table.loc[i-1,'DBR'])):
                                st.session_state[f'{i}dbr'] = str(csv_table.loc[i-1,'DBR'])
                            if not pd.isna(csv_table.loc[i-1,'Comment']) and len(str(csv_table.loc[i-1,'Comment'])):
                                st.session_state[f'{i}comment'] = str(csv_table.loc[i-1,'Comment'])
                                
                            if len(str(csv_table.loc[i-1,'WOBin'])) > 1:
                                if not pd.isna(csv_table.loc[i-1,'WOBin']):
                                    st.session_state[f'{i}wobrange'] = (int(csv_table.loc[i-1,'WOBin']),int(csv_table.loc[i-1,'WOBout']))
                                if not pd.isna(csv_table.loc[i-1,'RPMin']):
                                    st.session_state[f'{i}rpmrange'] = (int(csv_table.loc[i-1,'RPMin']),int(csv_table.loc[i-1,'RPMax']))
                                if not pd.isna(csv_table.loc[i-1,'Flowrate']):
                                    st.session_state[f'{i}flowrate'] = f"{csv_table.loc[i-1,'Flowrate']}"
                                
                                if not pd.isna(csv_table.loc[i-1,'Motorspecs']):
                                    csvmotorspecs = str(csv_table.loc[i-1,'Motorspecs'])
                                    csvmotor = float(csvmotorspecs.split(' ', 1)[0])
                                    st.session_state[f'{i}motor'] = csvmotor
                                    
                                    if len(csvmotorspecs.split(' ', 1)) > 2:
                                        csvmotorbend = float(csvmotorspecs.split(' ', 1)[3])
                                        st.session_state[f'{i}motorbend'] = csvmotorbend
                                                                        
                            if i == sindex:
                                # end csvoverwrite command
                                st.session_state.csvoverwrite = True
                            
                    except Exception as e:
                        print(e)
                        st.info(f"error csv_table: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                    # https://discuss.streamlit.io/t/increase-expanders-label-text-size/36227/4
                    # with st.expander(f'**Section {i}**',expanded=True):
                    with st.container():
                        sheader = st.write(f'Section {i}') 
                        # ,anchor=False) 
                        
                        try:
                            # if defbits is not None:                        
                            #     df_section = df.loc[df['PartNumber'].isin(defbits)]
                            #     df_section['Price'] = None
                            #     df_section['Backup'] = NoneIf using all scalar values, you must pass an index
                            defbit = None
                            defbittype = None
                            defbitbackup = ''
                            defbackuptype = ''
                            
                            scols = st.columns([3,2,3,3])
                            comps = st.container()
                            with scols[0]:
                                # with st.expander('',expanded=True):
                                # with st.container():
                                
                                # section = st.selectbox('Choose **Section**:',['Surface','Vertical','Drill Out','Intermediate','Curve','Lateral','Extended Lateral'],index=i,key=f'{i}section') 
                                try:
                                    # https://discuss.streamlit.io/t/multiselectbox-with-lists-as-options/46240/4
                                    # if sectioncount == 3:
                                    #     sectionpick = st.multiselect('Choose **Section** Name:',['Surface','Vertical','Drill Out','Intermediate','Curve','Lateral'],default=sectionlist3[i],key=f'{i}sectionpick3')   
                                    # elif sectioncount == 4:
                                    #     sectionpick = st.multiselect('Choose **Section** Name:',['Surface','Vertical','Drill Out','Intermediate','Curve','Lateral'],default=sectionlist4[i],key=f'{i}sectionpick4')                                                 
                                    # else:
                                    #     sectionpick = st.multiselect('Choose **Section** Name:',['Surface','Vertical','Drill Out','Intermediate','Curve','Lateral'],default=sectionlist[i],key=f'{i}sectionpick')  
                                                                                
                                    if uploaded_bpm:
                                        # sectionpick = st.multiselect('Choose **Section** Name:',['Surface','Vertical','Drill Out','Intermediate','Curve','Lateral'],key=f'{i}sectionpick',disabled=True,help='Single or combined. Select curve & lateral in order to create a Curve/Lateral. Incorrect order may create errors.')                                         
                                        # st.session_state[f'{i}sectionpick'] = ''
                                        section = ' '
                                    # elif st.session_state.csv is True:
                                        
                                    #     section = ' '
                                    else:
                                        sectionpick = st.multiselect('Choose **Section** Name:',['Surface','Vertical','Drill Out','Intermediate','Curve','Lateral',' '],key=f'{i}sectionpick',help='Single or combined. Select curve & lateral in order to create a Curve/Lateral. Incorrect order may create errors.')                                         
                                        
                                        section = '/'.join(str(e) for e in sectionpick)
                                        
                                    # st.write(section)
                                except Exception as e:
                                    print(e)
                                    # st.info(f"error sectionpicker: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                    
                            # if len(section) > 0:
                                colsec1, colsec2 = st.columns([1,2],gap="small")
                                with colsec1:
                                    sizelist_filtered=df_bits['BitSize'].sort_values(ascending=False).unique()
                                    try:
                                        if len(section) > 1:
                                            if section in sectionsize:                                            
                                                if uploaded_bpm:
                                                    # default_ix = sizelist_filtered.index(bpm_table.loc[i-1,'Size'])
                                                    default_ix = int(np.where(sizelist_filtered == bpm_table.loc[i-1,'Size'])[0].item())                
                                                elif st.session_state.csv is True:
                                                    default_ix = int(np.where(sizelist_filtered == csv_table.loc[i-1,'Size'])[0].item())                                                         
                                                else:      
                                                    default_ix = sectionsize_index[section]                                                         
                                                    if f'{i}size' not in st.session_state:
                                                        st.session_state[f'{i}size'] = sectionsize[section]
                                                        
                                                # defsize = st.selectbox(f'**{section}** Size:',options=sizelist_filtered,index=default_ix,key=f'{i}size',help='Type a number to quickfind your size.',)
                                            else:
                                                if uploaded_bpm:
                                                    # default_ix = sizelist_filtered.index(bpm_table.loc[i-1,'Size'])
                                                    try:
                                                        default_ix = int(np.where(sizelist_filtered == bpm_table.loc[i-1,'Size'])[0].item())
                                                    except Exception as e:
                                                        print(e)
                                                        default_ix = 34
                                                        st.info(f"error default_ix: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                                                    # st.write(int(np.where(sizelist_filtered == bpm_table.loc[i-1,'Size'])[0].item()))             
                                                elif st.session_state.csv is True:
                                                    default_ix = int(np.where(sizelist_filtered == csv_table.loc[i-1,'Size'])[0].item())   
                                                else:
                                                    # st.session_state[f'{i}size'] = 34
                                                    default_ix = 34
                                        else:
                                            if uploaded_bpm:
                                                # default_ix = sizelist_filtered.index(bpm_table.loc[i-1,'Size'])
                                                try:
                                                    default_ix = int(np.where(sizelist_filtered == bpm_table.loc[i-1,'Size'])[0].item())
                                                except Exception as e:
                                                    print(e)
                                                    default_ix = 34
                                                    st.info(f"error default_ix: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                                                # st.write(int(np.where(sizelist_filtered == bpm_table.loc[i-1,'Size'])[0].item()))       
                                                        
                                            elif st.session_state.csv is True:
                                                default_ix = int(np.where(sizelist_filtered == csv_table.loc[i-1,'Size'])[0].item())   
                                            else:
                                                # st.session_state[f'{i}size'] = 8.75
                                                default_ix = 34
                                        
                                        # # # If the session state has a value for "number_input", use that value, and *don't* set a default value.
                                        # # if st.session_state[f'{i}size'] >= 1:                                            
                                        # if uploaded_bpm or st.session_state.csv is True:
                                        #     defsize = st.selectbox(f'**{section}** Size:',options=sizelist_filtered,key=f'{i}size',help='Type a number to quickfind your size.',)
                                        # # Otherwise, set a default value.
                                        # else:
                                        #     defsize = st.selectbox(f'**{section}** Size:',options=sizelist_filtered,index=default_ix,key=f'{i}size',help='Type a number to quickfind your size.',) 
                                                                                    
                                        defsize = st.selectbox(f'**{section}** Size:',options=sizelist_filtered,key=f'{i}size',help='Type a number to quickfind your size.',) 
                                                                                                                                
                                                                                                                                
                                    except Exception as e:
                                        print(e)
                                        st.info(f"error sectionlist: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                
                                with colsec2:  
                                    try:                                                                                
                                        if uploaded_bpm or st.session_state.csv is True:
                                            bitoptions = df['PartNumber'].sort_values(ascending=False)
                                            # localbit = st.multiselect('Select **Bit** (Primary, Backup):',options=df['PartNumber'].sort_values(ascending=False), key=f'{i}bit', help='First selection is Primary, second is backup.')
                                        else:
                                            bitoptions = df['PartNumber'].sort_values(ascending=False).loc[df['BitSize'] == defsize]
                                            # localbit = st.multiselect('Select **Bit** (Primary, Backup):',options=df['PartNumber'].sort_values(ascending=False).loc[df['BitSize'] == defsize],max_selections=2, key=f'{i}bit', help='First selection is Primary, second is backup.')
                                        
                                        localbit = st.multiselect('Select **Bit** (Primary, Backup):', options=bitoptions, max_selections=2, key=f'{i}bit', help='First selection is Primary, second is backup.')
                                
                                    except Exception as e:
                                        print(e)
                                        # localbit = st.session_state[f'{i}bit']
                                        st.write(f"error localbit: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                        # st.write(localbit)
                                # defbitbackup = st.multiselect('Optional **Backup Bit**:',df['PartNumber'].sort_values(ascending=False).loc[df['BitSize'] == defsize],max_selections=1, key=f'{i}bitbackup')                                        
                                defbit = None
                                defbittype = None
                                defbitbackup = ''
                                defbackuptype = ''
                                for count, bit in enumerate(localbit): 
                                    if count == 0:
                                        defbit = bit
                                        defbittype = f"{df.loc[df[df['PartNumber'] == defbit].index, 'BitType'].values[0]}"
                                        defbitbackup = ''
                                        defbackuptype = ''
                                        
                                    elif count == 1:
                                        try:
                                            defbitbackup = f'{bit}' 
                                            defbackuptype = f"{df.loc[df[df['PartNumber'] == defbitbackup].index, 'BitType'].values[0]}"
                                        except Exception as e:
                                            print(e)
                                            # st.info(f"error backup: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                     
                                                        
                                coldep1, coldep2 = st.columns(2,gap='small')                                        
                                with coldep1:
                                    if i <= 1:
                                        depthdisabled = False 
                                        # defsectiondepthin = st.text_input('Depth In:', placeholder='ft or m', key=f'{i}depthin') 
                                        defsectiondepthin = st.text_input('Depth In:',placeholder='ft or m',disabled=False, key=f'{i}depthin',help='Required')
                                        if defsectiondepthin.isnumeric():
                                            defsectiondepthin = float(defsectiondepthin)
                                    else:
                                        if defsectiondepthin is not None:
                                            # prevdout = dfwell.loc[int(i-1), 'Dout']
                                            if uploaded_bpm or st.session_state.csv is True:
                                                defsectiondepthin = st.text_input('Depth In:',value=st.session_state[f'{i-1}depthout'],disabled=True, key=f'{i}depthin',)   
                                            else:
                                                defsectiondepthin = st.text_input('Depth In:',value=dfwell.loc[int(i-1), 'Dout'],disabled=True, key=f'{i}depthin',)   
                                            if defsectiondepthin.isnumeric():
                                                defsectiondepthin = float(defsectiondepthin)                                         
                                with coldep2:
                                    defsectiondepthout = st.text_input('Depth Out:', placeholder='ft or m', key=f'{i}depthout',help='Required')
                                    if defsectiondepthout.isnumeric():
                                        defsectiondepthout = float(defsectiondepthout)                                                        
                                
                                if uploaded_bpm or st.session_state.csv is True:
                                    defsectionrop = st.text_input('ROP:',placeholder='ft/hr', key=f'{i}rop')
                                elif section in sectionrop:
                                    defsectionrop = st.text_input('ROP:',value=f'{sectionrop[section]}',placeholder='ft/hr', key=f'{i}rop')  
                                else:
                                    defsectionrop = st.text_input('ROP:',value=0,placeholder='ft/hr', key=f'{i}rop') 
                                        
                                if defsectionrop.isnumeric():
                                        defsectionrop = float(defsectionrop)
                                
                        # if len(section) > 0:
                            with scols[1]:
                                colimg1, colimg2 = st.columns(2,gap="small")
                                with colimg1:
                                    if len(str(defsectiondepthin)) > 0 and len(str(defsectiondepthout)) > 0:
                                        if float(defsectiondepthout) > 0:
                                            figsec[i] = go.Figure()
                                            figsec[i].add_trace(go.Bar(x=[section],y=[defsectiondepthin],name='',marker_color='rgba(0,0,0,0)',text=f"{defsectiondepthin}",textposition='inside',insidetextanchor='end',hoverinfo = 'skip',)) # textfont=dict(color="blue"),
                                            figsec[i].add_trace(go.Bar(x=[section],y=[defsectiondepthout],name='Section',text=f'{float(defsectiondepthout) - float(defsectiondepthin)} ft',textposition='inside',insidetextanchor='middle',marker_color='rgba(138,207,221,1)',hoverinfo='text+x+y',))
                                            figsec[i].add_trace(go.Bar(x=[section],y=[100],name='',marker_color='rgba(0,0,0,0)',text=f"{defsectiondepthout}",textposition='inside',insidetextanchor='start',hoverinfo = 'skip',))
                                            # figsec[i].add_annotation(x=[section], y=[defsectiondepthin],text=f'{defsectiondepthin}',showarrow=False,)
                                            # figsec[i].add_annotation(x=[section], y=[defsectiondepthout],text=f'{defsectiondepthout}',showarrow=False,yshift=10)
                                            figsec[i].update_layout(xaxis=dict(visible=False))
                                            figsec[i].update_layout(yaxis=dict(title='Depth',autorange='reversed',))
                                            figsec[i].update_yaxes(ticklabelposition="outside top")      
                                            figsec[i].update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor='rgba(0,0,0,0)',margin=go.layout.Margin(l=0,r=0,b=0,t=45,))
                                            figsec[i].update_layout(uniformtext_minsize=8, uniformtext_mode='show')  # 'hide' to hide overflow text
                                            figsec[i].update_layout(title={'text': f'{defsize} {section}','y':1,'x':0.5,'xanchor': 'center','yanchor': 'top'})
                                            figsec[i].update_layout(height=340,barmode='stack', showlegend = False) # bargap=(1/(sectioncount)), bargap=0.7,
                                            st.plotly_chart(figsec[i], use_container_width=True)
                                        else:
                                            st.write(float(defsectiondepthout))
                                    else:
                                        st.write(f"{len(str(defsectiondepthin))} {len(str(defsectiondepthout))}")
                                with colimg2:
                                    for count, bit in enumerate(localbit):                                                 
                                        url = f'https://bithub.ulterra.com/adira_files_fetch/{bit}' 
                                        headers = {'Adira_API_Key': st.secrets['Adira_API_Key']}                                    
                                        zipper = get_bit_files(url,headers)     
                                        # zipper.extractall('.')                           
                                        # st.write(zipper.infolist())
                                        # "<ZipInfo filename='U04770 - 12.250 - XP716 - BCDEU.pdf' compress_type=deflate filemode='-rwxrwxrwx' file_size=324291 compress_size=264884>"                                                
                                        # 0:"<ZipInfo filename='U04775 - 12.250 - XP716 - BCDEU.pdf' compress_type=deflate filemode='-rwxrwxrwx' file_size=311693 compress_size=251144>"
                                        # 1:"<ZipInfo filename='U04775-1.jpg' compress_type=deflate filemode='-rwxrwxrwx' file_size=85793 compress_size=83831>"
                                        # 2:"<ZipInfo filename='U04775-2.jpg' compress_type=deflate filemode='-rwxrwxrwx' file_size=38449 compress_size=37536>"
                                        
                                        if count == 0 and bit is not None:   
                                            st.write(f'Primary: {bit}')                                                  
                                            
                                            # st.image('https://static.streamlit.io/examples/cat.jpg',use_column_width='auto')                                                    
                                            if os.path.exists(f'{bit}-1.jpg'):
                                                image = f'{bit}-1.jpg'                                                         
                                            elif os.path.exists(f'data/Bit Pictures/{bit}-1.jpg'):
                                                image = f'data/Bit Pictures/{bit}-1.jpg' 
                                            elif os.path.exists(f'data/Bit Pictures/{bit}-1.JPG'):
                                                image = f'data/Bit Pictures/{bit}-1.JPG'
                                            else:
                                                image = f'data/Bit Pictures/blank.jpg'
                                            # bitimagewidth = 150 if len(localbit) > 0 else 300
                                            
                                            imagepng = f'{bit}-1.png'
                                            convert_png_transparent(image, imagepng)
                                            
                                            st.image(imagepng,caption=f"{df.loc[df[df['PartNumber'] == bit].index, 'BitSize'].values[0]}  - {df.loc[df[df['PartNumber'] == bit].index, 'BitType'].values[0]}",width=150)
                                                                                        
                                            # use_column_width='auto'
                                        elif count == 1 and bit is not None: 
                                            st.write(f'Backup: {bit}')
                                                                
                                            if os.path.exists(f'{bit}-1.jpg'):
                                                image = f'{bit}-1.jpg'       
                                            elif os.path.exists(f'data/Bit Pictures/{bit}-1.jpg'):                                
                                                image = f'data/Bit Pictures/{bit}-1.jpg' 
                                            elif os.path.exists(f'data/Bit Pictures/{bit}-1.JPG'):                                
                                                image = f'data/Bit Pictures/{bit}-1.JPG'
                                            else:
                                                image = f'data/Bit Pictures/blank.jpg'     
                                                
                                            imagepng = f'{bit}-1.png'
                                            convert_png_transparent(image, imagepng)
                                            
                                            st.image(imagepng,caption=f"{df.loc[df[df['PartNumber'] == bit].index, 'BitSize'].values[0]}  - {df.loc[df[df['PartNumber'] == bit].index, 'BitType'].values[0]}",width=150)
                            
                            # if localbit:
                            with scols[2]:
                                colbit1, colbit2, colbit3 = st.columns([2,1,2],gap="small")
                                with colbit1:
                                    # price = st.text_input('Price:', placeholder='$', key=f'{i}price')
                                                                                
                                    if uploaded_bpm or st.session_state.csv is True:
                                        price = st.text_input('Price:', key=f'{i}price', help='Enter the flat bit price or the $/ft minimum price to be calculated into the total.', disabled=False, label_visibility="visible")
                                    else:
                                        price = st.text_input('Price:', 0, placeholder='price or base cost', key=f'{i}price', help='Enter the flat bit price or the $/ft minimum price to be calculated into the total.', disabled=False, label_visibility="visible")
                                    
                                    price.replace('$', '')
                                    if price.isnumeric():
                                        price = float(price)
                                with colbit2:
                                    # priceft = st.text_input('Price:', placeholder='$', key=f'{i}priceft')
                                                                
                                    if uploaded_bpm or st.session_state.csv is True:
                                        priceft = st.text_input('$/ft:', key=f'{i}priceft', help='Enter the flat bit price or the $/ft minimum price to be calculated into the total.', disabled=False, label_visibility="visible")
                                    else:
                                        priceft = st.text_input('$/ft:', 0, placeholder='$/ft rate', key=f'{i}priceft', help='Enter the flat bit price or the $/ft minimum price to be calculated into the total.', disabled=False, label_visibility="visible")
                                    priceft.replace('$', '')
                                    if priceft.isnumeric():
                                        priceft = float(priceft)
                                    # priceft = st.radio("Choose price method",options=['Flat', '$/ft'],horizontal=True,key=f'{i}pricetype',)                                            
                                with colbit3:
                                    st.write("##")
                                    try:
                                        if price is not None and defsectiondepthout is not None and defsectiondepthin is not None:
                                            pricetotal = float(price) + (float(priceft) * (float(defsectiondepthout)-float(defsectiondepthin)))
                                            # st.write(f'Total: ${int(price) + (int(priceft) * (int(defsectiondepthout)-int(defsectiondepthin)))}')
                                        elif price is not None:
                                            pricetotal = float(price)
                                        else:
                                            pricetotal = 0
                                        st.write(f'Total: ${pricetotal}')                                                   
                                    except Exception as e:
                                        print(e)
                                        # st.info(f"error sectionpicker: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                        st.write(f'Set Depth for Totals.')                                                    
                                # if len(defbit) > 1:
                                #     backup = True
                                # else:
                                #     backup = False                         
                                dbr = st.text_input('DBR/LIH:', placeholder='Negotiable or $', key=f'{i}dbr')       
                                comment = st.text_area('Comments:', placeholder='Additional details', key=f'{i}comment')
                                
                            # with scols[3]:
                            #     if  count == 1 and bit is not None:
                            #         defbitbackup = bit                                            
                            #         if showpartnumber == True:
                            #             st.caption(f'Backup: {bit}')   
                            #         else:                                                
                            #             st.caption(f'Backup:')                                                                     
                            #         if os.path.exists(f'data/Bit Pictures/{bit}-1.jpg'):                                
                            #             image = f'data/Bit Pictures/{bit}-1.jpg' 
                            #             # st.image(f'data/Bit Pictures/{defbitbackup}-1.jpg',width=250)
                            #         elif os.path.exists(f'data/Bit Pictures/{bit}-1.JPG'):                                
                            #             image = f'data/Bit Pictures/{bit}-1.JPG'
                            #             # st.image(f'data/Bit Pictures/{defbitbackup}-1.JPG',width=250)
                            #         else:
                            #             image = f'data/Bit Pictures/blank.jpg'     
                            #             # st.image(f'data/Bit Pictures/blank.jpg',width=250)  
                                                                                                    
                            #         st.image(image,caption=f"{df.loc[df[df['PartNumber'] == bit].index, 'BitSize'].values[0]}  - {df.loc[df[df['PartNumber'] == bit].index, 'BitType'].values[0]}",width=250)
                                    
                            with scols[3]:
                                # with st.expander('Parameters (optional)',expanded=True):
                                setparams = st.checkbox('Parameters', value=False, key=f'{i}setparams', help='Enable to give parameter recommendations',)
                                if setparams:
                                # with st.container():
                                    # defsectionwob = st.text_input('WOB:', placeholder='lbs', key=f'{i}wob')
                                                
                                    if uploaded_bpm or st.session_state.csv is True:
                                        wobin, wobout = st.slider('Weight Range (klb):',min_value=5,max_value=140,step=1,format='%d klb',key=f'{i}wobrange')
                                    else:
                                        wobin, wobout = st.slider('Weight Range (klb):',min_value=5,max_value=140, value=(30, 55),step=1,format='%d klb',key=f'{i}wobrange')
                                    
                                    # if wobin.isnumeric():
                                    #     wobin = int(wobin)
                                    # if wobout.isnumeric():
                                    #         wobout = int(wobout)
                                    # defsectionrpm = st.text_input('RPM:', placeholder='RPM', key=f'{i}rpm')
                                    
                                    if uploaded_bpm or st.session_state.csv is True:
                                        rpmin, rpmax = st.slider('RPM:',min_value=80,max_value=400,step=1,key=f'{i}rpmrange')
                                    else:
                                        rpmin, rpmax = st.slider('RPM:',min_value=80,max_value=400, value=(160, 180),step=1,key=f'{i}rpmrange')
                                    # if rpmin.isnumeric():
                                    #     rpmin = int(rpmin)
                                    # if rpmax.isnumeric():
                                    #         rpmax = int(rpmax)
                                    
                                    defsectionflowrate = st.text_input('GPM:', placeholder='GPM', key=f'{i}flowrate')
                                    defsectionflowrate.replace('GPM', '')
                                    if defsectionflowrate.isnumeric():
                                        defsectionflowrate = float(defsectionflowrate)
                                # with st.expander('Motor Specs (optional)',expanded=True):     
                                # with st.container():                                   
                                    colmtr1, colmtr2 = st.columns([2,1],gap="small")
                                    with colmtr1:
                                        # defsectionmotorspeed = st.text_input('Motor Speed:', placeholder='Revs/Gal', key=f'{i}motor')
                                        # defsectionmotorspeed = st.selectbox('Choose **Motor**:',pdmotors,key=f'{i}motor') 
                                        # if defsectionmotorspeed == 'Other':
                                        #     defsectionmotorspeed = st.text_input('Motor Speed:', placeholder='Revs/min', key=f'{i}motorspeed')                                                    
                                        
                                        if uploaded_bpm or st.session_state.csv is True:
                                            defsectionmotorspeed = st.slider('Motor Rev/Gal:',min_value=0.00,max_value=1.20, step=0.01,key=f'{i}motor')
                                        else:
                                            defsectionmotorspeed = st.slider('Motor Rev/Gal:',min_value=0.00,max_value=1.20, value=0.80, step=0.01,key=f'{i}motor')
                                    with colmtr2:
                                        defsectionmotorbend = st.text_input('Motor Bend(°):', placeholder='1.7°', key=f'{i}motorbend')
                                        defsectionmotorbend.replace('°', '')
                                        defsectionmotorbend.replace('deg', '')
                                        defsectionmotorbend.replace('d', '')
                                        if defsectionmotorbend.isnumeric():
                                            defsectionmotorbend = float(defsectionmotorbend)
                                        
                                    if len(str(defsectionmotorbend)) > 0:
                                        defsectionmotorspeed = f'{defsectionmotorspeed} Rev/Gal, {defsectionmotorbend}° Bend'
                                    else:
                                        defsectionmotorspeed = f'{defsectionmotorspeed} Rev/Gal'
                                else:
                                    wobin, wobout = 0,0
                                    rpmin, rpmax = 0,0
                                    defsectionflowrate = 0
                                    defsectionmotorspeed = 0
                            # st.write(f'{defbitbackup} - {defbackuptype}')
                            dfwell.loc[i] = [i,section,defsize,defbit,defbittype,price,priceft,pricetotal,dbr,defbitbackup,defbackuptype,defsectiondepthin,defsectiondepthout,wobin,wobout,defsectionrop,rpmin,rpmax,defsectionmotorspeed,defsectionflowrate,comment]
                            
                        
                            # compsbtn = st.button('Find Offsets', key=f'{i}getcomps')
                            # if compsbtn:
                            if smartoffset:                                 
                                offsetstitle = 'Offsets'   
                                # with comps.expander('Offsets', expanded=False):
                                if lat and long and defsize and defsectionrop and defsectiondepthin and defsectiondepthout:
                                    with st.spinner():
                                        try:
                                            target_lat = lat # st.number_input("Latitude:")
                                            target_lon = long # st.number_input("Longitude:")
                                            target_size = defsize # st.number_input("Size:")
                                            target_speed = defsectionrop # st.number_input("Speed:")
                                            target_start_mileage = defsectiondepthin # st.number_input("Start Mileage:")
                                            target_end_mileage = defsectiondepthout # st.number_input("End Mileage:")
                                            target_class = defsize # st.text_input("Vehicle Class:")
                                            try:
                                                recordsexist = df_br.groupby('ItemNo').size()[defbit]
                                            except Exception as e:
                                                recordsexist = 0
                                                # comps.info(f"error recordsexist: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                                
                                            if len(str(target_lat)) > 0 and len(str(target_lon)) > 0: # and recordsexist > 0:
                                                if -90 <= float(target_lat) <= 90 and -180 <= float(target_lon) <= 180:
                                                    target_record = {
                                                        'Latitude': target_lat,
                                                        'Longitude': target_lon,
                                                        'Size': target_size,
                                                        'ROP': target_speed,
                                                        'DepthIn': target_start_mileage,
                                                        'DepthOut': target_end_mileage,
                                                    }

                                                if target_class:
                                                    target_record['Size'] = defsize
                                                
                                                # df_br = get_bit_records()
                                                
                                                # vectored_records = vector_ranking_score(df_br, target_record)
                                                # vectored_records['Drilled'] = (vectored_records['DepthOut'] - vectored_records['DepthIn'])
                                                # # comps.write(f'Matchmaking Vectors: {vectored_records}' ) 
                                                # # comps.write(f'Matchmaking2 Vectors: ' ) 
                                                # # comps.write(vectored_records)                                                         
                                                # # comps.write(vectored_records[0:5])
                                                # # comps.write(vectored_records.sort_values(by=['total_score'], ascending=False))
                                                
                                                # vectored_records2[i]=vectored_records.sort_values(by=['total_score'], ascending=False,ignore_index=True)   
                                                vectored_records2[i], offset_rad[i], offset_logic[i] = vector_ranking_score(df_br, target_record)
                                                
                                                vectored_groups[i] = vectored_records2[i].groupby('BitMfgr', as_index=False).mean()
                                                # vectored_group_count[i] = vectored_records2[i].groupby('BitMfgr', as_index=False).agg(['mean','count'])
                                                
                                                
                                                offsetstitle = f"**C-rank Offsets**: ({len(vectored_records2[i])}) found, {len(vectored_groups[i]['BitMfgr'].unique())} Bit Vendors"
                                                
                                                        
                                            else:
                                                vectored_records2[i] = None
                                                comps.write(f'Matchmaking2 missing lat/long from user.' )
                                                    
                                        except Exception as e:
                                            print(e)
                                            comps.info(f"error Vectormain2: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                
                                with comps.expander(offsetstitle, expanded=False):                                        
                                    # colcompsopt1, colcompsopt2 = st.columns(2,gap="small")
                                    # with colcompsopt1:
                                    # pass exclusion values to offset maths
                                    # offsetexclusions = st.multiselect('Exclude:',['No TD','Other Rigs','Runs 90+ days Old'],key=f'{i}offsetexclusions', help='Exclusion options for C-Ranking')
                                    
                                    if lat and long and defsize and defsectionrop and defsectiondepthin and defsectiondepthout:
                                        if vectored_records2[i] is not None:
                                            colcomps1, colcomps2, colcomps3 = st.columns(3,gap="small")
                                            
                                            color_discrete_map  = {'SLB':'#0012d5','NOV':'#E51636','BH':'#00B040','HAL':'#ED171F','ULT':'#05929f','TRX':'#4B4848','DF':'#dc4d01','VAR':'#4A3B6A','ATL':'#999999 ',
                                                                'CAP':'#e94024','CHR':'#999999 ','DA':'#999999 ','DB':'#999999 ','DF':'#999999 ','DIA':'#999999 ','DP':'#999999 ','DRM':'#999999 ','ENC':'#999999 ','HAC':'#999999 ',
                                                                'JZ':'#999999 ','KD':'#999999 ','KLR':'#999999 ','MSN':'#999999 ','OTH':'#999999 ','OTS':'#999999 ','REV':'#999999 ','SHR':'#999999 ','TCL':'#999999 ','UNK':'#999999 ','VIP':'#999999 ','WEA':'#999999 ',
                                                                }
                                            color_gray_map  = {'SLB':'#999999','NOV':'#999999','BH':'#999999','HAL':'#999999','ULT':'#05929f','TRX':'#999999','DF':'#999999','VAR':'#999999','ATL':'#999999 ',
                                                                'CAP':'#999999','CHR':'#999999 ','DA':'#999999 ','DB':'#999999 ','DF':'#999999 ','DIA':'#999999 ','DP':'#999999 ','DRM':'#999999 ','ENC':'#999999 ','HAC':'#999999 ',
                                                                'JZ':'#999999 ','KD':'#999999 ','KLR':'#999999 ','MSN':'#999999 ','OTH':'#999999 ','OTS':'#999999 ','REV':'#999999 ','SHR':'#999999 ','TCL':'#999999 ','UNK':'#999999 ','VIP':'#999999 ','WEA':'#999999 ',
                                                                }
                                            
                                            
                                            # Rig History
                                            # rigname = st.selectbox("Select Rig Name:",df_br['Rig_Name_Full'].sort_values(ascending=False).unique())
                                            if len(str(rigname)) > 1:
                                                try:
                                                    # Get Rig runs
                                                    df_br_rig[i] = df_br[df_br['Rig_Name_Full']==rigname]
                                                    # st.write(df_br_rig[i])
                                                    df_br_rig[i] = df_br_rig[i][df_br_rig[i]['Size']==defsize]
                                                    if df_br_rig[i].shape[1] > 0:    
                                                        df_br_rig[i] = df_br_rig[i].sort_values(by=['RunDate'], ascending=False,ignore_index=True)                                              
                                                        df_br_rig[i] = df_br_rig[i][0:5]
                                                    else:
                                                        df_br_rig[i] = df_br[df_br['OperatorName']== customer]
                                                        df_br_rig[i] = df_br_rig[i][df_br_rig[i]['Size']==defsize]
                                                        df_br_rig[i] = df_br_rig[i].sort_values(by=['RunDate'], ascending=False,ignore_index=True)                                              
                                                        df_br_rig[i] = df_br_rig[i][0:5]
                                                        
                                                    
                                                    figrigdate[i] = go.Figure()
                                                    figrigdate[i].add_trace(go.Bar(x=df_br_rig[i].index,y=df_br_rig[i]['DepthIn'],name='',marker_color='rgba(0,0,0,0)',))                                            
                                                    # figrigdate[i].add_trace(go.Bar(x=df_br_rig[i].index,y=df_br_rig[i]['Drilled'],name='FTG',text=df_br_rig[i]['Drilled'],textposition='inside',insidetextanchor='middle',marker_color='rgba(0,147,159,1)', hoverinfo='text+x+y',))
                                                    
                                                    for r in df_br_rig[i]['BitMfgr'].unique():
                                                        dfoffa = df_br_rig[i][df_br_rig[i]['BitMfgr']==r]
                                                        if graycomps == True:
                                                            figrigdate[i].add_traces(go.Bar(x=dfoffa.index, y = dfoffa['Drilled'], name=r,marker_color=color_gray_map[r],text=dfoffa['Drilled'],textposition='inside',insidetextanchor='middle',))
                                                        else:
                                                            figrigdate[i].add_traces(go.Bar(x=dfoffa.index, y = dfoffa['Drilled'], name=r,marker_color=color_discrete_map[r],text=dfoffa['Drilled'],textposition='inside',insidetextanchor='middle',))
                                                        
                                                    figrigdate[i].add_trace(go.Scatter(x=df_br_rig[i].index,y=df_br_rig[i]['ROP'],name='ROP',text=df_br_rig[i]['ROP'],textposition='bottom center',textfont=dict(color='rgba(152,194,31,1)'),mode='lines+markers+text',marker_color='rgba(152,194,31,1)',hoverinfo='text',yaxis = 'y2',))
                                                    figrigdate[i].update_layout(yaxis=dict(title='Depth',titlefont=dict(color='black'),tickfont=dict(color='black'),autorange='reversed',rangemode='normal'))
                                                    figrigdate[i].update_layout(yaxis2=dict(title='ROP',titlefont=dict(color='rgba(152,194,31,1)'),tickfont=dict(color='rgba(152,194,31,1)'), autorange=True,anchor='free',overlaying='y',side='right',showgrid=False,position=1.0)) 
                                                    figrigdate[i].update_yaxes(ticklabelposition="outside top")      
                                                    figrigdate[i].update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor='rgba(0,0,0,0)')
                                                    figrigdate[i].update_layout(xaxis=dict(side='top',visible=False)) 
                                                    figrigdate[i].update_layout(uniformtext_minsize=8, uniformtext_mode='show')  # 'hide' to hide overflow text
                                                    figrigdate[i].update_layout(title={'text': f"Last 5 {defsize} on {rigname}",'y':0.9,'x':0.5,'xanchor': 'center','yanchor': 'top'})
                                                    figrigdate[i].update_layout(barmode='stack', bargap=0.7, showlegend = False) # bargap=(1/(sectioncount))
                                                    with colcomps1:
                                                        st.plotly_chart(figrigdate[i], use_container_width=True)                                                             
                                                        # st.write(df_br_rig[i]) 
                                                    
                                                except Exception as e:
                                                    print(e)
                                                    st.info(f"error figrigdate: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                            else:
                                                st.write('No Rig selected in header.')
                                                
                                            # Offsets by rank
                                            if len(str(defsectiondepthin)) > 0 and len(str(defsectiondepthout)) > 0:
                                                if float(defsectiondepthout) > 0:
                                                    figcompoffset[i] = go.Figure()
                                                    figcompoffset[i].add_trace(go.Bar(x=vectored_records2[i].index,y=vectored_records2[i]['DepthIn'],name='',marker_color='rgba(0,0,0,0)',))                                            
                                                    # figcompoffset[i].add_trace(go.Bar(x=vectored_records2[i].index,y=vectored_records2[i]['Drilled'],name='FTG',text=vectored_records2[i]['Drilled'],textposition='inside',insidetextanchor='middle',marker_color='rgba(0,147,159,1)', hoverinfo='text+x+y',))
                                                    
                                                    for r in vectored_records2[i]['BitMfgr'].unique():
                                                        dfoffb = vectored_records2[i][vectored_records2[i]['BitMfgr']==r]
                                                        if graycomps == True:
                                                            figcompoffset[i].add_traces(go.Bar(x=dfoffb.index, y = dfoffb['Drilled'], name=r,marker_color=color_gray_map[r],text=dfoffb['Drilled'],textposition='inside',insidetextanchor='middle',))
                                                        else:
                                                            # figcompoffset[i].add_traces(go.Bar(x=dfoffb.index, y = dfoffb['Drilled'], name=r,marker_color=color_discrete_map[r],text=vectored_records2[i]['Drilled'],textposition='inside',insidetextanchor='middle',))
                                                            figcompoffset[i].add_traces(go.Bar(x=dfoffb.index, y = dfoffb['Drilled'], name=r,marker_color=color_discrete_map[r],text=dfoffb['Drilled'],textposition='inside',insidetextanchor='middle',))
                                                        
                                                    figcompoffset[i].add_trace(go.Scatter(x=vectored_records2[i].index,y=vectored_records2[i]['ROP'],name='ROP',text=round(vectored_records2[i]['ROP'],2),textposition='bottom center',textfont=dict(color='rgba(152,194,31,1)'),mode='markers+text',marker_color='rgba(152,194,31,1)',hoverinfo='text',yaxis = 'y2',))
                                                                    
                                                    figcompoffset[i].update_layout(yaxis=dict(title='Depth',titlefont=dict(color='black'),tickfont=dict(color='black'),autorange='reversed',))
                                                    figcompoffset[i].update_layout(yaxis2=dict(title='ROP',titlefont=dict(color='rgba(152,194,31,1)'),tickfont=dict(color='rgba(152,194,31,1)'), autorange='reversed',anchor='free',overlaying='y',side='right',showgrid=False,position=1.0)) 
                                                    
                                                    figcompoffset[i].update_yaxes(ticklabelposition="outside top")      
                                                    figcompoffset[i].update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor='rgba(0,0,0,0)')
                                                    figcompoffset[i].update_layout(xaxis=dict(side='top',visible=False)) 
                                                    figcompoffset[i].update_layout(uniformtext_minsize=8, uniformtext_mode='show')  # 'hide' to hide overflow text
                                                    figcompoffset[i].update_layout(title={'text': f"C-Ranked {offset_rad[i]}-mile Offsets ({len(vectored_records2[i])})",'y':0.9,'x':0.5,'xanchor': 'center','yanchor': 'top'})
                                                    figcompoffset[i].update_layout(barmode='stack', bargap=0.6, showlegend = False) # bargap=(1/(sectioncount))
                                                    figcompoffset[i].add_annotation(showarrow=False,text=offset_logic[i],font=dict(size=10),xref='x domain',x=0.5,yref='y domain',y=0)
                                                    with colcomps2:
                                                        st.plotly_chart(figcompoffset[i], use_container_width=True)
                                                        
                                            # Comp Averages
                                            if len(str(defsectiondepthin)) > 0 and len(str(defsectiondepthout)) > 0:
                                                if float(defsectiondepthout) > 0:
                                                    # st.write(vectored_groups[i])
                                                    figcompavg[i] = go.Figure()
                                                    figcompavg[i].add_trace(go.Bar(x=vectored_groups[i]['BitMfgr'],y=vectored_groups[i]['DepthIn'],name='',marker_color='rgba(0,0,0,0)',))
                                                    try:
                                                        # figcompavg[i].add_trace(go.Bar(x=vectored_groups[i]['BitMfgr'],y=vectored_groups[i].groupby('BitMfgr')['Drilled'] ,name='FTG',text=vectored_groups[i]['BitMfgr'],textposition='inside',insidetextanchor='middle',marker=dict(color=colors[vectored_groups[i]['BitMfgr']],colorscale=colors,showscale=False), hoverinfo='text+x+y',))                                                                 
                                                        for t in vectored_groups[i]['BitMfgr'].unique():
                                                            dfoffc = vectored_groups[i][vectored_groups[i]['BitMfgr']==t]
                                                            # st.write(dfp)
                                                            if graycomps == True:
                                                                figcompavg[i].add_traces(go.Bar(x=dfoffc['BitMfgr'], y=dfoffc['Drilled'], name=t, marker_color=color_gray_map[t], width=3*(len(dfoffc)/len(vectored_records2[i])), text=round(dfoffc['Drilled'],1), textposition='inside',insidetextanchor='middle',))
                                                                # st.write(f"widths {(vectored_records2[i][vectored_groups[i]['BitMfgr']==t].value_counts()['ULT']/len(vectored_records2[i]))}")
                                                            else:
                                                                figcompavg[i].add_traces(go.Bar(x=dfoffc['BitMfgr'], y=dfoffc['Drilled'], name=t, marker_color=color_discrete_map[t],  width=3*(len(dfoffc)/len(vectored_records2[i])), text=f"{round(dfoffc['Drilled'],1)} ({(len(dfoffc)/len(vectored_records2[i]))})", textposition='inside',insidetextanchor='middle',))
                                                            
                                                            # figcompavg[i].add_traces(go.Bar(x=dfoffc['BitMfgr'], y=dfoffc['Drilled'], name=t, marker_color=color_discrete_map[t], text=round(dfoffc['Drilled'],1), textposition='inside',insidetextanchor='middle',))
                                                    except Exception as e:
                                                        # figcompavg[i].add_trace(go.Bar(x=vectored_groups[i]['BitMfgr'],y=vectored_groups[i].groupby('BitMfgr')['Drilled'] ,name='FTG',text=vectored_groups[i]['BitMfgr'],textposition='inside',insidetextanchor='middle',marker_color='rgba(0,147,159,1)', hoverinfo='text+x+y',))
                                                        comps.info(f"error markercolor3: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                                        
                                                    figcompavg[i].add_trace(go.Scatter(x=vectored_groups[i]['BitMfgr'],y=vectored_groups[i]['ROP'] ,name='ROP',text=round(vectored_groups[i]['ROP'],1),textposition='bottom center',textfont=dict(color='rgba(0,0,0,1)'),mode='markers+text',marker_color='rgba(0,0,0,1)',hoverinfo='text',yaxis = 'y2',))
                                                    # figproposal.update_traces(texttemplate='%{text:,.1f}')                    
                                                    figcompavg[i].update_layout(yaxis=dict(title='Depth',titlefont=dict(color='black'),tickfont=dict(color='black'),autorange='reversed',))
                                                    figcompavg[i].update_layout(yaxis2=dict(title='ROP',titlefont=dict(color='rgba(0,0,0,1)'),tickfont=dict(color='rgba(0,0,0,1)'), autorange='reversed',anchor='free',overlaying='y',side='right',showgrid=False,position=1.0)) 
                                                    
                                                    figcompavg[i].update_yaxes(ticklabelposition="outside top")      
                                                    figcompavg[i].update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor='rgba(0,0,0,0)')
                                                    figcompavg[i].update_layout(xaxis=dict(side='top',visible=True)) 
                                                    figcompavg[i].update_layout(uniformtext_minsize=8, uniformtext_mode='show')  # 'hide' to hide overflow text
                                                    figcompavg[i].update_layout(title={'text': f"C-Ranked Competitor Averages",'y':0.9,'x':0.5,'xanchor': 'center','yanchor': 'top'})
                                                    figcompavg[i].update_layout(barmode='stack', bargap=0.7, showlegend = False) # bargap=(1/(sectioncount))
                                                    figcompavg[i].add_annotation(showarrow=False,text=offset_logic[i],font=dict(size=10),xref='x domain',x=0.5,yref='y domain',y=0)
                                                    try:
                                                        figcompavg[i].add_annotation(showarrow=False,text=f'{t} {3*(len(dfoffc)/len(vectored_records2[i]))} {len(dfoffc)} / {len(vectored_records2[i])}',font=dict(size=10),xref='x domain',x=0.5,yref='y domain',y=0)
                                                    except Exception as e:
                                                        print(e)
                                                        
                                                    with colcomps3:
                                                        st.plotly_chart(figcompavg[i], use_container_width=True)
                                                        
                                                # try:
                                                #     comps.write(f'Matchmaking3 Results: ' ) 
                                                #     comps.write(matchmain(lat, long, defsize,defsectionrop,defsectiondepthin,defsectiondepthout,defsize)) 
                                                # except Exception as e:
                                                #     comps.info(f"error matchmain3: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                                    
                                                # try:
                                                #     comps.write(f'Matchmaking4 Results: ' ) 
                                                #     comps.write(matchmain2(lat, long, defsize,defsectionrop,defsectiondepthin,defsectiondepthout,defsize)) 
                                                # except Exception as e:
                                                #     comps.info(f"error matchmain4: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")       
                                                
                                                    st.write(vectored_records2[i])                                                        
                                        else:
                                            st.empty()
                                    else:
                                        st.empty()    
                                                
                                
                            if bitstats:
                                try:
                                    # st.write(f"Total Runs: {len(df_br.loc[df_br['ItemNo'] == bit])} **ROP:** {(df_br.loc[df_br['ItemNo'] == bit, 'ROP'].mean()):,.2f} **FTG:** {(df_br.loc[df_br['ItemNo'] == bit, 'DepthOut'].mean() - df_br.loc[df_br['ItemNo'] == bit, 'DepthIn'].mean()):,.2f}")
                                    df_opbr = df_br.loc[df_br['OperatorName'] == customer.upper()]
                                    # st.write(f"Total Runs (**{customer}**): {len(df_br.loc[df_br['ItemNo'] == bit])} (**{len(df_opbr.loc[df_opbr['ItemNo'] == bit]):,.0f}**) ROP: {(df_br.loc[df_br['ItemNo'] == bit, 'ROP'].mean()):,.1f} (**{(df_opbr.loc[df_opbr['ItemNo'] == bit, 'ROP'].mean()):,.1f}**) FTG: {(df_br.loc[df_br['ItemNo'] == bit, 'DepthOut'].mean() - df_br.loc[df_br['ItemNo'] == bit, 'DepthIn'].mean()):,.0f} (**{(df_opbr.loc[df_opbr['ItemNo'] == bit, 'DepthOut'].mean() - df_opbr.loc[df_opbr['ItemNo'] == bit, 'DepthIn'].mean()):,.0f}**)")
                                    bitstatstitle = f"Bit Stats: {defbit} **{customer}** (Total): {len(df_opbr.loc[df_opbr['ItemNo'] == defbit]):,.0f} (**{len(df_br.loc[df_br['ItemNo'] == defbit])}**) | ROP: {(df_opbr.loc[df_opbr['ItemNo'] == defbit, 'ROP'].mean()):,.1f} (**{(df_br.loc[df_br['ItemNo'] == defbit, 'ROP'].mean()):,.1f}**) | FTG: {(df_opbr.loc[df_opbr['ItemNo'] == defbit, 'DepthOut'].mean() - df_opbr.loc[df_opbr['ItemNo'] == defbit, 'DepthIn'].mean()):,.0f} (**{(df_br.loc[df_br['ItemNo'] == defbit, 'DepthOut'].mean() - df_br.loc[df_br['ItemNo'] == defbit, 'DepthIn'].mean()):,.0f}**)"                                        
                                    
                                except Exception as e:
                                    print(e)
                                    bitstatstitle = 'Bit Stats'
                                    # st.info(f"error bitrunstats: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                    
                                with st.expander(bitstatstitle, expanded=False):                                        
                                    try:
                                        if defbit is not None:
                                            colstat1, colstat2, colstat3 = st.columns(3,gap="small")
                                            df_bitrop[i] = df_br.loc[df_br['ItemNo'] == defbit]
                                            figbit[i] = go.Figure()
                                            figbit[i].add_trace(go.Violin(x=df_bitrop[i]['ItemNo'][ df_bitrop[i]['OperatorName'] == customer.upper() ],
                                                                    y=df_bitrop[i]['ROP'][ df_bitrop[i]['OperatorName'] == customer.upper() ],
                                                                    legendgroup=customer, scalegroup='Other', name=customer,
                                                                    side='negative',
                                                                    line_color='lightseagreen')
                                                        )
                                            figbit[i].add_trace(go.Violin(x=df_bitrop[i]['ItemNo'][ df_bitrop[i]['OperatorName'] != customer.upper() ],
                                                                    y=df_bitrop[i]['ROP'][ df_bitrop[i]['OperatorName'] != customer.upper() ],
                                                                    legendgroup='Other', scalegroup='Other', name='Other',
                                                                    side='positive',
                                                                    line_color='gray')
                                                        )
                                            ## loop through the values you want to label and add them as annotations
                                            # https://stackoverflow.com/questions/65463145/how-to-show-value-in-text-not-hover-in-boxplot-q1-q3-fences-with-px-box
                                            try:
                                                df_bitop = df_bitrop[i]['ROP'][df_bitrop[i]['OperatorName'] == customer.upper()]
                                                # st.write(df_bitop)
                                                for y in zip(["min","q1","med","q3","max"],df_bitop.quantile([0,0.25,0.5,0.75,1]).iloc[:,0].values):
                                                    figbit[i].add_annotation(x=-0.3,y=y[1],text=y[0] + ":" + str(y[1]),showarrow=False)
                                            except Exception as e:
                                                print(e)
                                                # st.info(f"error violin annotation: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                                
                                            figbit[i].update_layout(title={'text': f"{defbit} Run Stats: {customer} vs. Others",'y':0.9,'x':0.5,'xanchor': 'center','yanchor': 'top'})
                                            figbit[i].update_traces(meanline_visible=True,scalemode='count',)
                                            figbit[i].update_layout(legend=dict(orientation='h',xanchor='center',x=0.5)) # yanchor="top",y=0.01
                                            figbit[i].update_layout(violinmode='overlay')
                                            # violingap=0,orientation='h',
                                            with colstat1:
                                                st.plotly_chart(figbit[i], use_container_width=True)  
                                        
                                            df_bitrop[i] = df_br.loc[df_br['ItemNo'] == defbit]
                                            df_bitrop[i]['Drilled'] = (df_bitrop[i]['DepthOut'] - df_bitrop[i]['DepthIn'])
                                            df_bitrop[i] = df_bitrop[i].sort_values(by=['ROP'], ascending=False,ignore_index=True)                                              
                                            df_bitrop[i] = df_bitrop[i][0:5]
                                            figbitrop[i] = go.Figure()
                                            figbitrop[i].add_trace(go.Bar(x=df_bitrop[i].index,y=df_bitrop[i]['DepthIn'],name='',marker_color='rgba(0,0,0,0)',))                                            
                                            figbitrop[i].add_trace(go.Bar(x=df_bitrop[i].index,y=df_bitrop[i]['Drilled'],name='FTG',text=df_bitrop[i]['Drilled'],textposition='inside',insidetextanchor='middle',marker_color='rgba(0,147,159,1)', hoverinfo='text+x+y',))
                                            
                                            figbitrop[i].add_trace(go.Scatter(x=df_bitrop[i].index,y=df_bitrop[i]['ROP'],name='ROP',text=df_bitrop[i]['ROP'],textposition='bottom center',textfont=dict(color='rgba(152,194,31,1)'),mode='lines+markers+text',marker_color='rgba(152,194,31,1)',hoverinfo='text',yaxis = 'y2',))
                                            # figproposal.update_traces(texttemplate='%{text:,.1f}')                    
                                            figbitrop[i].update_layout(yaxis=dict(title='Depth',titlefont=dict(color='black'),tickfont=dict(color='black'),autorange='reversed',))
                                            figbitrop[i].update_layout(yaxis2=dict(title='ROP',titlefont=dict(color='rgba(152,194,31,1)'),tickfont=dict(color='rgba(152,194,31,1)'), autorange=True,anchor='free',overlaying='y',side='right',showgrid=False,position=1.0)) 
                                            # arrow_list1=[]
                                            # counter=0
                                            # if df_bitrop[i].shape[1] > 0 and len(str(df_bitrop[i].loc[1,'Section'])) > 1:
                                            #     try:
                                            #         for i in df_bitrop[i]['Section'].tolist():
                                            #             if i != "":
                                            #                 if int(df_bitrop[i]['ROP'].values[counter]) > 0:
                                            #                     arrow1=dict(x=df_bitrop[i].index.values[counter],y=df_bitrop[i]['Dout'].values[counter],xref="x",yref="y",text=f"{df_bitrop[i]['Section'].values[counter]} <br> {df_bitrop[i]['Din'].values[counter]}-{df_bitrop[i]['Dout'].values[counter]} <br> {(int(df_bitrop[i]['Drilled'].values[counter]) / int(df_bitrop[i]['ROP'].values[counter])):,.1f} hrs @ {df_bitrop[i]['ROP'].values[counter]:,.1f} ft/hr",ax=60,ay=0,align="left",showarrow=False,)
                                            #                     arrow_list1.append(arrow1)
                                            #                 counter+=1
                                            #             else:
                                            #                 counter+=1
                                            #     except Exception as e:
                                            #         print(e)
                                            #         st.info(f"error figproposal1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                                 
                                            # figbitrop[i].update_layout(annotations=arrow_list1)
                                            figbitrop[i].update_yaxes(ticklabelposition="outside top")      
                                            figbitrop[i].update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor='rgba(0,0,0,0)')
                                            figbitrop[i].update_layout(xaxis=dict(side='top',visible=False)) 
                                            figbitrop[i].update_layout(uniformtext_minsize=8, uniformtext_mode='show')  # 'hide' to hide overflow text
                                            figbitrop[i].update_layout(title={'text': f"{bit} Top 5 ROP",'y':0.9,'x':0.5,'xanchor': 'center','yanchor': 'top'})
                                            figbitrop[i].update_layout(barmode='stack', bargap=0.7, showlegend = False) # bargap=(1/(sectioncount))
                                            with colstat2:
                                                st.plotly_chart(figbitrop[i], use_container_width=True) 
                                                
                                            
                                            df_bitftg[i] = df_br.loc[df_br['ItemNo'] == defbit]
                                            df_bitftg[i]['Drilled'] = (df_bitftg[i]['DepthOut'] - df_bitftg[i]['DepthIn'])
                                            df_bitftg[i] = df_bitftg[i].sort_values(by=['Drilled'], ascending=False,ignore_index=True)                                              
                                            df_bitftg[i] = df_bitftg[i][0:5]
                                            figbitftg[i] = go.Figure()
                                            figbitftg[i].add_trace(go.Bar(x=df_bitftg[i].index,y=df_bitftg[i]['DepthIn'],name='',marker_color='rgba(0,0,0,0)',))                                            
                                            figbitftg[i].add_trace(go.Bar(x=df_bitftg[i].index,y=df_bitftg[i]['Drilled'],name='FTG',text=df_bitftg[i]['Drilled'],textposition='inside',insidetextanchor='middle',marker_color='rgba(0,147,159,1)', hoverinfo='text+x+y',))
                                            
                                            figbitftg[i].add_trace(go.Scatter(x=df_bitftg[i].index,y=df_bitftg[i]['ROP'],name='ROP',text=df_bitftg[i]['ROP'],textposition='bottom center',textfont=dict(color='rgba(152,194,31,1)'),mode='lines+markers+text',marker_color='rgba(152,194,31,1)',hoverinfo='text',yaxis = 'y2',))
                                            # figproposal.update_traces(texttemplate='%{text:,.1f}')                    
                                            figbitftg[i].update_layout(yaxis=dict(title='Depth',titlefont=dict(color='black'),tickfont=dict(color='black'),autorange='reversed',))
                                            figbitftg[i].update_layout(yaxis2=dict(title='ROP',titlefont=dict(color='rgba(152,194,31,1)'),tickfont=dict(color='rgba(152,194,31,1)'), autorange=True,anchor='free',overlaying='y',side='right',showgrid=False,position=1.0)) 
                                            # arrow_list1=[]
                                            # counter=0
                                            # if df_bitftg[i].shape[1] > 0 and len(str(df_bitftg[i].loc[1,'Section'])) > 1:
                                            #     try:
                                            #         for i in df_bitftg[i]['Section'].tolist():
                                            #             if i != "":
                                            #                 if int(df_bitftg[i]['ROP'].values[counter]) > 0:
                                            #                     arrow1=dict(x=df_bitftg[i].index.values[counter],y=df_bitftg[i]['Dout'].values[counter],xref="x",yref="y",text=f"{df_bitftg[i]['Section'].values[counter]} <br> {df_bitftg[i]['Din'].values[counter]}-{df_bitftg[i]['Dout'].values[counter]} <br> {(int(df_bitftg[i]['Drilled'].values[counter]) / int(df_bitftg[i]['ROP'].values[counter])):,.1f} hrs @ {df_bitftg[i]['ROP'].values[counter]:,.1f} ft/hr",ax=60,ay=0,align="left",showarrow=False,)
                                            #                     arrow_list1.append(arrow1)
                                            #                 counter+=1
                                            #             else:
                                            #                 counter+=1
                                            #     except Exception as e:
                                            #         print(e)
                                            #         st.info(f"error figproposal1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                                 
                                            # figbitftg[i].update_layout(annotations=arrow_list1)
                                            figbitftg[i].update_yaxes(ticklabelposition="outside top")      
                                            figbitftg[i].update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor='rgba(0,0,0,0)')
                                            figbitftg[i].update_layout(xaxis=dict(side='top',visible=False)) 
                                            figbitftg[i].update_layout(uniformtext_minsize=8, uniformtext_mode='show')  # 'hide' to hide overflow text
                                            figbitftg[i].update_layout(title={'text': f"{bit} Top 5 Footage",'y':0.9,'x':0.5,'xanchor': 'center','yanchor': 'top'})
                                            figbitftg[i].update_layout(barmode='stack', bargap=0.7, showlegend = False) # bargap=(1/(sectioncount))
                                            with colstat3:
                                                st.plotly_chart(figbitftg[i], use_container_width=True) 
                                        else:
                                            st.empty()
                                        
                                    except Exception as e:
                                        comps.info(f"error bitroptop5: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                                
                                            
                        except Exception as e:
                            print(e)
                            st.info(f"error priceloop2: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                                                        
                    st.divider()
                
                # # Call Section block
                # sectioncard()            
            
            # Call Sections
            definesections(sectioncount,df)                

            # Check well df
            st.divider()
            
            dfwell = dfwell.drop(dfwell.index[0])                
            dfchart = dfwell.replace('',np.nan,regex = True)
            dfchart.fillna(0)
            dfchart.apply(pd.to_numeric, errors='ignore')
            # dfchart = st.experimental_data_editor(dfwell)
            dfchart['Drilled'] = None
            
            if dfchart['Section'].isnull().values.any() == False and dfchart['Din'].isnull().values.any() == False and dfchart['Dout'].isnull().values.any() == False and dfchart['ROP'].isnull().values.any() == False:                    
                dfchart['Drilled'] = dfchart['Dout'].astype(float) - dfchart['Din'].astype(float)
                dfchart['ROP'] = dfchart['ROP'].astype(float)
                # # https://plotly.com/python/waterfall-charts/
                # figsidewell = go.Figure(go.Waterfall(name = "", orientation = "v",measure = ["relative", "relative", "relative", "relative", "total", "relative", "relative", "relative", "relative","total"],
                #     x = ["S","V","C","L","T","L1","C1","V1","S1"],
                #     y = [1500, 4000, 1000, 5000, 11500, -5000, -1000,-4000,-1500,None], connector = {"line":{"color":"rgb(63, 63, 63)"}},base = 100,                                    
                #     decreasing = {"marker":{"color":"gray"}},
                #     increasing = {"marker":{"color":"gray"}},
                #     totals = {"marker":{"color":"black"}}
                #     ))
                # figsidewell.update_yaxes(autorange="reversed",ticklabelposition="inside top",title=None)
                # figsidewell.update_xaxes(visible=False, showticklabels=False)
                # figsidewell.update_layout(title = "Well Diagram",showlegend = False)
                # # waterfallgap = 0.3
                # st.sidebar.plotly_chart(figsidewell, use_container_width=True)

                # https://plotly.com/python/bar-charts/
                        
                # dfchart = st.experimental_data_editor(dfchart)
                
                chtcol1, chtcol2, chtcol3 = st.columns([2,2,1],gap="small")
                # if showperfchart:
                try:
                    # st.write('dfchart')
                    figproposal = go.Figure()
                    try:
                        # text=f"{dfchart['Section'].tolist()} <br> {dfchart['Din'].tolist()}-{dfchart['Dout'].tolist()} <br> {set(dfchart['Drilled'].astype(int) / dfchart['ROP'].astype(int))} hrs @ {dfchart['ROP'].tolist()} ft/hr",textposition='inside',insidetextanchor='end',
                        figproposal.add_trace(go.Bar(x=dfchart.index,y=dfchart['Din'],name='',marker_color='rgba(0,0,0,0)',))
                    except Exception as e:
                        figproposal.add_trace(go.Bar(x=dfchart.index,y=dfchart['Din'],name='',marker_color='rgba(0,0,0,0)',))
                        st.info(f"error figproposal2a: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    # 'rgba(138,207,221,1)'
                    colors = {0: 'rgba(63,112,119,1)',
                                1: 'rgba(63,112,119,1)',
                                2: 'rgba(63,112,119,0.91)',
                                3: 'rgba(63,112,119,0.82)',
                                4: 'rgba(63,112,119,0.73)',
                                5: 'rgba(63,112,119,0.64)',
                                6: 'rgba(63,112,119,0.55)',
                                7: 'rgba(63,112,119,0.46)',
                                8: 'rgba(63,112,119,0.37)',
                                9: 'rgba(63,112,119,0.28)',
                                10: 'rgba(63,112,119,0.19)',}
                    
                    color_discrete_sequence = ['rgba(63,112,119,1)','rgba(63,112,119,1)','rgba(63,112,119,0.91)','rgba(63,112,119,0.82)','rgba(63,112,119,0.73)','rgba(63,112,119,0.64)','rgba(63,112,119,0.55)','rgba(63,112,119,0.46)','rgba(63,112,119,0.37)','rgba(63,112,119,0.28)','rgba(63,112,119,0.19)']
                    figproposal.add_trace(go.Bar(x=dfchart.index,y=dfchart['Drilled'],name='Sections',text=dfchart['Drilled'],textposition='inside',insidetextanchor='middle',marker_color=color_discrete_sequence, hoverinfo='text+x+y',))
                    # marker_color=colors[dfchart.index],
                    # figproposal.add_trace(go.Bar(x=dfchart.index,y=dfchart['Drilled'],name='Sections',text=dfchart['Drilled'],textposition='inside',insidetextanchor='middle',marker_color='rgba(138,207,221,1)',hoverinfo='text+x+y',))
                    
                    # hoemplate="<b>%{x}</b><br><br>" +"Depth in: %{y:.0f}ft<br>" +"Depth out: %{y1:.0f}<br>" +"<extra></extra>",
                    # color_discrete_sequence=['#1f77b4', '#17becf', '#d62728'],             hover_data={'Math':':.2f', 'Science':':.2f', 'English':':.2f'}
                    figproposal.add_trace(go.Scatter(x=dfchart.index,y=dfchart['ROP'],name='ROP',text=dfchart['ROP'],textposition='bottom center',textfont=dict(color='rgba(152,194,31,1)'),mode='lines+markers+text',marker_color='rgba(152,194,31,1)',hoverinfo='text',yaxis = 'y2',))
                    # figproposal.update_traces(texttemplate='%{text:,.1f}')                    
                    figproposal.update_layout(yaxis=dict(title='Depth',titlefont=dict(color='black'),tickfont=dict(color='black'),autorange='reversed',))
                    figproposal.update_layout(yaxis2=dict(title='ROP',titlefont=dict(color='rgba(152,194,31,1)'),tickfont=dict(color='rgba(152,194,31,1)'),range=[0,250], anchor='free',overlaying='y',side='right',showgrid=False,position=1.0)) 
                    
                    # if dfchart['Dout'].max() <= 25000:
                    #     figproposal.update_layout(yaxis=dict(title='Depth',range=[25000, 0],))
                    # else:
                    #     figproposal.update_layout(yaxis=dict(title='Depth',range=[dfchart['Dout'].max() + 500, 0],))
                        
                    # if dfchart['ROP'].max() <= 250:
                    #     figproposal.update_layout(yaxis2=dict(title='ROP',titlefont=dict(color='black'),tickfont=dict(color='black'), range=[250, 0],anchor='free',overlaying='y',side='right',showgrid=False,title_standoff=10,position=1.0)) 
                    # else:
                    #     figproposal.update_layout(yaxis2=dict(title='ROP',titlefont=dict(color='black'),tickfont=dict(color='black'), range=[(dfchart['ROP'].max() + 50), 0],anchor='free',overlaying='y',side='right',showgrid=False,title_standoff=10,position=1.0)) 
                    arrow_list1=[]
                    counter=0
                    if dfchart.shape[1] > 0 and len(str(dfchart.loc[1,'Section'])) > 1:
                        try:
                            for i in dfchart['Section'].tolist():
                                if i != "":
                                    if float(dfchart['ROP'].values[counter]) > 0:
                                        arrow1=dict(x=dfchart.index.values[counter],y=dfchart['Dout'].values[counter],xref="x",yref="y",text=f"{dfchart['Section'].values[counter]} <br> {dfchart['Din'].values[counter]}-{dfchart['Dout'].values[counter]} <br> {(float(dfchart['Drilled'].values[counter]) / float(dfchart['ROP'].values[counter])):,.1f} hrs @ {dfchart['ROP'].values[counter]:,.1f} ft/hr",ax=60,ay=0,align="left",showarrow=False,)
                                        arrow_list1.append(arrow1)
                                    counter+=1
                                else:
                                    counter+=1
                        except Exception as e:
                            print(e)
                            st.info(f"error figproposal1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                        
                    figproposal.update_layout(annotations=arrow_list1)
                    figproposal.update_yaxes(ticklabelposition="outside top")      
                    figproposal.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor='rgba(0,0,0,0)')
                    figproposal.update_layout(xaxis=dict(side='top',visible=False)) 
                    figproposal.update_layout(uniformtext_minsize=8, uniformtext_mode='show')  # 'hide' to hide overflow text
                    figproposal.update_layout(title={'text': "Well Plan",'y':0.9,'x':0.5,'xanchor': 'center','yanchor': 'top'})
                    # ,title=None
                    figproposal.update_layout(barmode='stack', bargap=0.7, showlegend = False) # bargap=(1/(sectioncount))
                    with chtcol1:
                        st.plotly_chart(figproposal, use_container_width=True)
                    # https://plotly.com/python/marker-style/
                    
                except Exception as e:
                    print(e)
                    st.info(f"error figproposal2: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                
                # if showdvdchart:
                try:
                    # Prop Run #	Prop Depth	Prop Hrs
                    #0 1	0	0
                    #1 1	5065	25
                    #2 2	5065	30.065
                    #3 2	9750	90.065
                    #4 3	9750	99.815
                    #5 3	11563	129.815
                    if dfchart.shape[1] > 0 and  (dfchart['ROP'] == 0).sum() == 0: #len(str(dfchart.loc[1,'Section'])) > 1 and
                        dvddata = []        
                        dfdvd = pd.DataFrame(dvddata, columns=['run','Depth','Hours','state','note'])
                        # Loop                    
                        dfdvd.loc[0,'run'] = 1      
                        dfdvd.loc[0,'Depth'] = float(dfchart.loc[1,'Din'] )                            
                        dfdvd.loc[0,'Hours'] = 0
                        dfdvd.loc[0,'state'] = 'base' 
                        dfdvd.loc[0,'note'] = '' 
                        if len(dfchart) >= 1: 
                            dfdvd.loc[1,'run'] = 1
                            dfdvd.loc[1,'Depth'] = float(dfchart.loc[1,'Dout']  )                           
                            dfdvd.loc[1,'Hours'] = (float(dfchart.loc[1,'Dout']) - float(dfchart.loc[1,'Din'])) / float(dfchart.loc[1,'ROP'])
                            dfdvd.loc[1,'state'] = 'drill'
                            dfdvd.loc[1,'note'] = f"{dfchart.loc[1,'Section']} <br> {dfchart.loc[1,'Size']} {dfchart.loc[1,'Type']}"
                        if len(dfchart) >= 2: 
                            dfdvd.loc[2,'run'] = 2     
                            dfdvd.loc[2,'Depth'] = float(dfchart.loc[2,'Din'] )                            
                            dfdvd.loc[2,'Hours'] = dfdvd.loc[1,'Hours'] + (dfdvd.loc[2,'Depth'] / float(triprate))
                            dfdvd.loc[2,'state'] = 'trip'
                            dfdvd.loc[2,'note'] = ''
                            dfdvd.loc[3,'run'] = 2      
                            dfdvd.loc[3,'Depth'] = float(dfchart.loc[2,'Dout'] )                            
                            dfdvd.loc[3,'Hours'] = dfdvd.loc[2,'Hours'] + ((float(dfchart.loc[2,'Dout']) - float(dfchart.loc[2,'Din'])) / float(dfchart.loc[2,'ROP']))
                            dfdvd.loc[3,'state'] = 'drill'
                            dfdvd.loc[3,'note'] = f"{dfchart.loc[2,'Section']} <br> {dfchart.loc[2,'Size']} {dfchart.loc[2,'Type']}"
                        if len(dfchart) >= 3: 
                            dfdvd.loc[4,'run'] = 3     
                            dfdvd.loc[4,'Depth'] = float(dfchart.loc[3,'Din'])                         
                            dfdvd.loc[4,'Hours'] = dfdvd.loc[3,'Hours'] + (dfdvd.loc[4,'Depth'] / float(triprate))
                            dfdvd.loc[4,'state'] = 'trip'
                            dfdvd.loc[4,'note'] = ''
                            dfdvd.loc[5,'run'] = 3      
                            dfdvd.loc[5,'Depth'] = float(dfchart.loc[3,'Dout'])
                            dfdvd.loc[5,'Hours'] = dfdvd.loc[4,'Hours'] + ((float(dfchart.loc[3,'Dout']) - float(dfchart.loc[3,'Din'])) / float(dfchart.loc[3,'ROP']))
                            dfdvd.loc[5,'state'] = 'drill'
                            dfdvd.loc[5,'note'] = f"{dfchart.loc[3,'Section']} <br> {dfchart.loc[3,'Size']} {dfchart.loc[3,'Type']}"
                        if len(dfchart) >= 4: 
                            dfdvd.loc[6,'run'] = 4     
                            dfdvd.loc[6,'Depth'] = float(dfchart.loc[4,'Din'])                         
                            dfdvd.loc[6,'Hours'] = dfdvd.loc[5,'Hours'] + (dfdvd.loc[6,'Depth'] / float(triprate))
                            dfdvd.loc[6,'state'] = 'trip'
                            dfdvd.loc[6,'note'] = ''
                            dfdvd.loc[7,'run'] = 4     
                            dfdvd.loc[7,'Depth'] = float(dfchart.loc[4,'Dout'])
                            dfdvd.loc[7,'Hours'] = dfdvd.loc[6,'Hours'] + ((float(dfchart.loc[4,'Dout']) - float(dfchart.loc[4,'Din'])) / float(dfchart.loc[4,'ROP']))
                            dfdvd.loc[7,'state'] = 'drill'
                            dfdvd.loc[7,'note'] = f"{dfchart.loc[4,'Section']} <br> {dfchart.loc[4,'Size']} {dfchart.loc[4,'Type']}"
                        if len(dfchart) >= 5: 
                            dfdvd.loc[8,'run'] = 5
                            dfdvd.loc[8,'Depth'] = float(dfchart.loc[5,'Din'])
                            dfdvd.loc[8,'Hours'] = dfdvd.loc[7,'Hours'] + (dfdvd.loc[8,'Depth'] / float(triprate))
                            dfdvd.loc[8,'state'] = 'trip'
                            dfdvd.loc[8,'note'] = ''
                            dfdvd.loc[9,'run'] = 5
                            dfdvd.loc[9,'Depth'] = float(dfchart.loc[5,'Dout'])
                            dfdvd.loc[9,'Hours'] = dfdvd.loc[8,'Hours'] + ((float(dfchart.loc[5,'Dout']) - float(dfchart.loc[5,'Din'])) / float(dfchart.loc[5,'ROP']))
                            dfdvd.loc[9,'state'] = 'drill'
                            dfdvd.loc[9,'note'] = f"{dfchart.loc[5,'Section']} <br> {dfchart.loc[5,'Size']} {dfchart.loc[5,'Type']}"
                        if len(dfchart) >= 6: 
                            dfdvd.loc[10,'run'] = 6
                            dfdvd.loc[10,'Depth'] = float(dfchart.loc[6,'Din'])
                            dfdvd.loc[10,'Hours'] = dfdvd.loc[9,'Hours'] + (dfdvd.loc[10,'Depth'] / float(triprate))
                            dfdvd.loc[10,'state'] = 'trip'
                            dfdvd.loc[10,'note'] = ''
                            dfdvd.loc[11,'run'] = 6
                            dfdvd.loc[11,'Depth'] = float(dfchart.loc[6,'Dout'])
                            dfdvd.loc[11,'Hours'] = dfdvd.loc[10,'Hours'] + ((float(dfchart.loc[6,'Dout']) - float(dfchart.loc[6,'Din'])) / float(dfchart.loc[6,'ROP']))
                            dfdvd.loc[11,'state'] = 'drill'
                            dfdvd.loc[11,'note'] = f"{dfchart.loc[6,'Section']} <br> {dfchart.loc[6,'Size']} {dfchart.loc[6,'Type']}"
                        if len(dfchart) >= 7: 
                            dfdvd.loc[12,'run'] = 7
                            dfdvd.loc[12,'Depth'] = float(dfchart.loc[7,'Din'])
                            dfdvd.loc[12,'Hours'] = dfdvd.loc[11,'Hours'] + (dfdvd.loc[12,'Depth'] / float(triprate))
                            dfdvd.loc[12,'state'] = 'trip'
                            dfdvd.loc[12,'note'] = ''
                            dfdvd.loc[13,'run'] = 7
                            dfdvd.loc[13,'Depth'] = float(dfchart.loc[7,'Dout'])
                            dfdvd.loc[13,'Hours'] = dfdvd.loc[12,'Hours'] + ((float(dfchart.loc[7,'Dout']) - float(dfchart.loc[7,'Din'])) / float(dfchart.loc[7,'ROP']))
                            dfdvd.loc[13,'state'] = 'drill'
                            dfdvd.loc[13,'note'] = f"{dfchart.loc[7,'Section']} <br> {dfchart.loc[7,'Size']} {dfchart.loc[7,'Type']}"
                        if len(dfchart) >= 8: 
                            dfdvd.loc[14,'run'] = 8
                            dfdvd.loc[14,'Depth'] = float(dfchart.loc[8,'Din'])
                            dfdvd.loc[14,'Hours'] = dfdvd.loc[13,'Hours'] + (dfdvd.loc[14,'Depth'] / float(triprate))
                            dfdvd.loc[14,'state'] = 'trip'
                            dfdvd.loc[14,'note'] = ''
                            dfdvd.loc[15,'run'] = 8
                            dfdvd.loc[15,'Depth'] = float(dfchart.loc[8,'Dout'])
                            dfdvd.loc[15,'Hours'] = dfdvd.loc[14,'Hours'] + ((float(dfchart.loc[8,'Dout']) - float(dfchart.loc[8,'Din'])) / float(dfchart.loc[8,'ROP']))
                            dfdvd.loc[15,'state'] = 'drill'
                            dfdvd.loc[15,'note'] = f"{dfchart.loc[8,'Section']} <br> {dfchart.loc[8,'Size']} {dfchart.loc[8,'Type']}"
                        if len(dfchart) >= 9: 
                            dfdvd.loc[16,'run'] = 9
                            dfdvd.loc[16,'Depth'] = float(dfchart.loc[9,'Din'])
                            dfdvd.loc[16,'Hours'] = dfdvd.loc[15,'Hours'] + (dfdvd.loc[16,'Depth'] / float(triprate))
                            dfdvd.loc[16,'state'] = 'trip'
                            dfdvd.loc[16,'note'] = ''
                            dfdvd.loc[17,'run'] = 9
                            dfdvd.loc[17,'Depth'] = float(dfchart.loc[9,'Dout'])
                            dfdvd.loc[17,'Hours'] = dfdvd.loc[16,'Hours'] + ((float(dfchart.loc[9,'Dout']) - float(dfchart.loc[9,'Din'])) / float(dfchart.loc[9,'ROP']))
                            dfdvd.loc[17,'state'] = 'drill'
                            dfdvd.loc[17,'note'] = f"{dfchart.loc[9,'Section']} <br> {dfchart.loc[9,'Size']} {dfchart.loc[9,'Type']}"
                        if len(dfchart) >= 10: 
                            dfdvd.loc[18,'run'] = 10
                            dfdvd.loc[18,'Depth'] = float(dfchart.loc[10,'Din'])
                            dfdvd.loc[18,'Hours'] = dfdvd.loc[17,'Hours'] + (dfdvd.loc[18,'Depth'] / float(triprate))
                            dfdvd.loc[18,'state'] = 'trip'
                            dfdvd.loc[18,'note'] = ''
                            dfdvd.loc[19,'run'] = 10
                            dfdvd.loc[19,'Depth'] = float(dfchart.loc[10,'Dout'])
                            dfdvd.loc[19,'Hours'] = dfdvd.loc[18,'Hours'] + ((float(dfchart.loc[10,'Dout']) - float(dfchart.loc[10,'Din'])) / float(dfchart.loc[10,'ROP']))
                            dfdvd.loc[19,'state'] = 'drill'
                            dfdvd.loc[19,'note'] = f"{dfchart.loc[10,'Section']} <br> {dfchart.loc[10,'Size']} {dfchart.loc[10,'Type']}"
                        
                        
                        # for i in range(1,(len(dfchart)+1)):
                        #     # createsectionpage(dfchart.loc[i,'Section'],dfchart.loc[i])   
                        #     dfdvd.loc[-1,'run'] = i 
                        #     dfdvd.loc[-1,'Depth'] = dfchart.loc[i,'Dout']
                        #     dfdvd.loc[-1,'Hours'] = dfdvd.loc[i-1,'Hours'] + ((dfchart.loc[i,'Dout'] - dfchart.loc[i,'Din']) / dfchart.loc[i,'ROP'])
                        #     dfdvd.loc[-1,'state'] = 'drill'
                                                            
                        #     dfdvd.loc[-1,'run'] = i+1 
                        #     dfdvd.loc[-1,'Depth'] = dfchart.loc[i,'Dout']
                        #     dfdvd.loc[-1,'Hours'] = dfchart.loc[i,'Dout'] / triprate + dfdvd.loc[-1,'Hours']
                        #     dfdvd.loc[-1,'state'] = 'trip'
                        #     # For each section, create entry and exit
                        #     # Account for trip time

                        # # st.write('dfchart')
                        # dfchart = st.experimental_data_editor(dfchart)
                        figdvd = go.Figure()
                        
                        figdvd.add_trace(go.Scatter(x=dfdvd['Hours'],y=dfdvd['Depth'],name='Hrs/Depth',marker_color='rgba(0,147,159,1)',))
                        
                        # figdvd.update_yaxes(ticklabelposition="outside top")
                        # figdvd.update_layout(xaxis=dict(side='top',range=[0, 100],ticksuffix='%')) 
                        figdvd.update_layout(title={'text':"Well Hours vs. Depth",'y':0.9,'x':0.5,'xanchor': 'center','yanchor': 'top'},showlegend = False)
                        # ,title=None
                        # figdvd.update_layout(title = "Well Hours vs. Depth",showlegend = False) # bargap=(1/(sectioncount))
                            
                        figdvd.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor='rgba(0,0,0,0)')
                        figdvd.update_layout(yaxis=dict(title='Depth',autorange="reversed",tick0=0))
                        figdvd.update_layout(xaxis=dict(title='Hours',tick0=0,dtick=24,minor=dict(ticklen=5, tickcolor="black", showgrid=True)))                            
                        # construct menus https://stackoverflow.com/questions/68943989/plotly-how-to-change-between-daily-and-hourly-candlesticks-using-a-dropdown-men
                        # updatemenus = [{
                        # #                 'active':1,
                        #                 'buttons': [{'method': 'update',
                        #                             'label': 'Hourly',
                        #                             'args': [
                        #                                     # 1. updates to the traces
                        #                                     {'open': [list(df.Open)],
                        #                                     'high': [list(df.High)],
                        #                                     'low': [list(df.Low)],
                        #                                     'close': [list(df.Close)],
                        #                                     'x':[list(df.Date)],
                        #                                     'visible': True}, 
                                                            
                        #                                     # 2. updates to the layout
                        #                                     {'title':'Hourly'},
                                                            
                        #                                     # 3. which traces are affected 
                        # #                                       [0, 1],
                                                            
                        #                                     ],  },
                        #                             {'method': 'update',
                        #                             'label': 'Daily', 
                        #                             'args': [
                        #                                     # 1. updates to the traces  
                        #                                     {'open': [list(dfd.Open)],
                        #                                         'high': [list(dfd.High)],
                        #                                         'low': [list(dfd.Low)],
                        #                                         'close': [list(dfd.Close)],
                        #                                         'x':[list(dfd.Date)],
                        #                                     'visible': True},
                                                            
                        #                                     # 2. updates to the layout
                        #                                     {'title':'Daily'},
                                                            
                        #                                     # 3. which traces are affected
                        # #                                        [0, 1]
                        #                                     ]
                                                    
                        #                             },],
                        #                 'type':'dropdown',
                        #                 'direction': 'down',
                        #                 'showactive': True,}]

                        # # update layout with buttons, and show the figure
                        # figdvd.update_layout(updatemenus=updatemenus)

                        # figdvd.update_yaxes(autorange="reversed",ticklabelposition="inside top",title=None)
                        # if dfchart['Dout'].max() <= 25000:
                        #     figdvd.update_layout(yaxis=dict(title='Depth',range=[25000, 0],))
                        # else:
                        #     figdvd.update_layout(yaxis=dict(title='Depth',range=[dfchart['Dout'].max() + 500, 0],))
                        
                        # https://stackoverflow.com/questions/70395238/add-multiple-annotations-at-once-to-plotly-line-chart
                        arrow_list=[]
                        counter=0
                        for i in dfdvd['note'].tolist():
                            if i != "":
                                arrow=dict(x=dfdvd['Hours'].values[counter],y=dfdvd['Depth'].values[counter],xref="x",yref="y",text=i,ax=20,align="left",arrowhead=0,arrowwidth=1,arrowcolor='rgb(0,0,0)',)
                                arrow_list.append(arrow)
                                counter+=1
                            else:
                                counter+=1

                        figdvd.update_layout(annotations=arrow_list)
                        # figdvd.add_annotation(x=2,y=5,xref="x",yref="y",text="max=5",showarrow=True,align="center",arrowhead=2,arrowcolor="#636363",ax=20,ay=-30,bgcolor="#ff7f0e",opacity=0.8)
                    
                        with chtcol2:
                            st.plotly_chart(figdvd, use_container_width=True)
                    
                    # figproposal.update_yaxes(ticklabelposition="outside top")
                    # figproposal.update_layout(xaxis=dict(side='top',range=[0, 100],ticksuffix='%')) 
                    # figproposal.update_layout(title={'text': "Plot Title",'y':0.9,'x':0.5,'xanchor': 'center','yanchor': 'top'})
                    # # ,title=None
                    # figproposal.update_layout(title = "Well Plan",barmode='stack', bargap=0.7, showlegend = False) # bargap=(1/(sectioncount))
                    # st.plotly_chart(figproposal, use_container_width=True)
                    # # https://plotly.com/python/marker-style/
                    
                except Exception as e:
                    print(e)
                    st.info(f"error figdvd: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                        
                # if showcostchart:
                # https://plotly.com/python/waterfall-charts/
                try:                            
                    welldata = []   
                    counter=0     
                    dfwellplan = pd.DataFrame(welldata, columns=['measure','x','y'])
                    # st.write(dfchart['Section'].tolist())
                    # Loop                              
                    for i in dfchart['Size'].tolist(): 
                        if i != "":        
                            counter+=1
                            # st.write(f"up {i}")         
                            dfwellplan.loc[counter-1,'measure'] = 'relative'
                            dfwellplan.loc[counter-1,'x'] = f"{dfchart.loc[counter,'Size']}in"
                            dfwellplan.loc[counter-1,'y'] = dfchart.loc[counter,'Drilled']
                            # st.write(f"up {i} {counter} , {dfwellplan.loc[counter-1,'x']} { dfwellplan.loc[counter-1,'y']}")
                            # counter+=1
                        else:
                            counter+=1
                    
                    # Add gap data entries in middle for chart spacing 
                    dfwellplan.loc[counter,'measure'] = 'total'
                    dfwellplan.loc[counter,'x'] = ' '
                    dfwellplan.loc[counter,'y'] = float(dfchart.loc[counter,'Drilled'])
                    dfwellplan.loc[counter+1,'measure'] = 'total'
                    dfwellplan.loc[counter+1,'x'] = ' . '
                    dfwellplan.loc[counter+1,'y'] = float(dfchart.loc[counter,'Drilled'])
                    dfwellplan.loc[counter+2,'measure'] = 'total'
                    dfwellplan.loc[counter+2,'x'] = '   '
                    dfwellplan.loc[counter+2,'y'] = float(dfchart.loc[counter,'Drilled'])
                    # dfwellplan.loc[counter+2,'measure'] = 'total'
                    # dfwellplan.loc[counter+2,'x'] = '   k'
                    # dfwellplan.loc[counter+2,'y'] = int(dfchart.loc[counter,'Drilled'])
                    
                    # st.write(dfwellplan)    
                    countup = counter +4
                    # st.write(f"up {counter}") 
                    for i in reversed(dfchart['Size'].tolist()): 
                        if i != "":                 
                            dfwellplan.loc[countup,'measure'] = 'relative'
                            dfwellplan.loc[countup,'x'] = f" {dfchart.loc[counter,'Size']} in"
                            dfwellplan.loc[countup,'y'] = -1 * float(dfchart.loc[counter,'Drilled'])
                            # st.write(f"down {i} {counter} , {dfwellplan.loc[counter-1,'x']} {dfwellplan.loc[counter-1,'y']}")
                            counter-=1
                            countup+=1
                        else:
                            counter-=1
                            countup+=1
                    
                    # st.write(dfwellplan)
                    # dfwellplan.loc[dfwellplan.shape[0]] = ['total',dfchart.loc[counter,'Section'],None]
                    # dfwellplan.loc[-1,'measure'] = 'total'
                    # dfwellplan.loc[-1,'x'] = dfchart.loc[counter,'Section']
                    # dfwellplan.loc[-1,'y'] = None
                    
                    # st.write(dfwellplan)
                    figwell = go.Figure(go.Waterfall(name = "", orientation = "v",measure = dfwellplan['measure'],
                        x = dfwellplan['x'],
                        y = dfwellplan['y'], connector = {"line":{"color":"rgb(63, 63, 63)"}},base = 100,                                    
                        decreasing = {"marker":{"color":"gray"}},
                        increasing = {"marker":{"color":"gray"}},
                        totals = {"marker":{"color":"rgba(63, 63, 63,0)"}},
                        # textposition = "inside",
                        # text = dfwellplan["x"].tolist(),
                        ))
                    
                    arrow_list3=[]
                    counter=0
                    for i in dfchart['Size'].tolist(): 
                        if i != "":
                            arrow=dict(x=' . ',y=(((float(dfchart['Dout'].values[counter]) - float(dfchart['Din'].values[counter])) / 2) + float(dfchart['Din'].values[counter])),xref="x",yref="y",text=f"""{dfchart['Size'].values[counter]}" {dfchart['Section'].values[counter]}<br>$/ft: {(dfchart['Pricetotal'].values[counter] / dfchart['Drilled'].values[counter]):,.2f}""",ax=0,align="left",showarrow=False,)
                            arrow_list3.append(arrow)
                            counter+=1
                        else:
                            counter+=1

                    figwell.update_layout(annotations=arrow_list3)
                    figwell.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor='rgba(0,0,0,0)')
                    figwell.update_yaxes(autorange="reversed",ticklabelposition="inside top",title=None)
                    figwell.update_xaxes(visible=False, showticklabels=False)
                    # figwell.update_layout(title = "Well Diagram",showlegend = False)
                    figwell.update_layout(title={'text':"Well Diagram",'y':0.9,'x':0.5,'xanchor': 'center','yanchor': 'top'},showlegend = False)
                    # waterfallgap = 0.3
                    with chtcol3:
                        st.plotly_chart(figwell, use_container_width=True)

                    # https://plotly.com/python/bar-charts/
                except Exception as e:
                    print(e)
                    st.info(f"error figwell: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
            
                
                config0 = dict({'displayModeBar': False})
                indcol1, indcol2 = st.columns(2,gap="small")
                try:                        
                    dfchart['Pricetotal'] = dfchart['Pricetotal'].astype('float')
                    # pricetotaltotal = dfchart['Pricetotal'].cumsum()
                    dfchart['Dout'] = dfchart['Dout'].astype('float')
                    # st.write(f'totaltotal: {pricetotaltotal}')
                    with indcol1:
                        figtotalcost = go.Figure(go.Indicator(
                            mode = 'number',value = dfchart['Pricetotal'].sum(),number = {'prefix': '$', 'valueformat':',.2f'},
                            # delta = {"reference": 512, "valueformat": ".0f", "prefix": "$", "suffix": "m"},
                            title = {'text': 'Total Bit Cost'},
                            domain = {'x': [0, 1], 'y': [0, 1]}))
                        figtotalcost.update_traces(number_font_size=36, title_font_size=16)                            
                        figtotalcost.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor='rgba(0,0,0,0)',height=100, margin=go.layout.Margin(l=0,r=0,b=0,t=0,))
                        st.plotly_chart(figtotalcost, use_container_width=True,config=config0)
                    
                    with indcol2:
                        figtotalcostft = go.Figure(go.Indicator(
                            mode = 'number',value = (dfchart['Pricetotal'].sum() / dfchart['Dout'].sum()),number = {'prefix': '$', 'suffix': '/ft', 'valueformat':',.2f'},
                            # delta = {"reference": 512, "valueformat": ".0f", "prefix": "$", "suffix": "m"},
                            title = {'text': 'Well Average $/ft'},
                            domain = {'x': [0, 1], 'y': [0, 1]}))
                        figtotalcostft.update_traces(number_font_size=36, title_font_size=16)  
                        figtotalcostft.update_layout(paper_bgcolor='rgba(0,0,0,0)',plot_bgcolor='rgba(0,0,0,0)',height=100, margin=go.layout.Margin(l=0,r=0,b=0,t=0,))
                        st.plotly_chart(figtotalcostft, use_container_width=True,config=config0)     
                                        
                except Exception as e:
                    print(e)
                    st.info(f"error totalcostindicator: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
            else:
                errorstring = ''
                errorstring += 'Input'
                errorstring += ' **Section Names**,' if dfchart['Section'].isnull().values.any() == True else ''
                errorstring += ' **Depth In**,' if dfchart['Din'].isnull().values.any() == True else ''
                errorstring += ' **Depth Out**,' if dfchart['Dout'].isnull().values.any() == True else ''
                errorstring += ' **ROP**,' if dfchart['ROP'].isnull().values.any() == True else ''
                errorstring += ' to display charts.'                    
                st.write(f'	:sparkles: :red[{errorstring}]')         
                                
                # original_title = '<p style="font-family:Courier; color:Blue; font-size: 20px;">Text to display here</p>'
                # st.markdown(original_title, unsafe_allow_html=True) 
            
                            
            # prop_exported = False
            # sendemailflag = 'Please provide your email address.'               
            # bsendemail = False
            # def submit_pptx():
            #     # Validate user entered email
            #     if len(str(sendemail)) > 0:
            #         if sendemail.partition('@')[-1]  == 'ulterra.com':
            #             bsendemail is True
            #         else:
            #             sendemailflag = 'Please use an Ulterra email address.'
            #     else:
            #         sendemailflag = 'Please provide your email address.'
                
                
            #     if bsendemail is True:    
            #         emailerror = st.empty()            
            #         save_bar = st.progress(0, text='Initializing..')                        
            #         df_record.insert(4,'email', sendemail)
                    
            #         with st.spinner('Constructing...'):    
            #             if reporttype == 'PowerPoint':
            #                 if reportsize == 'Landscape':
            #                     try:
            #                         reportppt = create_pptx_Landscape(pagefootage,sendemail,showpartnumber2,showgage2,showbody2)
            #                         try:                
            #                             save_bar.progress(95, text='Preparing Download')
            #                             st.download_button(
            #                                 label="Download PowerPoint",
            #                                 data=reportppt,
            #                                 file_name=f"Adira_Proposal_{customer}_{datetime.datetime.now()}.pptx",
            #                                 mime="application/octet-stream",
            #                             ) 
            #                             save_bar.progress(100, text='Ready')
            #                         except Exception as e:
            #                             st.info(f"error pptx dl1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                             
            #                     except Exception as e:
            #                         st.info(f"error pptx: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                    
                            
            #                 if reportsize == 'Portrait':
            #                     try:
            #                         reportppt = create_pptx_Portrait(pagefootage,sendemail,showpartnumber2,showgage2,showbody2)
            #                         try:
            #                             save_bar.progress(95, text='Preparing Download')
            #                             st.download_button(
            #                                 label="Download PowerPoint",
            #                                 data=reportppt,
            #                                 file_name=f"Adira_Proposal_{customer}_{datetime.datetime.now()}.pptx",
            #                                 mime="application/octet-stream",
            #                             ) 
            #                             save_bar.progress(100, text='Ready')
            #                         except Exception as e:
            #                             st.info(f"error pptx dl1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                             
            #                     except Exception as e:
            #                         st.info(f"error pptx: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
            #     else:
            #         emailerror = st.write(f':red[{sendemailflag}]')
                
                
            
            with st.form(key='template_form2'):
                fc1, fc2 = st.columns(2)
                with fc1:
                    reporttype = st.selectbox('Choose Format',['PowerPoint',],index=0)
                with fc2:
                    reportsize = st.selectbox('Choose Type',['Landscape','Portrait'],index=0)
                    # pagefootage = st.selectbox('Choose Ratio',[2000,1000,5000],index=0)
                    pagefootage = 5
                
                optcol1, optcol2 = st.columns(2,gap="small")
                with optcol1:
                    st.write('Display Options:')
                    showpartnumber2 = st.checkbox("Show Part Number",value=False, key="showpartnumber2")
                    showgage2 = st.checkbox("Show Gage Type",value=True, key="showgage2")  
                    showbody2 = st.checkbox("Show Bit Material",value=True, key="showbody2")  
                    emailerror = st.empty()
                with optcol2:
                    st.write('Export Options:')
                    # showperfchart = st.checkbox("Performance Chart",value=True, key="showperfchart")  
                    # showdvdchart = st.checkbox("Time/Depth Chart",value=True, key="showdvdchart")    
                    # showcostchart = st.checkbox("Value Chart",value=True, key="showcostchart")                         
                    showoffsetchart = st.checkbox("Well Plan Chart",value=True, key="showperfchart")  
                    showbitchart = st.checkbox("Hours vs Depth Chart",value=True, key="showdvdchart")    
                    showsummarychart = st.checkbox("Well Diagram Chart",value=True, key="showcostchart")
                    showpricetotal = st.checkbox("Price Totals",value=True, key="showpricetotal")
                
                fcfooter1, fcfooter2 = st.columns([1,5],gap="small")
                with fcfooter1:
                    sendemail = st.text_input("E-mail address:",placeholder='me@ulterra.com',autocomplete='email',help='Your full e-mail address (jdoe@ulterra.com) to receive the powerpoint.',key='sendemail',label_visibility='collapsed')
                with fcfooter2:
                    with st.spinner('Request...'): 
                        submit_pptx = st.form_submit_button('Request Proposal') # , on_click=submit_pptx)
            
            # format chart to numerics if not
            if st.session_state.bpm is True:
                dfchart = dfchart.apply(pd.to_numeric, errors='ignore') 
                dfchart = dfchart.astype({"Price": float, "Priceft": float}, errors='ignore')
            
            # imported old history foc https://thisismeetpatel.medium.com/read-csv-file-from-azure-blob-storage-to-directly-to-data-frame-using-python-83d34c4cbe57
            # df_oldhistory
            # create proposal record
            df_record = dfchart.copy()
            df_record.insert(0,'triprate', st.session_state.triprate)
            df_record.insert(0,'spreadrate', st.session_state.spreadrate)
            df_record.insert(0,'sectioncount', st.session_state.sectioncount)
            df_record.insert(0,'long', st.session_state.long)
            df_record.insert(0,'lat', st.session_state.lat)
            df_record.insert(0,'targetformation', str(st.session_state.targetformation))
            df_record.insert(0,'rigname', str(st.session_state.rigname))
            df_record.insert(0,'wellname', st.session_state.wellname)
            df_record.insert(0,'datecreated', f"{datetime.datetime.now().isoformat(timespec='milliseconds')}")
            df_record.insert(0,'preparedfor', st.session_state.preparedfor)
            df_record.insert(0,'preparedby', st.session_state.preparedby)
            df_record.insert(0,'operator', st.session_state.customer1)                    
            df_record.insert(0,'district', st.session_state.district)
            df_record.insert(0,'propkey', f"{st.session_state.customer1}_{st.session_state.wellname}_")
            df_record['propkey'] = df_record['propkey'].astype(str) + df_record['Run'].astype(str)
            df_record.insert(0,'pid', f"{int(df_history['pid'].max())+1}")
            
            
            def save_history(df_oldhistory,df_newrecord):
                # https://stackoverflow.com/questions/48226460/efficient-solutions-insert-or-update-row-pandas
                df_newhistory = (pd.concat([df_oldhistory, df_newrecord])
                    .drop_duplicates(['propkey'] , keep='last')
                    .sort_values('propkey' , ascending=False)
                    .reset_index(drop=True))
                # st.write(df_newhistory)
                # if saved, commit df to central csv for recall
                # if prop_exported:
                
                #  ADD A OVERWRITE COUNTER / VERSION NUMBER !!!!!!!!
                #  Test Azure Blob
                with st.spinner('Saving data...'):
                    try:
                        if not df_newhistory.empty: 
                            dbstatus = blob_upload(df_newhistory, 'Proposal_history.csv', True)  # Proposal_history.csv
                            # dbmsg.info(dbstatus)
                    except Exception as e:
                        print(e)
                        st.error(f"error historyupdate: {e}")
                return dbstatus
            
            
            def create_pptx_Portrait(pagefootage,useremail,showpartnumber2,showgage2,showbody2):
                
                save_bar.progress(1, text='Initializing..')
                pptx = 'data/Adira Proposal Portrait.pptx'
                prs = Presentation(pptx)
                
                # declare positional variables
                WIDTH = Inches(8.5)
                HEIGHT = Inches(11)
                
                try:
                    # function to replace text in pptx first slide with selected filters
                    def replace_text(replacements, shapes):
                        """function to replace text on a PowerPoint slide. Takes dict of {match: replacement, ... } and replaces all matches"""
                        for shape in shapes:
                            for match, replacement in replacements.items():
                                if shape.has_text_frame:
                                    if (shape.text.find(match)) != -1:
                                        text_frame = shape.text_frame
                                        for paragraph in text_frame.paragraphs:
                                            whole_text = "".join(run.text for run in paragraph.runs)
                                            whole_text = whole_text.replace(str(match), str(replacement))
                                            for idx, run in enumerate(paragraph.runs):
                                                if idx != 0:
                                                    p = paragraph._p
                                                    p.remove(run._r)
                                            if bool(paragraph.runs):
                                                paragraph.runs[0].text = whole_text

    
                    save_bar.progress(5, text='Creating template..') 
                    lyt=prs.slide_masters[0].slide_layouts[5] # choosing a slide layout
                    lyt9=prs.slide_masters[0].slide_layouts[7] # choosing a slide layout
                    lytTop5=prs.slide_masters[0].slide_layouts[8] # choosing a slide layout
                                    
                    # declare pptx variables
                    first_slide = prs.slides[0]
                    second_slide = prs.slides[2]
                    # fourth_slide = prs.slides[3]
                    shapes_1 = []
                    shapes_2 = []
                    shapes_4 = []
                    index_to_drop = []

                    # create lists with shape objects
                    for shape in first_slide.shapes:
                        shapes_1.append(shape)

                    for shape in second_slide.shapes:
                        shapes_2.append(shape)
                        
                    # for shape in fourth_slide.shapes:
                    #     shapes_4.append(shape)

                    # Add Ulterra Logo
                    first_slide.shapes.add_picture(f'data/customer logos/Ulterra.png', left=Inches(0.97), top=Inches(3.7), height=Inches(0.76)) 
                    line1=first_slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin_x=Inches(4.25), begin_y=Inches(3.35), end_x=Inches(4.25), end_y=Inches(4.4))
                    line1.line.fill.background()
                    line1.line.fill.solid()
                    line1.line.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    line1.width = Pt(3)
                    # https://stackoverflow.com/questions/58792955/changing-format-of-connector-to-an-arrow-one-in-python-pptx
                    # line_elem = connector.line._get_or_add_ln()
                    # line_elem.append(parse_xml("""<a:headEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>"""))                        
                    # # or
                    # line_elem = connector.line._get_or_add_ln()
                    # line_elem.append(parse_xml("""<a:tailEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>"""))

                    # Add Customer Logo
                    try:
                        first_slide.shapes.add_picture(customerlogo, left=Inches(4.86), top=Inches(3.55), height=Inches(1)) 
                    except Exception as e:
                        print(e)
                        st.info(f"error customerlogo: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                    # initiate a dictionary of placeholders and values to replace                        
                    pushwellname = f'Well: {wellname}' if len(wellname) > 1 else ''
                    pushrigname = f'Rig: {rigname}' if len(rigname) > 1 else ''
                    pushcounty = f'{gcounty}, {gstate}' if len(gcounty) > 1 else ''
                    pushlegals = f'{lat}, {long}' if len(lat) > 1 else ''
                    replaces_cover = {
                        '{company}': ' ',
                        '{document}':'PROPOSAL',
                        '{basin}': district,
                        '{year}': f'{datetime.date.today().strftime("%b")} {datetime.date.today().year}',
                        '{wellname}': pushwellname,
                        '{rigname}': pushrigname,
                        '{county}': pushcounty,
                        '{legals}': f'{pushlegals}\n{targetformation}',
                        '{preparedfor}': f'Prepared for:\n{preparedfor}',
                        '{preparedby}': f'Ulterra Rep:\n{preparedby}',
                        }

                    replaces_basin = {
                        '{region}': basins.loc[basins[basins['district'] == district].index,'region'].values[0],
                        '{district}': basins.loc[basins[basins['district'] == district].index,'district'].values[0],
                        '{description}': basins.loc[basins[basins['district'] == district].index,'description'].values[0],
                    }
                    
                    # if len(wellname) > 0 and len(rigname) > 0 and len(str(lat)) > 0 and len(str(long)) > 0:
                    #     replaces_county = {
                    #         '{county}': f'{district} - {wellname} {rigname} @ {lat}, {long}',
                    #     }
                    # elif len(wellname) > 0 and len(rigname) > 0:
                    #     replaces_county = {
                    #         '{county}': f'{district} - {wellname} {rigname}',
                    #     }                            
                    # else:
                    #     replaces_county = {
                    #         '{county}': f'{district} ',
                    #     }
                                            
                    countystring = f'{district}'
                    countystring += f' - {wellname}' if len(wellname) > 0 else ''
                    countystring += f' {rigname}' if len(rigname) > 0 else ''
                    countystring += f' @ {lat}, {long}' if len(str(lat)) > 0 and len(str(long)) > 0 else ''
                                                
                    # replaces_county = {
                    #     '{county}': countystring,
                    # }

                    # run the function to replace placeholders with values
                    replace_text(replaces_cover, shapes_1)
                    replace_text(replaces_basin, shapes_2)
                    # replace_text(replaces_county, shapes_4)
                    
                    
                    prs.slides[2].shapes.add_picture(f'data/maps/map-{district}.png', left=Inches(-0.02), top=Inches(-0.04), height=Inches(7)) 
                    
                except Exception as e:
                    print(e)
                    st.info(f"error ppt intro: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")

                try:
                    # slidelast=prs.slides.add_slide(lyt) # adding a slide   
                    def createbitcard(aslide, pn, lpos, tpos):
                        
                        bitsize = df.loc[df['PartNumber'] == pn,'BitSize'].values[0]
                        bittype = df.loc[df['PartNumber'] == pn,'BitType'].values[0]
                        gaugelen = df.loc[df['PartNumber'] == pn,'GaugeLength'].values[0]
                        gaugetype = df.loc[df['PartNumber'] == pn,'GaugeType'].values[0]
                        
                        shapes2 = aslide.shapes
                        bshape = shapes2.add_shape(MSO_SHAPE.RECTANGLE, left=Inches(lpos), top=Inches(tpos), width=Inches(5), height=Inches(2))
                        # MSO_SHAPE.ROUNDED_RECTANGLE
                        # try:
                        #     # https://python-pptx.readthedocs.io/en/latest/user/autoshapes.html#adjusting-an-autoshape
                        #     adjs = bshape.adjustments
                            
                        #     # # Size & Type
                        #     # txBox = aslide.shapes.add_textbox(left=Inches(0), top=Inches(0), width=Inches(2), height=Inches(0.22))            
                        #     # tf = txBox.text_frame
                        #     # tf.text = f"{bshape.adjustments}"
                        #     # tf.paragraphs[0].font.size = Pt(18)
                            
                        #     adjs[1].effective_value = 0.1                                       
                        # except Exception as e:
                        #     print(e)
                        #     # st.info(f"error shaperound: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                    
                        bshape.fill.solid()
                        # bshape.fill.fore_color.rgb = RGBColor(212, 230, 232)
                        bshape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        bshape.line.fill.background()
                        # D4E6E8
                        # set lighter or darker, -0.2 is 20% darker, 0.4 is 40% lighter
                        # shape.fill.fore_color.brightness = 0.4
                        # set fill to transparent (no fill)
                        # shape.fill.background()
                        
                        # if os.path.exists(f'data/Bit Pictures/{pn}-1.jpg'):
                        #     aslide.shapes.add_picture(f'data/Bit Pictures/{pn}-1.jpg', left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2))                            
                        #     # image = f'data/Bit Pictures/{pn}-1.jpg'
                        # elif os.path.exists(f'data/Bit Pictures/{pn}-1.JPG'):
                        #     aslide.shapes.add_picture(f'data/Bit Pictures/{pn}-1.JPG', left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2))                           
                        #     # image = f'data/Bit Pictures/{pn}-1.JPG'
                        # else:                           
                        #     # image = f'data/Bit Pictures/blank.jpg'
                        #     aslide.shapes.add_picture(f'data/Bit Pictures/blank.jpg', left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2)) 
                        # # aslide.shapes.add_picture(image, left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2)) 
                                                                                
                        if os.path.exists(f'{pn}-1.jpg'):
                            image = f'{pn}-1.jpg'       
                        elif os.path.exists(f'data/Bit Pictures/{pn}-1.jpg'):                                
                            image = f'data/Bit Pictures/{pn}-1.jpg' 
                        elif os.path.exists(f'data/Bit Pictures/{pn}-1.JPG'):                                
                            image = f'data/Bit Pictures/{pn}-1.JPG'
                        else:
                            image = f'data/Bit Pictures/blank.jpg'     
                            
                        imagepng = f'{pn}-1.png'
                        convert_png_transparent(image, imagepng)
                        aslide.shapes.add_picture(imagepng, left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2)) 
                        # st.image(imagepng,caption=f"{df.loc[df[df['PartNumber'] == pn].index, 'BitSize'].values[0]}  - {df.loc[df[df['PartNumber'] == bit].index, 'BitType'].values[0]}",width=150)
                        
                        
                        # Size & Type
                        txBox = aslide.shapes.add_textbox(left=Inches(lpos - 0.43), top=Inches(tpos-0.4), width=Inches(2), height=Inches(0.4))            
                        tf = txBox.text_frame
                        # bitsize = df.loc[df['PartNumber'] == pn,'BitSize']
                        tf.text = f"{df.loc[df['PartNumber'] == pn,'BitSize'].values[0]} {bittype}"
                        tf.paragraphs[0].font.size = Pt(18)
                        txBox.fill.solid()
                        txBox.fill.fore_color.rgb = RGBColor(0, 147, 159)                    
                        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  
                        
                        #  Features
                        if showgage2:
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+2.11), top=Inches(tpos), width=Inches(1.9), height=Inches(1))            
                            tf = txBox.text_frame
                            if gaugetype.lower().find('step') > 0 and gaugetype.lower().find('taper') > 0:
                                gaugetypesimple = 'Stepped & Tapered'
                            elif gaugetype.lower().find('step') > 0:
                                gaugetypesimple = 'Stepped'
                            elif gaugetype.lower().find('taper') > 0:
                                gaugetypesimple = 'Tapered'
                            else:
                                gaugetypesimple = 'Nominal'                                        
                            
                            feat1 = f"{df.loc[df['PartNumber'] == pn,'PartNumber'].values[0]}\v" if showpartnumber2 else ''
                            feat2 = f"{gaugelen} in Gauge\v" if showgage2 else ''
                            feat3 = f"{gaugetypesimple}\v" if showgage2 else ''
                            feat4 = f"{df.loc[df['PartNumber'] == pn,'BitMaterial'].values[0]} Body\v" if showbody2 else ''
                            # feat5 = f"{df.loc[df['PartNumber'] == pn,'PerformancePackage'].values[0]}"
                            tf.text = f"{feat1}{feat2}{feat3}{feat4}"                            
                            # tf.text = f"{gaugelen} in Gauge \v{gaugetypesimple}"
                            tf.paragraphs[0].font.size = Pt(14)
                        
                        # Technology Logo                            
                        if bittype.find('CF') >= 0:
                            techpic = 'CounterForce'
                        elif bittype.find('SPL') >= 0:
                            techpic = 'SplitBlade'
                        elif bittype.find('RPS') >= 0:
                            techpic = 'RipSaw'
                        elif bittype.find('WAV') >= 0:
                            techpic = 'WaveCut'
                        elif bittype.find('AIR') >= 0:
                            techpic = 'AirRaid'
                        elif bittype.find('XP') >= 0:
                            techpic = 'XP'
                        else:
                            techpic = None  
                        
                        if techpic is not None:
                            aslide.shapes.add_picture(f'data/tech logos/{techpic}_logo.png', left=Inches(lpos+1.72), top=Inches(tpos-.41), height=Inches(0.4))  
                            aslide.shapes.add_picture(f'data/tech logos/ad_{techpic}.png', left=Inches(2.61), top=Inches(8.64), height=Inches(1.75)) 
                            
                        # Price
                        if not pd.isna(dfchart.loc[i,'Price']) or not pd.isna(dfchart.loc[i,'Priceft']):
                            # Total Price
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+2.11), top=Inches(tpos+0.8), width=Inches(1.69), height=Inches(0.4))            
                            tf = txBox.text_frame
                            if float(dfchart.loc[i,'Priceft']) > 0:
                                tf.text = f"${float(dfchart.loc[i,'Pricetotal']):,.0f} est."                   
                            else:
                                tf.text = f"${float(dfchart.loc[i,'Pricetotal']):,.0f}"
                            tf.paragraphs[0].font.size = Pt(20)
                            txBox.fill.solid()
                            txBox.fill.fore_color.rgb = RGBColor(242, 242, 242)
                            # Price Breakdown
                            if not pd.isna(dfchart.loc[i,'Priceft'])  and float(dfchart.loc[i,'Priceft']) > 0:                            
                                txBox = aslide.shapes.add_textbox(left=Inches(lpos+3.89), top=Inches(tpos+0.8), width=Inches(1.36), height=Inches(0.44))            
                                tf = txBox.text_frame
                                tf.text = f"Flat: ${float(dfchart.loc[i,'Price']):,} \n$/Ft: {float(dfchart.loc[i,'Priceft']):,} @ {int(dfchart.loc[i,'Drilled']):,} ft" 
                                for paragraph in tf.paragraphs:
                                    paragraph.font.size = Pt(10)
                        
                        # DBR
                        if not pd.isna(dfchart.loc[i,'DBR']):                                
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+2.11), top=Inches(tpos+1.17), width=Inches(3.5), height=Inches(0.4))            
                            tf = txBox.text_frame
                            tf.text = f"DBR: {dfchart.loc[i,'DBR']}"  
                            tf.paragraphs[0].font.size = Pt(12)
                            # txBox.fill.solid()
                            # txBox.fill.fore_color.rgb = RGBColor(242, 242, 242)
                            
                        # Comments
                        if not pd.isna(dfchart.loc[i,'Comment']):                                
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+2.11), top=Inches(tpos+1.47), width=Inches(3.5), height=Inches(0.54))            
                            tf = txBox.text_frame
                            tf.text = f"Comment: \n{dfchart.loc[i,'Comment']}"  
                            for paragraph in tf.paragraphs:
                                paragraph.font.size = Pt(12) 
                            tf.paragraphs[0].font.size = Pt(12)
                            tf.paragraphs[0].font.bold = True
                        
                        # Backup
                        if not pd.isna(dfchart.loc[i,'Backup']):                                  
                            backuppn = dfchart.loc[i,'Backup']
                            backuptype = dfchart.loc[i,'Backuptype']
                            # if os.path.exists(f'data/Bit Pictures/{backuppn}-1.jpg'):
                            #     aslide.shapes.add_picture(f'data/Bit Pictures/{backuppn}-1.jpg', left=Inches(lpos+5.28), top=Inches(tpos+1.13), width=Inches(1))                            
                            #     # image = f'data/Bit Pictures/{pn}-1.jpg'
                            # elif os.path.exists(f'data/Bit Pictures/{backuppn}-1.JPG'):
                            #     aslide.shapes.add_picture(f'data/Bit Pictures/{backuppn}-1.JPG', left=Inches(lpos+5.28), top=Inches(tpos+1.13), width=Inches(1))                           
                            #     # image = f'data/Bit Pictures/{pn}-1.JPG'
                            # else:                           
                            #     # image = f'data/Bit Pictures/blank.jpg'
                            #     aslide.shapes.add_picture(f'data/Bit Pictures/blank.jpg', left=Inches(lpos+5.28), top=Inches(tpos+1.13), width=Inches(1)) 
                            # # aslide.shapes.add_picture(image, left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2)) 
                                        
                            if os.path.exists(f'{backuppn}-1.jpg'):
                                image = f'{backuppn}-1.jpg'       
                            elif os.path.exists(f'data/Bit Pictures/{backuppn}-1.jpg'):                                
                                image = f'data/Bit Pictures/{backuppn}-1.jpg' 
                            elif os.path.exists(f'data/Bit Pictures/{backuppn}-1.JPG'):                                
                                image = f'data/Bit Pictures/{backuppn}-1.JPG'
                            else:
                                image = f'data/Bit Pictures/blank.jpg'     
                                
                            imagepng = f'{backuppn}-1.png'
                            convert_png_transparent(image, imagepng)
                            aslide.shapes.add_picture(imagepng, left=Inches(lpos+5.28), top=Inches(tpos+0.6), width=Inches(1.5)) # 6.32" 2.99
                            
                            # Backup Size & Type
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+4.86), top=Inches(tpos+0.26), width=Inches(1.5), height=Inches(0.3))     # 5.9" 2.67       
                            tf = txBox.text_frame
                            # bitsize = df.loc[df['PartNumber'] == pn,'BitSize']
                            tf.text = f"Backup: {backuptype}"
                            tf.paragraphs[0].font.size = Pt(12)
                            txBox.fill.solid()
                            txBox.fill.fore_color.rgb = RGBColor(63, 112, 119)                    
                            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  
                        
                        # # # Status / Price Notes
                        # txBox = aslide.shapes.add_textbox(left=Inches(lpos+0.1), top=Inches(tpos+4.06), width=Inches(1.9), height=Inches(0.4))            
                        # tf = txBox.text_frame
                        # tf.text = f"{'Backup' if df.loc[df['PartNumber'] == pn,'Backup'].values[0] == 'True' else ''}"
                        # tf.paragraphs[0].font.size = Pt(10)
                        
                    
                    def createappcard(aslide, pn, lpos, tpos):                                                   
                        # aslide.shapes.add_picture(f'images/icon_interval.png', left=Inches(3.15), top=Inches(2.28), height=Inches(2.49))                              
                        filename = "figsec.png"
                        pio.write_image(figsec[i], filename, scale=1, width=284, height=604) 
                        # check pixel dimensions
                        placeholder = aslide.shapes.add_picture(filename, left=Inches(lpos), top=Inches(tpos), height=Inches(3.15))
                        
                        # Interval Definition                            
                        if len(str(dfchart.loc[i,'Din'])) > 1 and len(str(dfchart.loc[i,'Dout'])) > 1 and len(str(dfchart.loc[i,'Drilled'])) > 1:
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos-1.68), top=Inches(tpos+0.46), width=Inches(1.67), height=Inches(0.34))
                            tf = txBox.text_frame
                            # text_frame.word_wrap = False
                            tf.text = f'Interval Definition:'
                            tf.paragraphs[0].font.bold = True
                            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                            # .italic = None 
                            tf.paragraphs[0].font.size = Pt(14)
                            # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                            # Interval Info
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos-1.68), top=Inches(tpos+0.72), width=Inches(1.67), height=Inches(0.81))
                            tf = txBox.text_frame
                            tf.text = f"Depth In: {float(dfchart.loc[i,'Din']):,.0f}\n Depth Out: {float(dfchart.loc[i,'Dout']):,.0f}\n Footage: {float(dfchart.loc[i,'Drilled']):,.0f}"
                            for paragraph in tf.paragraphs:
                                paragraph.font.size = Pt(14)
                                paragraph.alignment = PP_ALIGN.RIGHT
                            # paragraph.font.color.rgb = RGBColor(0, 147, 159)
                        # tf.paragraphs[0].font.size = Pt(14)
                        # tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                                                
                        # Parameters
                        # ***********ADD VALUE DETECTION TO PROCESS OUTPUTS BELOW. DO THIS FOR ALL TEXT OUTPUTS
                        if float(dfchart.loc[i,'WOBout']) > 0:
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos-1.68), top=Inches(tpos+1.45), width=Inches(1.67), height=Inches(0.34))
                            tf = txBox.text_frame
                            # text_frame.word_wrap = False
                            tf.text = f'Parameters:'
                            tf.paragraphs[0].font.bold = True
                            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                            # .italic = None 
                            tf.paragraphs[0].font.size = Pt(14)
                            # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                            # Parameters Info
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos-1.68), top=Inches(tpos+1.71), width=Inches(1.67), height=Inches(0.81))
                            tf = txBox.text_frame
                            tf.text = f"WOB: {float(dfchart.loc[i,'WOBin']):,.0f}/{float(dfchart.loc[i,'WOBout']):,.0f}\n RPM: {float(dfchart.loc[i,'RPMin']):,.0f}/{float(dfchart.loc[i,'RPMax']):,.0f}\n GPM: {dfchart.loc[i,'Flowrate']}"
                            for paragraph in tf.paragraphs:
                                paragraph.font.size = Pt(14)
                                paragraph.alignment = PP_ALIGN.RIGHT
                            # tf.paragraphs[0].font.size = Pt(14)
                            # tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                            # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                        
                            
                        # Performance
                        txBox = aslide.shapes.add_textbox(left=Inches(lpos+1.4), top=Inches(tpos+0.46), width=Inches(1.25), height=Inches(0.34))
                        tf = txBox.text_frame
                        # text_frame.word_wrap = False
                        tf.text = f'Performance:'
                        tf.paragraphs[0].font.bold = True
                        # .italic = None 
                        tf.paragraphs[0].font.size = Pt(14)
                        # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                        # Performance Info
                        txBox = aslide.shapes.add_textbox(left=Inches(lpos+1.4), top=Inches(tpos+0.72), width=Inches(2.3), height=Inches(0.81))
                        tf = txBox.text_frame
                        tf.text = f"{(dfchart.loc[i,'Drilled'] / dfchart.loc[i,'ROP']):,.1f} Hrs @ {dfchart.loc[i,'ROP']:,.1f} ft/hr"
                        tf.paragraphs[0].font.size = Pt(14)
                        # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                        
                            
                        # Motor Recs
                        if len(str(dfchart.loc[i,'Motorspecs'])) > 1:
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+1.4), top=Inches(tpos+2.18), width=Inches(2.27), height=Inches(0.34))
                            tf = txBox.text_frame
                            # text_frame.word_wrap = False
                            tf.text = f'Motor Recommendations:'
                            tf.paragraphs[0].font.bold = True
                            # .italic = None 
                            tf.paragraphs[0].font.size = Pt(14)
                            # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                            # Motor Info
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+1.94), top=Inches(tpos+2.62), width=Inches(2.34), height=Inches(0.34))
                            tf = txBox.text_frame
                            tf.text = f"{dfchart.loc[i,'Motorspecs']}"
                            tf.paragraphs[0].font.size = Pt(14)
                            aslide.shapes.add_picture(f'images/icon_revgal.png', left=Inches(lpos+1.54), top=Inches(tpos+2.49), height=Inches(0.4))     
                    
                        
                    
                    def createsectionpage(section,df):
                        newslide=prs.slides.add_slide(lyt9) # adding a slide              
                        # Section Name
                        txBox = newslide.shapes.add_textbox(left=Inches(0.61), top=Inches(0.24), width=Inches(3), height=Inches(0.33))
                        tf = txBox.text_frame
                        tf.text = f'Bit Run {i} {section}'
                        tf.paragraphs[0].font.size = Pt(24)
                        tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                        # subtitle
                        txBox = newslide.shapes.add_textbox(left=Inches(0.61), top=Inches(0.67), width=Inches(5), height=Inches(0.33))
                        tf = txBox.text_frame
                        tf.text = f'BIT AVAILABILITY AND PRICING'
                        tf.paragraphs[0].font.size = Pt(22)
                        # description             
                        # txBox = newslide.shapes.add_textbox(left=Inches(0.61), top=Inches(1.1), width=Inches(5), height=Inches(1.35))            
                        # tf = txBox.text_frame
                        # tf.text = f"We combine our company's unique product design with a rapid \v prototyping manufacturing process, leading to reliable and \v high-end performance drill bits customers can count on in \v the {district}. The proof is in the numbers."
                        # tf.paragraphs[0].font.size = Pt(12)
                        
                        bitcount = 0
                        start_h = 1.04
                        start_v = 2.41
                        # createbitcard(newslide,dfchart.loc[i,'Bit'],(start_h + (2 * (bitcount-1))+ (0.25 * (bitcount-1))),start_v)
                        createbitcard(newslide,dfchart.loc[i,'Bit'],start_h,start_v)
                                                    
                        # Write section chart
                        # filename = "figproposal.png"
                        # pio.write_image(figproposal, filename, scale=2, width=1746, height=1000) 
                        # check pixel dimensions
                        # placeholder = newslide.shapes.add_picture(filename, left=Inches(1.33), top=Inches(2.23), width=Inches(4), height=Inches(4)) 
                                                            
                        # start_h = 9.05
                        # start_v = 1.68
                        start_h = 3.51 # 1.04
                        start_v = 4.95 # 6.8
                        createappcard(newslide,dfchart.loc[i,'Bit'],start_h,start_v)
                        
                        # # newslide.shapes.add_picture(f'images/icon_interval.png', left=Inches(3.15), top=Inches(2.28), height=Inches(2.49))                              
                        # filename = "figsec.png"
                        # pio.write_image(figsec[i], filename, scale=2, width=284, height=604) 
                        # # check pixel dimensions
                        # placeholder = newslide.shapes.add_picture(filename, left=Inches(3.51), top=Inches(2.32), height=Inches(3.15))
                        
                        
                        # # Interval Definition                            
                        # if str(dfchart.loc[i,'Din']).isnumeric() and str(dfchart.loc[i,'Dout']).isnumeric() and str(dfchart.loc[i,'Drilled']).isnumeric():
                        #     txBox = newslide.shapes.add_textbox(left=Inches(1.83), top=Inches(2.78), width=Inches(1.67), height=Inches(0.34))
                        #     tf = txBox.text_frame
                        #     # text_frame.word_wrap = False
                        #     tf.text = f'Interval Definition:'
                        #     tf.paragraphs[0].font.bold = True
                        #     tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        #     # .italic = None 
                        #     tf.paragraphs[0].font.size = Pt(14)
                        #     # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                        #     # Interval Info
                        #     txBox = newslide.shapes.add_textbox(left=Inches(1.99), top=Inches(3.04), width=Inches(1.52), height=Inches(0.81))
                        #     tf = txBox.text_frame
                        #     tf.text = f"Depth In: {float(dfchart.loc[i,'Din']):,.0f}\n Depth Out: {float(dfchart.loc[i,'Dout']):,.0f}\n Footage: {float(dfchart.loc[i,'Drilled']):,.0f}"
                        #     for paragraph in tf.paragraphs:
                        #         paragraph.font.size = Pt(14)
                        #         paragraph.alignment = PP_ALIGN.RIGHT
                        #     # paragraph.font.color.rgb = RGBColor(0, 147, 159)
                        # # tf.paragraphs[0].font.size = Pt(14)
                        # # tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        # # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                                                
                        # # Parameters
                        # # ***********ADD VALUE DETECTION TO PROCESS OUTPUTS BELOW. DO THIS FOR ALL TEXT OUTPUTS
                        # if float(dfchart.loc[i,'WOBout']) > 0:
                        #     txBox = newslide.shapes.add_textbox(left=Inches(1.83), top=Inches(3.77), width=Inches(1.67), height=Inches(0.34))
                        #     tf = txBox.text_frame
                        #     # text_frame.word_wrap = False
                        #     tf.text = f'Parameters:'
                        #     tf.paragraphs[0].font.bold = True
                        #     tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        #     # .italic = None 
                        #     tf.paragraphs[0].font.size = Pt(14)
                        #     # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                        #     # Parameters Info
                        #     txBox = newslide.shapes.add_textbox(left=Inches(1.99), top=Inches(4.03), width=Inches(1.52), height=Inches(0.81))
                        #     tf = txBox.text_frame
                        #     tf.text = f"WOB: {float(dfchart.loc[i,'WOBin']):,.0f}/{float(dfchart.loc[i,'WOBout']):,.0f}\n RPM: {float(dfchart.loc[i,'RPMin']):,.0f}/{float(dfchart.loc[i,'RPMax']):,.0f}\n GPM: {dfchart.loc[i,'Flowrate']}"
                        #     for paragraph in tf.paragraphs:
                        #         paragraph.font.size = Pt(14)
                        #         paragraph.alignment = PP_ALIGN.RIGHT
                        #     # tf.paragraphs[0].font.size = Pt(14)
                        #     # tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        #     # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                        
                            
                        # # Performance
                        # txBox = newslide.shapes.add_textbox(left=Inches(4.91), top=Inches(2.78), width=Inches(1.25), height=Inches(0.34))
                        # tf = txBox.text_frame
                        # # text_frame.word_wrap = False
                        # tf.text = f'Performance:'
                        # tf.paragraphs[0].font.bold = True
                        # # .italic = None 
                        # tf.paragraphs[0].font.size = Pt(14)
                        # # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                        # # Performance Info
                        # txBox = newslide.shapes.add_textbox(left=Inches(4.98), top=Inches(3.04), width=Inches(2.3), height=Inches(0.81))
                        # tf = txBox.text_frame
                        # tf.text = f"{(dfchart.loc[i,'Drilled'] / dfchart.loc[i,'ROP']):,.1f} Hrs @ {dfchart.loc[i,'ROP']:,.1f} ft/hr"
                        # tf.paragraphs[0].font.size = Pt(14)
                        # # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                        
                            
                        # # Motor Recs
                        # if len(str(dfchart.loc[i,'Motorspecs'])) > 1:
                        #     txBox = newslide.shapes.add_textbox(left=Inches(4.91), top=Inches(4.13), width=Inches(2.27), height=Inches(0.34))
                        #     tf = txBox.text_frame
                        #     # text_frame.word_wrap = False
                        #     tf.text = f'Motor Recommendations:'
                        #     tf.paragraphs[0].font.bold = True
                        #     # .italic = None 
                        #     tf.paragraphs[0].font.size = Pt(14)
                        #     # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                        #     # Motor Info
                        #     txBox = newslide.shapes.add_textbox(left=Inches(5.62), top=Inches(4.51), width=Inches(2.34), height=Inches(0.34))
                        #     tf = txBox.text_frame
                        #     tf.text = f"{dfchart.loc[i,'Motorspecs']}"
                        #     tf.paragraphs[0].font.size = Pt(14)
                        #     newslide.shapes.add_picture(f'images/icon_revgal.png', left=Inches(5.22), top=Inches(4.48), height=Inches(0.4))     
                    
                    
                    def createspecsheetpage(section,i):                                                 
                        # PRIMARY BIT SPEC SHEET IMPORT
                        pn = dfchart.loc[i,'Bit']
                        if len(str(pn)) == 6:                         
                            try:         
                                newslide=prs.slides.add_slide(lyt) # adding a slide 
                                bitsize = df.loc[df['PartNumber'] == pn,'BitSize'].values[0]
                                bittype = df.loc[df['PartNumber'] == pn,'BitType'].values[0]
                                bitfeatures = df.loc[df['PartNumber'] == pn,'Features'].values[0]                     
                                # Primary Spec                                
                                specpdfname = f'{pn} - {float(str(bitsize)):.3f} - {bittype} - {bitfeatures}.pdf'
                                if os.path.isfile(specpdfname):
                                    pages = convert_from_path(specpdfname, 500)
                                    for count, page in enumerate(pages):
                                        page.save(f'{pn}-{count}.jpg', 'JPEG')
                                                                    
                                    filename = f'{pn}-0.jpg'
                                else:                                        
                                    filename = f'data/Bit Pictures/blank.jpg'
                                    
                                placeholder = newslide.shapes.add_picture(filename, left=Inches(0), top=Inches(0), height=Inches(11))                                    
                                
                                # Cover part number
                                if showpartnumber2 is False:
                                    shapes2 = newslide.shapes
                                    bshape = shapes2.add_shape(MSO_SHAPE.RECTANGLE, left=Inches(0.23), top=Inches(1.68), width=Inches(1.46), height=Inches(0.12))                 
                                    bshape.fill.solid()
                                    bshape.fill.fore_color.rgb = RGBColor(81, 81, 81)
                                    bshape.line.fill.background()
                                    
                                shapes2 = newslide.shapes
                                bshape = shapes2.add_shape(MSO_SHAPE.RECTANGLE, left=Inches(0), top=Inches(0), width=Inches(4.35), height=Inches(0.96))                    
                                bshape.fill.solid()
                                bshape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                                bshape.line.fill.background()    
                                        
                                # Section Name
                                txBox = newslide.shapes.add_textbox(left=Inches(0), top=Inches(0.08), width=Inches(3), height=Inches(0.33))
                                tf = txBox.text_frame
                                tf.text = f'Bit Run {i}: {section} Spec Sheets'
                                tf.paragraphs[0].font.size = Pt(24)
                                tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)  
                                
                                txBox = newslide.shapes.add_textbox(left=Inches(0), top=Inches(0.47), width=Inches(3), height=Inches(0.33))
                                tf = txBox.text_frame
                                tf.text = f'Primary:'
                                tf.paragraphs[0].font.size = Pt(16)                         
                                                                
                            except Exception as e:
                                print(e)
                                st.info(f"error ziptest: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                        
                                                    
                        # BACKUP BIT SPEC SHEET IMPORT
                        pn = dfchart.loc[i,'Backup']                            
                        if len(str(pn)) == 6:                
                            try:
                                newslide=prs.slides.add_slide(lyt) # adding a slide    
                                bitsize = df.loc[df['PartNumber'] == pn,'BitSize'].values[0]
                                bittype = df.loc[df['PartNumber'] == pn,'BitType'].values[0]
                                bitfeatures = df.loc[df['PartNumber'] == pn,'Features'].values[0]                     
                                # BACKUP Spec                                
                                specpdfname = f'{pn} - {float(str(bitsize)):.3f} - {bittype} - {bitfeatures}.pdf'
                                if os.path.isfile(specpdfname):
                                    pages = convert_from_path(specpdfname, 500)
                                    for count, page in enumerate(pages):
                                        page.save(f'{pn}-{count}.jpg', 'JPEG')
                                                                    
                                    filename = f'{pn}-0.jpg'
                                else:                                        
                                    filename = f'data/Bit Pictures/blank.jpg'
                                    
                                placeholder = newslide.shapes.add_picture(filename, left=Inches(0), top=Inches(0), height=Inches(11))                                    
                                    
                                # Cover part number
                                if showpartnumber2 is False:
                                    shapes2 = newslide.shapes
                                    bshape = shapes2.add_shape(MSO_SHAPE.RECTANGLE, left=Inches(0.23), top=Inches(1.68), width=Inches(1.46), height=Inches(0.12))                 
                                    bshape.fill.solid()
                                    bshape.fill.fore_color.rgb = RGBColor(81, 81, 81)
                                    bshape.line.fill.background()
                                    
                                shapes2 = newslide.shapes
                                bshape = shapes2.add_shape(MSO_SHAPE.RECTANGLE, left=Inches(0), top=Inches(0), width=Inches(4.35), height=Inches(0.96))                    
                                bshape.fill.solid()
                                bshape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                                bshape.line.fill.background()    
                                        
                                # Section Name
                                txBox = newslide.shapes.add_textbox(left=Inches(0), top=Inches(0.08), width=Inches(3), height=Inches(0.33))
                                tf = txBox.text_frame
                                tf.text = f'Bit Run {i}: {section} Spec Sheets'
                                tf.paragraphs[0].font.size = Pt(24)
                                tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)  
                                
                                txBox = newslide.shapes.add_textbox(left=Inches(0), top=Inches(0.47), width=Inches(3), height=Inches(0.33))
                                tf = txBox.text_frame
                                tf.text = f'Backup:'
                                tf.paragraphs[0].font.size = Pt(16)                         
                                                                
                            except Exception as e:
                                print(e)
                                st.info(f"error ziptest: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                            
                    def createsummarytables(aslide, i, lpos, tpos):                                                        
                        # bitsize = df.loc[df['PartNumber'] == pn,'BitSize'].values[0]
                        # bittype = df.loc[df['PartNumber'] == pn,'BitType'].values[0]
                        # gaugelen = df.loc[df['PartNumber'] == pn,'GaugeLength'].values[0]
                        # gaugetype = df.loc[df['PartNumber'] == pn,'GaugeType'].values[0]                                                        
                        # # Size & Type
                        # txBox = aslide.shapes.add_textbox(left=Inches(lpos - 0.43), top=Inches(tpos-0.4), width=Inches(2), height=Inches(0.4)) 
                        
                        # Data Table https://python-pptx.readthedocs.io/en/latest/user/table.html 
                        try:
                            shape = newslide.shapes.add_table(4, 12, Inches(lpos), Inches(tpos+(0.43*i)), Inches(6.8), Inches(0.5))
                            table = shape.table
                            # https://stackoverflow.com/questions/61982333/how-to-change-default-table-style-using-pptx-python
                            # style_id lists https://github.com/scanny/python-pptx/issues/27#issuecomment-263076372                                
                            # MediumStyle1 = '{793D81CF-94F2-401A-BA57-92F5A7B2D0C5}'
                            # MediumStyle1Accent1 = '{B301B821-A1FF-4177-AEE7-76D212191A09}'
                            # MediumStyle1Accent2 = '{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}'
                            # MediumStyle1Accent3 = '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}'
                            # MediumStyle1Accent4 = '{1E171933-4619-4E11-9A3F-F7608DF75F80}'
                            # MediumStyle1Accent5 = '{FABFCF23-3B69-468F-B69F-88F6DE6A72F2}'
                            # MediumStyle1Accent6 = '{10A1B5D5-9B99-4C35-A422-299274C87663}'
                            tbl =  shape._element.graphic.graphicData.tbl
                            if 'Surface' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}'
                            elif 'Vertical' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{B301B821-A1FF-4177-AEE7-76D212191A09}'
                            elif 'Drill Out' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{FABFCF23-3B69-468F-B69F-88F6DE6A72F2}'
                            elif 'Intermediate' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{1E171933-4619-4E11-9A3F-F7608DF75F80}'
                            elif 'Curve' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{793D81CF-94F2-401A-BA57-92F5A7B2D0C5}'
                            elif 'Lateral' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}'                                
                            else:
                                style_id = '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}'
                            tbl[0][-1].text = style_id
                            
                            def setcell(r,c,isbold,celltext):
                                cell = table.cell(r, c)
                                cell.text = celltext                    
                                cell.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                                cell.margin_left = cell.margin_right = cell.margin_top = cell.margin_bottom = 0                                    
                                for paragraph in cell.text_frame.paragraphs:
                                    # paragraph = cell.text_frame.paragraphs[0]
                                    paragraph.font.size = Pt(8)
                                    paragraph.font.bold = isbold
                                    # paragraph.alignment = PP_ALIGN.CENTER
                                    # paragraph.alignment = PP_ALIGN.RIGHT
                                
                            def mergecell(r1,c1,r2,c2):
                                table.cell(r1, c1).merge(table.cell(r2, c2))
                                
                            # setcell(0,1,str(dfchart.iloc[i][j]))
                            setcell(0,0,True,'Bit Run')
                            setcell(0,1,True,'Size')
                            setcell(0,2,True,'Type')
                            setcell(0,3,True,'BackUp')
                            setcell(0,4,True,'Depth In')
                            setcell(0,5,True,'Depth Out')
                            setcell(0,6,True,'FTG')
                            setcell(0,7,True,'HRS')
                            setcell(0,8,True,'ROP')
                            setcell(0,9,True,'WOB')
                            setcell(0,10,True,'RPM')
                            setcell(0,11,True,'GPM')
                            
                            mergecell(1,0,3,0) # setcell(2,1, ) #'Section')
                            setcell(1,0,False,f"Bit Run-{i+1}\n{str(dfchart.iloc[i]['Section'])}") #'Section')
                            setcell(1,1,False,str(dfchart.iloc[i]['Size'])) #'Size')
                            setcell(1,2,False,str(dfchart.iloc[i]['Type'])) #'Type')
                            setcell(1,3,False,str(dfchart.iloc[i]['Backuptype'])) #'BackUp')
                            setcell(1,4,False,str(dfchart.iloc[i]['Din'])) #'Depth In')
                            setcell(1,5,False,str(dfchart.iloc[i]['Dout'])) #'Depth Out')
                            setcell(1,6,False,str(dfchart.iloc[i]['Drilled'])) #'FTG')
                            setcell(1,7,False,f"{(float(dfchart.iloc[i]['Drilled']) / float(dfchart.iloc[i]['ROP'])):,.1f}") #'HRS')
                            setcell(1,8,False,str(dfchart.iloc[i]['ROP'])) #'ROP')
                            setcell(1,9,False,f"{dfchart.iloc[i]['WOBin']}/{dfchart.iloc[i]['WOBout']}") #'WOB')
                            setcell(1,10,False,f"{dfchart.iloc[i]['RPMin']}/{dfchart.iloc[i]['RPMax']}") #'RPM')
                            setcell(1,11,False,str(dfchart.iloc[i]['Flowrate'])) #'GPM')
                            
                            setcell(2,1,True,'Cost Details') #Size')
                            mergecell(2,1,2,2) 
                            setcell(2,3,True,'Rotary:') #BackUp')
                            setcell(2,4,False,f"{float(dfchart.iloc[i]['Price']):,.0f}") #Depth In')
                            setcell(2,5,True,'$/Ft:') #Depth Out')
                            setcell(2,6,False,f"{float(dfchart.iloc[i]['Priceft']):,.0f}") #FTG')
                            setcell(2,7,True,'Total Cost:') #HRS')
                            mergecell(2,7,2,8) 
                            setcell(2,9,False,f"{float(dfchart.iloc[i]['Pricetotal']):,.0f}") #ROP')
                            setcell(2,10,True,'DBR/LIH:') #WOB')
                            setcell(2,11,False,str(dfchart.iloc[i]['DBR'])) #RPM')
                            # mergecell(2,10,2,11) # setcell(2,12, ) #GPM') 
                            
                            setcell(3,1,True,'Motors:') 
                            setcell(3,2,False,str(dfchart.iloc[i]['Motorspecs']))                                 
                            mergecell(3,2,3,4) 
                            setcell(3,5,True,'Comments:') 
                            setcell(3,6,False,str(dfchart.iloc[i]['Comment']))                                
                            mergecell(3,6,3,11) 
                            
                            cols = table.columns
                            cols[0].width = Inches(0.69) # section
                            cols[1].width = Inches(0.47) # size
                            cols[2].width = Inches(0.67) # type
                            cols[3].width = Inches(0.62) # backup
                            cols[4].width = Inches(0.52) # din
                            cols[5].width = Inches(0.52) # dout
                            cols[6].width = Inches(0.45) # ftg
                            cols[7].width = Inches(0.45) # hrs
                            cols[8].width = Inches(0.47) # rop
                            cols[9].width = Inches(0.5) # wob
                            cols[10].width = Inches(0.73) # rpm
                            cols[11].width = Inches(0.77) # gpm
                            
                            def remove_row(table, row):
                                tbl = table._tbl
                                tr = row._tr
                                tbl.remove(tr)
                            
                            row = table.rows[0]
                            remove_row(table, row)                                
                                
                        except Exception as e:
                            print(e)
                            st.info(f"error ppt-table: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                        
                        
                    # dfWell['Section','Size''Bit','Backup','Din','Dout','WOBin','WOBout','ROP','RPMin','RPMax','Motorspecs','Flowrate','Price','Priceft','DBR','Comment']   
                    for i in range(1,(len(dfchart)+1)):                            
                        save_bar.progress(25+(2*i), text=f'Creating run page {i}..') 
                        createsectionpage(dfchart.loc[i,'Section'],i)
                    
                    
                    save_bar.progress(50, text='Creating Overview..') 
                    # Create Overview page
                    newslide=prs.slides.add_slide(lyt9) # adding a slide              
                    # Section Name
                    txBox = newslide.shapes.add_textbox(left=Inches(0.61), top=Inches(0.24), width=Inches(3), height=Inches(0.33))
                    tf = txBox.text_frame
                    tf.text = 'Overview'
                    tf.paragraphs[0].font.size = Pt(24)
                    tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                                
                    
                    # Write section chart
                    if showoffsetchart:
                        filename = "figproposal.png"
                        pio.write_image(figproposal, filename, scale=2, width=1280, height=888) 
                        # check pixel dimensions
                        placeholder = newslide.shapes.add_picture(filename, left=Inches(0.77), top=Inches(1.04), width=Inches(6.67), height=Inches(4.63))                     
                                
                    # Write section chart
                    if showbitchart:
                        filename = "figdvd.png"
                        pio.write_image(figdvd, filename, scale=2, width=1280, height=888) 
                        # check pixel dimensions
                        placeholder = newslide.shapes.add_picture(filename, left=Inches(0.77), top=Inches(5.5), width=Inches(6.67), height=Inches(4.63)) 
                    
                    
                    # Create Summary page
                    newslide=prs.slides.add_slide(lyt9) # adding a slide              
                    # Section Name
                    txBox = newslide.shapes.add_textbox(left=Inches(0.61), top=Inches(0.24), width=Inches(3), height=Inches(0.33))
                    tf = txBox.text_frame
                    tf.text = 'Summary'
                    tf.paragraphs[0].font.size = Pt(24)
                    tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                    
                    # Write section chart
                    if showsummarychart:
                        filename = "figwell.png"
                        pio.write_image(figwell, filename, scale=2, width=384, height=1920) 
                        # check pixel dimensions
                        placeholder = newslide.shapes.add_picture(filename, left=Inches(-0.36), top=Inches(0), width=Inches(2.2), height=Inches(11)) 
                    
                    # Data Table https://python-pptx.readthedocs.io/en/latest/user/table.html                         
                    shape = newslide.shapes.add_table(2, 12, Inches(1.52), Inches(1.37), Inches(6.86), Inches(0.75))
                    table = shape.table
                    tbl =  shape._element.graphic.graphicData.tbl
                    style_id = '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}'
                    tbl[0][-1].text = style_id                        
                    def setcell(r,c,isbold,celltext):
                        cell = table.cell(r, c)
                        cell.text = celltext                    
                        cell.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(9)
                        paragraph.font.bold = isbold 
                        cell.margin_left = cell.margin_right = cell.margin_top = cell.margin_bottom = 0
                        paragraph.alignment = PP_ALIGN.CENTER                             
                    # setcell(0,1,str(dfchart.iloc[i][j]))
                    setcell(0,0,True,f'Section')
                    setcell(0,1,True,'Size')
                    setcell(0,2,True,'Type')
                    setcell(0,3,True,'BackUp')
                    setcell(0,4,True,'Depth In')
                    setcell(0,5,True,'Depth Out')
                    setcell(0,6,True,'FTG')
                    setcell(0,7,True,'HRS')
                    setcell(0,8,True,'ROP')
                    setcell(0,9,True,'WOB')
                    setcell(0,10,True,'RPM')
                    setcell(0,11,True,'GPM')
                    
                    cols = table.columns
                    cols[0].width = Inches(0.68) # section
                    cols[1].width = Inches(0.46) # size
                    cols[2].width = Inches(0.67) # type
                    cols[3].width = Inches(0.6) # backup
                    cols[4].width = Inches(0.52) # din
                    cols[5].width = Inches(0.63) # dout
                    cols[6].width = Inches(0.44) # ftg
                    cols[7].width = Inches(0.44) # hrs
                    cols[8].width = Inches(0.46) # rop
                    cols[9].width = Inches(0.49) # wob
                    cols[10].width = Inches(0.72) # rpm
                    cols[11].width = Inches(0.76) # gpm
                    
                    for i in range(0,dfchart.shape[0]):                             
                        save_bar.progress(60+i, text=f'Creating Summary {i}..')                                
                        start_h = 1.52
                        start_v = 1.8
                        # createbitcard(newslide,dfchart.loc[i,'Bit'],(start_h + (2 * (bitcount-1))+ (0.25 * (bitcount-1))),start_v)
                        createsummarytables(newslide,i,start_h,start_v)
                            
                    
                    try:
                        
                        save_bar.progress(75, text=f'Totalling..') 
                        if showpricetotal:
                            # Write totalcost indicator
                            filename = "figtotalcost.png"
                            pio.write_image(figtotalcost, filename, scale=2, width=300, height=200) 
                            # check pixel dimensions
                            placeholder = newslide.shapes.add_picture(filename, left=Inches(1.51), top=Inches(9.29), width=Inches(3.38))
                            # Write totalcostft indicator
                            filename = "figtotalcostft.png"
                            pio.write_image(figtotalcostft, filename, scale=2, width=300, height=200) 
                            # check pixel dimensions
                            placeholder = newslide.shapes.add_picture(filename, left=Inches(4.57), top=Inches(9.29), width=Inches(3.38))
                    except Exception as e:
                        print(e)
                        st.info(f"error figtotalcost: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                    
                    # Add in Top 5 reasons closer slide
                    newslide=prs.slides.add_slide(lytTop5) # adding a slide                          
                    
                    # Create spec sheet pages    
                    for i in range(1,(len(dfchart)+1)): 
                        
                        save_bar.progress(80+i, text=f'Importing Spec Sheet {i}..') 
                        createspecsheetpage(dfchart.loc[i,'Section'],i)                       
                                                
                except Exception as e:
                    print(e)
                    st.info(f"error ppt-slides: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                        
                    
                # Close out
                save_bar.progress(92, text='Saving..') 
                binary_output = BytesIO()
                prs.save(binary_output)
                try:                        
                    target_stream = f"Adira_Proposal_{customer}_{datetime.datetime.now()}.pptx"
                    prs.save(target_stream)
                    yag.send(to=useremail, bcc='ccasad@ulterra.com', subject=f'Adira: {customer} Proposal', contents=f"Proposal created by {useremail.split('@', 1)[0]} using Adira v8.01. Please review fully before presenting to a customer. Contact Chris Casad with any questions or help.", attachments=target_stream)       
                except Exception as e:
                    print(e)
                    st.info(f"error yag test: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                
                
                try:                
                    statusupdate = save_history(df_history, df_record)
                    # st.write(f'History {statusupdate}')
                    save_bar.progress(90, text=f'{statusupdate}..')  
                    prop_exported = True
                    return binary_output
                except Exception as e:
                    print(e)
                    st.info(f"error finish: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                
            
            def create_pptx_Landscape(pagefootage,useremail,showpartnumber2,showgage2,showbody2):
                
                save_bar.progress(1, text='Initializing..')
                pptx = 'data/Adira Proposal Template.pptx'
                prs = Presentation(pptx)
                
                # declare positional variables
                WIDTH = Inches(10)
                HEIGHT = Inches(7.5)
                # left = Inches(2.5)
                # top = Inches(1)
                                    
                try:
                    # function to replace text in pptx first slide with selected filters
                    def replace_text(replacements, shapes):
                        """function to replace text on a PowerPoint slide. Takes dict of {match: replacement, ... } and replaces all matches"""
                        for shape in shapes:
                            for match, replacement in replacements.items():
                                if shape.has_text_frame:
                                    if (shape.text.find(match)) != -1:
                                        text_frame = shape.text_frame
                                        for paragraph in text_frame.paragraphs:
                                            whole_text = "".join(run.text for run in paragraph.runs)
                                            whole_text = whole_text.replace(str(match), str(replacement))
                                            for idx, run in enumerate(paragraph.runs):
                                                if idx != 0:
                                                    p = paragraph._p
                                                    p.remove(run._r)
                                            if bool(paragraph.runs):
                                                paragraph.runs[0].text = whole_text

    
                    save_bar.progress(5, text='Creating template..') 
                    lyt=prs.slide_masters[0].slide_layouts[2] # choosing a slide layout
                    lyt7=prs.slide_masters[0].slide_layouts[6] # choosing a slide layout
                    lytTop5=prs.slide_masters[0].slide_layouts[11] # choosing a slide layout
                                    
                    # declare pptx variables
                    first_slide = prs.slides[0]
                    second_slide = prs.slides[2]
                    # fourth_slide = prs.slides[3]
                    shapes_1 = []
                    shapes_2 = []
                    shapes_4 = []
                    index_to_drop = []

                    # create lists with shape objects
                    for shape in first_slide.shapes:
                        shapes_1.append(shape)

                    for shape in second_slide.shapes:
                        shapes_2.append(shape)
                        
                    # for shape in fourth_slide.shapes:
                    #     shapes_4.append(shape)

                    # Add Ulterra Logo
                    first_slide.shapes.add_picture(f'data/customer logos/Ulterra.png', left=Inches(3.43), top=Inches(2.88), height=Inches(0.76)) 
                    line1=first_slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin_x=Inches(6.67), begin_y=Inches(2.73), end_x=Inches(6.67), end_y=Inches(3.78))
                    line1.line.fill.background()
                    line1.line.fill.solid()
                    line1.line.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    line1.width = Pt(2.25)
                    # https://stackoverflow.com/questions/58792955/changing-format-of-connector-to-an-arrow-one-in-python-pptx
                    # line_elem = connector.line._get_or_add_ln()
                    # line_elem.append(parse_xml("""<a:headEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>"""))                        
                    # # or
                    # line_elem = connector.line._get_or_add_ln()
                    # line_elem.append(parse_xml("""<a:tailEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>"""))

                    # Add Customer Logo
                    try:
                        first_slide.shapes.add_picture(customerlogo, left=Inches(7.38), top=Inches(2.59), height=Inches(1.25)) # old 7.38 2.76 1
                    except Exception as e:
                        print(e)
                        st.info(f"error customerlogo: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                    # initiate a dictionary of placeholders and values to replace
                    pushwellname = f'Well: {wellname}' if len(wellname) > 1 else ''
                    pushrigname = f'Rig: {rigname}' if len(rigname) > 1 else ''
                    pushcounty = f'{gcounty}, {gstate}' if len(gcounty) > 1 else ''
                    pushlegals = f'{lat}, {long}' if len(lat) > 1 else ''
                    replaces_cover = {
                        '{company}': ' ',
                        '{document}':'PROPOSAL',
                        '{basin}': district,
                        '{year}': f'{datetime.date.today().strftime("%b")} {datetime.date.today().year}',
                        '{wellname}': pushwellname,
                        '{rigname}': pushrigname,
                        '{county}': pushcounty,
                        '{legals}': f'{pushlegals}\n{targetformation}',
                        '{preparedfor}': f'Prepared for:\n{preparedfor}',
                        '{preparedby}': f'Ulterra Rep:\n{preparedby}',
                        }

                    replaces_basin = {
                        '{region}': basins.loc[basins[basins['district'] == district].index,'region'].values[0],
                        '{district}': basins.loc[basins[basins['district'] == district].index,'district'].values[0],
                        '{description}': basins.loc[basins[basins['district'] == district].index,'description'].values[0],
                    }
                    
                    # build location string
                    countystring = f'{district}'
                    countystring += f' - {wellname}' if len(wellname) > 0 else ''
                    countystring += f' {rigname}' if len(rigname) > 0 else ''
                    countystring += f' @ {lat}, {long}' if len(str(lat)) > 0 and len(str(long)) > 0 else ''
                                                
                    # replaces_county = {
                    #     '{county}': countystring,
                    # }

                    # run the function to replace placeholders with values
                    replace_text(replaces_cover, shapes_1)
                    replace_text(replaces_basin, shapes_2)
                    # replace_text(replaces_county, shapes_4)
                    
                    
                    prs.slides[2].shapes.add_picture(f'data/maps/map-{district}.png', left=Inches(-0.02), top=Inches(-0.04),height=Inches(7)) 
                    
                except Exception as e:
                    print(e)
                    st.info(f"error ppt intro: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")

                try:
                    # slidelast=prs.slides.add_slide(lyt) # adding a slide   
                    def createbitcard(aslide, pn, lpos, tpos):
                        
                        bitsize = df.loc[df['PartNumber'] == pn,'BitSize'].values[0]
                        bittype = df.loc[df['PartNumber'] == pn,'BitType'].values[0]
                        gaugelen = df.loc[df['PartNumber'] == pn,'GaugeLength'].values[0]
                        gaugetype = df.loc[df['PartNumber'] == pn,'GaugeType'].values[0]
                        
                        # shapes2 = aslide.shapes
                        # bshape = shapes2.add_shape(MSO_SHAPE.RECTANGLE, left=Inches(lpos), top=Inches(tpos), width=Inches(5), height=Inches(2))                 
                        # bshape.fill.solid()
                        # bshape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        # bshape.line.fill.background()
                        # # MSO_SHAPE.ROUNDED_RECTANGLE
                        # # try:
                        # #     # https://python-pptx.readthedocs.io/en/latest/user/autoshapes.html#adjusting-an-autoshape
                        # #     adjs = bshape.adjustments
                            
                        # #     # # Size & Type
                        # #     # txBox = aslide.shapes.add_textbox(left=Inches(0), top=Inches(0), width=Inches(2), height=Inches(0.22))            
                        # #     # tf = txBox.text_frame
                        # #     # tf.text = f"{bshape.adjustments}"
                        # #     # tf.paragraphs[0].font.size = Pt(18)
                            
                        # #     adjs[1].effective_value = 0.1                                       
                        # # except Exception as e:
                        # #     print(e)
                        # #     # st.info(f"error shaperound: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")   
                        # # D4E6E8
                        # # set lighter or darker, -0.2 is 20% darker, 0.4 is 40% lighter
                        # # shape.fill.fore_color.brightness = 0.4
                        # # set fill to transparent (no fill)
                        # # shape.fill.background()
                        
                        
                        # if os.path.exists(f'data/Bit Pictures/{pn}-1.jpg'):
                        #     aslide.shapes.add_picture(f'data/Bit Pictures/{pn}-1.jpg', left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2))                            
                        #     # image = f'data/Bit Pictures/{pn}-1.jpg'
                        # elif os.path.exists(f'data/Bit Pictures/{pn}-1.JPG'):
                        #     aslide.shapes.add_picture(f'data/Bit Pictures/{pn}-1.JPG', left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2))                           
                        #     # image = f'data/Bit Pictures/{pn}-1.JPG'
                        # else:                           
                        #     # image = f'data/Bit Pictures/blank.jpg'
                        #     aslide.shapes.add_picture(f'data/Bit Pictures/blank.jpg', left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2)) 
                        # # aslide.shapes.add_picture(image, left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2)) 
                                                                    
                        if os.path.exists(f'{pn}-1.jpg'):
                            image = f'{pn}-1.jpg'       
                        elif os.path.exists(f'data/Bit Pictures/{pn}-1.jpg'):                                
                            image = f'data/Bit Pictures/{pn}-1.jpg' 
                        elif os.path.exists(f'data/Bit Pictures/{pn}-1.JPG'):                                
                            image = f'data/Bit Pictures/{pn}-1.JPG'
                        else:
                            image = f'data/Bit Pictures/blank.jpg'     
                            
                        imagepng = f'{pn}-1.png'
                        convert_png_transparent(image, imagepng)
                        aslide.shapes.add_picture(imagepng, left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2)) 
                        # st.image(imagepng,caption=f"{df.loc[df[df['PartNumber'] == pn].index, 'BitSize'].values[0]}  - {df.loc[df[df['PartNumber'] == bit].index, 'BitType'].values[0]}",width=150)
                        
                        # Size & Type
                        txBox = aslide.shapes.add_textbox(left=Inches(lpos - 0.43), top=Inches(tpos-0.4), width=Inches(2), height=Inches(0.4))            
                        tf = txBox.text_frame
                        # bitsize = df.loc[df['PartNumber'] == pn,'BitSize']
                        tf.text = f"{df.loc[df['PartNumber'] == pn,'BitSize'].values[0]} {bittype}"
                        tf.paragraphs[0].font.size = Pt(18)
                        txBox.fill.solid()
                        txBox.fill.fore_color.rgb = RGBColor(0, 147, 159)                    
                        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  
                        
                        #  Features
                        txBox = aslide.shapes.add_textbox(left=Inches(lpos+2.11), top=Inches(tpos), width=Inches(1.9), height=Inches(1))            
                        tf = txBox.text_frame
                        # tf.paragraphs[0].font.size = Pt(12)
                        if gaugetype.lower().find('step') > 0 and gaugetype.lower().find('taper') > 0:
                            gaugetypesimple = 'Stepped & Tapered'
                        elif gaugetype.lower().find('step') > 0:
                            gaugetypesimple = 'Stepped'
                        elif gaugetype.lower().find('taper') > 0:
                            gaugetypesimple = 'Tapered'
                        else:
                            gaugetypesimple = 'Nominal'
                        
                        feat1 = f"{df.loc[df['PartNumber'] == pn,'PartNumber'].values[0]}\v" if showpartnumber2 else ''
                        feat2 = f"{gaugelen} in Gauge\v" if showgage2 else ''
                        feat3 = f"{gaugetypesimple}\v" if showgage2 else ''
                        feat4 = f"{df.loc[df['PartNumber'] == pn,'BitMaterial'].values[0]} Body\v" if showbody2 else ''
                        # feat5 = f"{df.loc[df['PartNumber'] == pn,'PerformancePackage'].values[0]}"
                        tf.text = f"{feat1}{feat2}{feat3}{feat4}"                            
                        # tf.text = f"{gaugelen} in Gauge \v{gaugetypesimple}"
                        tf.paragraphs[0].font.size = Pt(14)
                        
                        # Technology Logo                            
                        if bittype.find('CF') >= 0:
                            techpic = 'CounterForce'
                        elif bittype.find('SPL') >= 0:
                            techpic = 'SplitBlade'
                        elif bittype.find('RPS') >= 0:
                            techpic = 'RipSaw'
                        elif bittype.find('WAV') >= 0:
                            techpic = 'WaveCut'
                        elif bittype.find('AIR') >= 0:
                            techpic = 'AirRaid'
                        elif bittype.find('XP') >= 0:
                            techpic = 'XP'
                        else:
                            techpic = None  
                        
                        if techpic is not None:
                            aslide.shapes.add_picture(f'data/tech logos/{techpic}_logo.png', left=Inches(lpos+1.72), top=Inches(tpos-.41), height=Inches(0.4))  
                            aslide.shapes.add_picture(f'data/tech logos/ad_{techpic}.png', left=Inches(9.61), top=Inches(5.01), height=Inches(1.75)) 
                            
                        # Price
                        if not pd.isna(dfchart.loc[i,'Price']) or not pd.isna(dfchart.loc[i,'Priceft']):
                            # Total Price
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+2.11), top=Inches(tpos+0.8), width=Inches(1.77), height=Inches(0.4))            
                            tf = txBox.text_frame
                            if float(dfchart.loc[i,'Priceft']) > 0:
                                tf.text = f"${float(dfchart.loc[i,'Pricetotal']):,.0f} est."                   
                            else:
                                tf.text = f"${float(dfchart.loc[i,'Pricetotal']):,.0f}"
                            tf.paragraphs[0].font.size = Pt(20)
                            tf.paragraphs[0].font.bold = True
                            txBox.fill.solid()
                            txBox.fill.fore_color.rgb = RGBColor(242, 242, 242)
                            # Price Breakdown    
                            if not pd.isna(dfchart.loc[i,'Priceft'])  and float(dfchart.loc[i,'Priceft']) > 0:                            
                                txBox = aslide.shapes.add_textbox(left=Inches(lpos+3.89), top=Inches(tpos+0.8), width=Inches(1.36), height=Inches(0.44))
                                tf = txBox.text_frame
                                tf.text = f"Flat: ${float(dfchart.loc[i,'Price']):,} \n$/Ft: {float(dfchart.loc[i,'Priceft']):,} @ {int(dfchart.loc[i,'Drilled']):,} ft" 
                                for paragraph in tf.paragraphs:
                                    paragraph.font.size = Pt(10)
                        
                        # DBR
                        if not pd.isna(dfchart.loc[i,'DBR']):                                
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+2.11), top=Inches(tpos+1.17), width=Inches(3.5), height=Inches(0.4))            
                            tf = txBox.text_frame
                            tf.text = f"DBR: {dfchart.loc[i,'DBR']}"  
                            tf.paragraphs[0].font.size = Pt(12)
                            # txBox.fill.solid()
                            # txBox.fill.fore_color.rgb = RGBColor(242, 242, 242)
                            
                        # Comments
                        if not pd.isna(dfchart.loc[i,'Comment']):                                
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+2.11), top=Inches(tpos+1.47), width=Inches(3.5), height=Inches(0.54))            
                            tf = txBox.text_frame
                            tf.text = f"Comments: \n{dfchart.loc[i,'Comment']}"  
                            for paragraph in tf.paragraphs:
                                paragraph.font.size = Pt(12) 
                            tf.paragraphs[0].font.size = Pt(12)
                            tf.paragraphs[0].font.bold = True
                        
                        # Backup
                        if not pd.isna(dfchart.loc[i,'Backup']):                                   
                            backuppn = dfchart.loc[i,'Backup']
                            backuptype = dfchart.loc[i,'Backuptype']
                            # if os.path.exists(f'data/Bit Pictures/{backuppn}-1.jpg'):
                            #     aslide.shapes.add_picture(f'data/Bit Pictures/{backuppn}-1.jpg', left=Inches(lpos+5.28), top=Inches(tpos+1.13), width=Inches(1))                            
                            #     # image = f'data/Bit Pictures/{pn}-1.jpg'
                            # elif os.path.exists(f'data/Bit Pictures/{backuppn}-1.JPG'):
                            #     aslide.shapes.add_picture(f'data/Bit Pictures/{backuppn}-1.JPG', left=Inches(lpos+5.28), top=Inches(tpos+1.13), width=Inches(1))                           
                            #     # image = f'data/Bit Pictures/{pn}-1.JPG'
                            # else:                           
                            #     # image = f'data/Bit Pictures/blank.jpg'
                            #     aslide.shapes.add_picture(f'data/Bit Pictures/blank.jpg', left=Inches(lpos+5.28), top=Inches(tpos+1.13), width=Inches(1)) 
                            # # aslide.shapes.add_picture(image, left=Inches(lpos+0.02), top=Inches(tpos+0.03), width=Inches(2)) 
                                    
                            if os.path.exists(f'{backuppn}-1.jpg'):
                                image = f'{backuppn}-1.jpg'       
                            elif os.path.exists(f'data/Bit Pictures/{backuppn}-1.jpg'):                                
                                image = f'data/Bit Pictures/{backuppn}-1.jpg' 
                            elif os.path.exists(f'data/Bit Pictures/{backuppn}-1.JPG'):                                
                                image = f'data/Bit Pictures/{backuppn}-1.JPG'
                            else:
                                image = f'data/Bit Pictures/blank.jpg'     
                                
                            imagepng = f'{backuppn}-1.png'
                            convert_png_transparent(image, imagepng)
                            aslide.shapes.add_picture(imagepng, left=Inches(lpos+0.02), top=Inches(tpos+2.5), width=Inches(1.5))  
                            
                            
                            # Backup Size & Type
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos), top=Inches(tpos+2.3), width=Inches(2), height=Inches(0.3))            
                            tf = txBox.text_frame
                            # bitsize = df.loc[df['PartNumber'] == pn,'BitSize']
                            tf.text = f"Backup: {backuptype}"
                            tf.paragraphs[0].font.size = Pt(14)
                            txBox.fill.solid()
                            txBox.fill.fore_color.rgb = RGBColor(63, 112, 119)                    
                            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) 
                        
                        # # # Status / Price Notes
                        # txBox = aslide.shapes.add_textbox(left=Inches(lpos+0.1), top=Inches(tpos+4.06), width=Inches(1.9), height=Inches(0.4))            
                        # tf = txBox.text_frame
                        # tf.text = f"{'Backup' if df.loc[df['PartNumber'] == pn,'Backup'].values[0] == 'True' else ''}"
                        # tf.paragraphs[0].font.size = Pt(10)
                        
                    
                    # slidelast=prs.slides.add_slide(lyt) # adding a slide   
                    def createappcard(aslide, pn, lpos, tpos):
                                                
                        # aslide.shapes.add_picture(f'images/icon_interval.png', left=Inches(3.15), top=Inches(2.28), height=Inches(2.49))                              
                        filename = "figsec.png"
                        pio.write_image(figsec[i], filename, scale=1, width=284, height=604) 
                        # check pixel dimensions
                        placeholder = aslide.shapes.add_picture(filename, left=Inches(lpos), top=Inches(tpos), height=Inches(3.15))
                        
                        # Interval Definition                            
                        if len(str(dfchart.loc[i,'Din'])) > 1 and len(str(dfchart.loc[i,'Dout'])) > 1 and len(str(dfchart.loc[i,'Drilled'])) > 1:
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos-1.68), top=Inches(tpos+0.46), width=Inches(1.67), height=Inches(0.34))
                            tf = txBox.text_frame
                            # text_frame.word_wrap = False
                            tf.text = f'Interval Definition:'
                            tf.paragraphs[0].font.bold = True
                            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                            # .italic = None 
                            tf.paragraphs[0].font.size = Pt(14)
                            # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                            # Interval Info
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos-1.68), top=Inches(tpos+0.72), width=Inches(1.67), height=Inches(0.81))
                            tf = txBox.text_frame
                            tf.text = f"Depth In: {float(dfchart.loc[i,'Din']):,.0f}\n Depth Out: {float(dfchart.loc[i,'Dout']):,.0f}\n Footage: {float(dfchart.loc[i,'Drilled']):,.0f}"
                            for paragraph in tf.paragraphs:
                                paragraph.font.size = Pt(14)
                                paragraph.alignment = PP_ALIGN.RIGHT
                            # paragraph.font.color.rgb = RGBColor(0, 147, 159)
                        # tf.paragraphs[0].font.size = Pt(14)
                        # tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                                                
                        # Parameters
                        # ***********ADD VALUE DETECTION TO PROCESS OUTPUTS BELOW. DO THIS FOR ALL TEXT OUTPUTS
                        if float(dfchart.loc[i,'WOBout']) > 0:
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos-1.68), top=Inches(tpos+1.45), width=Inches(1.67), height=Inches(0.34))
                            tf = txBox.text_frame
                            # text_frame.word_wrap = False
                            tf.text = f'Parameters:'
                            tf.paragraphs[0].font.bold = True
                            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                            # .italic = None 
                            tf.paragraphs[0].font.size = Pt(14)
                            # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                            # Parameters Info
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos-1.68), top=Inches(tpos+1.71), width=Inches(1.67), height=Inches(0.81))
                            tf = txBox.text_frame
                            tf.text = f"WOB: {float(dfchart.loc[i,'WOBin']):,.0f}/{float(dfchart.loc[i,'WOBout']):,.0f}\n RPM: {float(dfchart.loc[i,'RPMin']):,.0f}/{float(dfchart.loc[i,'RPMax']):,.0f}\n GPM: {dfchart.loc[i,'Flowrate']}"
                            for paragraph in tf.paragraphs:
                                paragraph.font.size = Pt(14)
                                paragraph.alignment = PP_ALIGN.RIGHT
                            # tf.paragraphs[0].font.size = Pt(14)
                            # tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                            # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                        
                            
                        # Performance
                        txBox = aslide.shapes.add_textbox(left=Inches(lpos+1.4), top=Inches(tpos+0.46), width=Inches(1.25), height=Inches(0.34))
                        tf = txBox.text_frame
                        # text_frame.word_wrap = False
                        tf.text = f'Performance:'
                        tf.paragraphs[0].font.bold = True
                        # .italic = None 
                        tf.paragraphs[0].font.size = Pt(14)
                        # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                        # Performance Info
                        txBox = aslide.shapes.add_textbox(left=Inches(lpos+1.4), top=Inches(tpos+0.72), width=Inches(2.3), height=Inches(0.81))
                        tf = txBox.text_frame
                        tf.text = f"{(dfchart.loc[i,'Drilled'] / dfchart.loc[i,'ROP']):,.1f} Hrs @ {dfchart.loc[i,'ROP']:,.1f} ft/hr"
                        tf.paragraphs[0].font.size = Pt(14)
                        # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                        
                            
                        # Motor Recs
                        if len(str(dfchart.loc[i,'Motorspecs'])) > 1:
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+1.4), top=Inches(tpos+2.18), width=Inches(2.27), height=Inches(0.34))
                            tf = txBox.text_frame
                            # text_frame.word_wrap = False
                            tf.text = f'Motor Recommendations:'
                            tf.paragraphs[0].font.bold = True
                            # .italic = None 
                            tf.paragraphs[0].font.size = Pt(14)
                            # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                            # Motor Info
                            txBox = aslide.shapes.add_textbox(left=Inches(lpos+1.94), top=Inches(tpos+2.62), width=Inches(2.34), height=Inches(0.34))
                            tf = txBox.text_frame
                            tf.text = f"{dfchart.loc[i,'Motorspecs']}"
                            tf.paragraphs[0].font.size = Pt(14)
                            aslide.shapes.add_picture(f'images/icon_revgal.png', left=Inches(lpos+1.54), top=Inches(tpos+2.49), height=Inches(0.4))     
                    
                    
                    def createsectionpage(section,i):                               
                        newslide=prs.slides.add_slide(lyt) # adding a slide              
                        # Section Name
                        txBox = newslide.shapes.add_textbox(left=Inches(1.1), top=Inches(0.24), width=Inches(3), height=Inches(0.33))
                        tf = txBox.text_frame
                        # tf.text = f'{section}'
                        tf.text = f'Bit Run {i} {section}'
                        tf.paragraphs[0].font.size = Pt(24)
                        tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                        # subtitle
                        txBox = newslide.shapes.add_textbox(left=Inches(1.1), top=Inches(0.67), width=Inches(5), height=Inches(0.33))
                        tf = txBox.text_frame
                        tf.text = f'BIT AVAILABILITY AND PRICING'
                        tf.paragraphs[0].font.size = Pt(22)
                        # description             
                        # txBox = newslide.shapes.add_textbox(left=Inches(1.1), top=Inches(1.1), width=Inches(5), height=Inches(1.35))            
                        # tf = txBox.text_frame
                        # tf.text = f"We combine our company's unique product design with a rapid \v prototyping manufacturing process, leading to reliable and \v high-end performance drill bits customers can count on in \v the {district}. The proof is in the numbers."
                        # tf.paragraphs[0].font.size = Pt(12)
                        
                        bitcount = 0
                        start_h = 2
                        start_v = 2.5
                        # createbitcard(newslide,dfchart.loc[i,'Bit'],(start_h + (2 * (bitcount-1))+ (0.25 * (bitcount-1))),start_v)
                        if not dfchart.loc[i,'Bit'] is None:
                            createbitcard(newslide,dfchart.loc[i,'Bit'],start_h,start_v)
                                                    
                        # Write section chart
                        # filename = "figproposal.png"
                        # pio.write_image(figproposal, filename, scale=2, width=1746, height=1000) 
                        # check pixel dimensions
                        # placeholder = newslide.shapes.add_picture(filename, left=Inches(1.33), top=Inches(2.23), width=Inches(4), height=Inches(4)) 
                                
                        start_h = 9.05
                        start_v = 1.68
                        createappcard(newslide,dfchart.loc[i,'Bit'],start_h,start_v)
                                                    
                        # # newslide.shapes.add_picture(f'images/icon_interval.png', left=Inches(3.15), top=Inches(2.28), height=Inches(2.49))                              
                        # filename = "figsec.png"
                        # pio.write_image(figsec[i], filename, scale=2, width=284, height=604) 
                        # # check pixel dimensions
                        # placeholder = newslide.shapes.add_picture(filename, left=Inches(2.77), top=Inches(2.02), height=Inches(3.15))
                        
                        
                        # # Interval Definition
                        # if str(dfchart.loc[i,'Din']).isnumeric() and str(dfchart.loc[i,'Dout']).isnumeric() and str(dfchart.loc[i,'Drilled']).isnumeric():
                        #     txBox = newslide.shapes.add_textbox(left=Inches(1.09), top=Inches(2.48), width=Inches(1.67), height=Inches(0.34))
                        #     tf = txBox.text_frame
                        #     # text_frame.word_wrap = False
                        #     tf.text = f'Interval Definition:'
                        #     tf.paragraphs[0].font.bold = True
                        #     tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        #     # .italic = None 
                        #     tf.paragraphs[0].font.size = Pt(14)
                        #     # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                        #     # Interval Info
                        #     txBox = newslide.shapes.add_textbox(left=Inches(1.25), top=Inches(2.74), width=Inches(1.52), height=Inches(0.81))
                        #     tf = txBox.text_frame
                        #     tf.text = f"Depth In: {float(dfchart.loc[i,'Din']):,}\n Depth Out: {float(dfchart.loc[i,'Dout']):,}\n Footage: {float(dfchart.loc[i,'Drilled']):,}"
                        #     for paragraph in tf.paragraphs:
                        #         paragraph.font.size = Pt(14)
                        #         paragraph.alignment = PP_ALIGN.RIGHT
                        #         # paragraph.font.color.rgb = RGBColor(0, 147, 159)
                        #     # tf.paragraphs[0].font.size = Pt(14)
                        #     # tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        #     # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                                                
                        # # Parameters
                        # # ***********ADD VALUE DETECTION TO PROCESS OUTPUTS BELOW. DO THIS FOR ALL TEXT OUTPUTS
                        # if float(dfchart.loc[i,'WOBout']) > 0:
                        #     txBox = newslide.shapes.add_textbox(left=Inches(1.63), top=Inches(3.47), width=Inches(1.13), height=Inches(0.34))
                        #     tf = txBox.text_frame
                        #     # text_frame.word_wrap = False
                        #     tf.text = f'Parameters:'
                        #     tf.paragraphs[0].font.bold = True
                        #     tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        #     # .italic = None 
                        #     tf.paragraphs[0].font.size = Pt(14)
                        #     # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                        #     # Parameters Info
                        #     txBox = newslide.shapes.add_textbox(left=Inches(1.43), top=Inches(3.73), width=Inches(1.34), height=Inches(0.81))
                        #     tf = txBox.text_frame
                        #     tf.text = f"WOB: {float(dfchart.loc[i,'WOBin']):,}/{float(dfchart.loc[i,'WOBout']):,}\n RPM: {float(dfchart.loc[i,'RPMin']):,}/{float(dfchart.loc[i,'RPMax']):,}\n GPM: {dfchart.loc[i,'Flowrate']}"
                        #     for paragraph in tf.paragraphs:
                        #         paragraph.font.size = Pt(14)
                        #         paragraph.alignment = PP_ALIGN.RIGHT
                        #     # tf.paragraphs[0].font.size = Pt(14)
                        #     # tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
                        #     # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                        
                            
                        # # Performance
                        # txBox = newslide.shapes.add_textbox(left=Inches(4.17), top=Inches(2.48), width=Inches(1.25), height=Inches(0.34))
                        # tf = txBox.text_frame
                        # # text_frame.word_wrap = False
                        # tf.text = f'Performance:'
                        # tf.paragraphs[0].font.bold = True
                        # # .italic = None 
                        # tf.paragraphs[0].font.size = Pt(14)
                        # # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                        # # Performance Info
                        # txBox = newslide.shapes.add_textbox(left=Inches(4.24), top=Inches(2.74), width=Inches(1.52), height=Inches(0.81))
                        # tf = txBox.text_frame
                        # tf.text = f"{(dfchart.loc[i,'Drilled'] / dfchart.loc[i,'ROP']):,.1f} Hrs @ {dfchart.loc[i,'ROP']:,.1f} ft/hr"
                        # tf.paragraphs[0].font.size = Pt(14)
                        # # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                        
                            
                        # # Motor Recs
                        # if len(str(dfchart.loc[i,'Motorspecs'])) > 1:
                        #     txBox = newslide.shapes.add_textbox(left=Inches(1.63), top=Inches(5.16), width=Inches(2.27), height=Inches(0.34))
                        #     tf = txBox.text_frame
                        #     # text_frame.word_wrap = False
                        #     tf.text = f'Motor Recommendations:'
                        #     tf.paragraphs[0].font.bold = True
                        #     # .italic = None 
                        #     tf.paragraphs[0].font.size = Pt(14)
                        #     # tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)                                                                       
                        #     # Motor Info
                        #     txBox = newslide.shapes.add_textbox(left=Inches(3.07), top=Inches(5.52), width=Inches(2.34), height=Inches(0.34))
                        #     tf = txBox.text_frame
                        #     tf.text = f"{dfchart.loc[i,'Motorspecs']}"
                        #     tf.paragraphs[0].font.size = Pt(14)
                        #     newslide.shapes.add_picture(f'images/icon_revgal.png', left=Inches(2.67), top=Inches(5.49), height=Inches(0.4)) 
                        #     # newslide.shapes.add_picture(f'images/icon_flowrate.png', left=Inches(2.67), top=Inches(5.49), height=Inches(0.4)) 
                        #     # newslide.shapes.add_picture(f'images/icon_revmin.png', left=Inches(2.67), top=Inches(6), height=Inches(0.4))    
                    
                    
                    def createspecsheetpage(section,i):    
                        newslide=prs.slides.add_slide(lyt) # adding a slide          
                        # Section Name
                        txBox = newslide.shapes.add_textbox(left=Inches(1.1), top=Inches(0.24), width=Inches(3), height=Inches(0.33))
                        tf = txBox.text_frame
                        tf.text = f'Section {i}: {section} Spec Sheets'
                        tf.paragraphs[0].font.size = Pt(24)
                        tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                        
                        
                        # PRIMARY BIT SPEC SHEET IMPORT
                        pn = dfchart.loc[i,'Bit']
                        if len(str(pn)) == 6:                 
                            try:
                                bitsize = df.loc[df['PartNumber'] == pn,'BitSize'].values[0]
                                bittype = df.loc[df['PartNumber'] == pn,'BitType'].values[0]
                                bitfeatures = df.loc[df['PartNumber'] == pn,'Features'].values[0]                     
                                # Primary
                                txBox = newslide.shapes.add_textbox(left=Inches(1.1), top=Inches(0.63), width=Inches(5), height=Inches(0.33))
                                tf = txBox.text_frame
                                tf.text = f'Primary:'
                                tf.paragraphs[0].font.size = Pt(16)
                                
                                # https://www.python-engineer.com/posts/python-zip/
                                # with zipper.open("data/assets/index.txt") as fp:
                                #     print(fp.read().decode())   
                                
                                specpdfname = f'{pn} - {float(str(bitsize)):.3f} - {bittype} - {bitfeatures}.pdf'
                                if os.path.isfile(specpdfname):
                                    pages = convert_from_path(specpdfname, 500)
                                    for count, page in enumerate(pages):
                                        page.save(f'{pn}-{count}.jpg', 'JPEG')
                                                                    
                                    filename = f'{pn}-0.jpg'
                                else:
                                    
                                    
                                    filename = f'data/Bit Pictures/blank.jpg'
                                    
                                # newslide.shapes.add_picture(f'images/icon_interval.png', left=Inches(3.15), top=Inches(2.28), height=Inches(2.49))                              
                                # filename = f'{pn}-0.jpg'
                                # pio.write_image(figsec[i], filename, scale=2, width=284, height=604) 
                                placeholder = newslide.shapes.add_picture(filename, left=Inches(1.58), top=Inches(1), height=Inches(6))                            
                                
                                if showpartnumber2 is False:
                                    shapes2 = newslide.shapes
                                    bshape = shapes2.add_shape(MSO_SHAPE.RECTANGLE, left=Inches(1.71), top=Inches(1.89), width=Inches(0.8), height=Inches(0.12))                 
                                    bshape.fill.solid()
                                    bshape.fill.fore_color.rgb = RGBColor(81, 81, 81)
                                    bshape.line.fill.background()
                            except Exception as e:
                                print(e)
                                st.info(f"error ziptest: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                        
                        # BACKUP BIT SPEC SHEET IMPORT
                        pn = dfchart.loc[i,'Backup']
                        if len(str(pn)) == 6:
                            try:
                                bitsize = df.loc[df['PartNumber'] == pn,'BitSize'].values[0]
                                bittype = df.loc[df['PartNumber'] == pn,'BitType'].values[0]
                                bitfeatures = df.loc[df['PartNumber'] == pn,'Features'].values[0]                             
                                # Primary
                                txBox = newslide.shapes.add_textbox(left=Inches(6.82), top=Inches(0.63), width=Inches(5), height=Inches(0.33))
                                tf = txBox.text_frame
                                tf.text = f'Backup:'
                                tf.paragraphs[0].font.size = Pt(16)
                                
                                specpdfname = f'{pn} - {float(str(bitsize)):.3f} - {bittype} - {bitfeatures}.pdf'
                                if os.path.isfile(specpdfname):
                                    pages = convert_from_path(specpdfname, 500)
                                    for count, page in enumerate(pages):
                                        page.save(f'{pn}-{count}.jpg', 'JPEG')
                                                                    
                                    filename = f'{pn}-0.jpg'
                                else:
                                    # url = f'https://bithub.ulterra.com/adira_files_fetch/{pn}' 
                                    # headers = {'Adira_API_Key': st.secrets['Adira_API_Key']}                                    
                                    # zipper = get_bit_files(url,headers)     
                                    # zipper.extractall('.')    
                                    
                                    if os.path.isfile(specpdfname):
                                        pages = convert_from_path(specpdfname, 500)
                                        for count, page in enumerate(pages):
                                            page.save(f'{pn}-{count}.jpg', 'JPEG')                              
                                        
                                        filename = f'{pn}-0.jpg'
                                    
                                    else:
                                        filename = f'{pn}-0.jpg'
                                # newslide.shapes.add_picture(f'images/icon_interval.png', left=Inches(3.15), top=Inches(2.28), height=Inches(2.49))                              
                                filename = f'{pn}-0.jpg'
                                # pio.write_image(figsec[i], filename, scale=2, width=284, height=604) 
                                placeholder = newslide.shapes.add_picture(filename, left=Inches(7.3), top=Inches(1), height=Inches(6))                            
                                                                
                                if showpartnumber2 is False:
                                    shapes2 = newslide.shapes
                                    bshape = shapes2.add_shape(MSO_SHAPE.RECTANGLE, left=Inches(7.42), top=Inches(1.89), width=Inches(0.8), height=Inches(0.12))                 
                                    bshape.fill.solid()
                                    bshape.fill.fore_color.rgb = RGBColor(81, 81, 81)
                                    bshape.line.fill.background()
                            except Exception as e:
                                print(e)
                                st.info(f"error ziptest: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                            
                    def createsummarytables(aslide, i, lpos, tpos):                            
                        # bitsize = df.loc[df['PartNumber'] == pn,'BitSize'].values[0]
                        # bittype = df.loc[df['PartNumber'] == pn,'BitType'].values[0]
                        # gaugelen = df.loc[df['PartNumber'] == pn,'GaugeLength'].values[0]
                        # gaugetype = df.loc[df['PartNumber'] == pn,'GaugeType'].values[0]                                                        
                        # # Size & Type
                        # txBox = aslide.shapes.add_textbox(left=Inches(lpos - 0.43), top=Inches(tpos-0.4), width=Inches(2), height=Inches(0.4)) 
                        
                        # Data Table https://python-pptx.readthedocs.io/en/latest/user/table.html 
                        try:
                            shape = newslide.shapes.add_table(4, 12, Inches(lpos), Inches(tpos+(0.43*i)), Inches(10), Inches(0.5))
                            table = shape.table
                            # https://stackoverflow.com/questions/61982333/how-to-change-default-table-style-using-pptx-python
                            # style_id lists https://github.com/scanny/python-pptx/issues/27#issuecomment-263076372                                
                            # MediumStyle1 = '{793D81CF-94F2-401A-BA57-92F5A7B2D0C5}'
                            # MediumStyle1Accent1 = '{B301B821-A1FF-4177-AEE7-76D212191A09}'
                            # MediumStyle1Accent2 = '{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}'
                            # MediumStyle1Accent3 = '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}'
                            # MediumStyle1Accent4 = '{1E171933-4619-4E11-9A3F-F7608DF75F80}'
                            # MediumStyle1Accent5 = '{FABFCF23-3B69-468F-B69F-88F6DE6A72F2}'
                            # MediumStyle1Accent6 = '{10A1B5D5-9B99-4C35-A422-299274C87663}'
                            tbl =  shape._element.graphic.graphicData.tbl
                            if 'Surface' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}'
                            elif 'Vertical' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{B301B821-A1FF-4177-AEE7-76D212191A09}'
                            elif 'Drill Out' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{FABFCF23-3B69-468F-B69F-88F6DE6A72F2}'
                            elif 'Intermediate' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{1E171933-4619-4E11-9A3F-F7608DF75F80}'
                            elif 'Curve' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{793D81CF-94F2-401A-BA57-92F5A7B2D0C5}'
                            elif 'Lateral' in f"{dfchart.iloc[i]['Section']}":
                                style_id = '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}'                                
                            else:
                                style_id = '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}'
                            tbl[0][-1].text = style_id
                            
                            def setcell(r,c,isbold,celltext):
                                cell = table.cell(r, c)
                                cell.text = celltext                    
                                cell.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT   
                                cell.margin_left = cell.margin_right = cell.margin_top = cell.margin_bottom = 0                                                                   
                                for paragraph in cell.text_frame.paragraphs:
                                    # paragraph = cell.text_frame.paragraphs[0]
                                    paragraph.font.size = Pt(8)
                                    paragraph.font.bold = isbold
                                    # paragraph.alignment = PP_ALIGN.CENTER
                                    # paragraph.alignment = PP_ALIGN.RIGHT
                                
                            def mergecell(r1,c1,r2,c2):
                                table.cell(r1, c1).merge(table.cell(r2, c2))
                                
                            # setcell(0,1,str(dfchart.iloc[i][j]))
                            setcell(0,0,True,f'Bit Run')
                            setcell(0,1,True,'Size')
                            setcell(0,2,True,'Type')
                            setcell(0,3,True,'BackUp')
                            setcell(0,4,True,'Depth In')
                            setcell(0,5,True,'Depth Out')
                            setcell(0,6,True,'FTG')
                            setcell(0,7,True,'HRS')
                            setcell(0,8,True,'ROP')
                            setcell(0,9,True,'WOB')
                            setcell(0,10,True,'RPM')
                            setcell(0,11,True,'GPM')
                            
                            setcell(1,0,False,f"Bit Run-{i+1}\n{str(dfchart.iloc[i]['Section'])}") #'Section')
                            setcell(1,1,False,str(dfchart.iloc[i]['Size'])) #'Size')
                            setcell(1,2,False,str(dfchart.iloc[i]['Type'])) #'Type')
                            setcell(1,3,False,str(dfchart.iloc[i]['Backuptype'])) #'BackUp')
                            setcell(1,4,False,str(dfchart.iloc[i]['Din'])) #'Depth In')
                            setcell(1,5,False,str(dfchart.iloc[i]['Dout'])) #'Depth Out')
                            setcell(1,6,False,str(dfchart.iloc[i]['Drilled'])) #'FTG')
                            setcell(1,7,False,f"{(float(dfchart.iloc[i]['Drilled']) / float(dfchart.iloc[i]['ROP'])):,.1f}") #'HRS')
                            setcell(1,8,False,str(dfchart.iloc[i]['ROP'])) #'ROP')
                            setcell(1,9,False,f"{dfchart.iloc[i]['WOBin']}/{dfchart.iloc[i]['WOBout']}") #'WOB')
                            setcell(1,10,False,f"{dfchart.iloc[i]['RPMin']}/{dfchart.iloc[i]['RPMax']}") #'RPM')
                            setcell(1,11,False,str(dfchart.iloc[i]['Flowrate'])) #'GPM')                                
                            
                            mergecell(1,0,3,0) # setcell(2,1, ) #'Section')
                            setcell(2,1,True,'Cost Details') #Size')
                            mergecell(2,1,2,2) 
                            setcell(2,3,True,'Rotary:') #BackUp')
                            setcell(2,4,False,f"{float(dfchart.iloc[i]['Price']):,.0f}") #Depth In')
                            setcell(2,5,True,'$/Ft:') #Depth Out')
                            setcell(2,6,False,f"{float(dfchart.iloc[i]['Priceft']):,.0f}") #FTG')
                            setcell(2,7,True,'Total Cost:') #HRS')
                            mergecell(2,7,2,8) 
                            setcell(2,9,False,f"{float(dfchart.iloc[i]['Pricetotal']):,.0f}") #ROP')
                            setcell(2,10,True,'DBR/LIH:') #WOB')
                            setcell(2,11,False,str(dfchart.iloc[i]['DBR'])) #RPM')
                            # mergecell(2,10,2,11) # setcell(2,12, ) #GPM') 
                            
                            setcell(3,1,True,'Motors:') 
                            setcell(3,2,False,str(dfchart.iloc[i]['Motorspecs']))                                 
                            mergecell(3,2,3,4) 
                            setcell(3,5,True,'Comments:') 
                            setcell(3,6,False,str(dfchart.iloc[i]['Comment']))                                
                            mergecell(3,6,3,11) 
                            
                            cols = table.columns
                            cols[0].width = Inches(1.02) # section
                            cols[1].width = Inches(0.69) # size
                            cols[2].width = Inches(1.2) # type
                            cols[3].width = Inches(0.7) # backup
                            cols[4].width = Inches(0.72) # din
                            cols[5].width = Inches(0.89) # dout
                            cols[6].width = Inches(0.55) # ftg
                            cols[7].width = Inches(0.85) # hrs
                            cols[8].width = Inches(0.64) # rop
                            cols[9].width = Inches(0.75) # wob
                            cols[10].width = Inches(0.72) # rpm
                            cols[11].width = Inches(1.29) # gpm
                            
                            def remove_row(table, row):
                                tbl = table._tbl
                                tr = row._tr
                                tbl.remove(tr)
                            
                            row = table.rows[0]
                            remove_row(table, row)
                            
                            
                            
                            # k = 1
                            # for col in dfchart.columns:
                            #     cell = table.cell(0, k-1)
                            #     cell.text = f'{col.upper()}'
                            #     paragraph = cell.text_frame.paragraphs[0]
                            #     paragraph.font.size = Pt(10)
                            #     k +=1
                                
                            # for i in range(0,dfchart.shape[0]):
                            #     for j in range(dfchart.shape[1]):
                            #         cell = table.cell(i+1, j)
                            #         cell.text = str(dfchart.iloc[i][j])                    
                            #         cell.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                            #         # cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
                            #         # cell.text_frame.fit_text() 
                            #         paragraph = cell.text_frame.paragraphs[0]
                            #         paragraph.font.size = Pt(10)
                            #         # paragraph.font.color.rgb = RGBColor(12, 34, 56)  
                                    
                            # cols = table.columns
                            # cols[0].width = Inches(0.65) # section
                            # cols[1].width = Inches(0.54) # size
                            # cols[2].width = Inches(0.68) # pn
                            # cols[3].width = Inches(0.72) # type
                            # cols[4].width = Inches(0.59) # price
                            # cols[5].width = Inches(0.3) # priceft
                            # cols[6].width = Inches(0.59) # pricetotal
                            # cols[7].width = Inches(0.65) # dbr
                            # cols[8].width = Inches(0.68) # backup
                            # cols[9].width = Inches(0.72) # backuptype
                            # cols[10].width = Inches(0.59) # din
                            # cols[11].width = Inches(0.59) # dout
                            # cols[12].width = Inches(0.37) # wobin
                            # cols[13].width = Inches(0.37) # wobout
                            # cols[14].width = Inches(0.45) # rop
                            # cols[15].width = Inches(0.45) # rpmin
                            # cols[16].width = Inches(0.45) # rpmax
                            # cols[17].width = Inches(1.46) # motor
                            # cols[18].width = Inches(0.45) # gpm
                            # cols[19].width = Inches(1.7) # comment
                            # cols[20].width = Inches(0.59) # drilled
                            
                            # def remove(table,column):
                            #     col_idx = table._tbl.tblGrid.index(column._gridCol)

                            #     for tr in table._tbl.tr_lst:
                            #         tr.remove(tr.tc_lst[col_idx])
                            #     table._tbl.tblGrid.remove(column._gridCol)
                            
                            # remove(table, table.columns[19]) # test remove comments
                            
                            # for colno in range(cols):
                            #     cols[colno].width = int(shapeWidth * widths[colno] / widths_total)      
                        except Exception as e:
                            print(e)
                            st.info(f"error ppt-table: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                        
                        
                    # dfWell['Section','Size''Bit','Backup','Din','Dout','WOBin','WOBout','ROP','RPMin','RPMax','Motorspecs','Flowrate','Price','Priceft','DBR','Comment']   
                    for i in range(1,(len(dfchart)+1)):                            
                        save_bar.progress(25+(2*i), text=f'Creating run page {i}..')   
                        createsectionpage(dfchart.loc[i,'Section'],i)
                    
                    # Create Overview page                        
                    save_bar.progress(50, text='Creating Overview..')
                    newslide=prs.slides.add_slide(lyt7) # adding a slide              
                    # Section Name
                    txBox = newslide.shapes.add_textbox(left=Inches(1.1), top=Inches(0.24), width=Inches(3), height=Inches(0.33))
                    tf = txBox.text_frame
                    tf.text = 'Overview'
                    tf.paragraphs[0].font.size = Pt(24)
                    tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                    
                            
                    # Write section chart
                    if showoffsetchart:
                        filename = "figproposal.png"
                        pio.write_image(figproposal, filename, scale=1, width=1280, height=888) 
                        # check pixel dimensions
                        placeholder = newslide.shapes.add_picture(filename, left=Inches(0), top=Inches(1.71), width=Inches(6.67), height=Inches(4.63))                     
                                
                    # Write section chart
                    if showbitchart:
                        filename = "figdvd.png"
                        pio.write_image(figdvd, filename, scale=1, width=1280, height=888) 
                        # check pixel dimensions
                        placeholder = newslide.shapes.add_picture(filename, left=Inches(6.67), top=Inches(1.59), width=Inches(6.67), height=Inches(4.63)) 
                    
                    
                    # Create Summary page
                    newslide=prs.slides.add_slide(lyt7) # adding a slide              
                    # Section Name
                    txBox = newslide.shapes.add_textbox(left=Inches(1.1), top=Inches(0.24), width=Inches(3), height=Inches(0.33))
                    tf = txBox.text_frame
                    tf.text = 'Summary'
                    tf.paragraphs[0].font.size = Pt(24)
                    tf.paragraphs[0].font.color.rgb = RGBColor(0, 147, 159)
                    
                    # Write section chart
                    if showsummarychart:
                        filename = "figwell.png"
                        pio.write_image(figwell, filename, scale=2, width=422, height=1106) 
                        # check pixel dimensions
                        placeholder = newslide.shapes.add_picture(filename, left=Inches(0), top=Inches(0), width=Inches(2.6), height=Inches(7.5)) 
                    
                    # Data Table https://python-pptx.readthedocs.io/en/latest/user/table.html                         
                    shape = newslide.shapes.add_table(4, 12, Inches(2.5), Inches(1), Inches(10), Inches(0.27))
                    table = shape.table
                    tbl =  shape._element.graphic.graphicData.tbl
                    style_id = '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}'
                    tbl[0][-1].text = style_id                        
                    def setcell(r,c,isbold,celltext):
                        cell = table.cell(r, c)
                        cell.text = celltext                    
                        cell.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(10)
                        paragraph.font.bold = isbold                              
                    # setcell(0,1,str(dfchart.iloc[i][j]))
                    setcell(0,0,True,f'Section')
                    setcell(0,1,True,'Size')
                    setcell(0,2,True,'Type')
                    setcell(0,3,True,'BackUp')
                    setcell(0,4,True,'Depth In')
                    setcell(0,5,True,'Depth Out')
                    setcell(0,6,True,'FTG')
                    setcell(0,7,True,'HRS')
                    setcell(0,8,True,'ROP')
                    setcell(0,9,True,'WOB')
                    setcell(0,10,True,'RPM')
                    setcell(0,11,True,'GPM')
                    
                    for i in range(0,dfchart.shape[0]):                                
                        start_h = 2.5
                        start_v = 1.32
                        # createbitcard(newslide,dfchart.loc[i,'Bit'],(start_h + (2 * (bitcount-1))+ (0.25 * (bitcount-1))),start_v)                            
                        save_bar.progress(70+i, text=f'Creating Summary {i}..')    
                        createsummarytables(newslide,i,start_h,start_v)
                            
                    
                    try:                            
                        save_bar.progress(75, text='Totalling..')
                        if showpricetotal:
                            # Write totalcost indicator
                            filename = "figtotalcost.png"
                            pio.write_image(figtotalcost, filename, scale=1, width=300, height=200) 
                            # check pixel dimensions
                            placeholder = newslide.shapes.add_picture(filename, left=Inches(4.1), top=Inches(6.05), width=Inches(3))
                            # Write totalcostft indicator
                            filename = "figtotalcostft.png"
                            pio.write_image(figtotalcostft, filename, scale=1, width=300, height=200) 
                            # check pixel dimensions
                            placeholder = newslide.shapes.add_picture(filename, left=Inches(8.27), top=Inches(6.05), width=Inches(3))
                    except Exception as e:
                        print(e)
                        st.info(f"error figtotalcost: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                    
                    
                    # Add in Top 5 reasons closer slide
                    newslide=prs.slides.add_slide(lytTop5) # adding a slide                          
                    
                    # Create spec sheet pages    
                    for i in range(1,(len(dfchart)+1)):                             
                        save_bar.progress(80+i, text=f'Importing Spec Sheet {i}..') 
                        createspecsheetpage(dfchart.loc[i,'Section'],i)                       
                                                
                except Exception as e:
                    print(e)
                    st.info(f"error ppt-slides: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                        
                    
                # Close out
                # save_bar.progress(92, text='Saving..') 
                binary_output = BytesIO()
                prs.save(binary_output)
                try:
                    target_stream = f"Adira_Proposal_{customer}_{datetime.datetime.now()}.pptx"
                    prs.save(target_stream)
                    yag.send(to=useremail, bcc='ccasad@ulterra.com', subject=f'Adira: {customer} Proposal', contents=f"Proposal created by {useremail.split('@', 1)[0]} using Adira v8.01. Please review fully before presenting to a customer. Contact Chris Casad with any questions or help.", attachments=target_stream)       
                except Exception as e:
                    print(e)
                    st.info(f"error yag test: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
                
                try:
                    statusupdate = save_history(df_history, df_record)
                    
                    st.write(f'History {statusupdate}')
                    save_bar.progress(90, text=f'{statusupdate}..')  
                    prop_exported = True
                    return binary_output
                
                except Exception as e:
                    print(e)
                    st.info(f"error finish: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}") 
        
        
        
        # old placement for form setup and df_record definition
        
        
        
            prop_exported = False
            sendemailflag = 'Please provide your email address.' 
            bsendemail = False
            if submit_pptx:
                # Validate user entered email
                if len(str(sendemail)) > 0:
                    if sendemail.partition('@')[-1]  == 'ulterra.com':
                        bsendemail = True
                    else:
                        sendemailflag = f"Please use an Ulterra email address. Found: {sendemail.partition('@')[-1]}"
                else:
                    sendemailflag = f'Please provide your email address. {len(str(sendemail))}'
                
                
                if bsendemail is True:    
                    emailerror = st.empty()            
                    save_bar = st.progress(0, text='Initializing..')                        
                    df_record.insert(4,'email', sendemail)
                    
                    with st.spinner('Constructing...'):    
                        if reporttype == 'PowerPoint':
                            if reportsize == 'Landscape':
                                try:
                                    reportppt = create_pptx_Landscape(pagefootage,sendemail,showpartnumber2,showgage2,showbody2)
                                    try:                
                                        save_bar.progress(95, text='Preparing Download')
                                        st.download_button(
                                            label="Download PowerPoint",
                                            data=reportppt,
                                            file_name=f"Adira_Proposal_{customer}_{datetime.datetime.now()}.pptx",
                                            mime="application/octet-stream",
                                        ) 
                                        save_bar.progress(100, text='Ready')
                                    except Exception as e:
                                        st.info(f"error pptx dl1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                             
                                except Exception as e:
                                    st.info(f"error pptx: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                    
                            
                            if reportsize == 'Portrait':
                                try:
                                    reportppt = create_pptx_Portrait(pagefootage,sendemail,showpartnumber2,showgage2,showbody2)
                                    try:
                                        save_bar.progress(95, text='Preparing Download')
                                        st.download_button(
                                            label="Download PowerPoint",
                                            data=reportppt,
                                            file_name=f"Adira_Proposal_{customer}_{datetime.datetime.now()}.pptx",
                                            mime="application/octet-stream",
                                        ) 
                                        save_bar.progress(100, text='Ready')
                                    except Exception as e:
                                        st.info(f"error pptx dl1: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")                                             
                                except Exception as e:
                                    st.info(f"error pptx: {e} //// {traceback.print_exc()} //// {traceback.format_exc()}")
                else:
                    emailerror = st.write(f':red[{sendemailflag}]')
                
            
                
                    
            # contents = [f"error date time: {e}", 'user info:', userdata,
            #             f'File name: {uploaded_file.name}',
            #             f'//// {traceback.print_exc()} //// {traceback.format_exc()}',
            #             "You can find an audio file attached.", '/local/path/to/song.mp3'
            # ]
            # yag.send('ccasad@ulterra.com', 'Adira: error date time', contents, attachments=uploaded_file)
            
            
            with st.expander('Data Table', expanded=False):
                dfchart = st.experimental_data_editor(dfchart, key='editor_chart')  
                # st.write(dfchart.dtypes)
                                        
            # with st.expander('Record Table', expanded=True):
            #     df_record = st.experimental_data_editor(df_record, key='editor_record') 
            #     csv = df_record.to_csv(index=False).encode('utf-8')
            #     st.download_button("Download Data",csv,"Proposal_history.csv","text/csv",key='download-csv')
            
            st.divider
            
        # Call app
        WellProposal()


            # # Sample usage
            # if __name__ == "__main__":
            #     # Replace 'target_record' with the record you want to use as the target for matching
            #     # target_record = {
            #     #     'Latitude': 37.7749,
            #     #     'Longitude': -122.4194,
            #     #     'Size': 10,
            #     #     'speed': 50,
            #     #     'start_mileage': 0,
            #     #     'end_mileage': 1000
            #     # }
                
            #     target_lat = lat # st.number_input("Latitude:")
            #     target_lon = long # st.number_input("Longitude:")
            #     target_size = dfchart.loc[1,'Size'] # st.number_input("Size:")
            #     target_speed = dfchart.loc[1,'ROP'] # st.number_input("Speed:")
            #     target_start_mileage = dfchart.loc[1,'Din'] # st.number_input("Start Mileage:")
            #     target_end_mileage = dfchart.loc[1,'Dout'] # st.number_input("End Mileage:")
            #     target_class = dfchart.loc[1,'Size'] # st.text_input("Vehicle Class:")

            #     target_record = {
            #         'Latitude': target_lat,
            #         'Longitude': target_lon,
            #         'Size': target_size,
            #         'ROP': target_speed,
            #         'DepthIn': target_start_mileage,
            #         'DepthOut': target_end_mileage,
            #     }

            #     csv_file = 'records_USonly.csv'
            #     vehicle_class = 'SUV'  # Replace this with the desired vehicle class
            #     num_matches = 3  # The number of high ranking matches to find
            #     max_distance = 100  # The maximum distance for initial GPS location search

            #     high_ranking_matches = find_high_ranking_matches(csv_file, target_record, vehicle_class, num_matches, max_distance)
            #     print(high_ranking_matches)

            # st.markdown('#### ')


        

# import pandas as pd
# from geopy.distance import geodesic

# def calculate_distance(source_coords, target_coords):
#     return geodesic(source_coords, target_coords).miles

# def calculate_ranking_score(record, target_record, vehicle_class):
#     distance_score = calculate_distance((record['latitude'], record['longitude']), 
#                                         (target_record['latitude'], target_record['longitude']))
#     size_score = abs(record['Size'] - target_record['Size'])
#     speed_score = abs(record['speed'] - target_record['speed'])

#     # Check if the vehicle class is available in the record, otherwise consider it as a separate class
#     if 'class' in record:
#         class_score = 0 if record['class'] == vehicle_class else 1
#     else:
#         class_score = 1

#     # Assuming the target_record includes the start and end mileage points
#     start_mileage = target_record['start_mileage']
#     end_mileage = target_record['end_mileage']
#     segment_start = record['start_mileage']
#     segment_end = record['end_mileage']

#     # Calculate the absolute difference between the middle segment start and end mileage
#     segment_mileage_diff = abs((segment_end - segment_start) - (end_mileage - start_mileage))

#     # You can adjust the weights below based on how you want to prioritize the factors.
#     # For example, if distance is more important, increase its weight, and vice versa.
#     distance_weight = 0.25
#     size_weight = 0.15
#     speed_weight = 0.15
#     class_weight = 0.2
#     segment_weight = 0.25

#     total_score = (distance_weight * distance_score +
#                    size_weight * size_score +
#                    speed_weight * speed_score +
#                    class_weight * class_score +
#                    segment_weight * segment_mileage_diff)

#     return total_score

# def rank_records(csv_file, target_record, vehicle_class):
#     df = pd.read_csv(csv_file)

#     # Calculate ranking score for each record in the dataframe
#     df['ranking_score'] = df.apply(lambda row: calculate_ranking_score(row, target_record, vehicle_class), axis=1)

#     # Sort the dataframe based on ranking score in ascending order
#     ranked_df = df.sort_values(by='ranking_score')

#     return ranked_df

# def find_high_ranking_matches(csv_file, target_record, vehicle_class, num_matches=3):
#     ranked_records = rank_records(csv_file, target_record, vehicle_class)
#     high_ranking_matches = ranked_records.head(num_matches)
    
#     while len(high_ranking_matches) < num_matches:
#         # Gradually broaden the class pool and re-rank
#         vehicle_class = 'Any'
#         ranked_records = rank_records(csv_file, target_record, vehicle_class)
#         high_ranking_matches = ranked_records.head(num_matches)

#     return high_ranking_matches

# def find_high_ranking_matches2(csv_file, target_record, vehicle_class, num_matches=3, max_distance=100):
#     ranked_records = rank_records(csv_file, target_record, vehicle_class)
#     high_ranking_matches = ranked_records.head(num_matches)
#     current_distance = 0

#     while len(high_ranking_matches) < num_matches and current_distance <= max_distance:
#         current_distance += 10  # Increase the distance search by 10 miles
#         # Re-rank the records based on the increased distance
#         ranked_records = rank_records(csv_file, target_record, vehicle_class)

#         # Filter the records based on the current distance and vehicle class
#         filtered_records = ranked_records[
#             (ranked_records['ranking_score'] <= current_distance) & 
#             ((ranked_records['class'] == vehicle_class) | (vehicle_class == 'Any'))
#         ]

#         high_ranking_matches = filtered_records.head(num_matches)

#     if len(high_ranking_matches) < num_matches:
#         # If still not enough high ranking matches, expand the search to other class sizes
#         vehicle_class = 'Any'
#         ranked_records = rank_records(csv_file, target_record, vehicle_class)
#         high_ranking_matches = ranked_records.head(num_matches)

#     return high_ranking_matches

# def find_high_ranking_matches3(csv_file, target_record, vehicle_class, num_matches=3, max_distance=100):
#     df = pd.read_csv(csv_file)
#     key_column = 'matching_key'  # Name of the column to store the matching key or tag

#     # Check if there are already records with the same matching key
#     previous_matches = df[df[key_column].notnull()]

#     if len(previous_matches) >= num_matches:
#         # If we have enough previous matches, use them to expedite future matching
#         filtered_records = df[df[key_column].isin(previous_matches[key_column])]
#     else:
#         # Otherwise, perform a new search using the previous logic
#         ranked_records = rank_records(csv_file, target_record, vehicle_class)
#         high_ranking_matches = ranked_records.head(num_matches)
#         current_distance = 0

#         while len(high_ranking_matches) < num_matches and current_distance <= max_distance:
#             current_distance += 10  # Increase the distance search by 10 miles
#             # Re-rank the records based on the increased distance
#             ranked_records = rank_records(csv_file, target_record, vehicle_class)

#             # Filter the records based on the current distance and vehicle class
#             filtered_records = ranked_records[
#                 (ranked_records['ranking_score'] <= current_distance) & 
#                 ((ranked_records['class'] == vehicle_class) | (vehicle_class == 'Any'))
#             ]

#             high_ranking_matches = filtered_records.head(num_matches)

#         if len(high_ranking_matches) < num_matches:
#             # If still not enough high ranking matches, expand the search to other class sizes
#             vehicle_class = 'Any'
#             ranked_records = rank_records(csv_file, target_record, vehicle_class)
#             high_ranking_matches = ranked_records.head(num_matches)

#         # Store the matching key or tag in the DataFrame for future matching
#         matching_key = 'key_{}'.format(len(previous_matches) + 1)
#         df.at[high_ranking_matches.index, key_column] = matching_key
#         filtered_records = df[df[key_column].isin([matching_key])]

#     return filtered_records

# # Sample usage
# if __name__ == "__main__":
#     # Replace 'target_record' with the record you want to use as the target for matching
#     target_record = {
#         'latitude': 37.7749,
#         'longitude': -122.4194,
#         'Size': 10,
#         'speed': 50,
#         'start_mileage': 0,
#         'end_mileage': 1000
#     }

#     csv_file = 'path/to/your/csv/file.csv'
#     vehicle_class = 'SUV'  # Replace this with the desired vehicle class
#     num_matches = 3  # The number of high ranking matches to find

#     high_ranking_matches = find_high_ranking_matches(csv_file, target_record, vehicle_class, num_matches)
#     print(high_ranking_matches)





    # # import pandas as pd
    # # from geopy.distance import geodesic

    # def calculate_distance(source_coords, target_coords):
    #     return geodesic(source_coords, target_coords).miles

    # def calculate_ranking_score(record, target_record, vehicle_class):
    #     distance_score = calculate_distance((record['Latitude'], record['Longitude']), 
    #                                         (target_record['Latitude'], target_record['Longitude']))
    #     size_score = abs(record['Size'] - target_record['Size'])
    #     speed_score = abs(record['Speed'] - target_record['Speed'])

    #     # Check if the vehicle class is available in the record, otherwise consider it as a separate class
    #     if 'class' in record:
    #         class_score = 0 if record['class'] == vehicle_class else 1
    #     else:
    #         class_score = 1

    #     # Assuming the target_record includes the start and end mileage points
    #     start_mileage = target_record['start_mileage']
    #     end_mileage = target_record['end_mileage']
    #     segment_start = record['start_mileage']
    #     segment_end = record['end_mileage']

    #     # Calculate the absolute difference between the middle segment start and end mileage
    #     segment_mileage_diff = abs((segment_end - segment_start) - (end_mileage - start_mileage))

    #     # You can adjust the weights below based on how you want to prioritize the factors.
    #     # For example, if distance is more important, increase its weight, and vice versa.
    #     distance_weight = 0.25
    #     size_weight = 0.15
    #     speed_weight = 0.15
    #     class_weight = 0.2
    #     segment_weight = 0.25

    #     total_score = (distance_weight * distance_score +
    #                 size_weight * size_score +
    #                 speed_weight * speed_score +
    #                 class_weight * class_score +
    #                 segment_weight * segment_mileage_diff)

    #     return total_score

    # def rank_records(csv_file, target_record, vehicle_class):
    #     dfrr = pd.read_csv(csv_file)

    #     # Calculate ranking score for each record in the dataframe
    #     dfrr['ranking_score'] = dfrr.apply(lambda row: calculate_ranking_score(row, target_record, vehicle_class), axis=1)

    #     # Sort the dataframe based on ranking score in ascending order
    #     ranked_df = dfrr.sort_values(by='ranking_score')

    #     return ranked_df

    # def find_high_ranking_matches(csv_file, target_record, vehicle_class, num_matches=3, max_distance=100):
    #     dfhrr = pd.read_csv(csv_file)
    #     key_column = 'matching_key'  # Name of the column to store the matching key or tag

    #     # Check if there are already records with the same matching key
    #     previous_matches = dfhrr[dfhrr[key_column].notnull()]

    #     if len(previous_matches) >= num_matches:
    #         # If we have enough previous matches, use them to expedite future matching
    #         filtered_records = dfhrr[dfhrr[key_column].isin(previous_matches[key_column])]
    #     else:
    #         # Otherwise, perform a new search using the previous logic
    #         ranked_records = rank_records(csv_file, target_record, vehicle_class)
    #         high_ranking_matches = ranked_records.head(num_matches)
    #         current_distance = 0

    #         while len(high_ranking_matches) < num_matches and current_distance <= max_distance:
    #             current_distance += 10  # Increase the distance search by 10 miles
    #             # Re-rank the records based on the increased distance
    #             ranked_records = rank_records(csv_file, target_record, vehicle_class)

    #             # Filter the records based on the current distance and vehicle class
    #             filtered_records = ranked_records[
    #                 (ranked_records['ranking_score'] <= current_distance) & 
    #                 ((ranked_records['class'] == vehicle_class) | (vehicle_class == 'Any'))
    #             ]

    #             high_ranking_matches = filtered_records.head(num_matches)

    #         if len(high_ranking_matches) < num_matches:
    #             # If still not enough high ranking matches, expand the search to other class sizes
    #             vehicle_class = 'Any'
    #             ranked_records = rank_records(csv_file, target_record, vehicle_class)
    #             high_ranking_matches = ranked_records.head(num_matches)

    #         # Store the matching key or tag in the DataFrame for future matching
    #         matching_key = 'key_{}'.format(len(previous_matches) + 1)
    #         df.at[high_ranking_matches.index, key_column] = matching_key
    #         filtered_records = dfhrr[dfhrr[key_column].isin([matching_key])]

    #     return filtered_records

    # # Sample usage
    # if __name__ == "__main__":
    #     # Replace 'target_record' with the record you want to use as the target for matching
    #     # target_record = {
    #     #     'latitude': 37.7749,
    #     #     'longitude': -122.4194,
    #     #     'Size': 10,
    #     #     'speed': 50,
    #     #     'start_mileage': 0,
    #     #     'end_mileage': 1000
    #     # }
        
    #     target_lat = lat # st.number_input("Latitude:")
    #     target_lon = long # st.number_input("Longitude:")
    #     target_size = dfchart.loc[1,'Size'] # st.number_input("Size:")
    #     target_speed = dfchart.loc[1,'ROP'] # st.number_input("Speed:")
    #     target_start_mileage = dfchart.loc[1,'Din'] # st.number_input("Start Mileage:")
    #     target_end_mileage = dfchart.loc[1,'Dout'] # st.number_input("End Mileage:")
    #     target_class = dfchart.loc[1,'Size'] # st.text_input("Vehicle Class:")

    #     target_record = {
    #         'Latitude': target_lat,
    #         'Longitude': target_lon,
    #         'Size': target_size,
    #         'ROP': target_speed,
    #         'DepthIn': target_start_mileage,
    #         'DepthOut': target_end_mileage,
    #     }

    #     csv_file = 'records_USonly.csv'
    #     vehicle_class = 'SUV'  # Replace this with the desired vehicle class
    #     num_matches = 3  # The number of high ranking matches to find
    #     max_distance = 100  # The maximum distance for initial GPS location search

    #     high_ranking_matches = find_high_ranking_matches(csv_file, target_record, vehicle_class, num_matches, max_distance)
    #     print(high_ranking_matches)


    # def main():
    #     st.title("Matchmaking Web App")
    #     st.write("Upload a CSV data file and test the matchmaking algorithm.")

    #     # uploaded_file = st.file_uploader("Choose a CSV file", type=["csv"])
    #     uploaded_file = 'records_USonly.csv'
        
    #     if uploaded_file:
    #         df = pd.read_csv(uploaded_file)

    #         # Create input fields for the target record
    #         st.caption("Enter Target Record Information:")
    #         target_lat = lat # st.number_input("Latitude:")
    #         target_lon = long # st.number_input("Longitude:")
    #         target_size = dfchart.loc[1,'Size'] # st.number_input("Size:")
    #         target_speed = dfchart.loc[1,'ROP'] # st.number_input("Speed:")
    #         target_start_mileage = dfchart.loc[1,'Din'] # st.number_input("Start Mileage:")
    #         target_end_mileage = dfchart.loc[1,'Dout'] # st.number_input("End Mileage:")
    #         target_class = dfchart.loc[1,'Size'] # st.text_input("Vehicle Class:")

    #         target_record = {
    #             'Latitude': target_lat,
    #             'Longitude': target_lon,
    #             'Size': target_size,
    #             'ROP': target_speed,
    #             'DepthIn': target_start_mileage,
    #             'DepthOut': target_end_mileage,
    #         }

    #         if target_class:
    #             target_record['Size'] = target_class

    #         # Perform the matchmaking algorithm and display results
    #         st.caption("Matchmaking Results:")
    #         matched_records = find_high_ranking_matches(df, target_record, vehicle_class=target_class)
    #         st.write(matched_records)
